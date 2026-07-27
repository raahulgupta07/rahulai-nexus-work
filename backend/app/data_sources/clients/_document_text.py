"""Best-effort plain-text extraction for rich document formats.

Lets file sources make PDF / Word / PowerPoint content *searchable* and
*readable* (not just opaque bytes). Dependency-light: PDF via pypdf, PPTX via
python-pptx (both already project deps), DOCX by reading the OOXML zip directly
(no python-docx needed).

Every extractor is defensive — a corrupt/locked/oversized file or a missing
optional lib yields "" rather than raising, so search over a mixed directory
never dies on one bad file.
"""
from __future__ import annotations

import logging
import re
import zipfile
from typing import Optional

logger = logging.getLogger(__name__)

# pypdf is very chatty on real-world PDFs ("Ignoring wrong pointing object",
# "FloatObject invalid", …) — harmless recovery warnings that flood the logs
# when scanning a whole directory. Silence them; genuine errors still surface
# via our own try/except below.
logging.getLogger("pypdf").setLevel(logging.ERROR)

# Extensions this module can turn into text. Callers gate on this set.
DOC_EXTS = {"pdf", "docx", "pptx"}

# Below this many non-whitespace characters, a rich document's text extraction
# is treated as failed — the common case being a scanned / image-based PDF or
# one using CID fonts with no ToUnicode map, where pypdf returns nothing or a
# stray glyph. Callers fall back to raw bytes so the file can be rendered to
# images for a vision model instead of surfacing an empty/garbage "text" read.
MIN_USABLE_DOC_CHARS = 16

# Formats whose extractor is EXACT rather than heuristic. OOXML stores literal
# text runs (<w:t> / <a:t>), so whatever comes back IS the document's text —
# there is no "scanned docx" failure mode for the length floor to catch.
_EXACT_TEXT_EXTS = {"docx", "pptx"}


def doc_text_is_usable(text, ext: Optional[str] = None) -> bool:
    """True when extracted document text is substantive enough to use as-is.

    ``MIN_USABLE_DOC_CHARS`` exists for PDFs, where a near-empty extraction
    signals a scanned/image-based page and the caller should fall back to raw
    bytes for raster + vision. It must NOT be applied to OOXML: those
    extractors are exact, so a short result means a short *document*, not a
    failed read. Applying the floor there discarded a correct extraction of a
    legitimately brief file (a one-line memo) and handed the caller bytes it
    had no way to recover from — unlike PDF, docx/pptx have no image fallback,
    so the read dead-ended as an opaque "binary" result.

    Pass ``ext`` (the source file's extension) to get the format-aware rule;
    omitting it keeps the conservative PDF-style floor.
    """
    if not text or not str(text).strip():
        return False
    if ext and str(ext).lstrip(".").lower() in _EXACT_TEXT_EXTS:
        return True
    return len(str(text).strip()) >= MIN_USABLE_DOC_CHARS


# Garble detection samples at most this many characters — plenty of signal,
# bounded cost on huge extractions.
_GARBLE_SAMPLE_CHARS = 8000
# Below this many non-whitespace chars the length gate (doc_text_is_usable)
# is the authority; a garble verdict on a handful of chars is noise.
_GARBLE_MIN_CHARS = 40
# Real prose in any script runs 60-80% letters; glyph-remapped soup runs ~10-20%.
_GARBLE_MAX_ALPHA_RATIO = 0.35
# In real text nearly all letters sit inside multi-letter words. In remapped
# soup letters appear as isolated singletons between symbols.
_GARBLE_MAX_WORD_LETTER_RATIO = 0.05

_WORD_RUN_RE = re.compile(r"[^\W\d_]{3,}", re.UNICODE)


def doc_text_looks_garbled(text) -> bool:
    """True when a document extraction "succeeded" but produced glyph soup.

    The signature failure: a subset-font PDF with a missing/broken ToUnicode
    CMap renders pixel-perfect but extracts as raw glyph codes — low-ASCII
    punctuation/digit soup with no words (e.g. ``$ * Q P % 1 0 $ +``). Length
    alone can't catch it (hundreds of chars come back), so this checks *shape*:
    flag only when letters are rare AND almost none of them form multi-letter
    words. Both conditions must hold, so numeric tables with real header words
    and prose in any script (Latin, Hebrew, CJK…) pass through. Deliberately
    conservative — a miss here still has the model-side ``as_images`` escape
    hatch behind it, while a false positive would burn a vision render on a
    readable file."""
    if not text:
        return False
    s = str(text)[:_GARBLE_SAMPLE_CHARS]
    non_ws = [c for c in s if not c.isspace()]
    n = len(non_ws)
    if n < _GARBLE_MIN_CHARS:
        return False
    letters = sum(1 for c in non_ws if c.isalpha())
    if letters / n >= _GARBLE_MAX_ALPHA_RATIO:
        return False
    word_letters = sum(len(m) for m in _WORD_RUN_RE.findall(s))
    return word_letters < _GARBLE_MAX_WORD_LETTER_RATIO * n

# Default cap on extracted characters. Bounds both memory and how much of a big
# document we scan on every search. ~200k chars ≈ 40-50 pages of prose.
DEFAULT_MAX_CHARS = 200_000

# How many PDF pages to read before giving up (paired with the char cap).
_PDF_MAX_PAGES = 200


def _ext(name: str) -> str:
    return name.rsplit(".", 1)[-1].lower() if name and "." in name else ""


def sanitize_extracted_text(text: str) -> str:
    """Make extractor output safe to persist and serialize.

    pypdf passes lone UTF-16 surrogates (U+D800–DFFF) straight through from
    broken/hostile ToUnicode CMaps. Python strings hold them happily, so the
    text survives all the way into stored JSON — and then every later
    serialization (`json.dumps(...).encode("utf-8")` in the API response)
    raises `surrogates not allowed`, 500ing the report permanently. NUL is
    stripped for the same reason (Postgres JSONB rejects it). This is THE
    chokepoint: every document extractor returns through here.
    """
    if not text:
        return text
    # Fast path: already clean (the overwhelmingly common case).
    try:
        text.encode("utf-8")
        clean = text
    except UnicodeEncodeError:
        clean = text.encode("utf-8", errors="replace").decode("utf-8")
    return clean.replace("\x00", "") if "\x00" in clean else clean


def extract_document_text(path: str, name: Optional[str] = None,
                          max_chars: int = DEFAULT_MAX_CHARS) -> str:
    """Return extracted text for a pdf/docx/pptx file, or "" if unsupported /
    unreadable. `name` overrides the extension source (defaults to `path`)."""
    ext = _ext(name or path)
    try:
        if ext == "pdf":
            return sanitize_extracted_text(_pdf(path, max_chars))
        if ext == "docx":
            return sanitize_extracted_text(_docx(path, max_chars))
        if ext == "pptx":
            return sanitize_extracted_text(_pptx(path, max_chars))
    except Exception as e:  # never let one bad file break a search
        # DEBUG, not WARNING: on a real directory plenty of files legitimately
        # fail (encrypted, corrupt, Office temp stubs) — that's expected noise,
        # not an actionable problem. The file is simply skipped.
        logger.debug("extract_document_text skipped %s: %s", path, e)
    return ""


def extract_document_text_from_bytes(data: bytes, name: str,
                                     max_chars: int = DEFAULT_MAX_CHARS) -> str:
    """Adapt the path-based document extractor to in-memory bytes (S3 objects,
    Graph/Drive downloads) by writing to a NamedTemporaryFile with the right
    extension (the extractors dispatch on filename). Returns "" on any failure."""
    import os
    import tempfile

    leaf = name.rsplit("/", 1)[-1]
    suffix = "." + _ext(leaf) if _ext(leaf) else ""
    tmp = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as fh:
            fh.write(data)
            tmp = fh.name
        return extract_document_text(tmp, leaf, max_chars=max_chars) or ""
    except Exception:
        return ""
    finally:
        if tmp:
            try:
                os.unlink(tmp)
            except OSError:
                pass


def extract_pdf_pages_text(path: str, first: int, last: int,
                           max_chars: int = DEFAULT_MAX_CHARS) -> tuple:
    """Extract text for an inclusive 1-based page range of a PDF.

    Returns (text, pages_total). Pages outside the document are simply not
    read (an empty range yields ""). Raises on an unreadable/locked PDF so the
    caller can surface a real error — unlike the search-oriented
    extract_document_text, a targeted page read failing silently would strand
    the model in a retry loop.
    """
    from pypdf import PdfReader

    reader = PdfReader(path)
    if getattr(reader, "is_encrypted", False):
        reader.decrypt("")
    pages_total = len(reader.pages)
    first = max(1, int(first))
    last = min(pages_total, int(last))
    parts: list[str] = []
    total = 0
    for i in range(first - 1, last):
        try:
            t = reader.pages[i].extract_text() or ""
        except Exception:
            t = ""
        parts.append(t)
        total += len(t)
        if total >= max_chars:
            break
    return sanitize_extracted_text("\n".join(parts)[:max_chars]), pages_total


def _pdf(path: str, max_chars: int) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)
    # Encrypted PDFs: many use an empty owner password — try that before giving
    # up so we can still read them. If it's genuinely locked, decrypt raises and
    # the caller catches it (file skipped).
    if getattr(reader, "is_encrypted", False):
        try:
            reader.decrypt("")
        except Exception:
            return ""
    parts: list[str] = []
    total = 0
    for i, page in enumerate(reader.pages):
        if i >= _PDF_MAX_PAGES:
            break
        try:
            t = page.extract_text() or ""
        except Exception:
            continue
        parts.append(t)
        total += len(t)
        if total >= max_chars:
            break
    return "\n".join(parts)[:max_chars]


# Match the visible-text runs in an OOXML part: <w:t> (Word) / <a:t>
# (PowerPoint). The tag name is ANCHORED — `t` must be followed by whitespace
# or `>` — because an unanchored `t[^>]*` also matches <w:tbl>/<w:tr>/<w:tc>/
# <w:tab>, which made every docx containing a TABLE extract with raw XML
# markup interleaved into its text. The namespace prefix is left open (any
# `\w+:`) since non-Microsoft generators emit prefixes like <ns0:t>.
_OOXML_TEXT_RE = re.compile(r"<(?:\w+:)?t(?:\s[^>]*)?>(.*?)</(?:\w+:)?t>", re.DOTALL)

# Any leftover XML/OOXML tag — used as a belt-and-suspenders scrub on OOXML text
# by _strip_residual_ooxml() below (our defence-in-depth on docx/pptx output).
_RESIDUAL_TAG_RE = re.compile(r"<[^>]+>")

# Sniff for a WordprocessingML file that is NOT a zip: flat OPC / "Word 2003
# XML" documents — single-file XML that Word saves/exports with a .docx or
# .xml name. zipfile chokes on them, but the text runs are scrapeable with
# the same regex as the zipped parts.
_WORDML_MARKERS = (
    b"wordprocessingml",        # flat OPC + modern namespaces
    b"word/2003/wordml",        # Word 2003 XML namespace
    b"progid=\"Word.Document\"",  # mso-application processing instruction
)


def _unescape(s: str) -> str:
    return (s.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
             .replace("&quot;", '"').replace("&apos;", "'"))


def _strip_residual_ooxml(text: str) -> str:
    """Defensive scrub for OOXML (docx/pptx) extractor output.

    The tightened ``_OOXML_TEXT_RE`` should already yield clean text, but any
    residual ``<...>`` markup (from a malformed part or an unforeseen element
    layout) is removed here and whitespace is collapsed, so NO ``<w:``/``<a:``/
    ``<p:`` tags can ever survive into an instruction / knowledge row. Applied
    only to OOXML formats — never to PDF text, where ``<``/``>`` can be legit
    content (e.g. ``a < b``).
    """
    if not text or "<" not in text:
        return text
    cleaned = _RESIDUAL_TAG_RE.sub(" ", text)
    cleaned = re.sub(r"[ \t]+", " ", cleaned)
    cleaned = re.sub(r" *\n *", "\n", cleaned)
    return cleaned.strip()


def _docx(path: str, max_chars: int) -> str:
    if not zipfile.is_zipfile(path):
        return _flat_wordml(path, max_chars)
    with zipfile.ZipFile(path) as z:
        # Body plus headers/footers so contract text in any part is searchable.
        names = [n for n in z.namelist()
                 if n == "word/document.xml"
                 or (n.startswith("word/") and (n.startswith("word/header")
                     or n.startswith("word/footer")))]
        out: list[str] = []
        total = 0
        for n in names:
            xml = z.read(n).decode("utf-8", "ignore")
            for m in _OOXML_TEXT_RE.findall(xml):
                txt = _unescape(m)
                out.append(txt)
                total += len(txt)
            if total >= max_chars:
                break
    return _strip_residual_ooxml(" ".join(out))[:max_chars]


def _flat_wordml(path: str, max_chars: int) -> str:
    """Text of a single-file WordprocessingML document (flat OPC / Word 2003
    XML) that carries a .docx name. Returns "" for anything that doesn't sniff
    as Word XML — the caller then falls through to the usual empty-extraction
    handling."""
    with open(path, "rb") as fh:
        head = fh.read(4096)
        if not (head.lstrip()[:5] == b"<?xml" or head.lstrip()[:1] == b"<"):
            return ""
        data = head + fh.read()
    if not any(m in data for m in _WORDML_MARKERS):
        return ""
    xml = data.decode("utf-8", "ignore")
    out: list[str] = []
    total = 0
    for m in _OOXML_TEXT_RE.findall(xml):
        txt = _unescape(m)
        out.append(txt)
        total += len(txt)
        if total >= max_chars:
            break
    return " ".join(out)[:max_chars]


def _pptx(path: str, max_chars: int) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return _ooxml_zip_fallback(path, "ppt/slides/", max_chars)

    prs = Presentation(path)
    out: list[str] = []
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = shape.text_frame.text
                if t:
                    out.append(t)
                    total += len(t)
        if total >= max_chars:
            break
    return _strip_residual_ooxml("\n".join(out))[:max_chars]


def _ooxml_zip_fallback(path: str, prefix: str, max_chars: int) -> str:
    """python-pptx unavailable — scrape <a:t> runs straight from the slide parts."""
    with zipfile.ZipFile(path) as z:
        out: list[str] = []
        total = 0
        for n in sorted(z.namelist()):
            if not (n.startswith(prefix) and n.endswith(".xml")):
                continue
            xml = z.read(n).decode("utf-8", "ignore")
            for m in _OOXML_TEXT_RE.findall(xml):
                txt = _unescape(m)
                out.append(txt)
                total += len(txt)
            if total >= max_chars:
                break
    return _strip_residual_ooxml("\n".join(out))[:max_chars]
