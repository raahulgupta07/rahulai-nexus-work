"""Standalone, deterministic file-classification module.

Given an uploaded file, decide which ONE destination store it belongs in.
This module is PURE: stdlib + pandas + openpyxl + (optional) python-docx / PyMuPDF.
It has NO app or DB dependencies, opens NO network connections, and NEVER raises
from ``classify_file`` (worst case it returns a low-confidence "knowledge" fallback).

INSTRUCTIONS
============
Destinations (exactly one is chosen per file):

  "table"        A queryable dataset. Clean rectangular tabular data
                 (CSV / XLSX grids with >=2 typed columns and many rows).
                 -> ingested as a data-source table the agent can query.

  "instruction"  Definitions, glossaries, term->meaning maps, rules, policies,
                 "X means ...", "Always ...", "Never ...", "must ..." statements.
                 -> loaded as agent instructions / grounding definitions.

  "skill"        A repeatable procedure: SOP / runbook / playbook / how-to /
                 numbered or bulleted step lists with imperative verbs.
                 -> registered as a skill / procedure the agent can follow.

  "knowledge"    Reference prose, background docs, images, and anything that is
                 informative but not tabular, not a rule set, not a procedure.
                 -> stored in the knowledge base for retrieval.

Classification is layered and deterministic (no LLM here):

  L0  file type       csv/xlsx -> tabular-candidate; docx/pdf/pptx -> text-candidate;
                      images -> knowledge.
  L1  filename hints   glossary/definition/dictionary/meaning -> instruction;
                      sop/runbook/playbook/procedure/how-to/steps/workflow -> skill;
                      q&a/logic/rules/policy -> instruction/knowledge.
  L2  content probe    open the file and inspect structure/shape:
                        - "tabular"       -> table
                        - "term_meaning"  -> instruction
                        - "step_list"     -> skill
                        - "prose"         -> knowledge
                        - "unknown"       -> fall back to L1/L0

Confidence:
  > 0.85  L2 is unambiguous (high).
  0.6-0.85  medium.
  < 0.6   ``needs_llm`` is set True so a caller may run an optional L3 (LLM) pass.
"""

from __future__ import annotations

import os
import re
import csv as _csv
import io

# ---- optional dependencies (guarded) ---------------------------------------
try:
    import pandas as pd  # type: ignore
    _HAS_PANDAS = True
except Exception:  # pragma: no cover
    _HAS_PANDAS = False

try:
    import openpyxl  # type: ignore
    _HAS_OPENPYXL = True
except Exception:  # pragma: no cover
    _HAS_OPENPYXL = False

try:
    import docx  # python-docx  # type: ignore
    _HAS_DOCX = True
except Exception:  # pragma: no cover
    _HAS_DOCX = False

try:
    import fitz  # PyMuPDF  # type: ignore
    _HAS_FITZ = True
except Exception:  # pragma: no cover
    _HAS_FITZ = False

try:
    from pptx import Presentation  # python-pptx  # type: ignore
    _HAS_PPTX = True
except Exception:  # pragma: no cover
    _HAS_PPTX = False


# ---- extension groups ------------------------------------------------------
_TABULAR_EXT = {".csv", ".tsv", ".xlsx", ".xls", ".xlsm"}
_TEXT_EXT = {".docx", ".doc", ".pdf", ".pptx", ".ppt", ".txt", ".md", ".rtf"}
_IMAGE_EXT = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp", ".svg", ".heic"}

# ---- filename hint patterns ------------------------------------------------
_FN_INSTRUCTION = re.compile(
    r"(definition|glossary|dictionar|meaning|terms?|lexicon|taxonom)", re.I
)
_FN_SKILL = re.compile(
    r"(sop|runbook|playbook|procedure|how[\-_ ]?to|steps?|workflow|guide|checklist|instructions?)",
    re.I,
)
_FN_RULES = re.compile(r"(q\s*&\s*a|q_?and_?a|faq|logic|rules?|policy|policies)", re.I)

# ---- content signal patterns -----------------------------------------------
# imperative / step verbs at start of a line
_IMPERATIVE = re.compile(
    r"^\s*(do|run|check|then|first|next|finally|open|click|select|create|configure|"
    r"install|set|enter|verify|ensure|start|stop|add|remove|navigate|press|type|"
    r"execute|apply|save|review|confirm|repeat|go to)\b",
    re.I,
)
# an ordered/bulleted list marker at the start of a line
_LIST_MARKER = re.compile(r"^\s*(?:\d{1,3}[.)]\s+|step\s*\d+\b|[-*•‣▪●]\s+)", re.I)
# rule / definition statements
_RULE_STMT = re.compile(
    r"(\bmeans\b|\bis defined as\b|\brefers to\b|\balways\b|\bnever\b|\bmust\b|"
    r"\bshould\b|\bshall\b|\bdo not\b|\bdon't\b|\bmay not\b|\bis required\b|\bstands for\b)",
    re.I,
)

_TEXT_SAMPLE_CHARS = 8000
_MAX_ROWS_PROBE = 200


# ============================================================================
# Public API
# ============================================================================
def classify_file(path: str, filename: str = "", content_type: str = "") -> dict:
    """Classify an uploaded file into ONE destination store.

    Returns:
      {
        "destination": "table" | "instruction" | "skill" | "knowledge",
        "confidence": float,       # 0..1
        "reason": str,
        "signals": {"l0_type": str, "l1_filename": str|None, "l2_shape": str},
        "needs_llm": bool,         # True when confidence < 0.6
      }
    """
    signals = {"l0_type": "unknown", "l1_filename": None, "l2_shape": "unknown"}
    try:
        name = filename or (os.path.basename(path) if path else "")
        ext = _ext_of(name, path, content_type)

        # ---------------- L0: coarse type ----------------
        l0 = _l0_type(ext, content_type)
        signals["l0_type"] = l0

        # images short-circuit to knowledge
        if l0 == "image":
            return _result("knowledge", 0.8, "image file -> knowledge base", signals)

        # ---------------- L1: filename hints ----------------
        l1 = _l1_filename(name)
        signals["l1_filename"] = l1

        # ---------------- L2: content probe ----------------
        if l0 == "tabular":
            shape, dest, conf, reason = _probe_tabular(path, ext, l1)
        elif l0 == "text":
            shape, dest, conf, reason = _probe_text(path, ext, l1)
        else:
            shape, dest, conf, reason = "unknown", None, 0.0, "unknown type"
        signals["l2_shape"] = shape

        # ---------------- resolve when L2 was inconclusive ----------------
        if dest is None:
            dest, conf, reason = _fallback_from_hints(l0, l1)

        return _result(dest, conf, reason, signals)
    except Exception as e:  # never raise
        return {
            "destination": "knowledge",
            "confidence": 0.3,
            "reason": "fallback: %s" % (type(e).__name__,),
            "signals": signals,
            "needs_llm": True,
        }


# ============================================================================
# L3: LLM "Librarian" — decide the destination from CONTENT, not filename
# ============================================================================
_LLM_DESTINATIONS = ("table", "instruction", "skill", "knowledge")


def classify_file_llm(content: str, filename: str = "", content_type: str = "",
                      infer=None) -> dict | None:
    """Read a file's extracted content and let an LLM pick the destination store.

    This is the content-driven counterpart to :func:`classify_file`. The caller
    passes ``content`` (already-extracted sheet preview OR document text) plus an
    ``infer(prompt: str) -> str`` callable (a real LLM in production, a mock in
    tests). The model is told to IGNORE the filename and decide purely from what
    the content actually is.

    Returns ``{"destination", "confidence", "reason"}`` on success, or ``None``
    on any failure / unparseable reply so the caller can fall back to the
    deterministic classifier. Pure: imports nothing from the app, opens no
    network by itself, and never raises.
    """
    if infer is None or not content or not str(content).strip():
        return None

    sample = str(content)[:6000]
    prompt = (
        "You are a data librarian. Read the file CONTENT below and decide which "
        "ONE store it belongs in. Decide from the CONTENT ONLY — ignore the "
        "filename; a future file will not be conveniently named.\n\n"
        "Destinations:\n"
        "- table: a queryable dataset — clean rectangular rows/columns of data "
        "(sales, transactions, a records export) meant to be queried/aggregated.\n"
        "- instruction: a CONCENTRATED set of definitions, a glossary, "
        "term->meaning maps, metric formulas, or explicit business rules / "
        "policies ('X means ...', 'always ...', 'never ...', 'must ...') meant to "
        "directly constrain how the agent computes or answers.\n"
        "- skill: a repeatable multi-step procedure / SOP / runbook / playbook — "
        "an ordered how-to the agent can follow.\n"
        "- knowledge: background/reference material that is mostly explanatory "
        "PROSE — narratives, explanations, overviews, tutorials, 'how it works' "
        "or 'detailed explanation' write-ups, and Q&A / FAQ reference documents. "
        "Informative context to retrieve, not a dataset, rule set, or procedure.\n\n"
        "Tie-breakers (decide from the actual content):\n"
        "- If the document is mostly explanatory prose or narrative — even if it "
        "mentions a few rules in passing — choose knowledge, NOT instruction.\n"
        "- Choose instruction ONLY when it is a dense, list-like set of "
        "definitions/formulas/rules with little surrounding prose.\n"
        "- A Q&A / FAQ explaining a topic is knowledge; a compact glossary of "
        "term=meaning pairs is instruction.\n\n"
        "Return STRICT JSON and NOTHING else:\n"
        '{"destination": "table|instruction|skill|knowledge", '
        '"confidence": 0.0, "reason": "one short sentence"}\n\n'
        f"FILENAME (context only, do NOT rely on it): {filename or '(unknown)'}\n"
        "CONTENT:\n"
        f"{sample}\n"
    )

    try:
        raw = infer(prompt)
    except Exception:
        return None

    parsed = _llm_extract_json(raw)
    if not isinstance(parsed, dict):
        return None
    dest = str(parsed.get("destination") or "").strip().lower()
    if dest not in _LLM_DESTINATIONS:
        return None
    try:
        conf = float(parsed.get("confidence"))
    except (TypeError, ValueError):
        conf = 0.7
    conf = max(0.0, min(1.0, conf))
    reason = str(parsed.get("reason") or "").strip() or "LLM content classification"
    return {"destination": dest, "confidence": round(conf, 2), "reason": reason}


def _llm_extract_json(text) -> object | None:
    """Fence-tolerant, brace-matched JSON extraction (local, app-free)."""
    if not text:
        return None
    import json as _json
    t = str(text).strip()
    if t.startswith("```"):
        nl = t.find("\n")
        t = t[nl + 1:] if nl != -1 else t[3:]
        rf = t.rfind("```")
        if rf != -1:
            t = t[:rf]
        t = t.strip()
    try:
        return _json.loads(t)
    except Exception:
        pass
    start = t.find("{")
    while start != -1:
        depth = 0
        in_str = False
        esc = False
        for i in range(start, len(t)):
            ch = t[i]
            if in_str:
                if esc:
                    esc = False
                elif ch == "\\":
                    esc = True
                elif ch == '"':
                    in_str = False
                continue
            if ch == '"':
                in_str = True
            elif ch == "{":
                depth += 1
            elif ch == "}":
                depth -= 1
                if depth == 0:
                    try:
                        return _json.loads(t[start:i + 1])
                    except Exception:
                        break
        start = t.find("{", start + 1)
    return None


# ============================================================================
# L0 / L1 helpers
# ============================================================================
def _ext_of(name: str, path: str, content_type: str) -> str:
    ext = os.path.splitext(name or "")[1].lower()
    if not ext and path:
        ext = os.path.splitext(path)[1].lower()
    if not ext and content_type:
        ct = content_type.lower()
        if "csv" in ct:
            ext = ".csv"
        elif "spreadsheet" in ct or "excel" in ct:
            ext = ".xlsx"
        elif "wordprocessing" in ct or "msword" in ct:
            ext = ".docx"
        elif "pdf" in ct:
            ext = ".pdf"
        elif "presentation" in ct or "powerpoint" in ct:
            ext = ".pptx"
        elif ct.startswith("image/"):
            ext = ".png"
        elif "plain" in ct or "text" in ct:
            ext = ".txt"
    return ext


def _l0_type(ext: str, content_type: str) -> str:
    if ext in _TABULAR_EXT:
        return "tabular"
    if ext in _IMAGE_EXT:
        return "image"
    if ext in _TEXT_EXT:
        return "text"
    ct = (content_type or "").lower()
    if ct.startswith("image/"):
        return "image"
    if "csv" in ct or "spreadsheet" in ct or "excel" in ct:
        return "tabular"
    if any(k in ct for k in ("pdf", "word", "presentation", "text", "powerpoint")):
        return "text"
    return "unknown"


def _l1_filename(name: str):
    if not name:
        return None
    if _FN_INSTRUCTION.search(name):
        return "instruction"
    if _FN_SKILL.search(name):
        return "skill"
    if _FN_RULES.search(name):
        return "instruction"  # rules/policy/faq lean instruction, knowledge as tiebreak
    return None


def _fallback_from_hints(l0: str, l1):
    """When L2 could not decide, lean on filename hint / coarse type. Low conf."""
    if l1 == "instruction":
        return "instruction", 0.55, "filename hint (definition/glossary/rules)"
    if l1 == "skill":
        return "skill", 0.55, "filename hint (sop/runbook/procedure)"
    if l0 == "tabular":
        return "table", 0.55, "tabular type but shape unclear"
    return "knowledge", 0.4, "no strong signal -> knowledge"


# ============================================================================
# L2: tabular probe
# ============================================================================
def _probe_tabular(path: str, ext: str, l1):
    """Return (shape, destination|None, confidence, reason)."""
    df = _read_tabular_sample(path, ext)
    if df is None or df.shape[0] == 0 or df.shape[1] == 0:
        # could not read as a grid; treat as unknown, defer
        return "unknown", None, 0.0, "could not read tabular content"

    ncols = df.shape[1]
    nrows = df.shape[0]

    # count columns that actually carry data
    populated = [c for c in df.columns if df[c].notna().sum() > 0]
    npop = len(populated)

    # -- two-column term -> meaning check (glossary/definition) --
    if npop == 2:
        if _looks_term_meaning(df, populated):
            conf = 0.88 if l1 in (None, "instruction") else 0.82
            return ("term_meaning", "instruction", conf,
                    "two columns look like term -> meaning (glossary)")

    # -- clean rectangular typed grid -> table --
    if npop >= 2:
        typed_ratio = _typed_ratio(df, populated)
        rect = _is_rectangular(df, populated)
        if rect and typed_ratio >= 0.6 and nrows >= 3:
            conf = 0.9 if typed_ratio >= 0.75 else 0.8
            return ("tabular", "table", conf,
                    "clean rectangular grid, %d cols, typed ratio %.2f" % (npop, typed_ratio))
        # 2+ cols but messy -> still probably a table, medium
        return ("tabular", "table", 0.68,
                "multi-column grid (%d cols) but mixed typing" % npop)

    # single populated column -> a list of things, likely prose/steps not a table
    if npop == 1:
        col = df[populated[0]].dropna().astype(str).tolist()
        joined = "\n".join(col[:_MAX_ROWS_PROBE])
        shape2 = _text_shape(joined)
        if shape2 == "step_list":
            return ("step_list", "skill", 0.72, "single column of step-like rows")
        if shape2 == "rule":
            return ("term_meaning", "instruction", 0.66, "single column of rule/definition rows")
        return ("prose", "knowledge", 0.55, "single-column list -> knowledge")

    return "unknown", None, 0.0, "empty grid"


def _read_tabular_sample(path: str, ext: str):
    if not (path and os.path.exists(path)):
        return None
    try:
        if ext in (".csv", ".tsv", ".txt"):
            sep = "\t" if ext == ".tsv" else None
            if _HAS_PANDAS:
                return pd.read_csv(path, sep=sep, engine="python", nrows=_MAX_ROWS_PROBE,
                                   header=0, on_bad_lines="skip", dtype=str,
                                   keep_default_na=True)
            return _read_csv_stdlib(path, sep or ",")
        # excel family
        if _HAS_PANDAS:
            return pd.read_excel(path, nrows=_MAX_ROWS_PROBE, header=0, dtype=str)
        if _HAS_OPENPYXL:
            return _read_xlsx_openpyxl(path)
    except Exception:
        # last-ditch: try csv reader for text-ish files
        try:
            if ext in (".csv", ".tsv", ".txt"):
                return _read_csv_stdlib(path, "\t" if ext == ".tsv" else ",")
        except Exception:
            return None
    return None


def _read_csv_stdlib(path: str, sep: str):
    if not _HAS_PANDAS:
        rows = []
        with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
            reader = _csv.reader(f, delimiter=sep)
            for i, row in enumerate(reader):
                if i > _MAX_ROWS_PROBE:
                    break
                rows.append(row)
        # crude DataFrame-less shim is not supported downstream; require pandas
        return None
    return None


def _read_xlsx_openpyxl(path: str):
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    ws = wb[wb.sheetnames[0]]
    data = []
    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i > _MAX_ROWS_PROBE:
            break
        data.append(list(row))
    wb.close()
    if not data:
        return None
    header = [("" if v is None else str(v)) for v in data[0]]
    body = data[1:]
    if _HAS_PANDAS:
        return pd.DataFrame(body, columns=header if len(set(header)) == len(header) else None)
    return None


def _looks_term_meaning(df, cols) -> bool:
    """Left column = short distinct labels, right column = longer explanatory text."""
    try:
        left = df[cols[0]].dropna().astype(str)
        right = df[cols[1]].dropna().astype(str)
        if len(left) < 2 or len(right) < 2:
            return False
        left_len = left.str.len().mean()
        right_len = right.str.len().mean()
        left_words = left.str.split().str.len().mean()
        distinct_ratio = left.nunique() / max(1, len(left))
        # left short & mostly unique labels, right substantially longer
        return (
            left_len <= 40
            and left_words <= 5
            and right_len >= left_len * 1.4
            and right_len >= 12
            and distinct_ratio >= 0.8
        )
    except Exception:
        return False


def _typed_ratio(df, cols) -> float:
    """Fraction of populated cells that parse as numbers or dates (i.e. 'typed' data)."""
    total = 0
    typed = 0
    for c in cols:
        s = df[c].dropna().astype(str)
        for v in s.head(60):
            total += 1
            if _looks_typed(v):
                typed += 1
    return (typed / total) if total else 0.0


def _looks_typed(v: str) -> bool:
    v = v.strip()
    if not v:
        return False
    # number (with optional thousands sep / currency / percent)
    if re.fullmatch(r"[-+$€£]?\s?\d{1,3}(?:[,\s]\d{3})*(?:\.\d+)?%?", v):
        return True
    if re.fullmatch(r"[-+]?\d+(?:\.\d+)?", v):
        return True
    # date-ish
    if re.fullmatch(r"\d{4}[-/]\d{1,2}[-/]\d{1,2}", v):
        return True
    if re.fullmatch(r"\d{1,2}[-/]\d{1,2}[-/]\d{2,4}", v):
        return True
    # boolean-ish / short code
    if v.lower() in ("true", "false", "yes", "no", "y", "n"):
        return True
    return False


def _is_rectangular(df, cols) -> bool:
    """Most rows fill most columns (few ragged/empty cells)."""
    try:
        sub = df[cols]
        filled = sub.notna().sum().sum()
        total = sub.shape[0] * sub.shape[1]
        return total > 0 and (filled / total) >= 0.7
    except Exception:
        return True


# ============================================================================
# L2: text probe
# ============================================================================
def _probe_text(path: str, ext: str, l1):
    """Return (shape, destination|None, confidence, reason)."""
    text = _extract_text_sample(path, ext)
    if not text or not text.strip():
        # no text available (missing libs / scanned pdf) -> lean on filename hint
        if l1 == "skill":
            return "unknown", "skill", 0.5, "no text extracted; filename hints skill"
        if l1 == "instruction":
            return "unknown", "instruction", 0.5, "no text extracted; filename hints instruction"
        return "unknown", None, 0.0, "no text extracted"

    shape = _text_shape(text)
    if shape == "step_list":
        conf = 0.9 if l1 in (None, "skill") else 0.82
        return "step_list", "skill", conf, "numbered/bulleted imperative step list"
    if shape == "rule":
        conf = 0.88 if l1 in (None, "instruction") else 0.8
        return "term_meaning", "instruction", conf, "definition/rule statements (means/always/never/must)"
    # prose
    # nudge with filename hint but keep it knowledge-ish
    if l1 == "skill":
        return "prose", "skill", 0.6, "prose but filename hints procedure"
    if l1 == "instruction":
        return "prose", "instruction", 0.6, "prose but filename hints definitions"
    return "prose", "knowledge", 0.75, "reference prose -> knowledge base"


def _text_shape(text: str) -> str:
    """Classify a text sample into step_list / rule / prose."""
    lines = [ln for ln in text.splitlines() if ln.strip()]
    if not lines:
        return "prose"
    sample = lines[:120]
    n = len(sample)

    # strip any leading list/step marker before checking for an imperative verb,
    # so "1. Open the dashboard" counts as imperative ("Open"), not as non-verb.
    stripped = [_LIST_MARKER.sub("", ln, count=1) for ln in sample]

    list_lines = sum(1 for ln in sample if _LIST_MARKER.search(ln))
    imperative_lines = sum(1 for ln in stripped if _IMPERATIVE.search(ln))
    rule_hits = sum(1 for ln in sample if _RULE_STMT.search(ln))

    list_ratio = list_lines / n
    imp_ratio = imperative_lines / n
    rule_ratio = rule_hits / n

    # step list: many list markers AND some imperative verbs (a procedure)
    if list_ratio >= 0.4 and (imperative_lines >= 2 or imp_ratio >= 0.2):
        return "step_list"
    # ordered numbered steps even without strong verbs
    if list_lines >= 3 and imp_ratio >= 0.15:
        return "step_list"
    # rules / definitions dominate
    if rule_ratio >= 0.3:
        return "rule"
    # standalone: lots of imperative lines, few of them prose -> steps
    if imp_ratio >= 0.35 and list_ratio < 0.4:
        return "step_list"
    return "prose"


def _extract_text_sample(path: str, ext: str) -> str:
    if not (path and os.path.exists(path)):
        return ""
    try:
        if ext in (".txt", ".md", ".rtf"):
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read(_TEXT_SAMPLE_CHARS)
        if ext == ".docx" and _HAS_DOCX:
            d = docx.Document(path)
            parts = []
            for p in d.paragraphs:
                if p.text and p.text.strip():
                    # preserve list style so markers survive
                    parts.append(p.text)
                if sum(len(x) for x in parts) > _TEXT_SAMPLE_CHARS:
                    break
            return "\n".join(parts)[:_TEXT_SAMPLE_CHARS]
        if ext == ".pdf" and _HAS_FITZ:
            doc = fitz.open(path)
            buf = []
            for page in doc:
                buf.append(page.get_text())
                if sum(len(x) for x in buf) > _TEXT_SAMPLE_CHARS:
                    break
            doc.close()
            return "".join(buf)[:_TEXT_SAMPLE_CHARS]
        if ext in (".pptx",) and _HAS_PPTX:
            prs = Presentation(path)
            buf = []
            for slide in prs.slides:
                for shp in slide.shapes:
                    if shp.has_text_frame and shp.text_frame.text.strip():
                        buf.append(shp.text_frame.text)
                if sum(len(x) for x in buf) > _TEXT_SAMPLE_CHARS:
                    break
            return "\n".join(buf)[:_TEXT_SAMPLE_CHARS]
        # unsupported / lib missing: try raw bytes for text-like content
        with open(path, "rb") as f:
            raw = f.read(4096)
        try:
            return raw.decode("utf-8", errors="ignore")
        except Exception:
            return ""
    except Exception:
        return ""


# ============================================================================
# result assembly
# ============================================================================
def _result(dest: str, conf: float, reason: str, signals: dict) -> dict:
    conf = max(0.0, min(1.0, float(conf)))
    return {
        "destination": dest,
        "confidence": round(conf, 2),
        "reason": reason,
        "signals": dict(signals),
        "needs_llm": conf < 0.6,
    }


# ============================================================================
# self-test
# ============================================================================
if __name__ == "__main__":
    import tempfile
    import json as _json

    tmp = tempfile.mkdtemp(prefix="fclf_")

    def _p(nm):
        return os.path.join(tmp, nm)

    cases = []

    # 1) clean tabular CSV -> table
    sales = _p("sales.csv")
    with open(sales, "w") as f:
        f.write("date,region,units,revenue\n")
        for i in range(20):
            f.write("2026-01-%02d,North,%d,%d\n" % ((i % 28) + 1, 10 + i, 100 * (i + 1)))
    cases.append((sales, "sales.csv", "expect table"))

    # 2) definitions glossary xlsx (term -> meaning) -> instruction
    defs = _p("definitions.xlsx")
    if _HAS_OPENPYXL:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["Term", "Meaning"])
        ws.append(["WMAPE", "Weighted mean absolute percentage error used to score demand forecasts across SKUs."])
        ws.append(["SKU", "Stock keeping unit; a unique identifier for a distinct sellable product."])
        ws.append(["Lead time", "The elapsed time between placing a replenishment order and receiving the goods."])
        ws.append(["Safety stock", "Buffer inventory held to protect against demand and supply variability."])
        wb.save(defs)
        cases.append((defs, "definitions.xlsx", "expect instruction"))

    # 3) SOP runbook docx (step list) -> skill
    sop = _p("deploy_runbook.docx")
    if _HAS_DOCX:
        d = docx.Document()
        d.add_paragraph("Deployment Runbook")
        d.add_paragraph("1. Open the CI dashboard and confirm the build is green.")
        d.add_paragraph("2. Run the migration script against staging.")
        d.add_paragraph("3. Check the health endpoint returns 200.")
        d.add_paragraph("4. Then promote the image to production.")
        d.add_paragraph("5. Verify logs show no errors and confirm the rollout.")
        d.save(sop)
        cases.append((sop, "deploy_runbook.docx", "expect skill"))

    # 4) nutrition charter prose docx -> knowledge
    charter = _p("nutrition_charter.docx")
    if _HAS_DOCX:
        d = docx.Document()
        d.add_paragraph("Nutrition Charter")
        d.add_paragraph(
            "Our organization is committed to promoting balanced nutrition and wellbeing "
            "across all facilities. This charter describes the principles that guide our "
            "food service partners and the values we share with the communities we serve."
        )
        d.add_paragraph(
            "We believe that access to wholesome food is fundamental. Over the coming years "
            "we will continue to invest in sourcing, education, and transparency so that every "
            "meal reflects our shared aspirations for health and sustainability."
        )
        d.save(charter)
        cases.append((charter, "nutrition_charter.docx", "expect knowledge (or skill)"))

    # 5) rules/policy txt -> instruction
    rules = _p("access_policy.txt")
    with open(rules, "w") as f:
        f.write("Access Policy\n")
        f.write("Users must authenticate before accessing any dataset.\n")
        f.write("Admins should never share credentials over email.\n")
        f.write("A 'viewer' means a role that can read but not modify data.\n")
        f.write("Access tokens must be rotated every 90 days.\n")
    cases.append((rules, "access_policy.txt", "expect instruction"))

    # 6) image -> knowledge
    img = _p("chart.png")
    with open(img, "wb") as f:
        f.write(b"\x89PNG\r\n\x1a\n")
    cases.append((img, "chart.png", "expect knowledge"))

    print("=== synthetic self-test ===")
    for pth, nm, note in cases:
        r = classify_file(pth, nm)
        print("%-26s -> %-11s conf=%.2f llm=%-5s | %s  (%s)" % (
            nm, r["destination"], r["confidence"], r["needs_llm"], r["reason"], note))
        print("   signals:", _json.dumps(r["signals"]))

    # === real samples, if present ===
    real_dirs = [
        "/private/tmp/claude-501/-Users-rahulgupta/f5ca543e-912c-4716-94dd-90fbefc22ce9/scratchpad/xlsx-parser/samples",
        os.path.expanduser("~/Desktop/excel-parser-test"),
    ]
    printed_header = False
    for d in real_dirs:
        if not os.path.isdir(d):
            continue
        files = [os.path.join(d, x) for x in sorted(os.listdir(d))
                 if os.path.splitext(x)[1].lower() in (_TABULAR_EXT | _TEXT_EXT | _IMAGE_EXT)]
        if not files:
            continue
        if not printed_header:
            print("\n=== real samples ===")
            printed_header = True
        for fp in files:
            r = classify_file(fp, os.path.basename(fp))
            print("%-40s -> %-11s conf=%.2f llm=%-5s | %s" % (
                os.path.basename(fp)[:40], r["destination"], r["confidence"],
                r["needs_llm"], r["reason"]))
