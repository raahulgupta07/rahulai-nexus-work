"""Draw a deck theme's structural motifs onto a deck that is already saved.

``theme_layouts`` declares what a theme's furniture *is*; this module is what
draws it. It reopens a finished ``.pptx``, adds the ruling, margin rules,
masthead, tracker, footer, stamp and ornaments the theme asks for, and saves.

Contract, copied deliberately from ``normalize_chart_axes`` in
``app/ai/code_execution/pptx_executor.py``: open the deck once, mutate it,
save only if something actually changed, and **never raise**. A deck that
could not be decorated is still a valid deck; a deck that failed to export
because a decoration threw is not.

★ Z-order is the whole difficulty.

python-pptx appends every new shape to the end of the slide's shape tree,
which is the *front*. Ruling drawn that way covers the slide's own text, and a
deck whose text is hidden behind a background is far worse than a deck with no
motif at all. python-pptx has no send-to-back API, so this module moves the
new element in the XML itself: remove it from ``spTree`` and reinsert it ahead
of every existing drawable child. ``_Backdrop`` does that, and holds a cursor
so that several background shapes keep their own relative order (paper first,
then ruling, then the margin rule on top of the ruling) while all of them stay
behind everything the deck already had.

Not every motif belongs at the back. A footer, a page number, a corner mark or
a stamp is meant to read *over* the paper, and those are left where
python-pptx put them -- at the front. Which is which is decided per motif in
``_paint_slide``, never globally.

Geometry is native vector throughout -- rectangles, ovals, freeform polylines,
gradient fills, rotation. Pillow is used for exactly one motif, the paper
grain, because a per-pixel texture is the one thing shapes cannot express; it
is placed full-bleed at the very back and is skipped silently when Pillow is
not importable.
"""

from __future__ import annotations

import io
import math
import random
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from .theme_layouts import CANVAS_H, CANVAS_W, EMPTY_LAYOUT

__all__ = ["paint_theme_furniture", "convert_boxes_to_rules", "ORNAMENTS"]


#: Ornament name -> painter. A layout naming something absent here is ignored,
#: not an error: a theme file may describe a motif nobody has drawn yet.
ORNAMENTS = (
    "double_hairline_border",
    "monogram_seal",
    "punch_holes",
    "star_field",
    "orbit_arcs",
    "title_line_motif",
    "status_dot",
    "paper_grain",
)


# =============================================================================
# Small helpers
# =============================================================================

def _rgb(value: Any, default: str = "000000") -> RGBColor:
    """`'#1F3A5F'` / `'1f3a5f'` -> RGBColor. Anything unusable -> `default`."""
    text = value if isinstance(value, str) else ""
    text = text.strip().lstrip("#")
    if len(text) == 3:
        text = "".join(c * 2 for c in text)
    if len(text) != 6:
        text = default
    try:
        return RGBColor.from_string(text.upper())
    except Exception:
        return RGBColor.from_string(default)


def _is_hex(value: Any) -> bool:
    text = value if isinstance(value, str) else ""
    text = text.strip().lstrip("#")
    if len(text) == 3:
        text = "".join(c * 2 for c in text)
    if len(text) != 6:
        return False
    try:
        int(text, 16)
    except ValueError:
        return False
    return True


class _Scale:
    """Maps the prompts' 1280x720 pixel canvas onto the real slide."""

    def __init__(self, slide_width: int, slide_height: int) -> None:
        self.width = int(slide_width)
        self.height = int(slide_height)
        self._x = self.width / float(CANVAS_W)
        self._y = self.height / float(CANVAS_H)

    def x(self, px: float) -> int:
        return int(round(px * self._x))

    def y(self, px: float) -> int:
        return int(round(px * self._y))

    def pt(self, px: float) -> Pt:
        """A CSS pixel is 0.75pt at 96dpi, which is exactly this canvas."""
        return Pt(max(px * 0.75, 1))


def _first_shape_index(sp_tree) -> int:
    """Index of the first *drawable* child of a shape tree.

    A ``spTree`` opens with ``nvGrpSpPr`` and ``grpSpPr``, which are not
    shapes. Hardcoding 2 works on every file python-pptx writes and this does
    not depend on that being true.
    """
    for index, child in enumerate(sp_tree):
        tag = child.tag.rsplit("}", 1)[-1]
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            return index
    return len(sp_tree)


class _Backdrop:
    """Moves shapes behind everything the slide already had, in order.

    ``push`` is called immediately after a shape is created. The first pushed
    shape lands ahead of all pre-existing content; the second lands directly
    after the first, and so on -- so background layers keep the order they
    were painted in while none of them can cover the deck's own text.
    """

    def __init__(self, slide) -> None:
        self._tree = slide.shapes._spTree
        self._cursor = _first_shape_index(self._tree)

    def push(self, shape) -> None:
        element = shape._element
        self._tree.remove(element)
        self._tree.insert(self._cursor, element)
        self._cursor += 1


def _no_line(shape) -> None:
    try:
        shape.line.fill.background()
    except Exception:
        pass


def _no_fill(shape) -> None:
    try:
        shape.fill.background()
    except Exception:
        pass


def _no_shadow(shape) -> None:
    """Every theme here that mentions shadows forbids them; none asks for one."""
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _rect(shapes, x, y, cx, cy, color) -> Any:
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(max(cx, 1)), Emu(max(cy, 1)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    _no_line(shape)
    _no_shadow(shape)
    return shape


def _circle(shapes, x, y, diameter) -> Any:
    """A round circle, in EMU on both axes.

    Scaling a diameter through ``_Scale.x`` and ``_Scale.y`` separately gives
    an ellipse on any deck that is not 16:9 -- a punch hole that is an oval on
    a 4:3 deck reads as a mistake, so one axis sets both.
    """
    size = Emu(max(int(diameter), 1))
    shape = shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(x)), Emu(int(y)), size, size)
    _no_shadow(shape)
    return shape


def _textbox(shapes, x, y, cx, cy, text, color, size, align, font=None) -> Any:
    box = shapes.add_textbox(Emu(x), Emu(y), Emu(max(cx, 1)), Emu(max(cy, 1)))
    frame = box.text_frame
    frame.word_wrap = False
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    if font:
        run.font.name = font
    return box


def _theme_font(theme: Any) -> Optional[str]:
    """A body font name out of whatever shape the theme dict arrived in."""
    if theme is None:
        return None
    for key in ("body_font", "font_body", "font", "fonts", "body"):
        value = theme.get(key) if isinstance(theme, dict) else getattr(theme, key, None)
        if isinstance(value, str) and value.strip():
            return value.strip()
        if isinstance(value, (list, tuple)) and value:
            first = value[-1] if len(value) > 1 else value[0]
            if isinstance(first, str) and first.strip():
                return first.strip()
    return None


def _pick(layout: dict, theme: Any, key: str, role: str, default: str) -> str:
    """A colour from the layout, else the theme's palette role, else default."""
    value = layout.get(key)
    if _is_hex(value):
        return value  # type: ignore[return-value]
    palette = None
    if isinstance(theme, dict):
        palette = theme.get("palette") if isinstance(theme.get("palette"), dict) else theme
    else:
        palette = getattr(theme, "palette", None)
    if isinstance(palette, dict) and _is_hex(palette.get(role)):
        return palette[role]
    return default


# =============================================================================
# Grounds
# =============================================================================

def _slide_sets_own_ground(slide) -> bool:
    """The model gave this slide its own background (`<p:bg>` in the slide XML).

    Painting the theme's ground sheet over it buries that background under an
    opaque rectangle while the text keeps the colors chosen for it -- a dark
    cover with white text becomes white-on-white at 1.00:1. The slide's own
    ground wins; the theme keeps its overlays (footer, tracker, chip) which sit
    above the text and carry their own colors.
    """
    return slide._element.cSld.find(qn("p:bg")) is not None


def _paint_ground(shapes, backdrop, scale, layout, theme, painted) -> None:
    ground = layout.get("ground", "flat")
    color = layout.get("ground_color")
    if ground == "gradient":
        second = layout.get("ground_color_2")
        shape = _rect(shapes, 0, 0, scale.width, scale.height, _rgb(color, "FFFFFF"))
        try:
            shape.fill.gradient()
            stops = shape.fill.gradient_stops
            stops[0].color.rgb = _rgb(color, "FFFFFF")
            stops[1].color.rgb = _rgb(second, "000000")
            # python-pptx writes linear gradients only; a radial glow is not
            # expressible, so the theme's "centre to corners" becomes the
            # nearest honest thing, a vertical fade.
            shape.fill.gradient_angle = float(layout.get("gradient_angle", 270))
        except Exception:
            pass
        backdrop.push(shape)
        painted.append("ground:gradient")
        return

    if not _is_hex(color):
        return
    shape = _rect(shapes, 0, 0, scale.width, scale.height, _rgb(color))
    backdrop.push(shape)
    painted.append("ground:flat" if ground != "ruled" else "ground:paper")


def _paint_ruling(shapes, backdrop, scale, layout, painted) -> None:
    spacing = layout.get("rule_spacing_px")
    if not isinstance(spacing, (int, float)) or spacing <= 0:
        return
    color = _rgb(layout.get("rule_color"), "DDDDDD")
    thickness = max(scale.y(float(layout.get("rule_width_pt", 0.75)) / 0.75), 1)
    start = float(layout.get("rule_start_px", spacing))
    drawn = 0
    y = start
    while y < CANVAS_H and drawn < 200:  # a ceiling; 720/26 is 27 lines
        backdrop.push(_rect(shapes, 0, scale.y(y), scale.width, thickness, color))
        drawn += 1
        y += spacing
    if drawn:
        painted.append(f"ruling:{drawn}")


def _paint_margin_rule(shapes, backdrop, scale, layout, painted) -> None:
    rule = layout.get("margin_rule")
    if not isinstance(rule, dict):
        return
    x_px = rule.get("x_px")
    if not isinstance(x_px, (int, float)):
        return
    color = _rgb(rule.get("color"), "C75146")
    width = max(scale.x(float(rule.get("width_pt", 1.5)) / 0.75), 1)
    backdrop.push(_rect(shapes, scale.x(x_px), 0, width, scale.height, color))
    painted.append("margin_rule")


def _paint_masthead(shapes, backdrop, scale, layout, theme, painted) -> None:
    if not layout.get("masthead"):
        return
    color = _rgb(_pick(layout, theme, "rule_color", "heading_text", "#1A1A1A"), "1A1A1A")
    # "a 3px black rule, then a row ..., closed by a thin rule".
    backdrop.push(_rect(shapes, scale.x(48), scale.y(26), scale.width - scale.x(96), scale.y(3), color))
    backdrop.push(_rect(shapes, scale.x(48), scale.y(74), scale.width - scale.x(96), max(scale.y(1), 1), color))
    painted.append("masthead")


# =============================================================================
# Overlay furniture -- these read over the paper and stay at the front
# =============================================================================

def _paint_tracker(shapes, scale, layout, index, count_slides, painted) -> None:
    tracker = layout.get("tracker")
    if not isinstance(tracker, dict):
        return
    kind = tracker.get("kind")
    if kind not in ("squares", "band"):
        return
    count = tracker.get("count")
    if not isinstance(count, int) or count <= 0:
        return
    on = _rgb(tracker.get("on"), "1F3A5F")
    off = _rgb(tracker.get("off"), "D5DCE4")

    # Which cell is lit: the slide's position mapped onto the tracker's cells,
    # so a 20-slide deck still moves a 6-square tracker from end to end.
    span = max(count_slides - 1, 1)
    current = min(int(round(index * (count - 1) / span)), count - 1) if count > 1 else 0

    if kind == "band":
        cell_w = scale.x(120 / count)
        x = scale.width - scale.x(56) - cell_w * count
        for i in range(count):
            shape = _rect(shapes, x + i * cell_w, scale.y(38), cell_w - scale.x(2), scale.y(6),
                          on if i == current else off)
            _no_shadow(shape)
        painted.append("tracker:band")
        return

    size = scale.x(8)
    gap = scale.x(6)
    total = count * size + (count - 1) * gap
    x = scale.width - scale.x(56) - total
    for i in range(count):
        shape = _rect(shapes, x + i * (size + gap), scale.y(38), size, size, on)
        if i != current:
            _no_fill(shape)
            try:
                shape.line.color.rgb = off
                shape.line.width = Pt(0.75)
            except Exception:
                pass
    painted.append("tracker:squares")


#: The footer band, on the prompts' 1280x720 canvas. The theme's own footer sits
#: at 660; 648 leaves a small margin above it so a model-drawn source line that
#: aims for "the bottom" and lands a few pixels high is still inside the band.
_FOOTER_BAND_TOP_PX = 648


def _strip_model_footer_text(shapes, scale, layout, painted) -> None:
    """Remove text the generated code drew in the band the footer is about to use.

    The deck's own code sometimes draws its own source line, footer or page
    number across the bottom of the slide. The theme then paints ITS footer at
    660 and the two overlap -- measured on a real deck as two source footers on
    top of each other and a stray page number beside the theme's "02 / 04". The
    prompt asks the model not to; this is the half that does not depend on the
    model listening.

    ★ Ordering is what keeps the theme's own furniture safe. This runs
    IMMEDIATELY BEFORE ``_paint_footer``, so nothing the furniture pass draws
    into this band exists yet and cannot be seen, let alone removed. Keep it
    there; moving it after the paint would make it eat the footer it protects.

    Only TEXT is removed. A picture, a chart, a table or any other graphic
    frame is left exactly where it is -- a chart whose legend dips into the
    bottom of the slide is not a footer, and deleting it costs the slide its
    content to fix a cosmetic overlap.

    Runs only when the layout will actually paint something here: a theme with
    no footer and no footer rule leaves the model's footer as the only one on
    the slide, which is correct and must be left alone.
    """
    footer = layout.get("footer")
    rule = layout.get("footer_rule")
    if not isinstance(footer, dict) and rule not in ("hairline", "double"):
        return

    band_top = scale.y(_FOOTER_BAND_TOP_PX)
    for shape in list(shapes):
        try:
            if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
                continue
            if getattr(shape, "shapes", None) is not None:
                continue  # a group: its children are not reachable as text here
            if not _has_real_text(shape):
                continue  # a picture, a graphic frame, a rule, an empty box
            top = shape.top
            if top is None or int(top) < band_top:
                continue
            element = shape._element
            parent = element.getparent()
            if parent is None:
                continue
            # python-pptx has no shape.delete(); this is the same spTree
            # surgery ``_Backdrop.push`` does, one step shorter.
            parent.remove(element)
            painted.append("stripped:footer_text")
        except Exception:
            continue  # never raise -- a strip that fails must not cost the deck


def _paint_footer(shapes, scale, layout, theme, index, count_slides, painted) -> None:
    footer = layout.get("footer")
    rule = layout.get("footer_rule")
    if not isinstance(footer, dict) and rule not in ("hairline", "double"):
        return

    color_hex = footer.get("color") if isinstance(footer, dict) else None
    color = _rgb(color_hex or _pick(layout, theme, "_none", "muted_text", "#777777"), "777777")
    y = scale.y(660)

    if rule == "hairline":
        _rect(shapes, scale.x(56), y, scale.width - scale.x(112), max(scale.y(1), 1), color)
        painted.append("footer_rule:hairline")
    elif rule == "double":
        # "two 1px charcoal lines 1px apart".
        _rect(shapes, scale.x(56), y, scale.width - scale.x(112), max(scale.y(1), 1), color)
        _rect(shapes, scale.x(56), y + max(scale.y(2), 2), scale.width - scale.x(112),
              max(scale.y(1), 1), color)
        painted.append("footer_rule:double")

    if not isinstance(footer, dict):
        return

    size_px = footer.get("size_px")
    size = scale.pt(size_px if isinstance(size_px, (int, float)) else 9)
    font = _theme_font(theme)
    fields = {"page": index + 1, "pages": count_slides}
    slots = (
        ("left", PP_ALIGN.LEFT, scale.x(56), scale.x(420)),
        ("center", PP_ALIGN.CENTER, (scale.width - scale.x(300)) // 2, scale.x(300)),
        ("right", PP_ALIGN.RIGHT, scale.width - scale.x(56) - scale.x(300), scale.x(300)),
    )
    wrote = False
    for key, align, x, width in slots:
        text = footer.get(key)
        if not isinstance(text, str) or not text.strip():
            continue
        try:
            text = text.format(**fields)
        except (KeyError, IndexError, ValueError):
            pass  # a literal brace in a footer is not an error
        _textbox(shapes, x, y + scale.y(6), width, scale.y(22), text, color, size, align, font)
        wrote = True
    if wrote:
        painted.append("footer")


def _paint_chip(shapes, scale, layout, theme, painted) -> None:
    chip = layout.get("chip")
    if not isinstance(chip, dict) or layout.get("forbid_boxes"):
        return
    text = chip.get("text")
    if not isinstance(text, str) or not text.strip():
        return
    width, height = scale.x(96), scale.y(24)
    x, y = scale.width - scale.x(56) - width, scale.y(36)
    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(chip.get("fill"), "E4EDF8")
    _no_line(shape)
    _no_shadow(shape)
    frame = shape.text_frame
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.size = scale.pt(11)
    run.font.color.rgb = _rgb(chip.get("color"), "1D4F91")
    font = _theme_font(theme)
    if font:
        run.font.name = font
    painted.append("chip")


def _paint_corner_mark(shapes, scale, layout, theme, painted) -> None:
    mark = layout.get("corner_mark")
    if not isinstance(mark, dict):
        return
    text = mark.get("text")
    if not isinstance(text, str) or not text.strip():
        return
    width = scale.x(320)
    _textbox(shapes, scale.width - scale.x(56) - width, scale.y(28), width, scale.y(18),
             text, _rgb(mark.get("color"), "66707D"), scale.pt(8), PP_ALIGN.RIGHT,
             _theme_font(theme))
    painted.append("corner_mark")


def _paint_stamp(shapes, scale, layout, painted) -> None:
    """One rotated outline stamp. Deck-level: the caller calls this once."""
    stamp = layout.get("stamp")
    if not isinstance(stamp, dict):
        return
    text = stamp.get("text")
    if not isinstance(text, str) or not text.strip():
        return
    color = _rgb(stamp.get("color"), "C75146")
    width, height = scale.x(240), scale.y(76)
    x, y = scale.width - scale.x(120) - width, scale.height - scale.y(190)
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(width), Emu(height))
    _no_fill(shape)
    _no_shadow(shape)
    try:
        shape.line.color.rgb = color
        shape.line.width = Pt(2.25)  # "3px border"
    except Exception:
        pass
    try:
        shape.rotation = float(stamp.get("rotation", 0))
    except (TypeError, ValueError):
        pass
    para = shape.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.size = scale.pt(24)
    run.font.bold = True
    run.font.color.rgb = color
    painted.append("stamp")


# =============================================================================
# Ornaments
# =============================================================================

def _orn_double_hairline_border(shapes, backdrop, scale, layout, theme, index, painted) -> None:
    """"one 1px gold line inset 24px and a second 1px line inset 30px"."""
    color = _rgb(_pick(layout, theme, "accent", "primary_accent", "#B49A5B"), "B49A5B")
    for inset in (24, 30):
        x, y = scale.x(inset), scale.y(inset)
        w, h = scale.width - 2 * x, scale.height - 2 * y
        t = max(scale.y(1), 1)
        for rect in ((x, y, w, t), (x, y + h - t, w, t), (x, y, t, h), (x + w - t, y, t, h)):
            backdrop.push(_rect(shapes, *rect, color))
    painted.append("double_hairline_border")


def _orn_monogram_seal(shapes, backdrop, scale, layout, theme, index, painted) -> None:
    """"two thin concentric gold circles about 96px wide", top centre."""
    color = _rgb(_pick(layout, theme, "accent", "primary_accent", "#B49A5B"), "B49A5B")
    for diameter in (96, 82):
        size = scale.x(diameter)
        shape = _circle(shapes, (scale.width - size) // 2, scale.y(56), size)
        _no_fill(shape)
        try:
            shape.line.color.rgb = color
            shape.line.width = Pt(0.75)
        except Exception:
            pass
        backdrop.push(shape)
    painted.append("monogram_seal")


def _orn_punch_holes(shapes, backdrop, scale, layout, theme, index, painted) -> None:
    """"three gray punch holes about 26px wide spaced down the left edge"."""
    size = scale.x(26)
    for y_px in (150, 350, 550):
        shape = _circle(shapes, scale.x(22), scale.y(y_px), size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb("#E5E7EB")
        try:
            shape.line.color.rgb = _rgb("#D1D5DB")
            shape.line.width = Pt(0.75)
        except Exception:
            pass
        backdrop.push(shape)
    painted.append("punch_holes")


def _orn_star_field(shapes, backdrop, scale, layout, theme, index, painted) -> None:
    """"about 24 tiny white star dots of 1 to 2px at opacities 0.25 to 0.6".

    Opacity is not expressible on a shape fill in python-pptx, so brightness
    stands in for it -- a dimmer grey is what a 25%-opacity white over this
    ground actually looks like. The placement is seeded per slide so a deck
    redrawn twice has the same sky.
    """
    rng = random.Random(1426 + index)
    for _ in range(24):
        x = rng.uniform(0.02, 0.98) * scale.width
        y = rng.uniform(0.02, 0.98) * scale.height
        d = rng.choice((1, 2))
        level = int(120 + rng.random() * 100)
        shape = _circle(shapes, x, y, scale.x(d) + 1)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(level, level, level + 8 if level < 248 else level)
        _no_line(shape)
        backdrop.push(shape)
    painted.append("star_field")


def _orn_orbit_arcs(shapes, backdrop, scale, layout, theme, index, painted) -> None:
    """"one or two large partial orbit arcs ... bleeding off the corners".

    A freeform polyline, because an arc that leaves the slide is geometry no
    preset autoshape describes.
    """
    color = _rgb(_pick(layout, theme, "accent", "primary_accent", "#E3C77B"), "E3C77B")
    specs = (
        (-0.15, 0.35, 0.85, 200, 340),
        (1.10, 0.90, 0.60, 150, 260),
    )
    drawn = 0
    for cx_f, cy_f, r_f, start_deg, end_deg in specs:
        cx, cy = cx_f * scale.width, cy_f * scale.height
        radius = r_f * scale.width
        points = []
        steps = 48
        for i in range(steps + 1):
            angle = math.radians(start_deg + (end_deg - start_deg) * i / steps)
            points.append((int(cx + radius * math.cos(angle)), int(cy + radius * math.sin(angle))))
        try:
            builder = shapes.build_freeform(points[0][0], points[0][1], scale=1.0)
            builder.add_line_segments(points[1:], close=False)
            shape = builder.convert_to_shape()
        except Exception:
            continue
        _no_fill(shape)
        try:
            shape.line.color.rgb = color
            shape.line.width = Pt(0.75)
        except Exception:
            pass
        _no_shadow(shape)
        backdrop.push(shape)
        drawn += 1
    if drawn:
        painted.append("orbit_arcs")


def _orn_title_line_motif(shapes, backdrop, scale, layout, theme, index, painted) -> None:
    """"a motif of thin electric blue horizontal lines of varying length"."""
    color = _rgb(_pick(layout, theme, "accent", "primary_accent", "#2251FF"), "2251FF")
    rng = random.Random(51 + index)
    y = 470
    for _ in range(9):
        width = scale.x(rng.uniform(70, 380))
        backdrop.push(_rect(shapes, scale.x(80), scale.y(y), width, max(scale.y(2), 1), color))
        y += 22
    painted.append("title_line_motif")


def _orn_status_dot(shapes, backdrop, scale, layout, theme, index, painted) -> None:
    """"small 8px teal status dots with a soft glow meaning healthy".

    The glow is a second, larger, dimmer dot behind the first -- two shapes
    where CSS would use one blur.
    """
    color = _rgb(_pick(layout, theme, "accent", "primary_accent", "#2DD4BF"), "2DD4BF")
    cx = scale.width - scale.x(56) - scale.x(8)
    cy = scale.y(30)
    glow = _circle(shapes, cx - scale.x(5), cy - scale.x(5), scale.x(18))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(
        (color[0] + 13) // 3, (color[1] + 22) // 3, (color[2] + 31) // 3
    )
    _no_line(glow)
    backdrop.push(glow)
    dot = _circle(shapes, cx, cy, scale.x(8))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    _no_line(dot)
    backdrop.push(dot)
    painted.append("status_dot")


def _orn_paper_grain(shapes, backdrop, scale, layout, theme, index, painted) -> None:
    """A per-pixel paper texture -- the one motif shapes cannot express.

    Rendered with Pillow and dropped in full-bleed at the very back, behind
    the ground itself. Skipped in silence when Pillow is unavailable, because
    a missing texture must never cost a deck its export.
    """
    try:
        from PIL import Image
    except Exception:
        return
    base = layout.get("ground_color")
    if not _is_hex(base):
        return
    rgb = _rgb(base)
    try:
        width, height = 256, 144
        image = Image.new("RGB", (width, height), (rgb[0], rgb[1], rgb[2]))
        pixels = image.load()
        rng = random.Random(7 + index)
        for y in range(height):
            for x in range(width):
                n = rng.randint(-6, 6)
                pixels[x, y] = (
                    max(0, min(255, rgb[0] + n)),
                    max(0, min(255, rgb[1] + n)),
                    max(0, min(255, rgb[2] + n)),
                )
        buffer = io.BytesIO()
        image.save(buffer, format="PNG")
        buffer.seek(0)
        picture = shapes.add_picture(buffer, Emu(0), Emu(0), Emu(scale.width), Emu(scale.height))
    except Exception:
        return
    backdrop.push(picture)
    painted.append("paper_grain")


_ORNAMENT_PAINTERS = {
    "double_hairline_border": _orn_double_hairline_border,
    "monogram_seal": _orn_monogram_seal,
    "punch_holes": _orn_punch_holes,
    "star_field": _orn_star_field,
    "orbit_arcs": _orn_orbit_arcs,
    "title_line_motif": _orn_title_line_motif,
    "status_dot": _orn_status_dot,
    "paper_grain": _orn_paper_grain,
}


# =============================================================================
# The pass
# =============================================================================

def _paint_slide(slide, scale, layout, theme, index, count_slides, painted) -> None:
    """One slide's furniture. Background layers first, overlays after.

    The ``_Backdrop`` is created before anything is drawn, so its cursor is
    measured against the slide as the deck left it -- every shape it is handed
    lands behind all of that, in the order it was painted.
    """
    shapes = slide.shapes
    backdrop = _Backdrop(slide)

    is_cover = index == 0
    cover = layout.get("cover") if isinstance(layout.get("cover"), dict) else {}
    skip = bool(layout.get("skip_title_slide")) and is_cover

    # The cover may restate the ground -- "Cover: full-bleed navy" -- and the
    # themes that do say so are exactly the ones that skip the rest.
    ground_layout = layout
    if is_cover and _is_hex(cover.get("ground_color")):
        ground_layout = dict(layout)
        ground_layout["ground"] = "flat"
        ground_layout["ground_color"] = cover["ground_color"]

    ornaments = cover.get("ornament") if is_cover and "ornament" in cover else layout.get("ornament")
    if not isinstance(ornaments, (list, tuple)):
        ornaments = ()

    own_ground = _slide_sets_own_ground(slide)
    if own_ground:
        painted.append("ground:kept_slide_own")

    if "paper_grain" in ornaments and not own_ground:
        _orn_paper_grain(shapes, backdrop, scale, layout, theme, index, painted)

    if not own_ground:
        _paint_ground(shapes, backdrop, scale, ground_layout, theme, painted)

    if not skip:
        if not own_ground:
            _paint_ruling(shapes, backdrop, scale, layout, painted)
            _paint_margin_rule(shapes, backdrop, scale, layout, painted)
        _paint_masthead(shapes, backdrop, scale, layout, theme, painted)

    for name in ornaments:
        painter = _ORNAMENT_PAINTERS.get(name)
        if painter is None or name == "paper_grain":
            continue
        try:
            painter(shapes, backdrop, scale, layout, theme, index, painted)
        except Exception:
            continue

    if skip:
        return

    _paint_tracker(shapes, scale, layout, index, count_slides, painted)
    # ★ Immediately before the footer, never after: the strip must not be able
    # to see -- and so cannot remove -- anything the pass itself painted here.
    _strip_model_footer_text(shapes, scale, layout, painted)
    _paint_footer(shapes, scale, layout, theme, index, count_slides, painted)
    _paint_chip(shapes, scale, layout, theme, painted)
    _paint_corner_mark(shapes, scale, layout, theme, painted)


def paint_theme_furniture(
    pptx_path: Path,
    theme: dict,
    layout: dict,
    logger=None,
) -> dict:
    """Draw one theme's structural motifs onto a deck that is already saved.

    ``theme`` is the plain dict handed to generated deck code -- it supplies
    palette roles and a font name when the layout does not name a colour
    itself. ``layout`` is an entry from :mod:`theme_layouts`.

    Returns ``{"painted": [names], "slides": int}``. ``painted`` names every
    motif drawn, once per drawing, so a 10-slide ledger deck reports its
    ruling ten times and its stamp once.

    Never raises. A deck that could not be decorated is still a valid deck --
    same contract as ``normalize_chart_axes``, and for the same reason.
    """
    result: dict = {"painted": [], "slides": 0}
    if not isinstance(layout, dict):
        layout = EMPTY_LAYOUT

    try:
        prs = Presentation(str(pptx_path))
    except Exception:
        if logger:
            logger.warning("theme furniture: could not reopen %s", pptx_path)
        return result

    try:
        slides = list(prs.slides)
    except Exception:
        return result

    result["slides"] = len(slides)
    if not slides:
        return result

    try:
        scale = _Scale(prs.slide_width, prs.slide_height)
    except Exception:
        return result

    painted: list[str] = []
    for index, slide in enumerate(slides):
        try:
            _paint_slide(slide, scale, layout, theme, index, len(slides), painted)
        except Exception:
            if logger:
                logger.warning("theme furniture: slide %d skipped", index + 1)
            continue

    # "One rotated outline stamp per deck" -- so it is painted once, on the
    # first slide the theme does not exempt.
    if isinstance(layout.get("stamp"), dict):
        first = 1 if (layout.get("skip_title_slide") and len(slides) > 1) else 0
        try:
            _paint_stamp(slides[first].shapes, scale, layout, painted)
        except Exception:
            pass

    result["painted"] = painted
    if not painted:
        return result

    try:
        prs.save(str(pptx_path))
    except Exception:
        if logger:
            logger.warning("theme furniture: could not save %s", pptx_path)
        return {"painted": [], "slides": len(slides)}

    return result


# =============================================================================
# forbid_boxes -- a card becomes a left-rule row
# =============================================================================
#
# The measured gap between our decks and the design systems they claim is not
# palette and not type -- the model gets those right every time. It is the
# PROHIBITIONS, and the structural one is `forbid_boxes`. Ours: four filled,
# bordered, shadowed panels. The reference for the same theme: four rows
# separated by nothing but a thin left rule, and half the slide empty.
#
# `enforce_theme_rules` clears the mechanical prohibitions (shadows, corners,
# gradients) and deliberately only COUNTS `boxes`, saying in its own docstring
# that "turning a filled card into a left-rule row is layout, not enforcement,
# and belongs where the slide is composed". This is that place.
#
# ★ The danger here is not the conversion, it is the identification. A filled
# rectangle is also the takeaway panel several themes REQUIRE, a chart plot
# ground, a KPI tile, a decorative bleed, a table cell. Converting one of those
# breaks the deck, and a wrong conversion is worse than an unconverted card --
# so every test below is a REASON TO REFUSE, the shape is left where it stands
# on the first refusal, and the refusal is counted and named rather than
# swallowed. `cards_left` is expected to be non-zero on a healthy deck.

#: Rectangle geometries. A card is a rectangle; an oval, an arrow or a chevron
#: is something else and is never touched, whatever it is filled with.
_RECT_PRESETS = frozenset({
    "rect",
    "roundRect",
    "snip1Rect",
    "snip2SameRect",
    "snip2DiagRect",
    "snipRoundRect",
    "round1Rect",
    "round2SameRect",
    "round2DiagRect",
})

#: The size window a card sits in, on the prompts' 1280x720 canvas. Below it a
#: shape is a bullet mark, a swatch or a divider; above it a shape is a bleed,
#: a plot ground or a full-width band, none of which is a row.
_CARD_MIN_W_PX = 180
_CARD_MIN_H_PX = 36
_CARD_MAX_W_FRAC = 0.72
_CARD_MAX_AREA_FRAC = 0.42

#: A card's fill is a light tint of the paper. A takeaway panel is the loud
#: shape on the slide -- dark, or strongly coloured -- and this is what keeps
#: the two apart before anything else is measured.
_LIGHT_MIN = 0.62      # perceptual luminance, 0..1
_CHROMA_MAX = 60       # max channel minus min channel, 0..255

#: Two shapes are siblings when their left edge, width and height agree within
#: this many canvas pixels. Generated decks lay a row out by arithmetic, so the
#: agreement is usually exact; the tolerance is for the ones that round.
_SIBLING_TOL_PX = 10

#: "a thin vertical accent rule" -- 3px on the canvas, 2.25pt.
_RULE_WIDTH_PX = 3

#: How many refusals are named individually before the list is summarised. The
#: reasons repeat; an unbounded list would be a log, not a summary.
_WHY_LIMIT = 40


def _geom_preset(shape) -> Optional[str]:
    """The shape's preset geometry name, or None if it has no `prstGeom`."""
    try:
        sp_pr = shape._element.spPr
    except Exception:
        return None
    if sp_pr is None:
        return None
    try:
        prst_geom = sp_pr.find(qn("a:prstGeom"))
    except Exception:
        return None
    if prst_geom is None:
        return None
    return prst_geom.get("prst")


def _is_rect_like(shape) -> bool:
    """A rectangle, and something whose geometry we can actually read.

    A picture, a chart, a table and a group all fail here: a ``graphicFrame``
    has no ``spPr``, so a table CELL is not reachable from this walk at all,
    which is exactly the outcome wanted.
    """
    if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
        return False
    if getattr(shape, "shapes", None) is not None:
        return False  # a group; its children are deliberately out of scope
    return _geom_preset(shape) in _RECT_PRESETS


def _fill_type(shape):
    try:
        return shape.fill.type
    except Exception:
        return None


def _solid_fill_rgb(shape) -> Optional[RGBColor]:
    """The shape's flat fill colour, or None if it has no readable one.

    A gradient, a picture fill, a pattern and a theme-INDEXED colour all
    answer None. A theme-indexed fill is a real refusal and not a technicality:
    we cannot tell whether it resolves light or dark, and a card is identified
    by being light.
    """
    if _fill_type(shape) != MSO_FILL.SOLID:
        return None
    try:
        return shape.fill.fore_color.rgb
    except Exception:
        return None


def _is_light_tint(rgb) -> bool:
    """Light enough, and neutral enough, to be a card rather than a panel."""
    try:
        r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
    except Exception:
        return False
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255.0
    chroma = max(r, g, b) - min(r, g, b)
    return luminance >= _LIGHT_MIN and chroma <= _CHROMA_MAX


def _has_real_text(shape) -> bool:
    try:
        return bool(shape.has_text_frame and shape.text_frame.text.strip())
    except Exception:
        return False


def _box(shape):
    """`(left, top, width, height)` in EMU, or None if any of it is missing."""
    try:
        values = (shape.left, shape.top, shape.width, shape.height)
    except Exception:
        return None
    if any(v is None for v in values):
        return None
    return tuple(int(v) for v in values)


def _sits_on(card_box, shape) -> bool:
    """True when the shape's CENTRE lies inside the card.

    Centre, not intersection: a heading whose box merely grazes a panel below
    it is not sitting on it, and using intersection made a full-width title
    "belong" to every panel under it.
    """
    box = _box(shape)
    if box is None:
        return False
    left, top, width, height = box
    cx, cy = left + width // 2, top + height // 2
    c_left, c_top, c_width, c_height = card_box
    return c_left <= cx <= c_left + c_width and c_top <= cy <= c_top + c_height


def _in_card_size_window(box, scale) -> bool:
    _left, _top, width, height = box
    if width < scale.x(_CARD_MIN_W_PX) or height < scale.y(_CARD_MIN_H_PX):
        return False
    if width > scale.width * _CARD_MAX_W_FRAC:
        return False
    return width * height <= scale.width * scale.height * _CARD_MAX_AREA_FRAC


def _stack_key(box, tol) -> tuple:
    left, _top, width, height = box
    return (round(left / tol), round(width / tol), round(height / tol))


def _is_vertical_stack(boxes, tol) -> bool:
    """Rows in a column: distinct tops, and no two of them overlapping.

    This is the test that separates a row of KPI tiles from a stack of cards.
    Tiles share a top and differ in left; rows share a left and differ in top,
    and only the second shape is what a left rule can replace -- a rule down
    the side of a tile in a horizontal row rules off nothing.
    """
    ordered = sorted(boxes, key=lambda b: b[1])
    for lower, upper in zip(ordered, ordered[1:]):
        if upper[1] - lower[1] < tol:
            return False  # same row, side by side
        if upper[1] + tol < lower[1] + lower[3]:
            return False  # overlapping, so not laid out as rows
    return True


def _draw_left_rule(slide, card, color, width) -> None:
    """The row's rule, at the card's left edge, BEHIND the card.

    Behind, not in front: the card is transparent by the time this is called,
    so the rule reads through it, and anything the deck drew over the card --
    its heading, its body copy, its own text frame -- still covers the rule
    rather than being covered by it. A rule painted on top of the text it
    belongs beside is worse than no rule.
    """
    box = _box(card)
    if box is None:
        return
    left, top, _width, height = box
    rule = _rect(slide.shapes, left, top, width, height, color)
    tree = slide.shapes._spTree
    element = rule._element
    tree.remove(element)
    tree.insert(list(tree).index(card._element), element)


def _convert_slide(slide, scale, color, rule_width, refuse) -> int:
    """One slide. Returns how many cards became rows."""
    shapes = list(slide.shapes)
    texts = [s for s in shapes if _has_real_text(s)]

    kept: list = []
    for shape in shapes:
        if not _is_rect_like(shape):
            continue
        box = _box(shape)
        if box is None:
            continue
        if _fill_type(shape) in (None, MSO_FILL.BACKGROUND):
            continue  # an outline or an invisible shape is not a box at all
        if not _in_card_size_window(box, scale):
            continue  # too small to be a panel, or too big to be a row

        rgb = _solid_fill_rgb(shape)
        if rgb is None:
            refuse(shape, "its fill is not a plain readable solid colour")
            continue
        if not _is_light_tint(rgb):
            refuse(shape, "its fill is dark or strongly coloured, which is how a "
                          "takeaway panel and a decorative bleed both read")
            continue
        if not any(_sits_on(box, other) for other in texts if other is not shape) \
                and not _has_real_text(shape):
            refuse(shape, "no text sits on it, which is how a chart plot ground "
                          "and a decorative panel both read")
            continue
        kept.append((shape, box))

    groups: dict = {}
    for shape, box in kept:
        groups.setdefault(_stack_key(box, scale.x(_SIBLING_TOL_PX) or 1), []).append((shape, box))

    converted = 0
    for members in groups.values():
        boxes = [box for _shape, box in members]
        if len(members) < 2 or not _is_vertical_stack(boxes, scale.y(_SIBLING_TOL_PX) or 1):
            for shape, _box_ in members:
                refuse(shape, "it is not one of several near-identical rows stacked "
                              "down the slide, so it may be the one panel the theme wants")
            continue
        for shape, _box_ in members:
            try:
                _no_fill(shape)
                _no_line(shape)
                _no_shadow(shape)
                _draw_left_rule(slide, shape, color, rule_width)
                converted += 1
            except Exception:
                refuse(shape, "the conversion itself could not be applied")
    return converted


def convert_boxes_to_rules(
    pptx_path: Path,
    theme: dict,
    layout: dict,
    logger=None,
) -> dict:
    """Turn a stark theme's filled cards into the reference's left-rule rows.

    Runs only when the layout says ``forbid_boxes``. A card loses its fill and
    its border and gains a thin vertical accent rule at its left edge; nothing
    moves, so every word on the slide stays exactly where the deck put it.

    What is deliberately NEVER converted, because each one is a filled
    rectangle too and converting it breaks the deck:

    * the takeaway / bottom-line panel -- it is the one loud shape, so it fails
      the light-tint test and, being alone, the sibling test as well;
    * a chart plot ground -- no text sits on it, and a chart's own frame is not
      reachable from this walk;
    * a KPI tile row -- tiles are laid out ACROSS, and only a column of rows is
      converted;
    * a decorative bleed -- too large for the size window;
    * a table cell -- inside a ``graphicFrame``, which this never opens;
    * anything inside a group, anything whose fill is a gradient, a picture or
      a theme-indexed colour, and anything that is not a rectangle.

    Returns ``{"cards_converted": int, "cards_left": int, "why_left": [str]}``.
    ``cards_left`` counts filled rectangles that looked box-shaped and were
    left alone anyway; on a healthy deck it is normally the larger number.

    Never raises -- same contract as :func:`paint_theme_furniture` and
    ``normalize_chart_axes``. A deck this could not restyle is still a deck.
    """
    result: dict = {"cards_converted": 0, "cards_left": 0, "why_left": []}

    if not isinstance(layout, dict) or not layout.get("forbid_boxes"):
        return result

    try:
        prs = Presentation(str(pptx_path))
    except Exception:
        if logger:
            logger.warning("forbid_boxes: could not reopen %s", pptx_path)
        return result

    try:
        slides = list(prs.slides)
        scale = _Scale(prs.slide_width, prs.slide_height)
    except Exception:
        return result

    color = _rgb(_pick(layout, theme, "accent", "primary_accent", "#1F3A5F"), "1F3A5F")
    rule_width = max(scale.x(_RULE_WIDTH_PX), 1)

    left_count = 0
    reasons: list[str] = []

    converted = 0
    for index, slide in enumerate(slides):
        def refuse(_shape, why, _index=index):
            nonlocal left_count
            left_count += 1
            if len(reasons) < _WHY_LIMIT:
                reasons.append(f"slide {_index + 1}: {why}")

        try:
            converted += _convert_slide(slide, scale, color, rule_width, refuse)
        except Exception:
            if logger:
                logger.warning("forbid_boxes: slide %d skipped", index + 1)
            continue

    if len(reasons) < left_count:
        reasons.append(f"... and {left_count - len(reasons)} more left alone")

    result["cards_left"] = left_count
    result["why_left"] = reasons

    if not converted:
        return result

    try:
        prs.save(str(pptx_path))
    except Exception:
        if logger:
            logger.warning("forbid_boxes: could not save %s", pptx_path)
        # Nothing persisted, so nothing was converted -- and those cards are
        # still cards. Say both.
        result["cards_left"] = left_count + converted
        result["why_left"] = reasons + [
            f"the deck could not be saved, so {converted} converted cards were lost"
        ]
        return result

    result["cards_converted"] = converted
    return result
