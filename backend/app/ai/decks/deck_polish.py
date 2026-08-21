"""FIX-F / FIX-G / FIX-H (2026-08-21) — post-save polish on the shipped deck.

Why this module exists
----------------------
The phase-2 verification deck (fabric rerun, artifact dc7cc144) was numerically
correct — the data gate saw to that — and still carried three defects the gate
is structurally blind to, because none of them is a figure:

* **A fabricated attribution** — "Source: team analysis" on every slide. No
  team analysed anything; the model invents the credit line as slide furniture.
  The codegen prompt now forbids it, but a prohibition that lives in a prompt
  is followed only usually. (FIX-G)
* **Raw float data labels** — charts labelled 86.438351 and 1711.2989, the
  honest side effect of computing values from rows instead of typing rounded
  ones. Correct number, unpresentable notation. (FIX-H)
* **Invisible text** — the original RCA's cause 5: the renderer-based layout
  check measures geometry, not readability, so a white title on a white ground
  passes it at exactly 1.00:1. The furniture pass now keeps a slide's own dark
  background (phase 1), which fixes the measured cases, but nothing MEASURES
  contrast on the shipped file — a regression would be invisible twice over.
  (FIX-F)

FIX-G and FIX-H edit the saved .pptx deterministically, in the same post-save
slot as the furniture pass. FIX-F only reports — a contrast fix would mean
choosing new colours for the model's design, and a wrong automatic recolour is
worse than a warning the agent can act on.

Contract, same as every pass in this slot: never raises, the deck ships
whatever happens here, and every function returns a small dict so "ran and
found nothing", "ran and changed things" and "could not run" stay
distinguishable in the stored artifact content.
"""
from __future__ import annotations

import logging
import re
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

_C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"

# How many contrast issues a verdict carries — same cap and same reasoning as
# the data gate's _MAX_VIOLATIONS.
_MAX_ISSUES = 12

# WCAG large-text minimum. Slide text is large text; 3:1 is the floor below
# which it stops being readable, and the measured defect sat at 1.00:1.
_MIN_CONTRAST = 3.0


# ---------------------------------------------------------------------------
# FIX-G — strip fabricated attributions
# ---------------------------------------------------------------------------

# A whole paragraph that credits a source that does not exist. Deliberately
# requires a generic-analysis terminal word, so a REAL attribution ("Source:
# Microsoft Fabric", "Source: POS sales data") can never match.
_FABRICATED_SOURCE = re.compile(
    r"^\s*source\s*[:\-–—]\s*"
    r"(?:(?:the|our|company|internal|team|management|analyst|staff)\s+)*"
    r"(?:data\s+)?"
    r"(?:analysis|analytics|estimates?|research|calculations?)"
    r"\s*\.?\s*$",
    re.IGNORECASE,
)


def strip_fabricated_attribution(pptx_path: str, logger_=None) -> Dict[str, Any]:
    """Remove "Source: team analysis"-style credit lines from the saved deck.

    Deletion, not replacement: this code has no data-source name to offer (the
    tool sees query results, not connectors), and an absent credit line is
    honest where an invented one is not. The cover already carries the
    deterministic "Data through <period>" stamp for provenance.
    """
    log = logger_ or logger
    try:
        from pptx import Presentation

        prs = Presentation(str(pptx_path))
        removed = 0
        for slide in prs.slides:
            for shape in list(slide.shapes):
                if not getattr(shape, "has_text_frame", False):
                    continue
                tf = shape.text_frame
                matched = [
                    p for p in tf.paragraphs if _FABRICATED_SOURCE.match(p.text or "")
                ]
                if not matched:
                    continue
                if len(matched) == len(tf.paragraphs):
                    # The whole shape is the fabricated credit — drop the shape.
                    el = shape._element
                    el.getparent().remove(el)
                    removed += len(matched)
                    continue
                for p in matched:
                    p._p.getparent().remove(p._p)
                    removed += 1
        if removed:
            prs.save(str(pptx_path))
        return {"status": "checked", "removed": removed}
    except Exception as err:
        log.warning("attribution strip failed open: %s", err)
        return {"status": "unavailable", "reason": str(err)}


# ---------------------------------------------------------------------------
# FIX-H — round chart data labels
# ---------------------------------------------------------------------------

def round_chart_data_labels(pptx_path: str, logger_=None) -> Dict[str, Any]:
    """Give every shown data label a display number format.

    Only charts whose plotted values actually carry sub-tenth precision are
    touched, and only their LABELS — the plotted values stay exact. The format
    is picked from the chart's own magnitude: thousands read as integers,
    smaller scales keep one decimal.
    """
    log = logger_ or logger
    try:
        from pptx import Presentation
        from lxml import etree

        prs = Presentation(str(pptx_path))
        charts_fixed = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_chart", False):
                    continue
                try:
                    chart = shape.chart
                    needs, max_abs = False, 0.0
                    for plot in chart.plots:
                        for series in plot.series:
                            for v in series.values:
                                if v is None:
                                    continue
                                max_abs = max(max_abs, abs(float(v)))
                                if abs(float(v) - round(float(v), 1)) > 1e-9:
                                    needs = True
                    if not needs:
                        continue
                    fmt = "#,##0" if max_abs >= 1000 else "#,##0.0"
                    space = chart._chartSpace
                    dlbls_list = space.findall(f".//{{{_C_NS}}}dLbls")
                    changed_here = False
                    for dlbls in dlbls_list:
                        delete = dlbls.find(f"{{{_C_NS}}}delete")
                        if delete is not None and delete.get("val") in ("1", "true"):
                            continue
                        old = dlbls.find(f"{{{_C_NS}}}numFmt")
                        if old is not None:
                            dlbls.remove(old)
                        num_fmt = etree.SubElement(dlbls, f"{{{_C_NS}}}numFmt")
                        num_fmt.set("formatCode", fmt)
                        num_fmt.set("sourceLinked", "0")
                        # numFmt is first in dLbls' schema sequence.
                        dlbls.remove(num_fmt)
                        dlbls.insert(0, num_fmt)
                        changed_here = True
                    if changed_here:
                        charts_fixed += 1
                except Exception:
                    continue
        if charts_fixed:
            prs.save(str(pptx_path))
        return {"status": "checked", "charts_fixed": charts_fixed}
    except Exception as err:
        log.warning("data label rounding failed open: %s", err)
        return {"status": "unavailable", "reason": str(err)}


# ---------------------------------------------------------------------------
# FIX-F — contrast check (reports only)
# ---------------------------------------------------------------------------

def _luminance(rgb) -> float:
    channels = []
    for c in (rgb[0], rgb[1], rgb[2]):
        x = c / 255.0
        channels.append(x / 12.92 if x <= 0.04045 else ((x + 0.055) / 1.055) ** 2.4)
    return 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]


def _contrast(rgb_a, rgb_b) -> float:
    hi, lo = sorted((_luminance(rgb_a), _luminance(rgb_b)), reverse=True)
    return (hi + 0.05) / (lo + 0.05)


def _slide_bg_rgb(slide) -> Optional[tuple]:
    """The slide's own ``<p:bg>`` solid colour, if the deck code set one."""
    try:
        from pptx.oxml.ns import qn

        bg = slide._element.cSld.find(qn("p:bg"))
        if bg is None:
            return None
        srgb = bg.find(".//" + qn("a:srgbClr"))
        if srgb is None:
            return None
        val = srgb.get("val")
        return tuple(int(val[i:i + 2], 16) for i in (0, 2, 4))
    except Exception:
        return None


def _solid_fill_rgb(shape) -> Optional[tuple]:
    try:
        from pptx.enum.dml import MSO_FILL

        if shape.fill.type != MSO_FILL.SOLID:
            return None
        rgb = shape.fill.fore_color.rgb
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


def _run_rgb(run) -> Optional[tuple]:
    try:
        rgb = run.font.color.rgb
        if rgb is None:
            return None
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


def check_deck_contrast(
    pptx_path: str,
    theme: Optional[Dict[str, Any]] = None,
    logger_=None,
) -> Dict[str, Any]:
    """Measure text-vs-backdrop contrast on the SHIPPED file. Reports only.

    The renderer-based layout check measures geometry; this measures the one
    thing it cannot — whether the text is readable against whatever actually
    sits behind it. Deterministic python-pptx walk, no renderer, no font
    substitution noise.

    Only EXPLICIT colours are judged: a run whose colour is inherited or
    theme-indexed resolves differently per master and is skipped rather than
    guessed at — a false "invisible" verdict on a readable deck would teach
    the agent to distrust the check.
    """
    log = logger_ or logger
    try:
        from pptx import Presentation

        prs = Presentation(str(pptx_path))
        slide_w = int(prs.slide_width or 12192000)
        slide_h = int(prs.slide_height or 6858000)
        slack = int(slide_w * 0.01)

        base_default = (255, 255, 255)
        if isinstance(theme, dict):
            candidate = str((theme.get("palette") or {}).get("background") or "").lstrip("#")
            if len(candidate) == 6:
                try:
                    base_default = tuple(int(candidate[i:i + 2], 16) for i in (0, 2, 4))
                except Exception:
                    pass

        issues: List[Dict[str, Any]] = []
        for idx, slide in enumerate(prs.slides, start=1):
            base = _slide_bg_rgb(slide) or base_default
            painted: List[tuple] = []  # (left, top, right, bottom, rgb)
            for shape in slide.shapes:
                rect = None
                try:
                    if shape.left is not None and shape.top is not None:
                        rect = (
                            int(shape.left), int(shape.top),
                            int(shape.left) + int(shape.width or 0),
                            int(shape.top) + int(shape.height or 0),
                        )
                except Exception:
                    rect = None

                fill_rgb = _solid_fill_rgb(shape)

                if getattr(shape, "has_text_frame", False):
                    text = (shape.text_frame.text or "").strip()
                    if len(text) >= 3:
                        # The shape's own fill is its text's backdrop; else the
                        # topmost earlier solid shape containing it; else the
                        # slide ground.
                        backdrop = fill_rgb
                        if backdrop is None and rect is not None:
                            for (pl, pt_, pr, pb, prgb) in reversed(painted):
                                if (pl - slack <= rect[0] and pt_ - slack <= rect[1]
                                        and pr + slack >= rect[2] and pb + slack >= rect[3]):
                                    backdrop = prgb
                                    break
                        if backdrop is None:
                            backdrop = base
                        worst = None
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if not (run.text or "").strip():
                                    continue
                                rgb = _run_rgb(run)
                                if rgb is None:
                                    continue
                                ratio = _contrast(rgb, backdrop)
                                if worst is None or ratio < worst[0]:
                                    worst = (ratio, rgb)
                        if worst is not None and worst[0] < _MIN_CONTRAST:
                            if len(issues) < _MAX_ISSUES:
                                issues.append({
                                    "slide": idx,
                                    "text": text[:60],
                                    "ratio": round(worst[0], 2),
                                    "text_color": "#%02X%02X%02X" % worst[1],
                                    "backdrop": "#%02X%02X%02X" % backdrop,
                                })

                # A big solid shape becomes the backdrop for what draws later.
                if fill_rgb is not None and rect is not None:
                    painted.append((rect[0], rect[1], rect[2], rect[3], fill_rgb))

        return {
            "status": "checked",
            "slides": len(prs.slides),
            "issues": issues,
        }
    except Exception as err:
        log.warning("deck contrast check failed open: %s", err)
        return {"status": "unavailable", "reason": str(err)}
