"""FIX-D / FIX-E (2026-08-21) — a deck may only state figures its own data holds.

Why this module exists
----------------------
Measured on 0.0.543.18, four decks generated through the chat API: the system
ran real queries and got correct answers, and then the slide-building step wrote
python-pptx code containing the numbers TYPED FROM MEMORY of the conversation
rather than read from its `visualizations` parameter. When the memory was good
the deck was right by luck (the Power BI deck); when it was not, the deck
shipped 20+ wrong figures (the excel deck put Riverside at 1,995 against a real
1,459 — flipping a branch ranking) or invented 7 of 12 months of a revenue
chart outright, the generated code even carrying the comment "Sample seasonal
monthly values". The invented numbers were self-consistent — totals summed,
percentages recomputed — so nothing inside the deck ever looked wrong.

The dashboard insight panel and the chat narrative already have exactly this
guarantee, provided by `app.services.figure_grounding`: every figure is checked
against the run's own data and an ungrounded one cannot be published. This
module extends that same engine — same tokenizer, same precision-derived slack,
same whitelist of derivable aggregates — to the third surface, the deck code,
which until now was the only figure-bearing surface with no check at all.

What is checked (FIX-D)
-----------------------
Two places fabrications were actually observed to live in generated deck code:

* **String literals** — slide titles, takeaways, stat tiles ("Riverside 1,995",
  "$4.8B full-year target"). Every free-standing figure in every string constant
  is checked with `figure_grounding.is_grounded`.
* **Hardcoded numeric series** — list/tuple literals of 3+ numbers fed to
  `chart_data.add_series(...)` (the invented seasonal months). Chart series are
  routinely written in display scale (202.3 meaning 202.3M), so each element is
  also tried with K/M/B suffixes before being called ungrounded.

Deliberately NOT checked, so the gate cannot reject an honest deck:

* Percentages — shares and growth rates are computed FROM the values and appear
  nowhere among them; `is_grounded` passes them by design, and a whole series
  living inside [0, 100] is skipped as percent-scale for the same reason.
* Small structural integers — "top 10", slide counts, `Inches(1)`, `Pt(24)`,
  `RGBColor(15, 23, 42)`: scalar call arguments are never collected, and
  `is_grounded` passes integers under 1000 as structural anyway.
* Fractions in [0, 1] — shares, opacities.
* Code that COMPUTES its numbers from `visualizations` — the desired behaviour
  produces no literals for the gate to look at.

The gate reports; the caller decides. `check_deck_code` returns a three-state
verdict like the layout check — "checked" (with any violations), "skipped"
(no data to check against), "unavailable" (the code would not parse) — and
never raises.

What is stamped (FIX-E)
-----------------------
The same measurement found every deck presenting old data as current: "how is
this year going" (August 2026) answered with 2025 and 2023 rows under a big
"2026" banner, with no slide anywhere saying when the data ends. The model is
now told to disclose, but disclosure that depends on a prompt being followed is
disclosure that sometimes does not happen — so `data_coverage_label` reads the
newest period straight out of the visualization rows (reusing the insight
panel's period parser) and `stamp_data_coverage` writes "Data through <period>"
onto the cover deterministically, in a colour picked against the cover's actual
background.
"""
from __future__ import annotations

import ast
import logging
from typing import Any, Dict, List, Optional

import re

from app.services.figure_grounding import (
    canonical,
    data_magnitudes,
    is_grounded,
    numbers_in,
)

logger = logging.getLogger(__name__)

# How many violations a verdict carries. Enough to show the pattern and feed a
# repair prompt; a deck wrong 40 ways does not need all 40 listed to be wrong.
_MAX_VIOLATIONS = 12

# A hardcoded list shorter than this is not a chart series — pairs like
# (width, height) or (x, y) are geometry, not data claims.
_MIN_SERIES_LEN = 3


# ---------------------------------------------------------------------------
# FIX-D — the gate
# ---------------------------------------------------------------------------

_NUMERIC_STRING = re.compile(r"^\s*-?[\d,]+(\.\d+)?\s*$")


def _deck_magnitudes(visualizations: List[Dict[str, Any]]) -> List[float]:
    """`figure_grounding.data_magnitudes` plus what a DECK legitimately cites.

    Replaying the gate over the three RCA decks' real stored code found two
    honest-claim classes the shared pool misses:

    * **Identifiers stored as strings.** The Power BI deck names article
      "1000000345139" — verbatim from a cell, but that column is stored as
      text, and `data_magnitudes` only pools numeric cells. A value copied
      verbatim from a cell is grounded regardless of the column's storage
      type, so numeric-parseable string cells join the pool.
    * **Dataset sizes.** "300,000+ line records" cites the query's row_count,
      which lives on the visualization dict, not in any row.
    """
    magnitudes = data_magnitudes(visualizations or [])
    for viz in visualizations or []:
        for key in ("row_count", "rows_available", "sample_row_count"):
            n = (viz or {}).get(key)
            if isinstance(n, (int, float)) and not isinstance(n, bool) and n:
                magnitudes.append(float(n))
        for row in (viz or {}).get("rows") or []:
            if not isinstance(row, dict):
                continue
            for v in row.values():
                if isinstance(v, str) and len(v) <= 40 and _NUMERIC_STRING.match(v):
                    try:
                        magnitudes.append(float(v.replace(",", "")))
                    except ValueError:
                        pass
    return magnitudes


def _relaxed_grounded(token: str, magnitudes: List[float]) -> bool:
    """A rounded round-number claim, held to the precision it actually shows.

    `write_precision_slack` reads "300,000" as stated to the unit and allows
    ±0.5 — which rejects an honest "300,000+ records" against a true 301,245.
    Trailing zeros are ambiguous precision: a writer who says 300,000 about a
    six-digit count is rounding to the zeros, not counting to the unit. So
    when the strict check fails AND the token ends in zeros AND the value is
    large enough that rounding is plausible (≥1000), retry with slack of half
    the last NON-zero digit's place.

    The measured fabrications all keep failing under this rule: 12,827 /
    352,747 / 1,995 have no trailing zeros, and 610M / 714M / 337.6M carry
    their precision in the mantissa.
    """
    value = canonical(token)
    if value is None or abs(value) < 1000:
        return False
    t = token.strip().rstrip("%")
    mult = 1.0
    if t and t[-1] in "KkMmBb":
        mult = {"k": 1e3, "m": 1e6, "b": 1e9}[t[-1].lower()]
        t = t[:-1]
    t = t.replace(",", "").lstrip("+-")
    if "." in t or not t.endswith("0"):
        return False
    trailing = len(t) - len(t.rstrip("0"))
    slack = (10.0 ** trailing) * mult / 2.0
    return any(abs(value - m) <= slack for m in magnitudes)


def _series_value_grounded(value: float, magnitudes: List[float]) -> bool:
    """A hardcoded chart value, tried at the scales charts are drawn in.

    `is_grounded` reads a TOKEN, and a bare `202.3` in a series list carries no
    scale suffix even when it means 202.3M — deck code divides by 1e6 for the
    axis and the literal keeps display scale. So the raw token is tried first,
    then with K/M/B welded on, which reuses `write_precision_slack`'s rule that
    tolerance scales with the stated magnitude.
    """
    v = abs(float(value))
    if float(v).is_integer():
        tok = str(int(v))
    else:
        tok = repr(v)
    for suffix in ("", "K", "M", "B"):
        if is_grounded(tok + suffix, magnitudes) or _relaxed_grounded(tok + suffix, magnitudes):
            return True
    return False


def _claims_in_code(code: str) -> List[Dict[str, Any]]:
    """Every figure-bearing claim the generated code states as a literal.

    Returns dicts of ``{"token", "where", "context"}`` where `where` is
    ``"text"`` (a figure inside a string constant) or ``"series"`` (an element
    of a hardcoded numeric list/tuple). Raises SyntaxError if the code does not
    parse — the caller maps that to status "unavailable".
    """
    tree = ast.parse(code)
    claims: List[Dict[str, Any]] = []
    # Elements of collected series, so a value is not reported twice when the
    # same list literal appears in both a List node and its parent expression.
    seen: set = set()

    for node in ast.walk(tree):
        if isinstance(node, ast.Constant) and isinstance(node.value, str):
            text = node.value
            for tok in numbers_in(text):
                claims.append({
                    "token": tok,
                    "where": "text",
                    "context": text.strip()[:90],
                })
        elif isinstance(node, (ast.List, ast.Tuple)):
            values: List[float] = []
            for elt in node.elts:
                inner = elt
                # -5.2 parses as UnaryOp(USub, Constant(5.2))
                if isinstance(inner, ast.UnaryOp) and isinstance(inner.op, (ast.USub, ast.UAdd)):
                    inner = inner.operand
                if (
                    isinstance(inner, ast.Constant)
                    and isinstance(inner.value, (int, float))
                    and not isinstance(inner.value, bool)
                ):
                    values.append(float(inner.value))
                else:
                    values = []
                    break
            if len(values) < _MIN_SERIES_LEN:
                continue
            # A series living entirely inside [0, 100] is percent-scale (or
            # small structural counts) — shares are computed from the data and
            # will not appear among its magnitudes, so checking them would
            # reject honest decks.
            if all(abs(v) <= 100 for v in values):
                continue
            key = (node.lineno, node.col_offset)
            if key in seen:
                continue
            seen.add(key)
            preview = ", ".join(
                str(int(v)) if float(v).is_integer() else f"{v:g}" for v in values[:6]
            )
            for v in values:
                claims.append({
                    "token": str(int(v)) if float(v).is_integer() else f"{v:g}",
                    "where": "series",
                    "value": v,
                    "context": f"hardcoded series [{preview}{', …' if len(values) > 6 else ''}]",
                })
    return claims


def check_deck_code(code: str, visualizations: List[Dict[str, Any]]) -> Dict[str, Any]:
    """Verify every literal figure in generated deck code against the data.

    Three-state, never raises:
      * ``{"status": "skipped", "reason": ...}`` — no numeric data to check
        against (a narrative deck must not be blocked by its own emptiness).
      * ``{"status": "unavailable", "reason": ...}`` — the code would not parse.
      * ``{"status": "checked", "claims": N, "violations": [...]}`` — the
        verdict; an empty violations list is a clean deck.
    """
    try:
        magnitudes = _deck_magnitudes(visualizations or [])
        if not magnitudes:
            return {"status": "skipped", "reason": "no numeric data in visualizations"}
        claims = _claims_in_code(code or "")
        violations: List[Dict[str, Any]] = []
        for claim in claims:
            if claim["where"] == "series":
                ok = _series_value_grounded(claim["value"], magnitudes)
            else:
                ok = (
                    is_grounded(claim["token"], magnitudes)
                    or _relaxed_grounded(claim["token"], magnitudes)
                )
            if not ok:
                violations.append({
                    "token": claim["token"],
                    "where": claim["where"],
                    "context": claim["context"],
                })
                if len(violations) >= _MAX_VIOLATIONS:
                    break
        return {
            "status": "checked",
            "claims": len(claims),
            "violations": violations,
        }
    except SyntaxError as err:
        return {"status": "unavailable", "reason": f"code does not parse: {err.msg}"}
    except Exception as err:  # a gate must never be the reason a deck dies
        logger.warning("deck data gate failed open: %s", err)
        return {"status": "unavailable", "reason": str(err)}


# ---------------------------------------------------------------------------
# FIX-E — data coverage
# ---------------------------------------------------------------------------

_MONTH_NAMES = (
    "Jan", "Feb", "Mar", "Apr", "May", "Jun",
    "Jul", "Aug", "Sep", "Oct", "Nov", "Dec",
)


def data_coverage_label(visualizations: List[Dict[str, Any]]) -> Optional[str]:
    """The newest period any visualization's rows reach, e.g. "Dec 2025".

    Reuses the insight panel's period parser (`artifact_insights._period_key`
    and `_period_column`) so a period reads the same on every surface. None
    when no visualization carries a recognisable period column — a deck about
    a subject with no timeline has no coverage to state.
    """
    try:
        from app.services.artifact_insights import _period_column, _period_key

        best: Optional[tuple] = None
        for viz in visualizations or []:
            rows = [r for r in ((viz or {}).get("rows") or []) if isinstance(r, dict)]
            col = _period_column(rows)
            if not col:
                continue
            for row in rows:
                key = _period_key(row.get(col))
                if key is not None and (best is None or key > best):
                    best = key
        if best is None:
            return None
        year, sub = best
        if sub == 0:
            return str(year)
        # Month-bearing keys are month*100(+day); the quarter form stores the
        # quarter's closing month bare (3, 6, 9, 12).
        month = sub // 100 if sub >= 100 else sub
        if 1 <= month <= 12:
            return f"{_MONTH_NAMES[month - 1]} {year}"
        return str(year)
    except Exception as err:
        logger.warning("data coverage label failed open: %s", err)
        return None


def _cover_is_dark(slide, theme: Optional[Dict[str, Any]]) -> bool:
    """Whether the cover's ground is dark, so the stamp stays readable on it.

    The cover's own ``<p:bg>`` wins when the model set one (that is the slide
    the furniture pass now leaves alone); otherwise the theme's background.
    """
    hex_color = None
    try:
        from pptx.oxml.ns import qn

        bg = slide._element.cSld.find(qn("p:bg"))
        if bg is not None:
            srgb = bg.find(".//" + qn("a:srgbClr"))
            if srgb is not None:
                hex_color = srgb.get("val")
    except Exception:
        pass
    if hex_color is None and isinstance(theme, dict):
        palette = theme.get("palette") or {}
        candidate = str(palette.get("background") or "").lstrip("#")
        if len(candidate) == 6:
            hex_color = candidate
    if hex_color is None:
        return False  # python-pptx default deck ground is white
    try:
        r, g, b = (int(hex_color[i:i + 2], 16) for i in (0, 2, 4))
        return (0.2126 * r + 0.7152 * g + 0.0722 * b) / 255 < 0.5
    except Exception:
        return False


def stamp_data_coverage(
    pptx_path: str,
    label: str,
    theme: Optional[Dict[str, Any]] = None,
    logger_=None,
) -> Dict[str, Any]:
    """Write "Data through <label>" onto the cover slide, deterministically.

    Runs on the SAVED deck, after the theme passes and before previews — same
    slot and same contract as the furniture pass: never raises, and a deck that
    built must ship whether or not it could be stamped.
    """
    log = logger_ or logger
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Emu, Pt

        prs = Presentation(str(pptx_path))
        if not prs.slides:
            return {"status": "skipped", "reason": "deck has no slides"}
        cover = prs.slides[0]

        dark = _cover_is_dark(cover, theme)
        color = RGBColor(0xCB, 0xD5, 0xE1) if dark else RGBColor(0x64, 0x74, 0x8B)

        slide_w = prs.slide_width or Emu(12192000)
        slide_h = prs.slide_height or Emu(6858000)
        margin = int(slide_w * 0.034)          # ≈0.45in on a 13.33in canvas
        box_h = int(slide_h * 0.045)
        box = cover.shapes.add_textbox(
            margin, slide_h - box_h - int(slide_h * 0.012),
            int(slide_w * 0.4), box_h,
        )
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"Data through {label}"
        p.font.size = Pt(9)
        p.font.color.rgb = color

        prs.save(str(pptx_path))
        return {"status": "stamped", "label": label, "dark_cover": dark}
    except Exception as err:
        log.warning("data coverage stamp failed open: %s", err)
        return {"status": "unavailable", "reason": str(err)}
