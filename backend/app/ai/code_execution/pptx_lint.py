"""
Layout check for generated decks.

python-pptx has no font metrics. Generated deck code therefore sizes every text
box by guess, and a guess that is wrong by 3 inches raises nothing: the code
runs, the file writes, `pptx_success` is True, and the deck ships with a
paragraph sitting on top of the card below it. The existing retry loop in
`pptx_executor` only fires on exceptions, so it never sees this class of bug.

The only way to catch it is to lay the text out. OfficeCLI renders a .pptx to
HTML carrying real per-shape geometry (`data-path`, absolute pt coordinates);
Chromium — already in the image for dashboard preflight and PDF export — then
reports what each shape actually measures.

★ OfficeCLI's own `view <file> issues` does NOT report pptx text overflow. It
emits text boxes as `height:auto; min-height:<declared>`, so in its model
nothing ever overflows. It is accurate on off-slide geometry and we keep that,
but the overflow verdict has to be computed from the rendered box.

★ The overflowing box and the shape it lands on usually do NOT intersect as
rectangles — the generator leaves a correct gap between the two declared boxes
and only the rendered text crosses it. A rect-intersection test on the DECLARED
boxes finds nothing here. The test has to be run against the GROWN rect.

★ Growth on its own is not a fault. Measured on a real board deck, correct text
boxes render 1.4-1.5x taller than declared and nothing is wrong: the text grew
into empty space inside a larger white card. The height the browser reports is
not the height PowerPoint will lay out — the deck's fonts are not installed in
the container, so Chromium substitutes, wraps wider and overstates. That error
is systematic and it is not removable from here. So a height ratio alone can
only ever be a noise gate; the verdict has to come from asking whether the
overflow HIT anything.

Reports only. Never raises to the caller, never edits a deck. A missing binary,
a crash, or a timeout yields an empty list and the deck ships exactly as before.
"""

import json
import shutil
import asyncio
import logging
import tempfile
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

OFFICECLI_BIN = "officecli"

# Per-slide budgets. A deck that trips these is already broken; spending more
# wall-clock confirming it helps nobody.
RENDER_TIMEOUT_S = 30
MEASURE_TIMEOUT_MS = 20_000
MAX_SLIDES = 40

# Noise gate on the height ratio. NOT the verdict — a shape past this gate is
# only a candidate, and still has to collide with something to be reported.
#
# Measured over 222 text shapes on 18 slides of four decks (one real board deck,
# three generated fixtures): every visually-correct box came in at or below
# 1.52x, and the one genuine overflow came in at 5.16x. The benign ceiling is
# set by font substitution in the container, not by the deck. 2.00 leaves 32%
# headroom over the worst benign case measured and still sits 2.6x below the
# real failure. The old 1.25 was tuned against the synthetic fixture alone
# (which overshoots 5.2x) and fired on every one of those benign 1.4-1.5x boxes.
GROWTH_RATIO = 2.00
GROWTH_SLACK_PX = 4.0

# Off-slide tolerance, in rendered px. The page is measured at a pinned 1600px
# viewport (see check_deck_layout), so a px figure here is a fixed fraction of
# the slide: 24px is 0.2in of a 16:9 slide.
#
# Measured: the largest benign overshoot across those 18 slides is 5.96px — the
# bottom edge of a correct text box, inflated by the same font substitution —
# and the smallest genuine one is 62.38px. 24.0 sits 4x above the noise and
# 2.6x below the real fault. The old 2.0 reported that 5.96px as an issue.
EDGE_TOLERANCE_PX = 24.0

# A shape enclosing the grown rect this loosely is its container, not something
# the growth crashed into. Same slack used to recognise a label that already sat
# inside the declared box.
CONTAINER_SLACK_PX = 4.0

# An intersection thinner than this in either axis is two boxes touching, not
# text landing on text.
COLLISION_MIN_PX = 2.0

# Shape geometry is emitted at 960pt of slide width; the rendered page is scaled.
SLIDE_WIDTH_PT = 960.0

# Collected from the rendered page. Kept in one place so the measurement and the
# analysis cannot drift apart.
_MEASURE_JS = """() => {
  // ★ The HTML is a viewer page, not a bare slide: it ships a sidebar of
  // thumbnails, and each thumbnail contains a full COPY of every shape with the
  // same data-path, collapsed to zero size. Querying `.shape` document-wide
  // returns each shape twice, smallest first — measure the wrong copy and a
  // 557px overflow reads as 0px. Scope to the real slide.
  const slides = Array.from(document.querySelectorAll('.slide'))
    .filter(el => !el.closest('.thumb') && !el.closest('.sidebar'));
  const slide = slides[slides.length - 1] || document.body;
  const sb = slide.getBoundingClientRect();
  const shapes = [];
  slide.querySelectorAll('.shape').forEach(el => {
    const r = el.getBoundingClientRect();
    const raw = el.style.minHeight || el.style.height || '0';
    shapes.push({
      path: el.getAttribute('data-path'),
      fill: el.classList.contains('has-fill'),
      text: (el.innerText || '').trim(),
      left: r.left - sb.left,
      top: r.top - sb.top,
      w: r.width,
      h: r.height,
      declared_pt: parseFloat(raw) || 0,
    });
  });
  return {slide: {w: sb.width, h: sb.height}, shapes};
}"""


def officecli_available() -> bool:
    return shutil.which(OFFICECLI_BIN) is not None


def _render_slide_html(pptx_path: Path, page: int, out_dir: Path) -> Optional[Path]:
    """Render one slide to HTML. Returns None on any failure — never raises."""
    out = out_dir / f"slide_{page}.html"
    try:
        proc = subprocess.run(
            [OFFICECLI_BIN, "view", str(pptx_path), "html", "--page", str(page)],
            capture_output=True,
            timeout=RENDER_TIMEOUT_S,
        )
    except (OSError, subprocess.SubprocessError) as e:
        logger.info("pptx_lint: could not render slide %s of %s: %s", page, pptx_path.name, e)
        return None
    if proc.returncode != 0 or not proc.stdout:
        logger.info(
            "pptx_lint: officecli returned %s for slide %s of %s",
            proc.returncode, page, pptx_path.name,
        )
        return None
    out.write_bytes(proc.stdout)
    return out


def _slide_count(pptx_path: Path) -> int:
    try:
        from pptx import Presentation
        return len(Presentation(str(pptx_path)).slides)
    except Exception as e:
        logger.info("pptx_lint: could not read slide count from %s: %s", pptx_path.name, e)
        return 0


def _rect(s: Dict[str, Any], height: Optional[float] = None):
    """(x1, y1, x2, y2). Pass `height` to get the box the generator DECLARED
    rather than the one the browser rendered."""
    h = s["h"] if height is None else height
    return (s["left"], s["top"], s["left"] + s["w"], s["top"] + h)


def _intersection(a, b) -> Optional[tuple]:
    """Overlap of two rects as (width, height), or None if they miss."""
    w = min(a[2], b[2]) - max(a[0], b[0])
    h = min(a[3], b[3]) - max(a[1], b[1])
    return (w, h) if w > 0 and h > 0 else None


def _contains(outer, inner, slack: float = 0.0) -> bool:
    return (
        outer[0] <= inner[0] + slack
        and outer[1] <= inner[1] + slack
        and outer[2] >= inner[2] - slack
        and outer[3] >= inner[3] - slack
    )


def _analyse(slide_no: int, data: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Turn one slide's measurements into issues. Pure — no I/O, so it is the
    part the tests pin down."""
    sw = data["slide"]["w"]
    sh = data["slide"]["h"]
    issues: List[Dict[str, Any]] = []

    # Belt and braces against the thumbnail copies the selector already excludes,
    # and against the viewer's own furniture (slide-number and date placeholders
    # carry no data-path). A phantom shape here becomes a phantom issue, and one
    # false alarm is enough for someone to switch the whole check off.
    shapes = [
        s for s in data["shapes"]
        if s.get("path") and s["w"] > 0 and s["h"] > 0
    ]
    by_path: Dict[str, Dict[str, Any]] = {}
    for s in shapes:
        prev = by_path.get(s["path"])
        if prev is None or s["w"] * s["h"] > prev["w"] * prev["h"]:
            by_path[s["path"]] = s
    shapes = list(by_path.values())

    for s in shapes:
        x1, y1 = s["left"], s["top"]
        x2, y2 = x1 + s["w"], y1 + s["h"]
        for edge, over in (
            ("right", x2 - sw), ("bottom", y2 - sh), ("left", -x1), ("top", -y1),
        ):
            if over > EDGE_TOLERANCE_PX:
                issues.append({
                    "kind": "off_slide",
                    "slide": slide_no,
                    "path": s["path"],
                    "detail": f"extends {over:.0f}px past the slide {edge} edge",
                    "text": s["text"][:60],
                })

    off_slide_paths = {i["path"] for i in issues if i["kind"] == "off_slide"}

    # A box that renders taller than it was given is only a fault if the extra
    # height LANDED on something. Growth into empty space, or growth that stays
    # inside a bigger card, is what a deck looks like when it is correct.
    scale = (sw / SLIDE_WIDTH_PT) if sw else 0.0
    for s in shapes:
        if s["fill"] or not s["text"] or s["declared_pt"] <= 0 or scale <= 0:
            continue
        allotted = s["declared_pt"] * scale
        if s["h"] <= allotted * GROWTH_RATIO + GROWTH_SLACK_PX:
            continue

        grown = _rect(s)
        declared = _rect(s, allotted)

        # Inside a larger container/background shape: the growth had room.
        if any(
            _contains(_rect(o), grown, CONTAINER_SLACK_PX)
            for o in shapes if o["path"] != s["path"]
        ):
            continue

        hit = None
        for o in shapes:
            if o["path"] == s["path"] or not o["text"].strip():
                continue
            other = _rect(o)
            # A shape enclosing the grown rect is its container; a shape that
            # already sat inside the DECLARED box is a nested label. Neither is
            # something the growth crashed into.
            if _contains(other, grown, CONTAINER_SLACK_PX):
                continue
            if _contains(declared, other, CONTAINER_SLACK_PX):
                continue
            overlap = _intersection(grown, other)
            if overlap is None or min(overlap) <= COLLISION_MIN_PX:
                continue
            # Overlapping the declared box too means the two boxes were laid out
            # on top of each other to begin with — not caused by the growth.
            if _intersection(declared, other) is not None:
                continue
            hit = o
            break

        if hit is not None:
            landed = f"lands on {hit['text'][:40]!r}"
        elif not _contains((0.0, 0.0, sw, sh), grown, EDGE_TOLERANCE_PX) and (
            _contains((0.0, 0.0, sw, sh), declared, EDGE_TOLERANCE_PX)
        ):
            # The growth is what pushed the shape off the slide. `off_slide`
            # measures the same rect against the same tolerance, so it has
            # normally reported this already — defer to it rather than raise a
            # second issue for one physical fault.
            if s["path"] in off_slide_paths:
                continue
            landed = "pushes the shape off the slide"
        else:
            continue

        issues.append({
            "kind": "grew_past_box",
            "slide": slide_no,
            "path": s["path"],
            "detail": (
                f"text renders {s['h']:.0f}px tall in a box given "
                f"{allotted:.0f}px ({s['h'] / allotted:.1f}x) and {landed}"
            ),
            "text": s["text"][:60],
        })

    return issues


@dataclass
class LayoutCheckResult:
    """Result of a layout check. Populated even when nothing was measured — that
    is the point of this type: check_deck_layout()'s bare `[]` cannot tell "the
    deck is clean" apart from "the check never ran", and that ambiguity is a bug
    once this check gates anything. See check_deck_layout_detailed().

    status:
      "checked"     - every slide in scope was measured. reason is None.
      "partial"     - some slides measured, but not all — including hitting the
                       MAX_SLIDES cap. reason says how many and why.
      "unavailable" - nothing was measured: missing file, officecli absent,
                       playwright absent, chromium launch failed, slide count
                       unreadable, every slide failed to render, every rendered
                       slide failed to measure, or the broad outer exception
                       fired. slides_measured == 0. reason says which.
    """
    status: str
    reason: Optional[str]
    slides_total: int
    slides_measured: int
    issues: List[Dict[str, Any]]


async def check_deck_layout_detailed(
    pptx_path: Path,
    log: Optional[logging.Logger] = None,
) -> LayoutCheckResult:
    """Measure a saved .pptx and report layout issues plus WHY the check did or
    didn't run — see LayoutCheckResult. Never raises: this is advisory, so any
    failure downgrades to status="unavailable" rather than propagating."""
    lg = log or logger
    path = Path(pptx_path)

    def _unavailable(reason: str, slides_total: int = 0) -> LayoutCheckResult:
        return LayoutCheckResult(
            status="unavailable", reason=reason,
            slides_total=slides_total, slides_measured=0, issues=[],
        )

    if not path.is_file():
        return _unavailable(f"file not found: {path}")
    if not officecli_available():
        lg.info("pptx_lint: officecli not on PATH; skipping layout check")
        return _unavailable("officecli not on PATH")

    slides_total = _slide_count(path)
    if slides_total <= 0:
        return _unavailable(f"could not read slide count from {path.name}")

    slides_in_scope = slides_total
    if slides_total > MAX_SLIDES:
        lg.info("pptx_lint: %s has %s slides; checking the first %s",
                path.name, slides_total, MAX_SLIDES)
        slides_in_scope = MAX_SLIDES

    try:
        from playwright.async_api import async_playwright
    except ImportError:
        lg.info("pptx_lint: playwright unavailable; skipping layout check")
        return _unavailable("playwright not installed", slides_total)

    issues: List[Dict[str, Any]] = []
    slides_measured = 0
    try:
        with tempfile.TemporaryDirectory(prefix="bow_pptx_lint_") as tmp:
            tmp_dir = Path(tmp)
            rendered = {}
            for n in range(1, slides_in_scope + 1):
                html = await asyncio.to_thread(_render_slide_html, path, n, tmp_dir)
                if html is not None:
                    rendered[n] = html
            if not rendered:
                return _unavailable("officecli failed to render every slide", slides_total)

            async with async_playwright() as p:
                browser = await p.chromium.launch(args=["--no-sandbox"])
                try:
                    page = await browser.new_page(
                        viewport={"width": 1600, "height": 1200}
                    )
                    for n, html in rendered.items():
                        try:
                            await page.goto(
                                html.as_uri(),
                                wait_until="networkidle",
                                timeout=MEASURE_TIMEOUT_MS,
                            )
                            data = await page.evaluate(_MEASURE_JS)
                        except Exception as e:
                            lg.info("pptx_lint: slide %s did not measure: %s", n, e)
                            continue
                        slides_measured += 1
                        issues.extend(_analyse(n, data))
                finally:
                    await browser.close()
    except Exception as e:
        # Deliberately broad: this is advisory. Whatever went wrong here, the
        # deck itself is already written and must still be delivered.
        lg.warning("pptx_lint: layout check failed for %s: %s", path.name, e)
        return _unavailable(f"layout check raised: {e}", slides_total)

    if issues:
        lg.info("pptx_lint: %s layout issue(s) in %s", len(issues), path.name)

    if slides_measured == 0:
        # Every slide rendered but none could be measured (goto/evaluate failed
        # on all of them) — still "nothing ran", same bucket as the cases above.
        return _unavailable(
            f"rendered {len(rendered)} slide(s) but none could be measured",
            slides_total,
        )

    if slides_measured < slides_total:
        # Two distinct causes can land here, either or both: the MAX_SLIDES cap
        # (slides_in_scope < slides_total), and slides inside the scope that
        # individually failed to render or measure.
        parts = []
        if slides_total > MAX_SLIDES:
            parts.append(f"capped at {MAX_SLIDES} slides")
        failed_in_scope = slides_in_scope - slides_measured
        if failed_in_scope > 0:
            parts.append(f"{failed_in_scope} slide(s) failed to render or measure")
        detail = "; ".join(parts) if parts else "not all slides could be measured"
        return LayoutCheckResult(
            status="partial",
            reason=f"measured {slides_measured} of {slides_total} slides ({detail})",
            slides_total=slides_total,
            slides_measured=slides_measured,
            issues=issues,
        )

    return LayoutCheckResult(
        status="checked", reason=None,
        slides_total=slides_total, slides_measured=slides_measured, issues=issues,
    )


async def check_deck_layout(
    pptx_path: Path,
    log: Optional[logging.Logger] = None,
) -> List[Dict[str, Any]]:
    """Measure a saved .pptx and return layout issues. Always returns a list.

    Thin wrapper over check_deck_layout_detailed() — unchanged public behaviour
    for existing callers/tests. Prefer the detailed entry point for anything
    that needs to tell "clean" apart from "the check never ran".
    """
    result = await check_deck_layout_detailed(pptx_path, log)
    return result.issues
