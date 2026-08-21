"""
Secure executor for python-pptx code generation.

Allows LLM-generated python-pptx code to run in a sandboxed environment
with security validation reused from code_execution.py.
"""

import io
import ast
import tempfile
import subprocess
from pathlib import Path
from contextlib import redirect_stdout
from typing import Dict, Any, List, Tuple, Optional

# python-pptx imports
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData, ChartData

from app.ai.code_execution.code_execution import (
    CodeSecurityError,
    UnsafePythonError,
    FORBIDDEN_MODULES,
    FORBIDDEN_BUILTINS,
    FORBIDDEN_ATTRIBUTES,
)


# =============================================================================
# PPTX-Specific Security Validation
# =============================================================================

# Modules allowed for PPTX generation (extend the forbidden list with exceptions)
PPTX_ALLOWED_MODULES = frozenset({
    'pptx',
})


class PptxSecurityVisitor(ast.NodeVisitor):
    """AST visitor that checks for dangerous code patterns, allowing pptx imports."""

    def __init__(self):
        self.errors: List[str] = []

    def visit_Import(self, node: ast.Import):
        for alias in node.names:
            module_name = alias.name.split('.')[0]
            # Allow pptx imports, block everything else that's forbidden
            if module_name not in PPTX_ALLOWED_MODULES and module_name in FORBIDDEN_MODULES:
                self.errors.append(f"Forbidden import: '{alias.name}'")
        self.generic_visit(node)

    def visit_ImportFrom(self, node: ast.ImportFrom):
        if node.module:
            module_name = node.module.split('.')[0]
            # Allow pptx imports, block everything else that's forbidden
            if module_name not in PPTX_ALLOWED_MODULES and module_name in FORBIDDEN_MODULES:
                self.errors.append(f"Forbidden import: 'from {node.module}'")
        self.generic_visit(node)

    def visit_Call(self, node: ast.Call):
        # Check for forbidden built-in calls like eval(), exec(), open()
        if isinstance(node.func, ast.Name):
            if node.func.id in FORBIDDEN_BUILTINS:
                self.errors.append(f"Forbidden function call: '{node.func.id}()'")

        # Check for __import__('os') style calls
        if isinstance(node.func, ast.Name) and node.func.id == '__import__':
            self.errors.append("Forbidden function call: '__import__()'")

        self.generic_visit(node)

    def visit_Attribute(self, node: ast.Attribute):
        # Check for direct access to forbidden attributes like obj.__class__
        if node.attr in FORBIDDEN_ATTRIBUTES:
            self.errors.append(f"Forbidden attribute access: '{node.attr}'")
        self.generic_visit(node)


import re as _re

# python-pptx's Chart object exposes no `.plot_area` or `.chart_area` attribute.
# LLM-generated deck code frequently hallucinates them (e.g.
# `chart.plot_area.format.fill.solid()`), which raises AttributeError at exec
# time and fails the whole slide-generation run. Neutralize those statement
# lines so the rest of the deck still builds.
_INVALID_CHART_ATTR = _re.compile(r'^(\s*)\S.*\.(?:plot_area|chart_area)\b')


def sanitize_pptx_code(code: str) -> str:
    """Replace statements touching non-existent python-pptx chart attributes
    (`.plot_area` / `.chart_area`) with a no-op so generation doesn't crash."""
    out = []
    for line in code.splitlines():
        m = _INVALID_CHART_ATTR.match(line)
        if m:
            out.append(f"{m.group(1)}pass  # neutralized: python-pptx Chart has no .plot_area/.chart_area")
        else:
            out.append(line)
    return "\n".join(out)


def normalize_chart_axes(pptx_path: Path, logger=None) -> int:
    """Give every chart in a deck the same axis numbers the browser shows.

    Generated deck code hands python-pptx raw values and python-pptx writes
    them out in full, so a chart that reads `4.3B` on screen printed
    `70000000000` in the exported file. This walks the saved deck, reads each
    chart's own plotted values, and applies the shared abbreviation from
    `app.services.number_format` to its value axis — the same rule, from the
    same definition, as the browser and the Word export.

    Category labels are also made unambiguous: two categories that render the
    same string are two different things presented as one.

    Returns the number of charts changed. Never raises — a deck that could not
    be post-processed is still a valid deck.
    """
    from app.services.number_format import (
        pptx_axis_number_format,
        qualify_duplicate_labels,
    )

    changed = 0
    try:
        prs = Presentation(str(pptx_path))
    except Exception:
        if logger:
            logger.warning("pptx axis normalize: could not reopen %s", pptx_path)
        return 0

    dirty = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_chart", False):
                continue
            chart = shape.chart
            touched = False

            # Value axis: abbreviate to the magnitude of the largest plotted
            # point, so the axis reads the same as it does on screen.
            magnitudes = []
            try:
                for plot in chart.plots:
                    for series in plot.series:
                        for v in series.values:
                            if isinstance(v, (int, float)):
                                magnitudes.append(abs(v))
            except Exception:
                magnitudes = []
            if magnitudes:
                try:
                    axis = chart.value_axis
                    axis.tick_labels.number_format = pptx_axis_number_format(max(magnitudes))
                    axis.tick_labels.number_format_is_linked = False
                    touched = True
                except (ValueError, NotImplementedError, AttributeError):
                    # Pie/doughnut and some plot types have no value axis.
                    pass

            # Category axis: qualify labels that repeat.
            try:
                for plot in chart.plots:
                    labels = list(plot.categories)
                    flat = [
                        "" if l is None else str(getattr(l, "label", l))
                        for l in labels
                    ]
                    if len(set(flat)) != len(flat):
                        qualified = qualify_duplicate_labels(flat)
                        _rewrite_category_labels(plot, qualified)
                        touched = True
            except Exception:
                pass

            if touched:
                changed += 1
                dirty = True

    if dirty:
        try:
            prs.save(str(pptx_path))
        except Exception:
            if logger:
                logger.warning("pptx axis normalize: could not save %s", pptx_path)
            return 0
    return changed


# =============================================================================
# Deck palette — a property of the product, not of the generating code
# =============================================================================

_HEX_RE = _re.compile(r'^#?([0-9A-Fa-f]{3}|[0-9A-Fa-f]{6})$')

# How close an off-palette colour has to be before it is treated as a near-miss
# of a palette role rather than as a deliberate second colour. Euclidean
# distance in RGB; 32 is roughly "the same swatch, rounded differently" and
# nowhere near "a different hue". Anything further away is REPORTED, never
# guessed at — a deck that quietly recolours a deliberate red to the brand blue
# is worse than one that tells you the red is there.
_NEAR_ENOUGH = 32.0


def _normalize_hex(value: Any) -> Optional[str]:
    """`'#1f4e79'` / `'1F4E79'` / `'#abc'` -> `'1F4E79'`. Anything else -> None."""
    if not isinstance(value, str):
        return None
    m = _HEX_RE.match(value.strip())
    if not m:
        return None
    digits = m.group(1).upper()
    if len(digits) == 3:
        digits = "".join(c * 2 for c in digits)
    return digits


def palette_from_theme(theme: Optional[Dict[str, Any]]) -> Dict[str, str]:
    """Collect every colour a theme names, as `{HEX: role}`.

    A theme is a plain dict handed to generated code, so it carries more than
    colours (fonts, sizes, names). Only values that parse as a hex colour are
    taken; a list value contributes each of its hex members (a series palette),
    keyed `ROLE[i]`. Non-colour entries are ignored rather than rejected —
    nothing here should be able to fail a deck.
    """
    palette: Dict[str, str] = {}
    if not isinstance(theme, dict):
        return palette
    for role, value in theme.items():
        hexed = _normalize_hex(value)
        if hexed is not None:
            palette.setdefault(hexed, str(role))
            continue
        if isinstance(value, (list, tuple)):
            for i, item in enumerate(value):
                hexed = _normalize_hex(item)
                if hexed is not None:
                    palette.setdefault(hexed, f"{role}[{i}]")
    return palette


def _distance(a: str, b: str) -> float:
    ar, ag, ab = int(a[0:2], 16), int(a[2:4], 16), int(a[4:6], 16)
    br, bg, bb = int(b[0:2], 16), int(b[2:4], 16), int(b[4:6], 16)
    return ((ar - br) ** 2 + (ag - bg) ** 2 + (ab - bb) ** 2) ** 0.5


def _nearest_palette_role(found: str, palette: Dict[str, str]) -> Tuple[Optional[str], float]:
    """Closest palette colour to `found`, with its distance. `(None, inf)` if empty."""
    best: Optional[str] = None
    best_d = float("inf")
    for candidate in palette:
        d = _distance(found, candidate)
        if d < best_d:
            best, best_d = candidate, d
    return best, best_d


def _color_slots(shape) -> List[Tuple[str, Any]]:
    """Every readable colour on one shape, as `(hex, setter)`.

    Each entry's setter takes an `RGBColor`. Reading is eager and each read is
    guarded on its own: python-pptx raises on a fill that is not solid, on a
    theme-indexed colour, and on shape types that have no line at all, and none
    of those are errors — they are simply colours this pass cannot speak about.
    """
    slots: List[Tuple[str, Any]] = []

    def _add(color_obj):
        try:
            current = str(color_obj.rgb).upper()
        except Exception:
            return
        if _normalize_hex(current) is None:
            return
        slots.append((current, lambda v, c=color_obj: setattr(c, "rgb", v)))

    try:
        _add(shape.fill.fore_color)
    except Exception:
        pass

    try:
        _add(shape.line.color)
    except Exception:
        pass

    try:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    try:
                        _add(run.font.color)
                    except Exception:
                        pass
    except Exception:
        pass

    try:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    try:
                        _add(cell.fill.fore_color)
                    except Exception:
                        pass
                    try:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                _add(run.font.color)
                    except Exception:
                        pass
    except Exception:
        pass

    try:
        if getattr(shape, "has_chart", False):
            for plot in shape.chart.plots:
                for series in plot.series:
                    try:
                        _add(series.format.fill.fore_color)
                    except Exception:
                        pass
                    try:
                        for point in series.points:
                            _add(point.format.fill.fore_color)
                    except Exception:
                        pass
    except Exception:
        pass

    # Groups nest; their children carry the colours a reader actually sees.
    try:
        for child in shape.shapes:
            slots.extend(_color_slots(child))
    except Exception:
        pass

    return slots


def apply_theme_palette(
    pptx_path: Path,
    theme: Optional[Dict[str, Any]] = None,
    logger=None,
) -> Dict[str, Any]:
    """Report — and where it is safe, correct — the colours a saved deck uses.

    A deck's palette is a property of the product, not of whatever the
    generating code happened to write. This walks the saved file the same way
    `normalize_chart_axes` does: reopen, walk every slide/shape/run/fill/line,
    mutate, save.

    With a `theme` palette supplied, a colour that is a near-miss of a palette
    role (within `_NEAR_ENOUGH`) is snapped onto that role. A colour that is far
    from every role is **reported and left alone** — this pass never guesses at
    a colour that was chosen deliberately, and it never refuses a deck for
    carrying one. `off_palette` is the signal that the generated code ignored
    the injected palette and hardcoded literals; it is a number to read, not a
    gate to fail.

    Returns:
        {
          "remapped": int,                  # colour slots snapped onto a role
          "off_palette": int,               # distinct colours not in the palette
          "off_palette_colors": [hex, ...], # sorted, deduped — what to go look at
          "colors_seen": [hex, ...],        # every distinct colour in the deck
          "palette": [hex, ...],            # the roles this pass compared against
        }

    Never raises — a deck that could not be post-processed is still a valid deck.
    """
    summary: Dict[str, Any] = {
        "remapped": 0,
        "off_palette": 0,
        "off_palette_colors": [],
        "colors_seen": [],
        "palette": [],
    }

    palette = palette_from_theme(theme)
    summary["palette"] = sorted(palette)

    try:
        prs = Presentation(str(pptx_path))
    except Exception:
        if logger:
            logger.warning("pptx palette: could not reopen %s", pptx_path)
        return summary

    seen: set = set()
    off: set = set()
    remapped = 0

    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                for found, setter in _color_slots(shape):
                    seen.add(found)
                    if not palette:
                        continue
                    if found in palette:
                        continue
                    nearest, distance = _nearest_palette_role(found, palette)
                    if nearest is not None and distance <= _NEAR_ENOUGH:
                        try:
                            setter(RGBColor.from_string(nearest))
                            remapped += 1
                            seen.add(nearest)
                            continue
                        except Exception:
                            pass
                    # Too far to be a rounding of a role, or not writable.
                    # Say so; do not invent an intent.
                    off.add(found)
    except Exception:
        if logger:
            logger.warning("pptx palette: walk stopped early on %s", pptx_path, exc_info=True)

    summary["remapped"] = remapped
    summary["colors_seen"] = sorted(seen)
    summary["off_palette_colors"] = sorted(off)
    summary["off_palette"] = len(off)

    if remapped:
        try:
            prs.save(str(pptx_path))
        except Exception:
            if logger:
                logger.warning("pptx palette: could not save %s", pptx_path)
            summary["remapped"] = 0

    return summary


# =============================================================================
# Theme prohibitions — the half of a design system the model ignores
# =============================================================================

# Every theme ships a "strictly avoid" list and, measured across live deck
# generations, the model honours the palette and the typography and ignores the
# prohibitions: drop shadows and rounded corners on every deck, every panel
# tinted instead of one. Prompt text did not fix it — it was already in the
# prompt. So the ones that are mechanical are applied to the saved file, the
# same way `normalize_chart_axes` and `apply_theme_palette` do it.

# Spelling drift in the theme files ("shadow", "drop shadows", "rounded") is
# not a different rule; it is the same rule written by hand.
_AVOID_ALIASES = {
    "shadow": "shadows",
    "shadows": "shadows",
    "drop_shadow": "shadows",
    "drop_shadows": "shadows",
    "box_shadow": "shadows",
    "box_shadows": "shadows",
    "rounded": "rounded_corners",
    "rounded_corner": "rounded_corners",
    "rounded_corners": "rounded_corners",
    "rounded_rectangles": "rounded_corners",
    "border_radius": "rounded_corners",
    "gradient": "gradients",
    "gradients": "gradients",
    "gradient_fill": "gradients",
    "gradient_fills": "gradients",
    "multiple_accent": "multiple_accents",
    "multiple_accents": "multiple_accents",
    "multi_accent": "multiple_accents",
    "multiple_accent_colors": "multiple_accents",
    "box": "boxes",
    "boxes": "boxes",
    "boxed": "boxes",
    "cards": "boxes",
    "legend": "legends",
    "legends": "legends",
    "chart_legends": "legends",
}

# Preset geometries whose corners are round. `rect` is the square counterpart
# for all of them; the adjustment list goes with the old geometry.
_ROUNDED_PRESETS = frozenset({
    "roundRect",
    "round1Rect",
    "round2DiagRect",
    "round2SameRect",
    "roundedRectangle",
})

# What this pass will actually change in the file. Everything else in the
# theme's avoid list is counted and handed back, never acted on.
_ENFORCEABLE = frozenset({"shadows", "rounded_corners", "gradients"})


def _avoid_tokens(theme: Optional[Dict[str, Any]]) -> set:
    """The prohibitions a theme states, normalised to canonical tokens.

    Only what the theme actually forbids. A theme with no `avoid` list forbids
    nothing, and this pass must leave such a deck byte-identical — squaring the
    corners of a theme that permits rounded corners is a new bug, not a fix.
    """
    tokens: set = set()
    if not isinstance(theme, dict):
        return tokens
    raw = theme.get("avoid")
    if isinstance(raw, str):
        raw = [raw]
    if not isinstance(raw, (list, tuple, set, frozenset)):
        return tokens
    for item in raw:
        if not isinstance(item, str):
            continue
        key = item.strip().lower().replace("-", "_").replace(" ", "_")
        if not key:
            continue
        tokens.add(_AVOID_ALIASES.get(key, key))
    return tokens


def _walk_shapes(shapes):
    """Every shape in a container, groups flattened — a group's children carry
    the shadows and the corners a reader actually sees."""
    for shape in shapes:
        yield shape
        try:
            children = shape.shapes
        except Exception:
            continue
        for nested in _walk_shapes(children):
            yield nested


def _sp_pr(shape):
    """A shape's `<p:spPr>`, or None for the shape types that have none."""
    try:
        return shape._element.spPr
    except Exception:
        return None


def _clear_shadow(shape) -> bool:
    """Turn one shape's shadow off. True if the file changed.

    Two things make a shadow: an explicit `<a:effectLst>` the generating code
    wrote, and — far more often — no effectLst at all, which means the shape
    inherits the theme's. An empty effectLst is the way to say "no effect,
    and do not inherit one", so both cases end at the same XML.
    """
    from pptx.oxml.ns import qn

    spPr = _sp_pr(shape)
    if spPr is None:
        return False
    try:
        effect_lst = spPr.find(qn("a:effectLst"))
        if effect_lst is not None and len(effect_lst) == 0:
            return False  # already explicitly shadow-free
        if effect_lst is not None:
            for child in list(effect_lst):
                effect_lst.remove(child)
            return True
        # No effectLst: the shape inherits. python-pptx writes the empty one.
        shape.shadow.inherit = False
        return spPr.find(qn("a:effectLst")) is not None
    except Exception:
        return False


def _has_shadow(shape) -> bool:
    """Whether this shape would render a shadow — explicit or inherited."""
    from pptx.oxml.ns import qn

    spPr = _sp_pr(shape)
    if spPr is None:
        return False
    try:
        effect_lst = spPr.find(qn("a:effectLst"))
    except Exception:
        return False
    if effect_lst is None:
        return True  # inherits the theme's
    return len(effect_lst) > 0


def _square_corners(shape) -> bool:
    """Swap a rounded preset geometry for `rect`. True if the file changed."""
    from pptx.oxml.ns import qn

    spPr = _sp_pr(shape)
    if spPr is None:
        return False
    try:
        prst_geom = spPr.find(qn("a:prstGeom"))
        if prst_geom is None:
            return False
        if prst_geom.get("prst") not in _ROUNDED_PRESETS:
            return False
        prst_geom.set("prst", "rect")
        # The adjustment values belonged to the geometry that just left.
        av_lst = prst_geom.find(qn("a:avLst"))
        if av_lst is not None:
            for child in list(av_lst):
                av_lst.remove(child)
        return True
    except Exception:
        return False


def _is_rounded(shape) -> bool:
    from pptx.oxml.ns import qn

    spPr = _sp_pr(shape)
    if spPr is None:
        return False
    try:
        prst_geom = spPr.find(qn("a:prstGeom"))
    except Exception:
        return False
    return prst_geom is not None and prst_geom.get("prst") in _ROUNDED_PRESETS


def _gradient_stops(shape):
    """A shape's gradient stops, or None if its fill is not a gradient."""
    try:
        fill = shape.fill
    except Exception:
        return None
    try:
        from pptx.enum.dml import MSO_FILL

        if fill.type != MSO_FILL.GRADIENT:
            return None
    except Exception:
        return None
    try:
        return list(fill.gradient_stops)
    except Exception:
        return None


def _dominant_stop(stops):
    """The stop that covers the most of the gradient.

    Each stop owns the ground halfway to its neighbours; the one owning the
    widest band is the colour a reader would call the shape's colour. Ties and
    unreadable positions fall back to the first stop, which is what a flat
    reading of the fill would have given anyway.
    """
    if not stops:
        return None
    try:
        positioned = sorted(stops, key=lambda s: float(s.position))
    except Exception:
        return stops[0]
    if len(positioned) == 1:
        return positioned[0]
    best = positioned[0]
    best_span = -1.0
    last = len(positioned) - 1
    for i, stop in enumerate(positioned):
        try:
            pos = float(stop.position)
            # Ownership runs to the midpoint of each neighbour, and to the edge
            # of the shape where there is no neighbour.
            lower = 0.0 if i == 0 else (float(positioned[i - 1].position) + pos) / 2.0
            upper = 1.0 if i == last else (float(positioned[i + 1].position) + pos) / 2.0
        except Exception:
            continue
        span = upper - lower
        if span > best_span:
            best, best_span = stop, span
    return best


def _flatten_gradient(shape, stops) -> bool:
    """Repaint a gradient fill as its dominant stop. True if the file changed."""
    stop = _dominant_stop(stops)
    if stop is None:
        return False
    try:
        rgb = stop.color.rgb  # raises on a theme-indexed stop colour
    except Exception:
        return False
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
        return True
    except Exception:
        return False


def _is_card(shape) -> bool:
    """A filled panel carrying text — the "box" a stark theme forbids."""
    try:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            return False
    except Exception:
        return False
    try:
        from pptx.enum.dml import MSO_FILL

        return shape.fill.type in (MSO_FILL.SOLID, MSO_FILL.GRADIENT, MSO_FILL.PATTERNED)
    except Exception:
        return False


def _accent_colors(prs) -> set:
    """Distinct non-neutral colours in the deck.

    Neutrals — the greys, the near-white, the near-black — are the ground a
    stark deck is made of, not accents; counting them would report every deck
    as polychrome and the number would mean nothing.
    """
    accents: set = set()
    for slide in prs.slides:
        for shape in _walk_shapes(slide.shapes):
            for found, _setter in _color_slots(shape):
                try:
                    r, g, b = (
                        int(found[0:2], 16),
                        int(found[2:4], 16),
                        int(found[4:6], 16),
                    )
                except Exception:
                    continue
                if max(r, g, b) - min(r, g, b) < 24:
                    continue
                accents.add(found)
    return accents


def enforce_theme_rules(
    pptx_path: Path,
    theme: Optional[Dict[str, Any]] = None,
    logger=None,
) -> Dict[str, Any]:
    """Apply a theme's prohibitions to a saved deck.

    Every design system here ships a "strictly avoid" list, and across live
    generations the model honoured each theme's palette and typography and
    ignored its prohibitions every time. The ones that are mechanical —
    shadows, rounded corners, gradient fills — are therefore not asked for,
    they are applied to the file: reopen, walk every slide and shape, mutate,
    save, exactly as `normalize_chart_axes` and `apply_theme_palette` do.

    Only what the theme actually forbids is touched. A theme with no `avoid`
    list leaves the deck untouched.

    Two prohibitions are deliberately counted and not acted on:

    * `multiple_accents` — repainting here would fight `apply_theme_palette`,
      which owns colour and has its own rule about when a colour is a
      deliberate choice.
    * `boxes` — turning a filled card into a left-rule row is layout, not
      enforcement, and belongs where the slide is composed.

    Nothing here is a gate. A violation is a number to read; a deck is never
    rejected for carrying one.

    Returns:
        {
          "shadows_cleared": int,
          "corners_squared": int,
          "gradients_flattened": int,
          "violations": {token: count},   # what was found, per prohibition
          "reported_only": [token, ...],  # forbidden, counted, left alone
        }

    Never raises — a deck that could not be post-processed is still a valid deck.
    """
    summary: Dict[str, Any] = {
        "shadows_cleared": 0,
        "corners_squared": 0,
        "gradients_flattened": 0,
        "violations": {},
        "reported_only": [],
    }

    avoid = _avoid_tokens(theme)
    if not avoid:
        return summary

    summary["reported_only"] = sorted(avoid - _ENFORCEABLE)

    try:
        prs = Presentation(str(pptx_path))
    except Exception:
        if logger:
            logger.warning("pptx theme rules: could not reopen %s", pptx_path)
        return summary

    violations = {token: 0 for token in avoid}
    shadows_cleared = 0
    corners_squared = 0
    gradients_flattened = 0

    try:
        for slide in prs.slides:
            for shape in _walk_shapes(slide.shapes):
                if "shadows" in avoid:
                    if _has_shadow(shape):
                        violations["shadows"] += 1
                        if _clear_shadow(shape):
                            shadows_cleared += 1

                if "rounded_corners" in avoid and _is_rounded(shape):
                    violations["rounded_corners"] += 1
                    if _square_corners(shape):
                        corners_squared += 1

                if "gradients" in avoid:
                    stops = _gradient_stops(shape)
                    if stops:
                        violations["gradients"] += 1
                        if _flatten_gradient(shape, stops):
                            gradients_flattened += 1

                if "boxes" in avoid and _is_card(shape):
                    violations["boxes"] += 1

                if "legends" in avoid and getattr(shape, "has_chart", False):
                    try:
                        if shape.chart.has_legend:
                            violations["legends"] += 1
                    except Exception:
                        pass

        if "multiple_accents" in avoid:
            accents = _accent_colors(prs)
            # One accent is the rule; each extra one is the violation.
            violations["multiple_accents"] = max(0, len(accents) - 1)
    except Exception:
        if logger:
            logger.warning(
                "pptx theme rules: walk stopped early on %s", pptx_path, exc_info=True
            )

    summary["violations"] = violations
    summary["shadows_cleared"] = shadows_cleared
    summary["corners_squared"] = corners_squared
    summary["gradients_flattened"] = gradients_flattened

    if shadows_cleared or corners_squared or gradients_flattened:
        try:
            prs.save(str(pptx_path))
        except Exception:
            if logger:
                logger.warning("pptx theme rules: could not save %s", pptx_path)
            # Nothing persisted, so nothing was done. Say that.
            summary["shadows_cleared"] = 0
            summary["corners_squared"] = 0
            summary["gradients_flattened"] = 0

    return summary


def _rewrite_category_labels(plot, labels: List[str]) -> None:
    """Replace a plot's category strings in the underlying chart XML.

    python-pptx exposes categories read-only, so the cached string points are
    edited directly. Guarded by the caller.
    """
    from pptx.oxml.ns import qn

    for ser in plot._element.iter(qn("c:ser")):
        # Scope to <c:cat> only. A series also carries a <c:strCache> for its
        # own NAME under <c:tx>, so an unscoped search returns one point too
        # many, the length check fails and the rewrite is skipped in silence —
        # which is what happened, and it took looking at a rendered slide to
        # notice.
        cat = ser.find(qn("c:cat"))
        if cat is None:
            continue
        pts = cat.findall(f'.//{qn("c:strCache")}/{qn("c:pt")}')
        if len(pts) != len(labels):
            continue
        for pt, text in zip(pts, labels):
            v = pt.find(qn("c:v"))
            if v is not None:
                v.text = text


def validate_pptx_code(code: str) -> None:
    """
    Validate Python code for security issues, allowing pptx imports.

    Raises:
        UnsafePythonError: If the code contains dangerous constructs.
    """
    try:
        tree = ast.parse(code)
    except SyntaxError:
        # Let syntax errors pass through - they'll fail at exec() time
        return

    visitor = PptxSecurityVisitor()
    visitor.visit(tree)

    if visitor.errors:
        raise UnsafePythonError(
            f"Code contains forbidden constructs: {'; '.join(visitor.errors)}"
        )


# =============================================================================
# PPTX Code Executor
# =============================================================================

def _make_image_opener(images: Dict[str, bytes]):
    """Build the `image(file_id)` helper handed to generated code.

    Returns a fresh BytesIO per call: python-pptx reads the stream it is given
    to exhaustion, so handing out one shared buffer would make the second
    placement of the same image silently render empty.
    """

    def image(file_id: str) -> io.BytesIO:
        raw = images.get(str(file_id))
        if raw is None:
            raise ValueError(
                f"Unknown image id {file_id!r}. Available: {sorted(images) or 'none'}"
            )
        return io.BytesIO(raw)

    return image


class PptxCodeExecutor:
    """
    Secure executor for python-pptx code generation.

    Reuses security patterns from StreamingCodeExecutor but with a namespace
    tailored for PPTX generation.
    """

    def __init__(self, logger=None):
        self.logger = logger

    def execute_pptx_code(
        self,
        *,
        code: str,
        visualizations: List[Dict[str, Any]],
        report: Dict[str, Any],
        output_path: Path,
        images: Optional[Dict[str, bytes]] = None,
        theme: Optional[Dict[str, Any]] = None,
    ) -> Tuple[Path, str]:
        """
        Execute python-pptx code and save the resulting presentation.

        Args:
            code: The python-pptx code to execute
            visualizations: List of visualization dicts with rows/columns
            report: Report info dict with id, title, theme
            output_path: Path where the PPTX file should be saved
            images: Raw bytes per embeddable file id. Resolved by the caller —
                generated code never touches the filesystem, so the sandbox
                keeps its no-IO guarantee while still being able to place art.
            theme: Optional palette handed to generated code as `theme` — a
                plain dict, so `theme['PRIMARY']` and `theme.get('ACCENT')` both
                work. When omitted the name is not bound at all and the
                namespace is exactly what it has always been.

        Returns:
            Tuple of (output_path, stdout_log)

        Raises:
            UnsafePythonError: If code contains forbidden imports, calls, or attributes
        """
        # Neutralize hallucinated python-pptx chart APIs (.plot_area/.chart_area)
        # that would otherwise crash the whole run with AttributeError.
        code = sanitize_pptx_code(code)

        # Security: Validate code before execution
        validate_pptx_code(code)

        output_log = ""

        # Build the namespace with pptx utilities and data
        local_namespace = {
            # python-pptx classes
            'Presentation': Presentation,
            'Inches': Inches,
            'Pt': Pt,
            'Emu': Emu,
            'RGBColor': RGBColor,
            'PP_ALIGN': PP_ALIGN,
            'MSO_ANCHOR': MSO_ANCHOR,
            'MSO_SHAPE': MSO_SHAPE,
            'XL_CHART_TYPE': XL_CHART_TYPE,
            'XL_LEGEND_POSITION': XL_LEGEND_POSITION,
            'CategoryChartData': CategoryChartData,
            'ChartData': ChartData,

            # Data access
            'visualizations': visualizations,
            'report': report,

            # Embeddable art. `image(file_id)` hands back a FRESH stream each
            # call — python-pptx consumes the stream it is given, so a shared
            # BytesIO would silently produce an empty picture the second time
            # the same image is placed.
            'image': _make_image_opener(images or {}),
            'image_ids': list((images or {}).keys()),

            # Output target (set by executor, not user code)
            '_pptx_output_path': str(output_path),
        }

        # The palette, when the caller has one. Bound only when supplied: a run
        # with no theme must see the namespace it has always seen.
        if theme is not None:
            local_namespace['theme'] = theme

        if self.logger:
            self.logger.debug(f"Executing PPTX code:\n{code}")

        # Capture via the per-thread stdout router (not redirect_stdout,
        # which swaps the process-global sys.stdout and cross-talks with
        # any concurrently-running sandboxed code execution).
        from app.ai.code_execution.code_execution import _stdout_router
        router = _stdout_router()
        stdout_capture = io.StringIO()
        router.bind(stdout_capture)
        try:
            exec(code, local_namespace)
            output_log = stdout_capture.getvalue()
        finally:
            router.unbind()
            stdout_capture.close()

        # Verify the file was created
        if not output_path.exists():
            raise RuntimeError(
                f"PPTX code executed but no file was created at {output_path}. "
                "Ensure the code calls prs.save(_pptx_output_path)"
            )

        # Axis numbers and category labels are a property of the product, not
        # of whatever the generating code happened to write. Applied after the
        # save so it covers every chart in the deck regardless of how it was
        # built. Never fatal — see normalize_chart_axes.
        try:
            normalize_chart_axes(output_path, logger=self.logger)
        except Exception:
            if self.logger:
                self.logger.warning("pptx axis normalize skipped", exc_info=True)

        return output_path, output_log


# =============================================================================
# PPTX to Image Preview Conversion
# =============================================================================

class PptxPreviewService:
    """
    Service for generating preview images from PPTX files.

    Uses LibreOffice (headless) to convert PPTX to PDF,
    then pdf2image to convert PDF pages to PNG images.
    """

    def __init__(self, preview_dir: Optional[Path] = None, logger=None):
        self.logger = logger
        if preview_dir:
            self.preview_dir = preview_dir
        else:
            # Default to uploads/pptx_previews relative to backend root
            backend_root = Path(__file__).parent.parent.parent.parent
            self.preview_dir = backend_root / "uploads" / "pptx_previews"

        # Ensure preview directory exists
        self.preview_dir.mkdir(parents=True, exist_ok=True)

    def generate_previews(
        self,
        pptx_path: Path,
        artifact_id: str,
        dpi: int = 220,
    ) -> List[str]:
        """
        Convert PPTX to PNG preview images.

        Args:
            pptx_path: Path to the PPTX file
            artifact_id: Artifact ID for organizing previews
            dpi: Resolution for preview images (default 220 — 150 produced
                 2000px-wide PNGs that retina displays upsample into a soft,
                 "blurry" preview of a perfectly sharp vector deck)

        Returns:
            List of relative paths to preview images (e.g., ["pptx_previews/{id}/slide-1.png", ...])
        """
        from pdf2image import convert_from_path
        from pdf2image.exceptions import PDFInfoNotInstalledError, PDFPageCountError

        # Create artifact-specific preview directory
        artifact_preview_dir = self.preview_dir / artifact_id
        artifact_preview_dir.mkdir(parents=True, exist_ok=True)

        # Step 1: Convert PPTX to PDF using LibreOffice
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)

            # LibreOffice convert to PDF
            try:
                result = subprocess.run(
                    [
                        'soffice',
                        '--headless',
                        '--convert-to', 'pdf',
                        '--outdir', str(tmp_path),
                        str(pptx_path),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=60,
                )
                if result.returncode != 0:
                    raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")
            except FileNotFoundError:
                raise RuntimeError(
                    "LibreOffice not found. Install with: apt-get install libreoffice-impress"
                )

            # Find the generated PDF
            pdf_files = list(tmp_path.glob("*.pdf"))
            if not pdf_files:
                raise RuntimeError("LibreOffice did not produce a PDF file")
            pdf_path = pdf_files[0]

            # Step 2: Convert PDF pages to PNG using pdf2image
            try:
                images = convert_from_path(pdf_path, dpi=dpi)
            except PDFInfoNotInstalledError:
                raise RuntimeError(
                    "poppler not found. Install with: apt-get install poppler-utils"
                )
            except PDFPageCountError as e:
                raise RuntimeError(f"Failed to read PDF: {e}")

            # Save images
            for i, image in enumerate(images):
                image_path = artifact_preview_dir / f"slide-{i + 1:02d}.png"
                image.save(str(image_path), "PNG")

        # Collect generated image paths
        preview_images = sorted(artifact_preview_dir.glob("slide-*.png"))

        # Return relative paths from uploads directory
        relative_paths = [
            f"pptx_previews/{artifact_id}/{img.name}" for img in preview_images
        ]

        if self.logger:
            self.logger.info(f"Generated {len(relative_paths)} preview images for artifact {artifact_id}")

        return relative_paths

    def get_preview_paths(self, artifact_id: str) -> List[str]:
        """Get existing preview image paths for an artifact."""
        artifact_preview_dir = self.preview_dir / artifact_id
        if not artifact_preview_dir.exists():
            return []

        preview_images = sorted(artifact_preview_dir.glob("slide-*.png"))
        return [
            f"pptx_previews/{artifact_id}/{img.name}" for img in preview_images
        ]

    def cleanup_previews(self, artifact_id: str) -> None:
        """Remove all preview images for an artifact."""
        artifact_preview_dir = self.preview_dir / artifact_id
        if artifact_preview_dir.exists():
            import shutil
            shutil.rmtree(artifact_preview_dir)
