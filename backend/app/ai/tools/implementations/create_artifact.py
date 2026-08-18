import asyncio
import base64
import json
import logging
import os
import re
import time
from pathlib import Path
from typing import AsyncIterator, Dict, Any, Type, List, Optional

import aiofiles
from pydantic import BaseModel
from sqlalchemy import select
from sqlalchemy.orm import selectinload

from app.ai.tools.base import Tool
from app.models.file import File
from app.models.report_file_association import report_file_association
from app.core.feature_flags import setting_enabled

logger = logging.getLogger(__name__)


# DEF-008: a model can answer a "build the dashboard" prompt with a sentence
# describing what it is about to do, and stop. `slides` mode already guards this
# (it compiles the python and retries on SyntaxError — added after grok-4.5 was
# seen doing exactly that), but `page` mode had no equivalent, so the prose went
# on to be wrapped in `function App() { … }` and shipped as a "dashboard" that
# only the browser could reject.
#
# JSX cannot be checked with `compile()`, and this image has no node, no npx and
# no Python JS parser — so a real Babel parse is not available here. This is a
# STRUCTURAL check instead: does the reply contain any of the things every React
# component must have? It answers "is this code at all", not "is this code
# valid". Syntactically broken JSX still gets through; prose does not.
#
# ★It used to be a substring scan over a bare marker list, and three of those
# markers match ordinary English: "I'll **return** the top 5 banners", "revenue
# **<** 1M", "create a **function that** aggregates". So the gate the whole
# DEF-008 recovery hangs on accepted the exact replies it exists to catch, and
# the strict retry never fired for them. A scan for words over free text cannot
# separate code from a sentence about code; these look at SHAPE instead.

# A JSX tag: `<` immediately followed by a tag name, then whitespace, `/` or
# `>`. `< 1M` and `a < b` do not match; `<div>`, `</Chart>`, `<Bar />` do.
_JSX_TAG = re.compile(r"</?[A-Za-z][\w.$-]*(?:\s|/|>)")

# A keyword standing alone as a token — not `functional`, not `returns`.
_CODE_KEYWORD = re.compile(
    r"(?<![\w$])(?:const|let|var|function|class|return|import|export|async|await)(?![\w$])"
)

# Punctuation that carries structure. Deliberately excludes `(` and `)`, which
# appear in ordinary prose ("two KPI cards (total and average)") far too often
# to mean anything on their own.
_CODE_STRUCTURE = re.compile(r"[{};]|=>")

# The data hook every generated dashboard calls, as a call rather than a
# mention of the name.
_DATA_HOOK = re.compile(r"useArtifactData\s*\(")


def _read_bool_setting(name: str, default: bool = True) -> bool:
    """Read a boolean runtime flag defensively.

    ★ This codebase has been bitten three times by `if not flag:` letting a deny
    state through — the string "off" is TRUTHY in Python. So type-check the value
    and fall back to the default explicitly rather than by truthiness accident.
    """
    try:
        from app.settings.config import settings as _settings  # lazy: no import cycle

        value = getattr(_settings, name, default)
    except Exception:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.strip().lower() in ("1", "true", "yes", "on")
    return default


def _completeness_gate_enabled() -> bool:
    """PHASE 1: refuse to build an artifact from truncated data. Default ON."""
    return _read_bool_setting("hybrid_artifact_completeness_gate", True)


def _render_preflight_enabled() -> bool:
    """PHASE 3: render the artifact headlessly before storing it. Default ON."""
    return _read_bool_setting("hybrid_artifact_render_preflight", True)


def _insights_enabled() -> bool:
    """PHASE 4: generate a grounded insight summary for each dashboard. Default ON."""
    return _read_bool_setting("hybrid_artifact_insights", True)


def _deck_layout_check_enabled() -> bool:
    """DEF-011: measure the saved deck's layout. Default OFF (settings/config.py).

    Reports only — it never blocks or rewrites a deck, and it runs after the
    .pptx has been written and its path recorded, so nothing it does can cost a
    deck that built fine.
    """
    return _read_bool_setting("hybrid_deck_layout_check", False)


def _deck_theme_furniture_enabled() -> bool:
    """Paint the theme's own furniture onto the saved deck. Default ON.

    Corrective, not measuring — unlike the layout check above it exists to make
    the deck right rather than to report on it, so it ships on. No entry in
    `config.py` is required: `_read_bool_setting` returns the default for a name
    it cannot find, so HYBRID_DECK_THEME_FURNITURE=false only bites once it is
    declared there.
    """
    return _read_bool_setting("hybrid_deck_theme_furniture", True)


def _deck_theme_enforcement_enabled() -> bool:
    """Enforce the theme's prohibitions on the saved deck. Default ON.

    Same reasoning as `_deck_theme_furniture_enabled`, and it must be able to
    run when furniture is off: the two are separate switches because a theme's
    prohibitions (no shadows, no rounded corners) are worth keeping even where
    an org does not want the decorative furniture painted.
    """
    return _read_bool_setting("hybrid_deck_theme_enforcement", True)


def _load_pptx_themes():
    """The deck theme registry, or None when it is not present.

    ★Imported defensively on purpose. The registry is a separate module; if it
    is absent — or raises while importing — this tool must build a deck exactly
    as it did before rather than fail. Hence `Exception`, not `ImportError`: a
    module that blows up at import time must not take deck generation with it.
    """
    try:
        from app.ai.decks import pptx_themes  # noqa: PLC0415 (optional dependency)

        return pptx_themes
    except Exception:  # pragma: no cover - exercised only when the module is absent
        return None


def _org_brand_hint(organization_settings: Any) -> Optional[str]:
    """A brand name/colour the theme registry can weigh, or None.

    ★`getattr` against a dict MISSES rather than raising, so the config blob is
    read both ways — the same bug class that shipped `{'name': …}:None` into a
    prompt from `mention_context_builder`.
    """
    try:
        cfg = getattr(organization_settings, "config", None)
        if cfg is None:
            return None
        branding = cfg.get("branding") if isinstance(cfg, dict) else getattr(cfg, "branding", None)
        if branding is None:
            return None
        if isinstance(branding, dict):
            hint = branding.get("accent_color") or branding.get("product_name")
        else:
            hint = getattr(branding, "accent_color", None) or getattr(branding, "product_name", None)
        return str(hint) if hint else None
    except Exception:
        return None


def _resolve_deck_theme(
    user_text: str = "",
    report: Any = None,
    organization_settings: Any = None,
) -> Optional[Any]:
    """Resolve the theme this deck should be built in, or None when the registry
    is unavailable. Never raises — a theme is a design improvement, not a
    precondition for producing a deck."""
    themes = _load_pptx_themes()
    if themes is None:
        return None
    try:
        return themes.resolve(
            user_text=user_text or "",
            report_theme_name=getattr(report, "theme_name", None) if report is not None else None,
            org_brand=_org_brand_hint(organization_settings),
            agent_default=getattr(themes, "DEFAULT_THEME_ID", None),
        )
    except Exception:
        return None


def _norm_theme_words(text: Any) -> str:
    """Lowercase, collapse everything non-alphanumeric to single spaces, pad.

    Padded so a caller can test for a WHOLE token with a plain `in` and not
    match inside a longer word — 'boardroom' must not be found inside
    'boardroomsomething'. Mirrors the registry's own normalisation so the two
    agree about what counts as the same name.
    """
    if not isinstance(text, str):
        return " "
    return " " + re.sub(r"[^a-z0-9]+", " ", text.lower()).strip() + " "


def _norm_theme_key(text: Any) -> str:
    """'McKinsey Style' / 'MCKINSEY_STYLE' / ' mckinsey  style ' -> 'mckinsey-style'."""
    return _norm_theme_words(text).strip().replace(" ", "-")


def _theme_by_id(themes: Any, requested: Any) -> Optional[Any]:
    """Find the theme the model named, accepting what models actually send.

    ★The registry's `get()` is an EXACT id lookup, and models do not send exact
    ids. They send the display name off the index ('McKinsey Style'), the id
    with the wrong separator ('mckinsey_style'), the family without its suffix
    ('mckinsey'), or the id wrapped in a sentence. Every one of those is an
    unambiguous statement of which theme was wanted, and rejecting them would
    silently cost the model its choice — the same class of failure `_lenient`
    exists for, where 79% of live `clarify` calls were thrown away over shape.

    Returns None when nothing matched. None is not a failure: the caller falls
    back to deterministic resolution and records that it did.
    """
    if not isinstance(requested, str) or not requested.strip():
        return None
    try:
        raw = requested.strip()
        key = _norm_theme_key(raw)
        if not key:
            return None

        # Tiers 1-4: the id itself, however it was spelled or suffixed.
        for probe in (raw, key, f"{key}-style", key[: -len("-style")] if key.endswith("-style") else None):
            if probe:
                found = themes.get(probe)
                if found is not None:
                    return found

        # Tier 5: the display name exactly as the index prints it.
        for theme in themes.all_themes():
            for label in (getattr(theme, "id", ""), getattr(theme, "name", "")):
                label_key = _norm_theme_key(label)
                if not label_key:
                    continue
                if label_key == key or label_key.removesuffix("-style") == key:
                    return theme

        # Tier 6: a theme named INSIDE a longer string ('use the boardroom
        # look'). `resolve` is used for its alias table, then the answer is
        # VERIFIED to actually appear in the text — resolve() never returns
        # None, so without this check its default fallback would be
        # indistinguishable from a real match and every unknown id would look
        # like a deliberate choice of the default theme.
        candidate = themes.resolve(user_text=raw)
        if candidate is not None:
            haystack = _norm_theme_words(raw)
            for label in (getattr(candidate, "id", ""), getattr(candidate, "name", "")):
                for probe in (label, str(label).removesuffix("-style")):
                    token = _norm_theme_words(probe).strip()
                    if token and f" {token} " in haystack:
                        return candidate
    except Exception:
        # A theme is a design improvement, never a precondition for a deck.
        return None
    return None


#: "in the ledger style", "use McKinsey Style", "make it art deco"
_STYLE_PHRASE = re.compile(
    r"(?:in|use|using|make\s+it|styled?\s+(?:as|like))\s+(?:the\s+)?"
    r"([A-Za-z][A-Za-z0-9 &'-]{2,28}?)\s*(?:style|theme|look)?\s*[.,;!]",
    re.I,
)


def _named_theme_in(text: str, themes: Any) -> Optional[Any]:
    """The style a person NAMED, or None.

    ★Deliberately a phrase match, not a word match. The conversation render
    handed to this contains every prior assistant turn and tool summary, and a
    bare-token search over it resolved a request for the "Atelier" style to
    `christmas` — a theme name that merely appeared somewhere in the noise.
    Requiring naming grammar ("in the X style") is what separates a request
    from a mention. Scans from the END so the latest instruction wins.
    """
    if not text:
        return None
    for m in reversed(list(_STYLE_PHRASE.finditer(text))):
        candidate = (m.group(1) or "").strip()
        if not candidate:
            continue
        hit = _theme_by_id(themes, candidate)
        if hit is not None:
            return hit
    return None


def _select_deck_theme(
    requested_theme_id: Any = None,
    user_text: str = "",
    report: Any = None,
    organization_settings: Any = None,
) -> tuple[Optional[Any], Optional[Dict[str, Any]]]:
    """Pick the deck's design system, and say HOW it was picked.

    Order, highest first:

      1. `theme_id` the MODEL named, validated against the registry.
      2. `resolve()` — this turn's request, then the report's saved theme, then
         the org brand, then the agent default.

    An id the registry does not know NEVER fails the deck: selection falls to
    (2) and the choice record says so, so the console can tell "the model chose
    this" from "the model asked for something that does not exist" from "no one
    chose, it was resolved". Returns (None, None) only when the registry itself
    is unavailable, which is exactly the pre-registry behaviour.
    """
    themes = _load_pptx_themes()
    if themes is None:
        return None, None

    requested = requested_theme_id.strip() if isinstance(requested_theme_id, str) else None
    requested = requested or None

    chosen = _theme_by_id(themes, requested) if requested else None
    method = "model" if chosen is not None else None

    # ★A style the person NAMED outranks anything inferred. Checked before
    # resolve() because resolve() searches loose text, and loose text over a
    # whole conversation is how "Atelier" became `christmas`.
    if chosen is None:
        chosen = _named_theme_in(user_text or "", themes)
        if chosen is not None:
            method = "named_by_user"

    if chosen is None:
        chosen = _resolve_deck_theme(
            user_text=user_text or "",
            report=report,
            organization_settings=organization_settings,
        )
        method = "resolved_unknown_id" if requested else "resolved"

    if chosen is None:
        return None, None

    choice: Dict[str, Any] = {
        "id": getattr(chosen, "id", None),
        "name": getattr(chosen, "name", None),
        "method": method,
    }
    if requested:
        choice["requested_id"] = requested
    return chosen, choice


def _theme_as_dict(theme: Any) -> Dict[str, Any]:
    """The registry's dataclass as the plain dict the post-passes are typed for.

    Their contract is `theme: dict`; handing them the dataclass would make every
    lookup inside them a `getattr` on an object that answers with `None` instead
    of raising — the exact miss that shipped `{'name': …}: None` into a prompt
    from `mention_context_builder`.
    """
    if isinstance(theme, dict):
        return dict(theme)
    out: Dict[str, Any] = {}
    for attr in ("id", "name", "category", "when_to_use"):
        out[attr] = getattr(theme, attr, None)
    # ★`avoid` is what `enforce_theme_rules` acts on. This function is a
    # field-by-field rebuild, which is precisely how a new field goes missing —
    # 0.0.542.7 shipped with the enforcement pass running against themes that
    # carried no prohibitions and honestly reporting zero every time.
    _avoid = getattr(theme, "avoid", None)
    out["avoid"] = list(_avoid) if isinstance(_avoid, (list, tuple)) else []
    fonts = getattr(theme, "fonts", None)
    out["fonts"] = list(fonts) if isinstance(fonts, (list, tuple)) else fonts
    palette = getattr(theme, "palette", None)
    out["palette"] = dict(palette) if isinstance(palette, dict) else {}
    return out


def _theme_layout_for(theme_id: Optional[str]) -> Dict[str, Any]:
    """The furniture layout for a theme id, or {} when there is none.

    Defensive on the import for the same reason the passes themselves are: the
    layouts module is optional, and a deck must build without it.
    """
    if not theme_id:
        return {}
    try:
        from app.ai.decks.theme_layouts import LAYOUTS  # noqa: PLC0415 (optional)

        layout = LAYOUTS.get(theme_id)
        return layout if isinstance(layout, dict) else {}
    except Exception:
        return {}


def _report_theme_payload(report: Any, resolved_theme: Any = None) -> Optional[Dict[str, Any]]:
    """What generated code receives as `report['theme']`.

    ★There is no `Report.theme` anywhere in the tree — the model carries
    `theme_name` and `theme_overrides` (models/report.py:21-22). Every call site
    here read `getattr(report, "theme", None)`, which is decided when you TYPE it:
    it returned None for every artifact ever generated, while the slides prompt
    documented a `theme` field that was therefore always empty.
    """
    payload: Dict[str, Any] = {}
    if report is not None:
        name = getattr(report, "theme_name", None)
        overrides = getattr(report, "theme_overrides", None)
        if name:
            payload["name"] = name
        if isinstance(overrides, dict) and overrides:
            payload["overrides"] = overrides
    if resolved_theme is not None:
        # The report's own stored name wins as the label a person chose; the
        # resolved theme supplies the parts generated code can actually use.
        payload.setdefault("name", getattr(resolved_theme, "name", None))
        payload["id"] = getattr(resolved_theme, "id", None)
        payload["fonts"] = getattr(resolved_theme, "fonts", None)
        payload["palette"] = getattr(resolved_theme, "palette", None)
    return payload or None


def _looks_like_component_code(inner: str) -> bool:
    """True when `inner` plausibly contains component code rather than prose.

    Deliberately generous: a JSX tag or a hook call is enough on its own, and a
    keyword only needs one piece of structure beside it. The aim is to catch the
    unmistakable case — a reply that is one English sentence — without
    second-guessing unusual but genuine code.

    ★The error is directed on purpose. A false positive costs one render retry;
    a false negative refuses a real dashboard and burns an LLM round-trip to
    ask for what it already had. So `<Chart>` mentioned in a sentence passes,
    and that is the intended trade rather than an oversight.
    """
    if not inner or not inner.strip():
        return False
    if _JSX_TAG.search(inner) or _DATA_HOOK.search(inner):
        return True
    # A keyword on its own is a word; with structure beside it, it is a token.
    return bool(_CODE_KEYWORD.search(inner) and _CODE_STRUCTURE.search(inner))


from app.ai.tools.metadata import ToolMetadata
from app.ai.tools.schemas import (
    ToolEvent,
    ToolStartEvent,
    ToolProgressEvent,
    ToolEndEvent,
)
from app.ai.tools.schemas.create_artifact import CreateArtifactInput, CreateArtifactOutput
from app.ai.llm import LLM
from app.ai.llm.types import ImageInput, Message, TextDeltaEvent
from app.models.artifact import Artifact
from app.models.visualization import Visualization
from app.dependencies import async_session_maker
from app.services.thumbnail_service import ThumbnailService
from app.services.artifact_insights_html import (
    INSIGHTS_CSS,
    insights_section_html,
)
from app.services.artifact_libs import get_inline_scripts
# DEF-010 — the one accessor. Every path that reads or writes the data an
# artifact is built on goes through this module; see its docstring.
from app.services.artifact_data import resolve_artifact_rows, store_artifact_dataset
from app.ai.code_execution.pptx_executor import PptxCodeExecutor, PptxPreviewService
from sqlalchemy import desc
from app.ai.tools.implementations._sandbox_context import SANDBOX_RUNTIME_PROMPT
from app.ai.tools.implementations._artifact_images import load_image_bytes
from app.ai.prompt_language import build_language_directive

# DEF-002 — row budgets, per consumer. Rows are carried in FULL through the
# visualization list; each consumer slices for its own reason, and any slice it
# takes is declared to whoever reads it.
#
# _PROMPT_STATS_ROWS: how many rows the codegen prompt's column stats are
#   computed from. A sample, to bound prompt size. Declared as
#   `stats_from_sample` in the profile so the model doesn't read sampled
#   min/max as whole-dataset ranges. `row_count` stays the TRUE count.
# _RENDER_ROW_LIMIT: upper bound on rows injected into the headless preview /
#   thumbnail HTML. Generous, because a preview whose KPI tiles disagree with
#   the live dashboard is worse than a slow one — that image is fed back to the
#   model by read_artifact and by the self-heal path. Only very large datasets
#   hit it, and when they do the render is told (`rows_truncated`).
_PROMPT_STATS_ROWS = 100
_RENDER_ROW_LIMIT = 20000

# DEF-009 — a large input must not be a dead end.
#
# The Phase 1 completeness gate below refuses to build a dashboard from a
# truncated dataset, and that was right: a 1,000-row prefix of 28,592 rows
# shipped KPI tiles that were simply wrong. But refusing is not the same as
# handling it. Live, the agent met the refusal, went away, wrote a GROUP BY that
# returned 960 rows and came back — the correct outcome, reached one full
# generation and ~2 minutes later, every time.
#
# The rows in hand cannot be repaired: they are a PREFIX in the query's own sort
# order, so the missing rows are genuinely missing and no amount of arithmetic
# over what was persisted recovers them. But the QUERY is still there
# (`steps.code`) and the clients that ran it are still in this run's context. So
# re-run it — no LLM, seconds — and reduce the full result here:
#
#   * the whole result fits the artifact cap -> use every row. Complete.
#   * it does not                            -> GROUP BY the low-cardinality
#                                               dimensions and SUM the measures,
#                                               which is exactly what the agent
#                                               wrote by hand. Complete, at a
#                                               coarser grain.
#
# Either way every figure is computed over the WHOLE dataset, so nothing here
# trades correctness for a faster answer. What was done is then declared in
# three places that a reader cannot miss: the codegen profile (so the model
# labels an aggregate as an aggregate), the stored artifact content, and the
# tool's own returned observation. If re-running fails, or the result cannot be
# reduced honestly, the Phase 1 refusal stands unchanged — this path can only
# ever turn a refusal into a correct artifact, never into a wrong one.
_RECOVERY_MAX_GROUP_COLUMNS = 4
# A "dimension" with more distinct values than this is an identifier, not
# something to group by — grouping on it returns roughly the input row count and
# calls it an aggregate.
_RECOVERY_MAX_CANDIDATE_CARDINALITY = 5000
# Fallback when the org has no `artifact_row_limit` (an older settings row).
# Matches the setting's own default.
_DEFAULT_ARTIFACT_ROW_CAP = 10000
# ★ The re-run deliberately carries no LIMIT — reading the whole result is the
# entire point. But the whole result then lands in THIS worker's memory, and the
# truncation this path exists to undo was also what kept it out. A result this
# far past the artifact cap has no honest reduction anyway (grouping millions of
# rows on their low-cardinality dimensions still exceeds the cap, and
# `_aggregate_dataframe` would refuse it a moment later), so stop before paying
# for the DataFrame rather than after. Env DASH_RECOVERY_MAX_SOURCE_ROWS.
_RECOVERY_MAX_SOURCE_ROWS = int(os.getenv("DASH_RECOVERY_MAX_SOURCE_ROWS", "500000"))


def _recovery_enabled() -> bool:
    """DEF-009: re-read and reduce a truncated dataset instead of refusing.

    Default ON. No entry in `config.py` is required — `_read_bool_setting`
    returns the default for a name it cannot find — so setting
    HYBRID_ARTIFACT_DATA_RECOVERY=false only has an effect once it is declared
    there; until then this is on and the Phase 1 refusal remains the fallback.
    """
    return _read_bool_setting("hybrid_artifact_data_recovery", True)


def _artifact_row_cap(organization_settings: Any) -> int:
    """The org's artifact row cap. 0 means "no cap" (the setting says so)."""
    try:
        value = int(organization_settings.get_config("artifact_row_limit").value)
    except Exception:
        # Also covers a settings row predating the setting: get_config returns
        # None -> AttributeError.
        return _DEFAULT_ARTIFACT_ROW_CAP
    return value if value > 0 else 0


def _df_to_rows(df: Any) -> List[Dict[str, Any]]:
    """Serialize a DataFrame the same way the widget formatter does.

    Same arguments as `format_df_for_widget` on purpose: dates, Decimals, NaT
    and numpy scalars all have to survive `json.dumps` later, and this is the
    call that has already been proven to make them.
    """
    return json.loads(df.to_json(orient="records", date_format="iso", default_handler=str))


def _aggregate_dataframe(df: Any, cap: int) -> tuple[Any, Optional[Dict[str, Any]]]:
    """Group a too-large result down to something an artifact can carry whole.

    Returns ``(aggregated_df, meta)``, or ``(None, None)`` when no honest
    reduction exists — no measures to add up, or no dimension coarse enough to
    group by. Refusing here is correct: the Phase 1 gate then refuses the build,
    which is what happens today.

    The aggregate is COMPLETE — every source row is counted in exactly one
    group — so sums, counts and shares taken from it match the full dataset.
    Averages do not survive a plain SUM, which is why the source row count per
    group is carried alongside as a column: mean = sum / that count.
    """
    import pandas as pd
    from pandas.api import types as ptypes

    if cap <= 0:
        return None, None

    work = df.copy()

    numeric_cols = [
        str(c) for c in work.columns
        if ptypes.is_numeric_dtype(work[c]) and not ptypes.is_bool_dtype(work[c])
    ]
    if not numeric_cols:
        # Nothing to add up. A COUNT-only aggregate would answer a question
        # nobody asked, so leave it to the refusal.
        return None, None

    # Bin datetimes to month before measuring cardinality — a timestamp column is
    # near-unique by nature and would be discarded as an identifier, when it is
    # usually the single most useful axis on the dashboard.
    binned: Dict[str, str] = {}
    for c in list(work.columns):
        if ptypes.is_datetime64_any_dtype(work[c]):
            try:
                work[c] = work[c].dt.to_period("M").astype(str)
                binned[str(c)] = "month"
            except Exception:
                pass

    candidates = [str(c) for c in work.columns if str(c) not in numeric_cols]
    if not candidates:
        return None, None

    cardinality = {}
    for c in candidates:
        try:
            cardinality[c] = int(work[c].nunique(dropna=False))
        except TypeError:
            # Unhashable cell values (a list/dict in a column) — not a dimension.
            continue
    ordered = sorted(
        (c for c, n in cardinality.items() if n <= _RECOVERY_MAX_CANDIDATE_CARDINALITY),
        key=lambda c: cardinality[c],
    )

    # Coarsest dimension first, adding finer ones while the result still fits.
    # Taking them in this order keeps as MANY dimensions as possible, which is
    # what makes the reduced dataset still worth charting.
    chosen: List[str] = []
    for c in ordered[:_RECOVERY_MAX_GROUP_COLUMNS]:
        trial = chosen + [c]
        try:
            groups = int(work.groupby(trial, dropna=False).ngroups)
        except Exception:
            continue
        if groups > cap:
            break
        chosen = trial
    if not chosen:
        return None, None

    grouped = work.groupby(chosen, dropna=False)
    agg = grouped[numeric_cols].sum().reset_index()

    count_col = "source_row_count"
    while count_col in agg.columns:
        count_col = "_" + count_col
    agg[count_col] = grouped.size().reset_index(drop=True)

    meta = {
        "method": "aggregate",
        "group_columns": chosen,
        "binned_columns": binned,
        "measures": numeric_cols,
        "row_count_column": count_col,
        "dropped_columns": [
            str(c) for c in df.columns
            if str(c) not in chosen and str(c) not in numeric_cols
        ],
    }
    return agg, meta


def _reduction_notice(reduction: Dict[str, Any]) -> str:
    """One sentence a person (and the model) can act on. No jargon, no IDs."""
    source = reduction.get("source_row_count")
    used = reduction.get("rows_used")
    if reduction.get("method") == "aggregate":
        groups = ", ".join(reduction.get("group_columns") or [])
        binned = reduction.get("binned_columns") or {}
        bin_note = (
            f" ({', '.join(f'{c} binned by {g}' for c, g in binned.items())})" if binned else ""
        )
        return (
            f"{source:,} rows were re-read in full and pre-aggregated to {used:,} rows, "
            f"grouped by {groups}{bin_note}; every numeric column is a SUM over the "
            f"WHOLE dataset, and '{reduction.get('row_count_column')}' is how many source "
            f"rows each group covers (use it as the denominator for averages)."
        )
    return (
        f"the stored copy held only {reduction.get('stored_row_count'):,} rows, so all "
        f"{used:,} rows were re-read from the source; the artifact now has the complete dataset."
    )


def _build_executor(organization_settings: Any, context_hub: Any) -> Any:
    """The code executor used to re-run a step's query.

    Its own function so a test can replace it without a database, a data source
    or a sandbox. Imported lazily — `code_execution` pulls in the whole
    execution stack and this module is imported at tool-registry build time.
    """
    from app.ai.code_execution.code_execution import StreamingCodeExecutor

    return StreamingCodeExecutor(
        organization_settings=organization_settings,
        logger=None,
        context_hub=context_hub,
    )


async def recover_truncated_visualizations(
    visualizations: List[Dict[str, Any]],
    step_code_by_viz: Dict[str, str],
    runtime_ctx: Dict[str, Any],
    organization_settings: Any,
) -> List[Dict[str, Any]]:
    """DEF-009. Re-read each truncated visualization and reduce it honestly.

    Mutates the entries in `visualizations` in place: a recovered entry carries
    the complete data, so its `rows_truncated` / `rows_available` markers are
    REMOVED — they described rows that are no longer what the artifact holds,
    and leaving them would make the Phase 1 gate refuse complete data.

    Returns one disclosure record per recovered visualization. A visualization
    that could not be recovered is left exactly as it was, so the Phase 1
    refusal still fires for it: this function can turn a refusal into a correct
    artifact and can do nothing else.
    """
    reductions: List[Dict[str, Any]] = []
    cap = _artifact_row_cap(organization_settings)

    for viz in visualizations:
        if not viz.get("rows_truncated"):
            continue
        viz_id = str(viz.get("id"))
        code = step_code_by_viz.get(viz_id)
        if not code:
            logger.warning("DEF-009: visualization %s has no stored query to re-run", viz_id)
            continue

        try:
            executor = _build_executor(organization_settings, runtime_ctx.get("context_hub"))
            result = await executor.execute_code_async(
                code=code,
                ds_clients=runtime_ctx.get("ds_clients") or {},
                excel_files=runtime_ctx.get("excel_files") or [],
            )
            df = result[0] if isinstance(result, tuple) else result
        except Exception as e:
            logger.warning("DEF-009: re-running the query for %s failed (%s)", viz_id, e)
            continue

        if df is None or len(df) == 0:
            logger.warning("DEF-009: re-run for %s returned no data — leaving it refused", viz_id)
            continue

        stored_rows = len(viz.get("rows") or [])
        source_rows = int(len(df))

        # See `_RECOVERY_MAX_SOURCE_ROWS`. Dropping the reference here matters:
        # the loop continues to the next visualization, and holding a
        # multi-million-row frame while re-running ANOTHER query is how one
        # oversized result takes the whole worker down with it.
        if source_rows > _RECOVERY_MAX_SOURCE_ROWS:
            logger.warning(
                "DEF-009: re-run for %s returned %d rows, past the %d ceiling — "
                "leaving it refused",
                viz_id, source_rows, _RECOVERY_MAX_SOURCE_ROWS,
            )
            df = None
            continue

        reduction: Dict[str, Any]

        if cap and source_rows > cap:
            agg, meta = _aggregate_dataframe(df, cap)
            if agg is None:
                logger.warning(
                    "DEF-009: %s (%d rows) has no honest reduction — leaving it refused",
                    viz_id, source_rows,
                )
                continue
            df = agg
            reduction = dict(meta)
        else:
            reduction = {"method": "full_reread"}

        try:
            rows = _df_to_rows(df)
        except Exception as e:
            logger.warning("DEF-009: could not serialize the reduced result for %s (%s)", viz_id, e)
            continue

        reduction.update({
            "visualization_id": viz_id,
            "title": viz.get("title"),
            "source_row_count": source_rows,
            "stored_row_count": stored_rows,
            "rows_used": len(rows),
        })
        reduction["notice"] = _reduction_notice(reduction)

        viz["rows"] = rows
        viz["columns"] = [{"headerName": str(c), "field": str(c)} for c in df.columns]
        # The rows in hand are now the whole dataset (or a complete aggregate of
        # it), so the count of them IS the count — and the stale column_info,
        # computed over the old prefix, would describe data that is gone.
        viz["row_count"] = len(rows)
        viz["column_info"] = {}
        viz["data_reduction"] = reduction
        viz.pop("rows_truncated", None)
        viz.pop("rows_total", None)
        viz.pop("rows_available", None)

        logger.info(
            "DEF-009: recovered '%s' — %s (%d stored -> %d source -> %d used)",
            viz.get("title") or viz_id, reduction["method"],
            stored_rows, source_rows, len(rows),
        )
        reductions.append(reduction)

    return reductions


class CreateArtifactTool(Tool):
    """Tool for generating React-based artifact code for dashboards.

    This tool generates standalone React/JSX code that renders visualizations
    using ECharts, styled with Tailwind CSS, and transpiled in-browser via Babel.

    The generated code runs in a sandboxed iframe and receives visualization
    data via window.ARTIFACT_DATA.
    """

    @property
    def metadata(self) -> ToolMetadata:
        return ToolMetadata(
            name="create_artifact",
            description=(
                "Create or fully rebuild artifacts (dashboards, pages, slide presentations) from visualizations. "
                "Use for: new dashboards, full redesigns, large layout changes, or when edit_artifact cannot handle the scope. "
                "Modes: 'page' for interactive dashboards with KPI cards, charts, and responsive grids; "
                "'slides' for presentation decks (exportable to PPTX). "
                "IMPORTANT: for 'page' mode visualization_ids are required - find them in previous create_data tool results "
                "shown as 'viz_id: <uuid>' in the conversation history. For 'slides' mode they are optional: "
                "a deck may include title, agenda and narrative slides that carry no chart. "
                "Do NOT ask the user for URLs or IDs - extract them from the conversation context. "
                "Only visualizations with successful step status are included."
            ),
            category="action",
            version="1.0.0",
            input_schema=CreateArtifactInput.model_json_schema(),
            output_schema=CreateArtifactOutput.model_json_schema(),
            max_retries=1,
            timeout_seconds=120,
            idempotent=False,
            required_permissions=[],
            is_active=True,
            tags=["artifact",  "dashboard", "slides"],
            allowed_modes=["chat"],
        )

    @property
    def input_model(self) -> Type[BaseModel]:
        return CreateArtifactInput

    @property
    def output_model(self) -> Type[BaseModel]:
        return CreateArtifactOutput

    # Path to the sandbox HTML file (relative to project root)
    # __file__ -> implementations -> tools -> ai -> app -> backend -> project_root
    SANDBOX_HTML_PATH = Path(__file__).parent.parent.parent.parent.parent.parent / "frontend" / "public" / "artifact-sandbox.html"

    async def _take_preview_screenshot(
        self,
        html_content: str,
    ) -> tuple[Optional[str], list[str]]:
        """Take a quick screenshot for planner reflection and capture JS errors.

        Returns (base64-encoded PNG string or None, list of JS error messages).
        """
        try:
            from playwright.async_api import async_playwright
        except ImportError:
            return None, []

        js_errors: list[str] = []

        try:
            import tempfile, os
            async with async_playwright() as p:
                # Optional executable override for deployments where the
                # Playwright-managed browser download is unavailable but a
                # compatible Chromium exists on disk.
                _exe = os.environ.get("BOW_CHROMIUM_EXECUTABLE") or None
                browser = await p.chromium.launch(headless=True, executable_path=_exe)
                page = await browser.new_page(viewport={"width": 1280, "height": 720})

                # Capture JS errors during render. Both channels matter:
                # pageerror catches thrown/uncaught exceptions (incl. Babel
                # transform errors), console 'error' catches failures that
                # libraries report without throwing (React render errors).
                def _on_pageerror(err):
                    if len(js_errors) < 10:
                        js_errors.append(str(err))

                def _on_console(msg):
                    try:
                        if msg.type == "error" and len(js_errors) < 10:
                            text = msg.text or ""
                            # Skip noise that isn't a code defect
                            if "favicon" in text.lower():
                                return
                            entry = f"[console.error] {text}"
                            if not any(text in e for e in js_errors):
                                js_errors.append(entry)
                    except Exception:
                        pass

                page.on("pageerror", _on_pageerror)
                page.on("console", _on_console)

                # Write HTML to a temp file and navigate via file:// URL.
                # This allows vendored scripts (e.g. Tailwind runtime) that use
                # document.write() to work correctly — document.write fails on
                # about:blank pages used by set_content().
                tmp = tempfile.NamedTemporaryFile(suffix=".html", delete=False, mode="w", encoding="utf-8")
                try:
                    tmp.write(html_content)
                    tmp.close()
                    await page.goto(f"file://{tmp.name}", wait_until="networkidle")

                    # Wait for React to mount and charts to render (short timeout)
                    try:
                        await page.wait_for_function(
                            "window.__ARTIFACT_RENDER_COMPLETE__ === true",
                            timeout=8000,
                        )
                    except Exception:
                        pass  # Take screenshot anyway — partial render is still useful

                    await asyncio.sleep(0.3)
                    screenshot_bytes = await page.screenshot(type="png", full_page=False)
                    await browser.close()
                    return base64.b64encode(screenshot_bytes).decode("utf-8"), js_errors
                finally:
                    os.unlink(tmp.name)
        except Exception as e:
            logger.warning(f"Preview screenshot failed: {e}")
            return None, js_errors

    async def _generate_thumbnail_background(
        self,
        artifact_id: str,
        html_content: str,
        mode: str = "page",
    ) -> None:
        """Generate thumbnail in background and update artifact.

        Runs independently with its own database session.
        """
        try:
            thumbnail_service = ThumbnailService()
            thumbnail_path = await thumbnail_service.generate_thumbnail(
                artifact_id=artifact_id,
                html_content=html_content,
                mode=mode,
            )
            if thumbnail_path:
                # Use a fresh database session for the background update
                async with async_session_maker() as db:
                    from sqlalchemy import update
                    from app.models.artifact import Artifact
                    stmt = update(Artifact).where(Artifact.id == artifact_id).values(thumbnail_path=thumbnail_path)
                    await db.execute(stmt)
                    await db.commit()
        except Exception as e:
            logger.warning(f"Failed to generate thumbnail for artifact {artifact_id}: {e}")

    async def _load_completion_images(
        self,
        db: Any,
        head_completion_id: Optional[str],
    ) -> List[ImageInput]:
        """Load images attached to the head completion as ImageInput objects.

        Args:
            db: Database session
            head_completion_id: The completion ID to load images for

        Returns:
            List of ImageInput objects ready for vision-capable LLM
        """
        if not head_completion_id:
            return []

        images: List[ImageInput] = []
        try:
            # Query files associated with this completion that are images
            result = await db.execute(
                select(File)
                .join(report_file_association, report_file_association.c.file_id == File.id)
                .where(report_file_association.c.completion_id == head_completion_id)
                .where(File.content_type.startswith("image/"))
            )
            image_files = result.scalars().all()

            for f in image_files:
                if not f.path:
                    continue
                try:
                    async with aiofiles.open(f.path, 'rb') as file:
                        content = await file.read()
                    images.append(ImageInput(
                        data=base64.b64encode(content).decode('utf-8'),
                        media_type=f.content_type or 'image/png',
                        source_type='base64'
                    ))
                except Exception as e:
                    logger.warning(f"Failed to load image file {f.id}: {e}")

        except Exception as e:
            logger.warning(f"Failed to query completion images: {e}")

        return images

    def _render_visualizations(self, visualizations: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Rows for the headless preview/thumbnail render (DEF-002).

        The live frontend hands the artifact every row of the step, so the
        preview must do the same or its KPI tiles disagree with what the user
        opens. Only a very large dataset is cut, and when it is, the injected
        payload says so — a truncated render that claims to be complete is the
        defect this replaced.
        """
        out: List[Dict[str, Any]] = []
        for viz in visualizations:
            rows = viz.get("rows") or []
            # DEF-010: carry any cap/aggregation disclosure into the rendered
            # payload, so the page that shows the numbers can also show what
            # shaped them.
            reduction = viz.get("data_reduction") or {}
            notice = reduction.get("notice") if isinstance(reduction, dict) else None
            if len(rows) <= _RENDER_ROW_LIMIT:
                if notice:
                    viz = dict(viz)
                    viz["dataNotice"] = str(notice)
                out.append(viz)
                continue
            logger.warning(
                "Artifact preview render: visualization %s has %d rows, capping at %d",
                viz.get("id"), len(rows), _RENDER_ROW_LIMIT,
            )
            capped = dict(viz)
            capped["rows"] = rows[:_RENDER_ROW_LIMIT]
            capped["rows_truncated"] = True
            capped["rows_total"] = len(rows)
            capped["dataNotice"] = (
                f"Preview shows {_RENDER_ROW_LIMIT:,} of {len(rows):,} rows."
                + (f" {notice}" if notice else "")
            )
            out.append(capped)
        return out

    def _build_thumbnail_html(
        self,
        artifact_data: dict,
        code: str,
        mode: str = "page",
        insights: Optional[dict] = None,
    ) -> str:
        """Build HTML for thumbnail generation in headless browser.

        Args:
            artifact_data: The data to inject as window.ARTIFACT_DATA
            code: The LLM-generated artifact code
            mode: 'page' for React dashboards, 'slides' for pure HTML presentations

        Returns:
            Complete HTML string ready for headless browser rendering
        """
        data_json = json.dumps(artifact_data, default=str)

        # Slides mode: pure HTML + Tailwind (no React/Babel)
        if mode == "slides":
            slides_scripts = get_inline_scripts(mode="slides")
            slides_template = """<!DOCTYPE html>
<html>
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  __SLIDES_SCRIPTS__
  <style>
    html, body { height: 100%; margin: 0; padding: 0; }
    body { font-family: system-ui, -apple-system, sans-serif; }
    .slide { transition: opacity 0.3s ease-in-out; }
  </style>
</head>
<body class="bg-slate-900">
  <script>
    window.ARTIFACT_DATA = __ARTIFACT_DATA_JSON__;
    window.__ARTIFACT_RENDER_COMPLETE__ = false;
    setTimeout(function() {
      window.__ARTIFACT_RENDER_COMPLETE__ = true;
    }, 500);
  </script>

  __LLM_GENERATED_CODE__
</body>
</html>"""
            return slides_template.replace("__SLIDES_SCRIPTS__", slides_scripts).replace("__ARTIFACT_DATA_JSON__", data_json).replace("__LLM_GENERATED_CODE__", code)

        # Page mode: Build self-contained HTML mirroring ArtifactFrame.vue's approach.
        # get_inline_scripts("page") already includes all vendored libs + artifact-globals.js
        # so we only need to inject ARTIFACT_DATA, the LLM code, and render-complete detection.
        page_scripts = get_inline_scripts(mode="page")
        insights_html = insights_section_html(
            insights, (artifact_data or {}).get("visualizations") or []
        )
        SC = '</' + 'script>'  # Avoid parser issues in this Python string too

        html = f"""<!DOCTYPE html>
<html>
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  {page_scripts}
  <style>
    {INSIGHTS_CSS}
    body {{ font-family: system-ui, -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; }}
  </style>
</head>
<body>
  <div id="root"></div>
  {insights_html}

  <script>
    window.ARTIFACT_DATA = {data_json};
    window.__ARTIFACT_RENDER_COMPLETE__ = false;
    window.__DASH_INFO = false;
  {SC}

  {code}

  <script>
    (function detectRenderComplete() {{
      var startTime = Date.now();
      var MAX_WAIT = 15000;
      function check() {{
        if (Date.now() - startTime > MAX_WAIT) {{
          window.__ARTIFACT_RENDER_COMPLETE__ = true;
          return;
        }}
        var root = document.getElementById('root');
        if (!root || root.children.length === 0) {{
          setTimeout(check, 200);
          return;
        }}
        var hasCharts = root.querySelectorAll('canvas').length > 0 ||
                        root.querySelectorAll('[_echarts_instance_]').length > 0;
        setTimeout(function() {{
          window.resizeAllCharts && window.resizeAllCharts();
          window.__ARTIFACT_RENDER_COMPLETE__ = true;
        }}, hasCharts ? 1500 : 300);
      }}
      setTimeout(check, 200);
    }})();
  {SC}
</body>
</html>"""
        return html

    # Maximum in-tool repair LLM calls per artifact operation. Bounded so a
    # stubborn defect can't consume the whole tool timeout — after this the
    # tool returns a structured failure and the planner decides.
    MAX_RENDER_REPAIR_ATTEMPTS = 2

    @staticmethod
    def fatal_render_errors(errors: List[str]) -> List[str]:
        """Errors that gate completion: thrown/uncaught exceptions.

        `[console.error]` entries are advisory — React dev builds log
        non-fatal warnings through console.error, so they inform repair but
        never fail an otherwise working artifact.
        """
        return [e for e in (errors or []) if not e.startswith("[console.error]")]

    async def _fix_code(
        self,
        code: str,
        errors: List[str],
        mode: str,
        runtime_ctx: Dict[str, Any],
    ) -> Optional[str]:
        """Compact in-tool repair: current code + exact errors → corrected code.

        Deliberately does NOT rebuild the full generation prompt — the model
        already produced this code; it needs the error, not the whole context.
        Returns the corrected code, or None if repair was unavailable/failed.
        """
        sigkill_event = runtime_ctx.get("sigkill_event")
        if sigkill_event and sigkill_event.is_set():
            return None

        error_text = "\n".join(f"- {e}" for e in errors[:5])

        fix_prompt = f"""You previously wrote the React dashboard code below. It runs in a sandboxed iframe (React 18 + Babel standalone + Tailwind + ECharts via <EChart>, data via useArtifactData()). When rendered, it produced these errors:

{error_text}

Current code:
```
{code}
```

Fix ONLY what the errors require — do not redesign, restyle, or restructure anything else. Common causes: using a variable/component before its declaration, duplicate declarations, undefined symbols, unguarded nullish values before string methods, malformed JSX.

Output the FULL corrected code wrapped in <script type="text/babel"> ... </script>. No explanations, no diff markers, no markdown fences."""

        llm = LLM(runtime_ctx.get("model"), usage_session_maker=async_session_maker)
        try:
            chunks: list[str] = []
            async for evt in llm.inference_stream_v2(
                messages=[Message(role="user", content=fix_prompt)],
                usage_scope="create_artifact_fix",
            ):
                if sigkill_event and sigkill_event.is_set():
                    return None
                if isinstance(evt, TextDeltaEvent):
                    chunks.append(evt.text)
            fixed = self._extract_code("".join(chunks), mode=mode)
            return fixed if fixed and fixed.strip() else None
        except Exception:
            logger.exception("In-tool code repair failed")
            return None

    async def _fix_pptx_code(
        self,
        code: str,
        error: str,
        runtime_ctx: Dict[str, Any],
    ) -> Optional[str]:
        """Compact in-tool repair for slides: current code + exact error → corrected code.

        Slides twin of _fix_code: the model already produced this code; it
        needs the executor's error, not the whole generation context.
        Returns the corrected code, or None if repair was unavailable/failed.
        """
        sigkill_event = runtime_ctx.get("sigkill_event")
        if sigkill_event and sigkill_event.is_set():
            return None

        fix_prompt = f"""You previously wrote the python-pptx code below. It runs in a sandboxed namespace that already provides: Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE, XL_CHART_TYPE, XL_LEGEND_POSITION, CategoryChartData, ChartData, plus the data variables `visualizations` (list of dicts with 'title', 'columns', 'rows'), `report`, `image`/`image_ids`, and `_pptx_output_path` (the code must call prs.save(_pptx_output_path)). Note that each entry of viz['columns'] is a DICT like {{'field': 'Revenue', 'headerName': 'Revenue'}} — use col['field'] to get the row key, never pass the dict itself where a string is expected. When executed, it failed with this error:

{error}

Current code:
```python
{code}
```

Fix ONLY what the error requires — do not redesign, restyle, or restructure anything else. Common causes: passing a non-string (e.g. a column dict) to a text property, indexing a row with the wrong key, referencing an undefined name, wrong python-pptx API usage.

Output the FULL corrected code in a ```python code block. No explanations, no diff markers."""

        llm = LLM(runtime_ctx.get("model"), usage_session_maker=async_session_maker)
        try:
            chunks: list[str] = []
            async for evt in llm.inference_stream_v2(
                messages=[Message(role="user", content=fix_prompt)],
                usage_scope="create_artifact_fix",
            ):
                if sigkill_event and sigkill_event.is_set():
                    return None
                if isinstance(evt, TextDeltaEvent):
                    chunks.append(evt.text)
            fixed = self._extract_code("".join(chunks), mode="slides")
            return fixed if fixed and fixed.strip() else None
        except Exception:
            logger.exception("In-tool PPTX code repair failed")
            return None

    @staticmethod
    def _pptx_error_text(exc: Exception) -> str:
        """Trimmed executor error for repair prompts and observations.

        Keeps the frames inside the generated code (reported as <string>) and
        the final exception line — the sandbox internals above them are noise
        the repair model can't act on.
        """
        import traceback as _tb

        lines = _tb.format_exception(type(exc), exc, exc.__traceback__)
        flat = "".join(lines).splitlines()
        kept = [l for l in flat if '<string>' in l]
        # Final line always carries the exception type + message
        if flat:
            kept.append(flat[-1])
        text = "\n".join(kept[-12:]) if kept else str(exc)
        return text[:2000]

    @staticmethod
    def _first_preview_base64(preview_images: List[str]) -> Optional[str]:
        """Base64 of the first slide preview PNG, or None if unavailable."""
        if not preview_images:
            return None
        import base64

        path = Path(__file__).parent.parent.parent.parent.parent / "uploads" / preview_images[0]
        try:
            return base64.b64encode(path.read_bytes()).decode("ascii")
        except OSError:
            return None

    async def _execute_and_repair_pptx(
        self,
        code: str,
        visualizations: List[Dict[str, Any]],
        report_data: Dict[str, Any],
        output_path: Path,
        images: Dict[str, bytes],
        runtime_ctx: Dict[str, Any],
        deadline_monotonic: Optional[float] = None,
    ) -> AsyncIterator[Any]:
        """Execute slides code and repair it in-tool when it fails.

        Slides twin of _validate_and_repair_stream. Async generator: yields
        ToolProgressEvent items for the UI, then a final dict result:
          {"code", "ok", "error", "repair_attempts"}

        - Only executor exceptions gate; a successful run ends the loop.
        - If repair can't produce a working deck, the ORIGINAL code and its
          error are returned so what's persisted matches what was executed.
        """
        import time as _time

        def _time_left() -> bool:
            return deadline_monotonic is None or _time.monotonic() < deadline_monotonic

        sigkill_event = runtime_ctx.get("sigkill_event")
        executor = PptxCodeExecutor(logger=logger)

        def _try(candidate_code: str) -> Optional[str]:
            """Run once; returns None on success, trimmed error text on failure."""
            try:
                executor.execute_pptx_code(
                    code=candidate_code,
                    visualizations=visualizations,
                    report=report_data,
                    output_path=output_path,
                    images=images,
                )
                return None
            except Exception as e:
                logger.error(f"PPTX execution failed: {e}")
                return self._pptx_error_text(e)

        yield ToolProgressEvent(type="tool.progress", payload={"stage": "executing_pptx_code"})
        error = _try(code)

        original_code = code
        original_error = error
        candidate = code
        attempts = 0

        while (
            error
            and attempts < self.MAX_RENDER_REPAIR_ATTEMPTS
            and _time_left()
            and not (sigkill_event and sigkill_event.is_set())
        ):
            attempts += 1
            yield ToolProgressEvent(
                type="tool.progress",
                payload={"stage": "repairing_code", "attempt": attempts, "error_count": 1},
            )
            fixed = await self._fix_pptx_code(candidate, error, runtime_ctx)
            if not fixed or fixed == candidate:
                break

            yield ToolProgressEvent(
                type="tool.progress",
                payload={"stage": "executing_pptx_code", "attempt": attempts},
            )
            error = _try(fixed)
            candidate = fixed

        if not error:
            yield {
                "code": candidate,
                "ok": True,
                "error": None,
                "repair_attempts": attempts,
            }
        else:
            # Repair didn't converge — return the original, verified-broken
            # state rather than an unverified intermediate, and drop any deck
            # a partial attempt may have left behind.
            try:
                if output_path.exists():
                    output_path.unlink()
            except OSError:
                pass
            yield {
                "code": original_code,
                "ok": False,
                "error": original_error,
                "repair_attempts": attempts,
            }

    async def _validate_and_repair_stream(
        self,
        code: str,
        artifact_data: Dict[str, Any],
        mode: str,
        runtime_ctx: Dict[str, Any],
        deadline_monotonic: Optional[float] = None,
    ) -> AsyncIterator[Any]:
        """Render-validate page code and repair it in-tool when it fails.

        Async generator: yields ToolProgressEvent items for the UI, then a
        final dict result:
          {"code", "clean", "screenshot", "errors", "repair_attempts"}

        - Validation always runs (it needs Playwright, not a vision model or
          the allow_llm_see_data flag — errors are not data).
        - Only fatal (thrown) errors gate; console errors ride along as
          repair context.
        - If repair can't produce a clean render, the ORIGINAL code and its
          errors are returned so what's persisted matches what was verified.
        """
        import time as _time

        def _time_left() -> bool:
            return deadline_monotonic is None or _time.monotonic() < deadline_monotonic

        sigkill_event = runtime_ctx.get("sigkill_event")

        yield ToolProgressEvent(type="tool.progress", payload={"stage": "validating_render"})
        try:
            html = self._build_thumbnail_html(artifact_data, code, mode=mode)
        except Exception as e:
            # Validation infrastructure unavailable (e.g. vendored libs not
            # downloaded) must not fail artifact creation — pass through
            # unvalidated, exactly like the missing-Playwright fallback.
            logger.warning(f"Render validation unavailable (thumbnail HTML build failed): {e}")
            yield {
                "code": code,
                "clean": True,
                "screenshot": None,
                "errors": [],
                "repair_attempts": 0,
                "validation_skipped": True,
            }
            return
        screenshot, errors = await self._take_preview_screenshot(html)
        fatal = self.fatal_render_errors(errors)

        original_code = code
        original_screenshot = screenshot
        original_errors = errors
        candidate = code
        attempts = 0

        while (
            fatal
            and attempts < self.MAX_RENDER_REPAIR_ATTEMPTS
            and _time_left()
            and not (sigkill_event and sigkill_event.is_set())
        ):
            attempts += 1
            yield ToolProgressEvent(
                type="tool.progress",
                payload={"stage": "repairing_code", "attempt": attempts, "error_count": len(fatal)},
            )
            fixed = await self._fix_code(candidate, errors, mode, runtime_ctx)
            if not fixed or fixed == candidate:
                break

            yield ToolProgressEvent(
                type="tool.progress",
                payload={"stage": "validating_render", "attempt": attempts},
            )
            html = self._build_thumbnail_html(artifact_data, fixed, mode=mode)
            screenshot, errors = await self._take_preview_screenshot(html)
            candidate = fixed
            fatal = self.fatal_render_errors(errors)

        if not fatal:
            yield {
                "code": candidate,
                "clean": True,
                "screenshot": screenshot,
                "errors": errors,
                "repair_attempts": attempts,
            }
        else:
            # Repair didn't converge — return the original, verified-broken
            # state rather than an unverified intermediate.
            yield {
                "code": original_code,
                "clean": False,
                "screenshot": original_screenshot,
                "errors": original_errors,
                "repair_attempts": attempts,
            }

    def _build_viz_profile(self, viz: Dict[str, Any], allow_llm_see_data: bool) -> Dict[str, Any]:
        """Build a privacy-aware profile of a visualization's data."""
        # Enrich columns with dtype/unique_count/min/max from column_info (always — not sensitive)
        column_info = viz.get("column_info") or {}
        raw_columns = viz.get("columns", [])
        enriched_columns = []
        for c in raw_columns:
            col = dict(c) if isinstance(c, dict) else {"field": c}
            field = col.get("field") or col.get("headerName") or col.get("name")
            if field and field in column_info:
                meta = column_info[field]
                col["dtype"] = meta.get("dtype")
                col["unique_count"] = meta.get("unique_count")
                if meta.get("min") is not None:
                    col["min"] = meta["min"]
                if meta.get("max") is not None:
                    col["max"] = meta["max"]
            enriched_columns.append(col)

        profile: Dict[str, Any] = {
            "id": viz.get("id"),
            "title": viz.get("title"),
            "chart_type": viz.get("data_model_type") or "table",
            # True dataset size; sample_row_count is how many rows are shown
            # to generation/preview (capped). At runtime the dashboard
            # receives ALL row_count rows.
            "row_count": viz.get("row_count", 0),
            "sample_row_count": viz.get("sample_row_count", len(viz.get("rows") or [])),
            "columns": enriched_columns,
        }

        # DEF-004: if the artifact only holds a prefix of the dataset, say so in the
        # profile the codegen prompt is built from — and say what it means, because
        # "rows_available < row_count" alone reads as a detail rather than as a
        # correctness constraint on every KPI tile the model is about to write.
        if viz.get("rows_truncated"):
            profile["rows_available"] = viz.get("rows_available")
            profile["rows_truncated"] = True
            profile["data_completeness_warning"] = (
                f"Only {viz.get('rows_available')} of {viz.get('row_count')} rows are "
                f"available here, taken in the query's own sort order — so the most "
                f"recent periods or lowest-ranked groups may be missing entirely. Any "
                f"SUM/COUNT/AVG over this data is PARTIAL: do not label it a total, and "
                f"do not describe the period it spans as the full period. Either state "
                f"it as partial in the tile or subtitle, or ask for a coarser-grained "
                f"dataset that fits the row limit."
            )

        # DEF-009: this dataset was re-read and reduced. The model MUST know —
        # writing `SUM(amount)` over a pre-aggregated table is right, but calling
        # a row "an order" when it is a month × branch total is not, and neither
        # is averaging a column of sums.
        reduction = viz.get("data_reduction")
        if reduction:
            profile["data_reduction"] = reduction
            if reduction.get("method") == "aggregate":
                profile["data_reduction_notice"] = (
                    f"This data is PRE-AGGREGATED and COMPLETE: all "
                    f"{reduction.get('source_row_count')} source rows were grouped by "
                    f"{', '.join(reduction.get('group_columns') or [])} and every numeric "
                    f"column is a SUM over the whole dataset. One row is a GROUP, not a "
                    f"record. Sums, counts and shares are exact; for an average divide by "
                    f"'{reduction.get('row_count_column')}', never take the mean of a summed "
                    f"column. Label tiles by what they are (a total), and note in the "
                    f"subtitle that the view is aggregated at this grain."
                )
            else:
                profile["data_reduction_notice"] = (
                    f"This dataset was re-read in full ({reduction.get('source_row_count')} "
                    f"rows) because the stored copy was truncated. It is COMPLETE — totals "
                    f"over it are real totals."
                )

        # Include data model hints
        data_model = viz.get("dataModel") or {}
        if data_model:
            series = data_model.get("series", [])
            if series:
                profile["series_config"] = series[:3]  # First 3 series configs
            if data_model.get("group_by"):
                profile["group_by"] = data_model.get("group_by")

        # Include view configuration hints
        view = viz.get("view") or {}
        if view:
            inner_view = view.get("view") or view
            profile["view_config"] = {
                "type": inner_view.get("type"),
                "x": inner_view.get("x"),
                "y": inner_view.get("y"),
                "category": inner_view.get("category"),
                "value": inner_view.get("value"),
            }
            # Surface aggregation (top-level) + per-series aggregations so the
            # artifact can honor granular-data handling rather than reading
            # the first row.
            if inner_view.get("aggregation"):
                profile["view_config"]["aggregation"] = inner_view.get("aggregation")
            series_styles = inner_view.get("seriesStyles") or []
            series_aggs = [
                {"key": s.get("key"), "aggregation": s.get("aggregation")}
                for s in series_styles
                if isinstance(s, dict) and s.get("aggregation")
            ]
            if series_aggs:
                profile["view_config"]["series_aggregations"] = series_aggs
            default_filters = inner_view.get("defaultFilters") or []
            if default_filters:
                profile["view_config"]["default_filters"] = default_filters
            # Include palette if present
            palette = inner_view.get("palette") or {}
            if palette.get("colors"):
                profile["colors"] = palette.get("colors")[:5]

        # Include sample data if allowed
        if allow_llm_see_data:
            all_rows = viz.get("rows", [])
            # DEF-002: stats are a SAMPLE by design (prompt size), so say so
            # rather than letting the model read them as whole-dataset ranges.
            rows = all_rows[:_PROMPT_STATS_ROWS]
            if len(all_rows) > len(rows):
                profile["stats_from_sample"] = len(rows)
            if rows:
                profile["sample_rows"] = rows[:5]  # First 5 rows
                # Compute basic stats for numeric columns
                if rows and isinstance(rows[0], dict):
                    stats = {}
                    for col in viz.get("columns", []):
                        col_name = col if isinstance(col, str) else col.get("field", col.get("name"))
                        if col_name:
                            values = [r.get(col_name) for r in rows if r.get(col_name) is not None]
                            numeric_values = [v for v in values if isinstance(v, (int, float))]
                            if numeric_values:
                                stats[col_name] = {
                                    "min": min(numeric_values),
                                    "max": max(numeric_values),
                                    "sample_values": numeric_values[:3]
                                }
                            elif values:
                                unique = list(set(str(v) for v in values[:20]))
                                stats[col_name] = {
                                    "unique_count": len(unique),
                                    "sample_values": unique[:5]
                                }
                    if stats:
                        profile["column_stats"] = stats

        return profile

    async def run_stream(self, tool_input: Dict[str, Any], runtime_ctx: Dict[str, Any]) -> AsyncIterator[ToolEvent]:
        data = CreateArtifactInput(**tool_input)
        # Repair budget: leave headroom under the runner's 300s hard timeout
        # for the final persist + observation after the last repair round.
        _repair_deadline = time.monotonic() + 210

        # Early validation: require at least one visualization OR at least one file
        # (an image/PDF-only artifact is allowed when file_ids are provided).
        # Slides are exempt: a deck legitimately opens with a title, agenda or
        # narrative slide that carries no chart, and a whole deck may be
        # narrative-only.
        if (
            (not data.visualization_ids or len(data.visualization_ids) == 0)
            and not getattr(data, "file_ids", None)
            and data.mode != "slides"
        ):
            yield ToolStartEvent(type="tool.start", payload={"title": data.title or "Artifact"})
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {
                        "success": False,
                        "error": "No visualization_ids provided. At least one visualization is required to create an artifact.",
                    },
                    "observation": {
                        "summary": "Failed to create artifact: no visualization_ids provided",
                        "error": {
                            "type": "validation_error",
                            "message": "visualization_ids is required and must contain at least one visualization ID. Create visualizations using create_data first, then use their IDs here.",
                        },
                    },
                },
            )
            return

        yield ToolStartEvent(type="tool.start", payload={"title": data.title or "Artifact"})
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "init"})

        # Get runtime context
        sigkill_event = runtime_ctx.get("sigkill_event")
        report = runtime_ctx.get("report")
        user = runtime_ctx.get("user")
        organization = runtime_ctx.get("organization")
        db = runtime_ctx.get("db")
        context_hub = runtime_ctx.get("context_hub")
        organization_settings = runtime_ctx.get("settings")

        # Check privacy setting
        allow_llm_see_data = True
        if organization_settings:
            try:
                allow_llm_see_data = setting_enabled(organization_settings, "allow_llm_see_data", default=True)
            except Exception:
                allow_llm_see_data = True

        instruction_context_builder = runtime_ctx.get("instruction_context_builder") or (
            getattr(context_hub, "instruction_builder", None) if context_hub else None
        )

        # Get conversation history context (similar to create_data.py)
        context_view = runtime_ctx.get("context_view")
        messages_context = ""
        try:
            _messages_section_obj = getattr(context_view.warm, "messages", None) if context_view else None
            messages_context = _messages_section_obj.render() if _messages_section_obj else ""
        except Exception as e:
            logger.warning(f"Failed to extract messages context: {e}")
            messages_context = ""

        # Load images attached to the head completion for vision-capable models
        head_completion = runtime_ctx.get("head_completion")
        head_completion_id = str(head_completion.id) if head_completion else None
        completion_images = await self._load_completion_images(db, head_completion_id)

        # Validate model supports vision if images are present
        model = runtime_ctx.get("model")
        if completion_images and not getattr(model, "supports_vision", False):
            logger.info(f"Model doesn't support vision, skipping {len(completion_images)} completion images")
            completion_images = []

        # Note: Previous artifacts are now available via observation context (from create_artifact/read_artifact)
        # No need to fetch from DB - the planner can call read_artifact if needed

        # Fetch visualizations by ID from database
        visualizations: List[Dict[str, Any]] = []
        warnings: List[str] = []
        included_viz_ids: List[str] = []
        # DEF-009: viz id -> the step's stored query code, so a truncated result
        # can be re-read at full width instead of failing the build.
        step_code_by_viz: Dict[str, str] = {}
        # DEF-010: viz id -> the step it was read from, so a dataset the build
        # actually used can be written back and every render path can see it.
        step_by_viz: Dict[str, Any] = {}
        data_reductions: List[Dict[str, Any]] = []

        # Fetch all visualizations in a single batched query
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "loading_visualizations"})
        from app.models.query import Query
        from app.models.step import Step
        report_id = str(report.id) if report else None
        try:
            # populate_existing=True forces SQLAlchemy to refresh objects from DB
            # rather than returning stale identity-map copies (e.g. query.steps or
            # query.default_step may have been loaded before the step was created/updated)
            result = await db.execute(
                select(Visualization)
                .options(
                    selectinload(Visualization.query).selectinload(Query.default_step),
                    selectinload(Visualization.query).selectinload(Query.steps),
                )
                .where(Visualization.id.in_(data.visualization_ids))
                .execution_options(populate_existing=True)
            )
            fetched_vizs = {str(v.id): v for v in result.scalars().all()}
        except Exception as e:
            logger.exception("Failed to batch-fetch visualizations")
            fetched_vizs = {}
            warnings.append(f"Error fetching visualizations: {str(e)}")

        # Process each requested viz in order, validating and building entries
        for viz_id in data.visualization_ids:
            viz = fetched_vizs.get(viz_id)
            if viz is None:
                warnings.append(f"Visualization {viz_id} not found")
                continue

            # Validate viz belongs to the report
            if report_id and str(viz.report_id) != report_id:
                warnings.append(f"Visualization {viz_id} does not belong to this report")
                continue

            # Get the step with data (prefer default_step, fallback to latest step)
            step = None
            if viz.query and viz.query.default_step:
                step = viz.query.default_step
            elif viz.query and viz.query.steps:
                step = viz.query.steps[-1] if viz.query.steps else None

            # Check if the associated step is successful
            step_status = step.status if step else None
            if step_status != "success":
                _has_query = viz.query is not None
                _has_default = viz.query.default_step is not None if _has_query else False
                _steps_len = len(viz.query.steps) if _has_query and viz.query.steps else 0
                _default_step_id = getattr(viz.query, 'default_step_id', None) if _has_query else None
                logger.warning(
                    f"Visualization {viz_id} skipped: step_status='{step_status}', "
                    f"has_query={_has_query}, default_step_id={_default_step_id}, "
                    f"has_default_step={_has_default}, steps_count={_steps_len}"
                )
                warnings.append(f"Visualization {viz_id} skipped: step status is '{step_status or 'unknown'}' (not success)")
                continue

            # ★Upstream 0.0.526 rewrote this block to `rows = _all_rows[:100]`
            # plus a `total_row_count`, aiming at the same defect from the other
            # side: it keeps the truncation and merely stops mislabeling the
            # sample as the whole dataset. That is a REGRESSION here. DEF-002 and
            # DEF-010 below removed the 100-row cut precisely because it also
            # corrupted the column stats and the preview screenshot / stored
            # thumbnail, which rendered a 100-row prefix while the live dashboard
            # rendered the full set. Reporting an accurate total does not fix a
            # thumbnail drawn from the wrong data. Ours is kept; only upstream's
            # new `sample_row_count` field is adopted (below).
            #
            # Get data directly from step (like frontend does)
            #
            # DEF-002: this used to truncate to rows[:100] here, which silently
            # corrupted three separate things downstream — row_count (reported
            # to the model as the FULL count), the column stats, and the preview
            # screenshot / stored thumbnail, which rendered a 100-row prefix of
            # the dataset while the live dashboard rendered all of it. A 360-row
            # dataset previewed as the first 10 of 36 months and the KPI tiles
            # read 1.4B instead of 5.1B.
            #
            # Rows are now carried in full; each consumer takes its own slice:
            #   - the codegen prompt      -> _PROMPT_STATS_ROWS (_build_viz_profile)
            #   - the headless render     -> _RENDER_ROW_LIMIT (artifact_data)
            step_data = step.data if step else {}
            # DEF-010: resolved through the ONE accessor every render path uses.
            # `rows` is the display prefix (limit_row_count, default 1000);
            # `rows_artifact` is the same dataset under the larger artifact cap.
            # Reading it here and nowhere else is what stops the build being
            # judged on one dataset and rendered from another.
            _resolved = resolve_artifact_rows(step_data)
            rows = _resolved.rows
            if _resolved.source != "rows":
                logger.info(
                    "DEF-010: visualization %s using artifact-width rows (%d, preview held %d)",
                    viz_id, len(rows), len((step_data or {}).get("rows") or []),
                )
            raw_columns = _resolved.columns
            data_model = step.data_model if step else {}
            step_info = step_data.get("info") or {} if step_data else {}
            column_info = step_info.get("column_info") or {}

            # Keep raw column objects (with field/headerName) — matches the prompt contract
            columns = raw_columns

            # Extract field names for internal use (filterable columns, logging)
            column_fields = []
            for c in raw_columns:
                if isinstance(c, str):
                    column_fields.append(c)
                elif isinstance(c, dict):
                    col_name = c.get("field") or c.get("colId") or c.get("headerName") or c.get("name")
                    if col_name:
                        column_fields.append(col_name)

            # Build visualization entry
            view_dict = viz.view or {}
            query_id = str(viz.query_id) if viz.query_id else None

            # DEF-004: the STEP may already hold only a prefix (limit_row_count,
            # default 1000). `len(rows)` is then the prefix length, not the dataset,
            # so reporting it as row_count tells the model a truncated series is
            # complete — the same lie DEF-002 fixed one layer up. Take the true count
            # from the step's own info and carry the truncation forward explicitly.
            # DEF-010: both come from the accessor, which computes them from the
            # rows actually in hand. `rows_truncated` on the step describes the
            # PREVIEW copy and is stale whenever the artifact copy is in use —
            # trusting it made the gate refuse data that was complete.
            rows_total = _resolved.rows_total
            step_truncated = _resolved.truncated

            ventry = {
                "id": str(viz.id),
                "title": viz.title,
                "query_id": query_id,
                "view": self._trim_none(view_dict),
                "data_model_type": (view_dict.get("view") or {}).get("type") or view_dict.get("type"),
                "columns": columns,
                "column_info": column_info,
                # row_count is the TRUE dataset size, from `resolve_artifact_rows`
                # rather than upstream's `len(_all_rows)` — same meaning, our
                # accessor. `sample_row_count` is upstream 0.0.526's addition and
                # is adopted: it tells the model how many rows are actually in
                # hand, so it stops generating truncation workarounds for data it
                # can already see in full.
                "row_count": rows_total,
                "sample_row_count": len(rows),
                "rows": rows,
                # The profile the planner reads calls this same list
                # `sample_rows` (see _build_viz_profile), so generated code
                # reaches for either name. Expose both: picking the wrong one
                # used to yield [] and fail the whole deck with
                # "chart data contains no categories".
                "sample_rows": rows,
                "dataModel": data_model or {},
            }
            if step_truncated:
                ventry["rows_truncated"] = True
                ventry["rows_total"] = rows_total
                ventry["rows_available"] = len(rows)
                # DEF-009: the warning is NOT raised here any more. Recovery runs
                # before the gate and may make this dataset complete, and a
                # warning that the data is partial next to data that is whole is
                # its own kind of wrong. Whatever is still truncated after
                # recovery is warned about below.

            # DEF-009: keep the query that produced these rows. Re-running it is
            # the only way to get the rows the row limit dropped.
            if step is not None and getattr(step, "code", None):
                step_code_by_viz[str(viz.id)] = step.code
            if step is not None:
                step_by_viz[str(viz.id)] = step

            # Debug logging
            logger.info(f"Visualization {viz.title}: {len(rows)} rows, {len(column_fields)} columns: {column_fields[:5] if column_fields else 'none'}")
            if rows:
                logger.info(f"  Sample row keys: {list(rows[0].keys())[:5] if isinstance(rows[0], dict) else 'not a dict'}")

            visualizations.append(ventry)
            included_viz_ids.append(str(viz.id))

        # Resolve any embedded files (generated images / uploaded images or PDFs).
        # Scoped to the org; stored on the artifact content so the frontend can
        # fetch + inject them into the sandbox for the <BowFile> component.
        included_files: List[Dict[str, Any]] = []
        requested_file_ids = getattr(data, "file_ids", None) or []
        if requested_file_ids:
            try:
                file_result = await db.execute(
                    select(File).where(
                        File.id.in_([str(f) for f in requested_file_ids]),
                        File.organization_id == str(organization.id) if organization else File.organization_id.is_(None),
                    )
                )
                fetched_files = {str(f.id): f for f in file_result.scalars().all()}
            except Exception as e:
                logger.warning(f"create_artifact: failed to fetch files: {e}")
                fetched_files = {}
            for fid in requested_file_ids:
                f = fetched_files.get(str(fid))
                if f is None:
                    warnings.append(f"File {fid} not found or not in this organization")
                    continue
                included_files.append({
                    "id": str(f.id),
                    "content_type": f.content_type or "application/octet-stream",
                    "filename": f.filename,
                })

        # Early failure: if no valid visualizations AND no files were resolved,
        # fail like create_data does with tables. Slides may be narrative-only
        # (see the mode exemption in the early validation above).
        if not visualizations and not included_files and data.mode != "slides":
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {
                        "success": False,
                        "error": "No valid visualizations found. All requested visualization_ids were either not found, don't belong to this report, or have non-success step status.",
                    },
                    "observation": {
                        "summary": "Failed to create artifact: no valid visualizations resolved",
                        "error": {
                            "type": "no_valid_visualizations",
                            "message": "None of the requested visualization_ids could be used. Ensure visualizations exist, belong to this report, and have successful step status.",
                            "requested_ids": data.visualization_ids,
                            "warnings": warnings,
                        },
                    },
                },
            )
            return

        # ═══════════════════════════════════════════════════════════════════════
        # PHASE 1 — completeness gate. A dashboard may not be built on a prefix.
        #
        # DEF-004 made truncation VISIBLE (the step declares `rows_truncated`, and
        # the profile carries a warning). Visibility was not enough: the model read
        # the warning, built the dashboard anyway, rendered it, noticed the totals
        # were partial, threw it away and rebuilt with compact aggregates. That is
        # the correct outcome reached the expensive way — it happened on all three
        # live runs, costing a full generation cycle each time, and any run where
        # the model did NOT notice shipped a dashboard whose KPI tiles were the
        # sum of a 1,000-row prefix (56.4B against a true 98.9B, 10 of 17 months).
        #
        # So refuse instead of warning. The agent is told, before it spends a
        # generation, to aggregate the data down to something complete. The reason
        # reaches it intact because a self-declared failure now keeps its own
        # message (DEF-003).
        # ═══════════════════════════════════════════════════════════════════════
        # ═══════════════════════════════════════════════════════════════════════
        # DEF-009 — recover before refusing. See the note beside
        # `_RECOVERY_MAX_GROUP_COLUMNS` for why this exists and why it cannot
        # make anything worse: it either produces a COMPLETE dataset (re-read in
        # full, or aggregated over every source row) or it leaves the
        # visualization exactly as it was for the gate below to refuse.
        # ═══════════════════════════════════════════════════════════════════════
        if _recovery_enabled() and any(v.get("rows_truncated") for v in visualizations):
            yield ToolProgressEvent(type="tool.progress", payload={"stage": "recovering_data"})
            try:
                data_reductions = await recover_truncated_visualizations(
                    visualizations, step_code_by_viz, runtime_ctx, organization_settings
                )
            except Exception as e:
                logger.warning("DEF-009: recovery pass failed (%s) — falling back to the gate", e)
                data_reductions = []
            for _red in data_reductions:
                warnings.append(
                    f"Visualization {_red.get('visualization_id')}: {_red.get('notice')}"
                )

            # ★ DEF-010: write the recovered dataset back to the step it came
            # from. Recovery used to mutate the tool's in-memory copy only, so
            # the dashboard was WRITTEN against a complete re-read or an
            # aggregate and RENDERED against the untouched 1,000-row prefix —
            # which, after an aggregation, does not even carry the same columns
            # the generated code refers to. An artifact must be rendered from
            # what it was built on; this is what makes that true.
            for _red in data_reductions:
                _viz_id = str(_red.get("visualization_id") or "")
                _step = step_by_viz.get(_viz_id)
                _entry = next((v for v in visualizations if str(v.get("id")) == _viz_id), None)
                if _step is None or _entry is None:
                    continue
                try:
                    _step.data = store_artifact_dataset(
                        _step.data,
                        _entry.get("rows") or [],
                        _entry.get("columns") or [],
                        _red,
                    )
                    db.add(_step)
                    await db.commit()
                except Exception as _persist_exc:
                    # The build can still proceed — but it would then render from
                    # data it was not built on, which is the whole defect. Drop
                    # the recovery for this visualization so the gate decides on
                    # what a renderer can actually see.
                    logger.warning(
                        "DEF-010: could not persist the recovered dataset for %s (%s) — "
                        "reverting it so the completeness gate judges the stored data",
                        _viz_id, _persist_exc,
                    )
                    try:
                        await db.rollback()
                    except Exception:
                        pass
                    _resolved_again = resolve_artifact_rows(getattr(_step, "data", None))
                    _entry["rows"] = _resolved_again.rows
                    _entry["columns"] = _resolved_again.columns
                    _entry["row_count"] = _resolved_again.rows_total
                    _entry.pop("data_reduction", None)
                    if _resolved_again.truncated:
                        _entry["rows_truncated"] = True
                        _entry["rows_total"] = _resolved_again.rows_total
                        _entry["rows_available"] = _resolved_again.rows_used
            # A reverted visualization is no longer reduced — do not go on
            # disclosing a reduction that was rolled back.
            _reverted = {
                str(v.get("id")) for v in visualizations if v.get("rows_truncated")
            }
            data_reductions = [
                _red for _red in data_reductions
                if str(_red.get("visualization_id") or "") not in _reverted
            ]

        # Anything STILL truncated after recovery is genuinely partial — say so
        # before the gate refuses it, so the reason survives into the result.
        for _v in visualizations:
            if _v.get("rows_truncated"):
                warnings.append(
                    f"Visualization {_v.get('id')}: only {_v.get('rows_available')} of "
                    f"{_v.get('rows_total')} rows are available to this artifact (row limit). "
                    f"Totals computed over them are PARTIAL."
                )

        if _completeness_gate_enabled():
            partial = [v for v in visualizations if v.get("rows_truncated")]
            if partial:
                detail = "; ".join(
                    f"'{v.get('title') or v.get('id')}' has {v.get('rows_available')} of "
                    f"{v.get('rows_total')} rows"
                    for v in partial
                )
                logger.warning(
                    "PHASE1: refusing artifact build — %d of %d visualizations are truncated (%s)",
                    len(partial), len(visualizations), detail,
                )
                yield ToolEndEvent(
                    type="tool.end",
                    payload={
                        "output": {
                            "success": False,
                            "error": (
                                f"Cannot build this artifact: {detail}. Only part of the data "
                                f"reaches an artifact, taken in the query's own sort order, so "
                                f"the most recent periods or lowest-ranked groups are missing "
                                f"entirely and every total would be WRONG."
                            ),
                            "truncated_visualizations": [
                                {
                                    "id": v.get("id"),
                                    "title": v.get("title"),
                                    "rows_available": v.get("rows_available"),
                                    "rows_total": v.get("rows_total"),
                                }
                                for v in partial
                            ],
                        },
                        "observation": {
                            "summary": (
                                "Artifact not created: the data is incomplete "
                                f"({len(partial)} of {len(visualizations)} visualizations truncated)."
                            ),
                            "error": {
                                "type": "incomplete_visualization_data",
                                "message": (
                                    f"{detail}. Do NOT build the artifact from this data and do "
                                    f"NOT retry unchanged. Call create_data first to produce "
                                    f"PRE-AGGREGATED datasets that fit the row limit — group by "
                                    f"the dimensions the dashboard needs (month, branch, "
                                    f"category, channel) and let the database compute the SUM / "
                                    f"COUNT / AVG, instead of returning granular rows for the "
                                    f"artifact to add up. Then build the artifact from those."
                                ),
                            },
                        },
                    },
                )
                return

        # Build visualization profiles (privacy-aware)
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "building_profiles"})
        viz_profiles = [self._build_viz_profile(v, allow_llm_see_data) for v in visualizations]

        # Emit visualizations_resolved
        yield ToolProgressEvent(type="tool.progress", payload={
            "stage": "visualizations_resolved",
            "tool_name": "create_artifact",
            "visualizations": [
                {"id": v["id"], "title": v["title"], "type": v.get("data_model_type", "")}
                for v in visualizations
            ],
        })

        # Build instruction context
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "building_context"})
        instructions_context = ""
        try:
            if instruction_context_builder is not None:
                inst_section = await instruction_context_builder.build(categories=["dashboard", "visualization", "general"])
                instructions_context = inst_section.render() or ""
        except Exception:
            pass

        # Create artifact early with pending status so frontend can show it.
        # One deliverable per run per mode: if this run already built an
        # artifact of this mode, this build supersedes it in place (same id,
        # next version) instead of starting a second independent chain that
        # would leave an orphan the user can open. See _artifact_run_scope.
        from app.ai.tools.implementations._artifact_run_scope import (
            find_run_artifact,
            next_run_version,
        )

        artifact = await find_run_artifact(
            db,
            report_id=str(report.id) if report else None,
            completion_id=head_completion_id,
            mode=data.mode,
        )
        if artifact is not None:
            artifact.title = data.title or artifact.title or "Untitled Artifact"
            artifact.content = {}  # Empty content initially
            artifact.generation_prompt = data.prompt
            artifact.version = next_run_version(artifact)
            artifact.status = "pending"
            artifact.screenshot_base64 = None
            artifact.render_errors = None
        else:
            artifact = Artifact(
                report_id=str(report.id) if report else None,
                user_id=str(user.id) if user else None,
                organization_id=str(organization.id) if organization else None,
                title=data.title or "Untitled Artifact",
                mode=data.mode,
                content={},  # Empty content initially
                generation_prompt=data.prompt,
                completion_id=head_completion_id,
                version=1,
                status="pending",
            )
        db.add(artifact)
        await db.commit()
        await db.refresh(artifact)

        # Notify frontend that artifact is created (pending)
        yield ToolProgressEvent(
            type="tool.progress",
            payload={
                "stage": "artifact_created",
                "artifact_id": str(artifact.id),
                "status": "pending",
                "timing": False,
            }
        )

        # Build the prompt for generating React code
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "building_prompt"})

        # The deck's design system, resolved once: the same object is injected
        # into the slides prompt below and handed to the generated code as
        # `report['theme']`, so the prompt and the runtime cannot disagree about
        # which theme this deck is in. None when the registry is unavailable.
        #
        # ★The MODEL gets first refusal. It has just been shown the whole theme
        # index in the slides prompt, and `theme_id` is where it names its pick;
        # only when it names nothing — or names something the registry does not
        # know — does deterministic resolution decide. `theme_choice` records
        # WHICH of those happened so the console never has to guess whether a
        # deck's look was chosen or defaulted into.
        deck_theme, theme_choice = (
            _select_deck_theme(
                requested_theme_id=getattr(data, "theme_id", None),
                # ★`data.prompt` is the PLANNER's brief for this tool, not what
                # the person typed. A user who asked for "the ledger style" had
                # that paraphrased away and the deck resolved to McKinsey.
                # The conversation carries their actual words, so it is offered
                # first and the tool brief is the fallback.
                user_text=f"{messages_context or ''}\n{data.prompt or ''}",
                report=report,
                organization_settings=organization_settings,
            )
            if data.mode == "slides"
            else (None, None)
        )

        prompt = self._build_prompt(
            user_prompt=data.prompt,
            title=data.title,
            mode=data.mode,
            viz_profiles=viz_profiles,
            instructions_context=instructions_context,
            report_title=getattr(report, 'title', None) if report else None,
            allow_llm_see_data=allow_llm_see_data,
            messages_context=messages_context,
            image_count=len(completion_images),
            organization_settings=organization_settings,
            files=included_files,
            report=report,
            resolved_theme=deck_theme,
        )
        # Static reference goes in the system prompt so provider-side prompt
        # caching reuses it across artifact calls (page mode only — slides
        # keeps its single-prompt path).
        system_prompt = self._build_page_system_prompt() if data.mode == "page" else None

        # Stream from LLM
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "llm_generating"})
        llm = LLM(runtime_ctx.get("model"), usage_session_maker=async_session_maker)
        buffer = ""
        slides_detected = 0  # Track number of slides detected during streaming

        async for evt in llm.inference_stream_v2(
            messages=[Message(role="user", content=prompt)],
            system=system_prompt,
            images=completion_images if completion_images else None,
            usage_scope="create_artifact",
            usage_scope_ref_id=str(report.id) if report else None,
        ):
            if sigkill_event and sigkill_event.is_set():
                break
            if isinstance(evt, TextDeltaEvent):
                buffer += evt.text

            # For slides mode, detect new slides as they're generated
            if data.mode == "slides":
                # Count slide sections in buffer
                current_slides = buffer.count('<section class="slide"')
                if current_slides > slides_detected:
                    # New slide detected
                    for i in range(slides_detected, current_slides):
                        yield ToolProgressEvent(
                            type="tool.progress",
                            payload={
                                "stage": "slide_generated",
                                "slide_index": i,
                                "total_slides": current_slides,
                                "timing": False,
                            }
                        )
                    slides_detected = current_slides

            # Stream partial updates
            if len(buffer) % 100 == 0:  # Throttle updates
                yield ToolProgressEvent(
                    type="tool.progress",
                    payload={"stage": "generating", "chars": len(buffer), "timing": False}
                )

        # Check sigkill after LLM generation
        if sigkill_event and sigkill_event.is_set():
            # Update artifact to stopped status
            artifact.status = "stopped"
            await db.commit()
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {"success": False, "artifact_id": str(artifact.id), "error": "Stopped by user"},
                    "observation": {"summary": "Artifact creation stopped by user", "artifact_id": str(artifact.id), "stopped": True},
                },
            )
            return

        # Extract the code from the response
        code = self._extract_code(buffer, mode=data.mode)

        # DEF-008 — page guard, the exact counterpart of the slides guard below.
        # Same failure, same models: a one-line prose reply instead of code. For
        # slides that ends as "invalid syntax"; for a dashboard it ends as a
        # 289-character artifact stored `completed` whose whole body is an
        # English sentence, and the user meets it as "Dashboard failed to render:
        # Missing semicolon (3:8)". One strict retry demanding code only —
        # cheaper than a wasted render round-trip, and the same shape the slides
        # path has used successfully.
        if data.mode == "page" and not _looks_like_component_code(code):
            logger.warning(
                "DEF-008: page reply contained no component code (%d chars) — one strict retry",
                len(code or ""),
            )
            yield ToolProgressEvent(type="tool.progress", payload={"stage": "llm_generating", "retry": True})
            strict_prompt = (
                prompt
                + "\n\nIMPORTANT: Your previous reply contained no component code — it was a "
                  "description of what you were going to do. Respond with ONLY the complete "
                  "dashboard inside a single <script type=\"text/babel\"> block, starting with "
                  "`function App() {`. No prose, no explanation, no plan — code only."
            )
            retry_buffer = ""
            async for evt in llm.inference_stream_v2(
                messages=[Message(role="user", content=strict_prompt)],
                usage_scope="create_artifact",
                usage_scope_ref_id=str(report.id) if report else None,
            ):
                if sigkill_event and sigkill_event.is_set():
                    break
                if isinstance(evt, TextDeltaEvent):
                    retry_buffer += evt.text
            retry_code = self._extract_code(retry_buffer, mode=data.mode)
            if _looks_like_component_code(retry_code):
                code = retry_code
                buffer = retry_buffer
            else:
                # Still prose. Storing it would produce another artifact that only
                # the browser can reject, so fail here with the real reason — which
                # now survives to the model (see DEF-003).
                logger.warning("DEF-008: retry also returned no component code — failing the tool")
                yield ToolEndEvent(
                    type="tool.end",
                    payload={
                        "output": {
                            "success": False,
                            "error": (
                                "The model replied with a description instead of dashboard code, "
                                "twice. No artifact was created. Retry with an explicit instruction "
                                "to emit only the <script type=\"text/babel\"> block."
                            ),
                        },
                        "observation": {
                            "summary": "Dashboard not created: the model returned prose, not component code.",
                            "error": {
                                "type": "no_component_code",
                                "message": (
                                    "Two generations returned narration instead of a React component. "
                                    "Nothing was stored. Ask again for code only."
                                ),
                            },
                        },
                    },
                )
                return

        # Slides guard: some models (observed: grok-4.5) reply with a one-line
        # ═══════════════════════════════════════════════════════════════════════
        # PHASE 3 — parse the dashboard before storing it.
        #
        # DEF-008's guard above catches PROSE. It cannot catch malformed JSX,
        # because JSX is not Python and compile() will not read it. So run the
        # SAME Babel the browser runs (shipped in the image for the artifact
        # sandbox) inside the Chromium that is already there for PDF export.
        # A parse error is fed back verbatim — the model gets the real message,
        # at the same line and column the user would have seen.
        # ═══════════════════════════════════════════════════════════════════════
        if data.mode == "page" and _render_preflight_enabled():
            from app.services.artifact_preflight import check_artifact_code

            for _attempt in range(2):
                _ok, _perr = await check_artifact_code(code)
                if _ok:
                    break
                logger.warning("PHASE3: artifact failed to parse (attempt %d): %s", _attempt + 1, _perr)
                yield ToolProgressEvent(
                    type="tool.progress",
                    payload={"stage": "llm_generating", "retry": True},
                )
                fix_prompt = (
                    prompt
                    + "\n\nIMPORTANT: your previous dashboard did not compile. Babel reported:\n\n"
                    + f"    {_perr}\n\n"
                      "Fix that error and return the COMPLETE corrected dashboard inside a single "
                      "<script type=\"text/babel\"> block. Code only — no prose, no explanation."
                )
                fix_buffer = ""
                async for evt in llm.inference_stream_v2(
                    messages=[Message(role="user", content=fix_prompt)],
                    usage_scope="create_artifact",
                    usage_scope_ref_id=str(report.id) if report else None,
                ):
                    if sigkill_event and sigkill_event.is_set():
                        break
                    if isinstance(evt, TextDeltaEvent):
                        fix_buffer += evt.text
                fixed = self._extract_code(fix_buffer, mode=data.mode)
                if _looks_like_component_code(fixed):
                    code = fixed
                    buffer = fix_buffer
            else:
                # Both attempts still fail to parse. Storing this would put another
                # "Dashboard failed to render" in front of the user, so stop here
                # and say why — the message survives to the model (DEF-003).
                _ok, _perr = await check_artifact_code(code)
                if not _ok:
                    logger.warning("PHASE3: artifact still does not parse after 2 attempts — failing")
                    yield ToolEndEvent(
                        type="tool.end",
                        payload={
                            "output": {
                                "success": False,
                                "error": (
                                    f"The generated dashboard does not compile: {_perr}. "
                                    f"Two attempts were made. Nothing was stored."
                                ),
                            },
                            "observation": {
                                "summary": "Dashboard not created: the generated code does not compile.",
                                "error": {
                                    "type": "artifact_parse_error",
                                    "message": (
                                        f"Babel rejected the generated dashboard: {_perr}. "
                                        f"Rewrite it as a complete, self-contained component and "
                                        f"check every JSX tag is closed."
                                    ),
                                },
                            },
                        },
                    )
                    return

        # prose description and stop instead of emitting python-pptx code — the
        # extraction fallback then returns prose that later dies with
        # "invalid syntax (<string>, line 1)". Validate the code actually
        # compiles; if not, ONE strict retry demanding code-only output.
        if data.mode == "slides":
            def _compiles(src: str) -> bool:
                try:
                    compile(src, "<slides>", "exec")
                    return True
                except SyntaxError:
                    return False
            if not _compiles(code):
                yield ToolProgressEvent(type="tool.progress", payload={"stage": "llm_generating", "retry": True})
                strict_prompt = (
                    prompt
                    + "\n\nIMPORTANT: Your previous reply contained no Python code. "
                      "Respond with ONLY the complete python-pptx code inside a single "
                      "```python fenced block. No prose, no explanation, no plan — code only."
                )
                retry_buffer = ""
                async for evt in llm.inference_stream_v2(
                    messages=[Message(role="user", content=strict_prompt)],
                    usage_scope="create_artifact",
                    usage_scope_ref_id=str(report.id) if report else None,
                ):
                    if sigkill_event and sigkill_event.is_set():
                        break
                    if isinstance(evt, TextDeltaEvent):
                        retry_buffer += evt.text
                retry_code = self._extract_code(retry_buffer, mode=data.mode)
                if _compiles(retry_code):
                    code = retry_code
                    buffer = retry_buffer

        # ═══════════════════════════════════════════════════════════════════════
        # Mode-specific processing: slides uses python-pptx, page skips to save
        # ═══════════════════════════════════════════════════════════════════════

        pptx_path: Optional[str] = None
        pptx_success: bool = True
        pptx_error: Optional[str] = None
        pptx_repair_attempts: int = 0
        preview_images: List[str] = []
        layout_issues: List[Dict[str, Any]] = []
        # The check's own verdict (status/reason/counts), kept separate from
        # `layout_issues` above so "we could not check" (status="unavailable")
        # is never indistinguishable from "we checked and it was clean"
        # (status="checked", issues=[]). None when the flag is off or the
        # check never ran.
        layout_check_result: Optional[Any] = None
        # The two theme post-passes, three-state exactly like the layout check:
        #   None                        -> did not run (flag off, no theme, no deck)
        #   {"status": "unavailable"}   -> ran and could not finish
        #   the pass's own summary dict -> ran
        # "did not run" and "ran clean" must never look the same; a deck whose
        # prohibitions were never enforced is not a deck that had none to break.
        furniture_result: Optional[Dict[str, Any]] = None
        enforcement_result: Optional[Dict[str, Any]] = None

        if data.mode == "slides":
            # ═══════════════════════════════════════════════════════════════════
            # SLIDES MODE: Execute python-pptx code (repairing in-tool on
            # failure, mirroring page-mode render validation) and generate
            # previews.
            # ═══════════════════════════════════════════════════════════════════
            report_data = {
                "id": str(report.id) if report else None,
                "title": getattr(report, "title", None) if report else None,
                # ★NOT `getattr(report, "theme", …)` — see _report_theme_payload.
                # That attribute does not exist, so this key was None on every
                # deck this product has ever generated.
                "theme": _report_theme_payload(report, deck_theme),
            }

            uploads_dir = Path(__file__).parent.parent.parent.parent.parent / "uploads" / "pptx"
            uploads_dir.mkdir(parents=True, exist_ok=True)
            output_path = uploads_dir / f"{artifact.id}.pptx"

            _pptx_result: Optional[Dict[str, Any]] = None
            async for _item in self._execute_and_repair_pptx(
                code,
                visualizations,
                report_data,
                output_path,
                await load_image_bytes(db, included_files),
                runtime_ctx,
                deadline_monotonic=_repair_deadline,
            ):
                if isinstance(_item, dict):
                    _pptx_result = _item
                else:
                    yield _item
            if _pptx_result is not None:
                code = _pptx_result["code"]
                pptx_success = bool(_pptx_result["ok"])
                pptx_error = _pptx_result["error"]
                pptx_repair_attempts = int(_pptx_result["repair_attempts"] or 0)
            if pptx_success:
                pptx_path = str(output_path)

            # ═══════════════════════════════════════════════════════════════════
            # THEME POST-PASSES — paint the theme's furniture onto the saved
            # deck, then enforce its prohibitions.
            #
            # ORDER IS LOAD-BEARING, twice over.
            #
            # Painter before enforcer: the painter adds shapes of its own, and
            # enforcement is what squares their corners and strips their
            # shadows. Run the other way round, the pass that guarantees the
            # theme's rules would be the one thing on the deck exempt from them.
            #
            # Both before previews and before the layout check: LibreOffice
            # renders the file as it stands, so a deck painted after its
            # previews would ship images of a deck nobody gets, and the layout
            # check would measure a file two passes out of date.
            #
            # ★OFF THE CRITICAL PATH, same three ways as the layout check
            # below: they run only once the .pptx is written AND `pptx_path` is
            # recorded, so nothing here can unmake a deck that built; both
            # functions are contracted never to raise; and both are wrapped
            # anyway, because an ImportError on the module comes from the
            # import, not from the call. A theme is a design improvement — it
            # must never be the reason a deck that built does not ship.
            # ═══════════════════════════════════════════════════════════════════
            if pptx_success and pptx_path and deck_theme is not None:
                _theme_payload = _theme_as_dict(deck_theme)
                _theme_layout = _theme_layout_for(_theme_payload.get("id"))

                if _deck_theme_furniture_enabled():
                    try:
                        from app.ai.decks.motifs import (  # noqa: PLC0415 (optional)
                            paint_theme_furniture,
                        )

                        furniture_result = paint_theme_furniture(
                            pptx_path, _theme_payload, _theme_layout, logger=logger
                        )
                    except Exception as e:
                        # Recorded, not silent: "we tried and could not" is a
                        # different fact from "we never tried", and the deck is
                        # unaffected either way.
                        logger.warning(
                            f"Deck theme furniture pass failed; the deck is unaffected: {e}"
                        )
                        furniture_result = {"status": "unavailable", "reason": str(e)}

                if _deck_theme_enforcement_enabled():
                    try:
                        from app.ai.code_execution.pptx_executor import (  # noqa: PLC0415
                            enforce_theme_rules,
                        )

                        enforcement_result = enforce_theme_rules(
                            pptx_path, _theme_payload, logger=logger
                        )
                    except Exception as e:
                        logger.warning(
                            f"Deck theme enforcement pass failed; the deck is unaffected: {e}"
                        )
                        enforcement_result = {"status": "unavailable", "reason": str(e)}

            # Previews are rendered by LibreOffice, which is a separate concern
            # from building the deck: it can be missing an import filter or be
            # misconfigured while the .pptx itself is perfectly valid. Failing
            # the artifact here would also make the export endpoint refuse to
            # serve a deck the user can open, so a preview failure only costs
            # the preview.
            if pptx_success and pptx_path:
                yield ToolProgressEvent(
                    type="tool.progress",
                    payload={"stage": "generating_previews"}
                )
                try:
                    preview_service = PptxPreviewService(logger=logger)
                    preview_images = preview_service.generate_previews(
                        pptx_path=Path(pptx_path),
                        artifact_id=str(artifact.id),
                    )
                except Exception as e:
                    logger.warning(
                        f"PPTX preview generation failed; deck is still downloadable: {e}"
                    )

            # ═══════════════════════════════════════════════════════════════════
            # DEF-011: measure the SAVED deck.
            #
            # python-pptx has no font metrics, so generated code sizes every text
            # box by guess and a guess that is wrong by three inches raises
            # nothing — the code runs, the file writes, `pptx_success` is True and
            # the deck ships with a paragraph sitting on the card below it. The
            # repair loop above only ever sees exceptions, so it cannot see this.
            #
            # ★OFF THE CRITICAL PATH, three ways over: it runs only after the
            # .pptx is written AND `pptx_path` is already recorded, so nothing
            # here can unmake a deck that built; `check_deck_layout_detailed`
            # never raises and degrades to status="unavailable"; and it is still
            # wrapped, because an ImportError on the module itself would be
            # raised by the import, not by the function. Whatever happens, the
            # deck is delivered — at worst without a verdict.
            # ═══════════════════════════════════════════════════════════════════
            if pptx_success and pptx_path and _deck_layout_check_enabled():
                yield ToolProgressEvent(
                    type="tool.progress",
                    payload={"stage": "checking_layout"},
                )
                try:
                    from app.ai.code_execution.pptx_lint import (
                        check_deck_layout_detailed,
                    )

                    layout_check_result = await check_deck_layout_detailed(
                        Path(pptx_path), log=logger
                    )
                    layout_issues = list(layout_check_result.issues or [])
                except Exception as e:
                    # Advisory only. Leave `layout_check_result` as None — the
                    # downstream reporting treats None as "the check never ran",
                    # which is exactly what happened, and is deliberately NOT the
                    # same as "checked and clean".
                    logger.warning(
                        f"Deck layout check failed; the deck is unaffected: {e}"
                    )
                    layout_check_result = None
                    layout_issues = []

        # ═══════════════════════════════════════════════════════════════════════
        # Page mode: render-validate (and repair in-tool) BEFORE persisting.
        # "completed" must mean "renders without fatal errors" — the old flow
        # committed completed first and only then discovered render errors,
        # pushing every repair through a full outer planner iteration.
        # Validation runs regardless of vision support or allow_llm_see_data:
        # errors are not data. Only the screenshot ATTACHMENT is gated below.
        # ═══════════════════════════════════════════════════════════════════════
        screenshot_base64: Optional[str] = None
        render_errors: list[str] = []
        render_clean = True
        repair_attempts = 0
        thumbnail_html: Optional[str] = None
        artifact_data: Optional[Dict[str, Any]] = None

        if data.mode == "page":
            artifact_data = {
                "report": {
                    "id": str(report.id) if report else None,
                    "title": getattr(report, "title", None) if report else None,
                    # Same fix as the slides branch above — the attribute read
                    # here never existed. Page mode takes no pptx theme.
                    "theme": _report_theme_payload(report),
                },
                # ★`_render_visualizations`, not raw `visualizations`. Upstream
                # has no such helper and its 0.0.526 hunk applied CLEANLY over
                # our call, silently dropping it — a clean apply is not a safe
                # apply. The helper is DEF-002/DEF-010: it hands the headless
                # preview every row (so the thumbnail cannot disagree with the
                # live dashboard, which gets them all) and carries any
                # cap/aggregation notice into the rendered payload.
                "visualizations": self._render_visualizations(visualizations),
            }
            # Inline embedded files as data URIs so the headless render
            # (which has no auth context) can show images/PDFs via <BowFile>.
            if included_files:
                artifact_data["files"] = await self._build_file_datauris(db, included_files)

            _validate_result: Optional[Dict[str, Any]] = None
            async for _item in self._validate_and_repair_stream(
                code, artifact_data, data.mode, runtime_ctx, deadline_monotonic=_repair_deadline,
            ):
                if isinstance(_item, dict):
                    _validate_result = _item
                else:
                    yield _item
            if _validate_result is not None:
                code = _validate_result["code"]
                render_clean = bool(_validate_result["clean"])
                screenshot_base64 = _validate_result["screenshot"]
                render_errors = list(_validate_result["errors"] or [])
                repair_attempts = int(_validate_result["repair_attempts"] or 0)
            try:
                thumbnail_html = self._build_thumbnail_html(artifact_data, code, mode=data.mode)
            except Exception as e:
                logger.warning(f"Thumbnail HTML build failed: {e}")
                thumbnail_html = None

        yield ToolProgressEvent(type="tool.progress", payload={"stage": "saving_artifact"})

        # Build content object (code is final — post-repair when repair ran)
        content: Dict[str, Any] = {
            "code": code,
            "visualization_ids": included_viz_ids,
        }

        # Embedded files (generated images / uploaded images/PDFs) referenced by
        # <BowFile id=...>. Stored as {id, content_type, filename}; the frontend
        # resolves the bytes and injects them into ARTIFACT_DATA.files.
        if included_files:
            content["files"] = included_files

        # DEF-009: what the dashboard was actually built from, stored WITH the
        # dashboard. A reduction that lives only in a tool result is invisible to
        # anyone who opens the artifact tomorrow.
        if data_reductions:
            content["data_reduction"] = data_reductions

        # DEF-010: if ANY cap or aggregation shaped what this artifact shows, say
        # so on the artifact itself, in one sentence per visualization. Stored on
        # the artifact (so it survives), and handed to the renderer below (so it
        # can be shown). A figure computed over less than the whole dataset — or
        # over an aggregate of it — must never be presented as if it were not.
        _data_notices = []
        for _v in visualizations:
            _red = _v.get("data_reduction")
            if _red and _red.get("notice"):
                _data_notices.append({
                    "visualization_id": _v.get("id"),
                    "title": _v.get("title"),
                    "notice": str(_red["notice"]),
                    "kind": "aggregated",
                })
            elif _v.get("rows_truncated"):
                _data_notices.append({
                    "visualization_id": _v.get("id"),
                    "title": _v.get("title"),
                    "notice": (
                        f"Showing {int(_v.get('rows_available') or 0):,} of "
                        f"{int(_v.get('rows_total') or 0):,} rows — figures here are PARTIAL."
                    ),
                    "kind": "truncated",
                })
        if _data_notices:
            content["data_notice"] = _data_notices

        # Add slides-specific content
        if data.mode == "slides" and preview_images:
            content["preview_images"] = preview_images

        # ★Upstream has no counterpart for the block below — DEF-011's layout
        # verdict and the Phase 4 insight panel are fork-only, absent from both
        # v0.0.525 and v0.0.526 (verified by grep on both tags). The 3-way merge
        # therefore produced an EMPTY 'theirs' side, which is positional drift,
        # not an upstream deletion. Taking 'theirs' here would have silently
        # removed 145 lines of working code with no conflict to warn anyone.
        # DEF-011: the layout check's whole verdict, not just its issue list —
        # stored WITH the deck for the same reason DEF-009 stores data_reduction
        # with the dashboard (a tool-result note is invisible to anyone who
        # opens the artifact tomorrow). `status` distinguishes "checked, clean"
        # from "could not check" from "partial" so neither is ever presented as
        # the other. `layout_issues` is kept populated on its own too — Phase 2
        # feeds it back to the model so it can reflow the offending slide, and
        # anything already reading that key for backwards compatibility still
        # finds it.
        # The deck's design system, stored WITH the deck — same reason DEF-009
        # stores data_reduction with the dashboard: which theme a deck is in,
        # and whether the model chose it, is invisible to anyone who opens the
        # artifact tomorrow if it lives only in a tool result. Keys are present
        # only for passes that RAN, so absence keeps meaning "did not run".
        if theme_choice is not None:
            _deck_theme_record = dict(theme_choice)
            if furniture_result is not None:
                _deck_theme_record["furniture"] = furniture_result
            if enforcement_result is not None:
                _deck_theme_record["enforcement"] = enforcement_result
            content["deck_theme"] = _deck_theme_record

        if layout_issues:
            content["layout_issues"] = layout_issues
        if layout_check_result is not None:
            content["layout_check"] = {
                "status": layout_check_result.status,
                "reason": layout_check_result.reason,
                "slides_total": layout_check_result.slides_total,
                "slides_measured": layout_check_result.slides_measured,
                "issues": layout_issues,
            }

        # ═══════════════════════════════════════════════════════════════════════
        # PHASE 4 — the insight panel. A dashboard is a wall of numbers; this is
        # the sentence a person reads first.
        #
        # Generated from the FINAL data, after the completeness gate has already
        # guaranteed it is whole — so the summary can never describe a prefix.
        # Every figure it claims is then checked against that same data, because
        # a summariser will occasionally produce a number from nowhere: during
        # testing one reported an average order value of 11,499 against a true
        # 11,488.57, present in no query result, no code and no tool output.
        # A finding citing an ungrounded figure is DROPPED, not published.
        #
        # Best-effort throughout: a dashboard without a summary is a fine
        # dashboard, so nothing here may fail the tool.
        # ═══════════════════════════════════════════════════════════════════════
        if data.mode == "page" and _insights_enabled() and visualizations:
            try:
                from app.services.artifact_insights import (
                    build_prompt as _ins_prompt,
                    parse_response as _ins_parse,
                    verify_findings as _ins_verify,
                )

                yield ToolProgressEvent(type="tool.progress", payload={"stage": "generating_insights"})
                _ibuf = ""
                async for _evt in llm.inference_stream_v2(
                    messages=[Message(role="user", content=_ins_prompt(data.title or "Dashboard", visualizations))],
                    usage_scope="create_artifact_insights",
                    usage_scope_ref_id=str(report.id) if report else None,
                ):
                    if sigkill_event and sigkill_event.is_set():
                        break
                    if isinstance(_evt, TextDeltaEvent):
                        _ibuf += _evt.text

                _parsed = _ins_parse(_ibuf)

                # ★★★One retry when findings were rejected. Rejection used to be
                # silent and final: a run that lost four of five findings shipped
                # a four-chart dashboard with a single bullet, and the model was
                # never told what failed. Telling it the exact rejected sentences
                # costs one call and usually recovers most of them.
                if _parsed:
                    _kept, _rejected = _ins_verify(_parsed.get("findings") or [], visualizations)
                    if _rejected and not (sigkill_event and sigkill_event.is_set()):
                        logger.info("PHASE4: %d finding(s) rejected — asking once more", len(_rejected))
                        _retry_prompt = (
                            _ins_prompt(data.title or "Dashboard", visualizations)
                            + "\n\nYOUR PREVIOUS ANSWER WAS PARTLY REJECTED.\n"
                            + "These sentences cited a figure that appears nowhere in the data above:\n"
                            + "\n".join(f"  - {r}" for r in _rejected[:5])
                            + "\n\nWrite the summary again. Keep what was accepted, and either correct "
                              "the rejected points using only figures present above, or replace them "
                              "with different observations. Do not repeat a rejected figure."
                        )
                        try:
                            _rbuf = ""
                            async for _evt in llm.inference_stream_v2(
                                messages=[Message(role="user", content=_retry_prompt)],
                                usage_scope="create_artifact_insights",
                                usage_scope_ref_id=str(report.id) if report else None,
                            ):
                                if sigkill_event and sigkill_event.is_set():
                                    break
                                if isinstance(_evt, TextDeltaEvent):
                                    _rbuf += _evt.text
                            _reparsed = _ins_parse(_rbuf)
                            if _reparsed:
                                _kept2, _rejected2 = _ins_verify(
                                    _reparsed.get("findings") or [], visualizations
                                )
                                # Keep the retry only if it actually did better.
                                if len(_kept2) > len(_kept):
                                    logger.info(
                                        "PHASE4: retry recovered %d finding(s) (%d -> %d)",
                                        len(_kept2) - len(_kept), len(_kept), len(_kept2),
                                    )
                                    _parsed, _kept, _rejected = _reparsed, _kept2, _rejected2
                        except Exception as _re:
                            logger.warning("PHASE4: retry failed (%s) — keeping the first answer", _re)

                    if _rejected:
                        logger.warning(
                            "PHASE4: dropped %d ungrounded finding(s): %s",
                            len(_rejected), " | ".join(_rejected[:3]),
                        )

                    # ★★★The headline is verified too. It used to be published as
                    # written, on the reasoning that it was "prose about direction
                    # and shape rather than a figure-bearing claim". In practice it
                    # carries figures AND a period: one real run opened with "In
                    # October 2025, in-store sales made up roughly 84% of net
                    # sales" above a dashboard spanning 36 months. Whatever the
                    # findings are held to, the sentence at the top — the one most
                    # people read and no one else checks — must be held to as well.
                    _headline = (_parsed.get("headline") or "").strip()
                    if _headline:
                        _hkept, _hrej = _ins_verify([{"text": _headline}], visualizations)
                        if not _hkept:
                            logger.warning(
                                "PHASE4: headline dropped as ungrounded: %s",
                                (_hrej[0] if _hrej else _headline)[:160],
                            )
                            _headline = None

                    from datetime import datetime as _dt, timezone as _tz

                    content["insights"] = {
                        "headline": _headline,
                        "findings": _kept,
                        "rejected_count": len(_rejected),
                        # The panel shows when the summary was written. It matters
                        # because the summary describes the data AS BUILT: if the
                        # dashboard is later rebuilt on fresher data, a stale
                        # narrative next to current numbers would be its own kind
                        # of wrong.
                        "generated_at": _dt.now(_tz.utc).isoformat(),
                    }
                    logger.info("PHASE4: stored %d grounded finding(s)", len(_kept))
                else:
                    logger.warning("PHASE4: summariser returned no parseable JSON — no panel")
            except Exception as _ie:
                logger.warning("PHASE4: insight generation failed (%s) — continuing without a panel", _ie)

        # Update the pending artifact with content and mark as completed
        artifact.content = content
        if data.mode == "slides":
            artifact.status = "completed" if pptx_success else "failed"
            # Persist the executor error like page mode persists render errors,
            # so read_artifact and later repairs can see WHY the deck failed.
            if pptx_error:
                artifact.render_errors = [pptx_error]
        else:
            artifact.status = "completed" if render_clean else "failed"

        # Set pptx_path for slides mode
        if pptx_path:
            artifact.pptx_path = pptx_path

        # Persist screenshot and render errors for later retrieval (read_artifact)
        if screenshot_base64 or render_errors:
            artifact.screenshot_base64 = screenshot_base64
            artifact.render_errors = render_errors or None

        await db.commit()
        await db.refresh(artifact)

        # ★Our 30-line block here was DELETED, not lost. Upstream 0.0.526 moved
        # this work ~250 lines earlier into _validate_and_repair_stream, which
        # already builds artifact_data, inlines file data URIs, captures the
        # screenshot and sets render_errors/repair_attempts. Keeping ours would
        # have re-initialised screenshot_base64=None and render_errors=[],
        # WIPING that validation result, and then run a second headless capture.
        if data.mode == "page" and thumbnail_html is not None and render_clean:
            # Generate thumbnail in background (for stored thumbnail, non-blocking)
            asyncio.create_task(
                self._generate_thumbnail_background(
                    artifact_id=str(artifact.id),
                    html_content=thumbnail_html,
                    mode=data.mode,
                )
            )
        elif preview_images:
            # For slides mode, use the first preview image as thumbnail
            first_preview = Path(__file__).parent.parent.parent.parent.parent / "uploads" / preview_images[0]
            if first_preview.exists():
                artifact.thumbnail_path = preview_images[0]
                await db.commit()

        # Slides mode that failed pptx execution (after bounded in-tool
        # repair): return a structured failure so the planner sees the real
        # error — mirror of the page-mode branch below. The artifact row is
        # persisted as status="failed" for debugging, but it is not presented
        # as a working deck.
        if data.mode == "slides" and not pptx_success:
            _error_msg = pptx_error or "unknown pptx execution error"
            slides_failure_observation: Dict[str, Any] = {
                "summary": (
                    f"Artifact '{data.title or 'Untitled'}' failed to build the presentation "
                    f"after {pptx_repair_attempts} in-tool repair attempt(s). "
                    f"Error: {_error_msg}"
                ),
                "error": {
                    "type": "pptx_execution_failed",
                    "message": _error_msg,
                    "repair_attempts": pptx_repair_attempts,
                    "remediation": (
                        "The generated python-pptx code does not execute. Call edit_artifact with an "
                        "edit_prompt that quotes the exact error above, or create_artifact to rebuild "
                        "with a simpler deck if the error persists."
                    ),
                },
                "artifact_id": str(artifact.id),
                "mode": data.mode,
                "visualization_count": len(visualizations),
                "visualization_ids": included_viz_ids,
                "render_errors": [_error_msg],
            }
            if warnings:
                slides_failure_observation["warnings"] = warnings
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {
                        "success": False,
                        "artifact_id": str(artifact.id),
                        "error": f"PPTX execution failed: {_error_msg}",
                        "code_preview": {
                            "language": "python",
                            "code": code,
                            "collapsed_default": True,
                        },
                    },
                    "observation": slides_failure_observation,
                },
            )
            return

        # Page mode that failed render validation (after bounded in-tool
        # repair): return a structured failure so the planner sees the real
        # errors. The artifact row is persisted as status="failed" for
        # debugging, but it is not presented as a working dashboard.
        if data.mode == "page" and not render_clean:
            fatal_errors = self.fatal_render_errors(render_errors)
            first_error = fatal_errors[0] if fatal_errors else (render_errors[0] if render_errors else "unknown render error")
            failure_observation: Dict[str, Any] = {
                "summary": (
                    f"Artifact '{data.title or 'Untitled'}' failed render validation with "
                    f"{len(fatal_errors)} fatal error(s) after {repair_attempts} in-tool repair attempt(s). "
                    f"First error: {first_error}"
                ),
                "error": {
                    "type": "render_validation_failed",
                    "message": first_error,
                    "render_errors": render_errors,
                    "repair_attempts": repair_attempts,
                    "remediation": (
                        "The generated code does not render. Call edit_artifact with an edit_prompt "
                        "that quotes the exact error(s) above, or create_artifact to rebuild with a "
                        "simpler layout if the errors persist."
                    ),
                },
                "artifact_id": str(artifact.id),
                "mode": data.mode,
                "visualization_count": len(visualizations),
                "visualization_ids": included_viz_ids,
                "render_errors": render_errors,
            }
            if warnings:
                failure_observation["warnings"] = warnings
            _model = runtime_ctx.get("model")
            if screenshot_base64 and allow_llm_see_data and _model and getattr(_model, "supports_vision", False):
                failure_observation["images"] = [{
                    "data": screenshot_base64,
                    "media_type": "image/png",
                    "source_type": "base64",
                }]
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {
                        "success": False,
                        "artifact_id": str(artifact.id),
                        "error": f"Render validation failed: {first_error}",
                        "code_preview": {
                            "language": "jsx",
                            "code": code,
                            "collapsed_default": True,
                        },
                    },
                    "observation": failure_observation,
                },
            )
            return

        output = CreateArtifactOutput(
            artifact_id=str(artifact.id),
            code=code,
            mode=data.mode,
            title=data.title,
            version=artifact.version,
        ).model_dump()

        # Add UI preview fields (similar to read_artifact)
        code_lines = code.count('\n') + 1 if code else 0
        output["artifact_preview"] = {
            "artifact_id": str(artifact.id),
            "title": data.title or "Untitled",
            "mode": data.mode,
            "version": artifact.version,
            "code_stats": {
                "chars": len(code),
                "lines": code_lines,
            },
            "visualization_ids": included_viz_ids,
            "visualization_count": len(visualizations),
        }
        # Code for collapsible toggle (collapsed by default in UI)
        output["code_preview"] = {
            "language": "jsx",
            "code": code,
            "collapsed_default": True,
        }

        # DEF-011: turn the layout check's verdict into one plain sentence, the
        # same job DEF-009 does for a data reduction — a fact that is true of
        # the artifact but invisible unless it rides the tool result back to
        # the model. Never claims the deck is clean when it is only unchecked.
        layout_notice: Optional[str] = None
        if layout_check_result is not None:
            if layout_check_result.status == "unavailable":
                reason = layout_check_result.reason or "an internal error"
                layout_notice = (
                    f"Slide layout could not be checked ({reason}) — the deck may "
                    "still have text running off a slide; this was not verified."
                )
            else:
                by_slide: Dict[int, int] = {}
                for issue in layout_issues:
                    slide_no = issue.get("slide")
                    if slide_no is not None:
                        by_slide[slide_no] = by_slide.get(slide_no, 0) + 1
                coverage = ""
                if layout_check_result.status == "partial":
                    coverage = (
                        f" (checked {layout_check_result.slides_measured} of "
                        f"{layout_check_result.slides_total} slides"
                    )
                    if layout_check_result.reason:
                        coverage += f", {layout_check_result.reason}"
                    coverage += ")"
                if by_slide:
                    slide_list = ", ".join(str(s) for s in sorted(by_slide))
                    n_slides = len(by_slide)
                    plural = "s" if n_slides != 1 else ""
                    verb = "have" if n_slides != 1 else "has"
                    layout_notice = (
                        f"{n_slides} slide{plural} {verb} text running past the "
                        f"slide edge or its box (slides {slide_list}){coverage}."
                    )
                elif layout_check_result.status == "partial":
                    layout_notice = (
                        "Slide layout check only covered "
                        f"{layout_check_result.slides_measured} of "
                        f"{layout_check_result.slides_total} slides"
                    )
                    if layout_check_result.reason:
                        layout_notice += f" ({layout_check_result.reason})"
                    layout_notice += "; no issues found in what was checked."
                # status == "checked" and no issues → nothing to say, same as
                # DEF-010 staying silent when no reduction shaped a visualization.

        # Build observation message
        summary_msg = f"Created artifact '{data.title or 'Untitled'}' with {len(code)} characters of code"
        # DEF-009: the row count actually used, in the first sentence the agent
        # reads. Reduction that is not stated is reduction that gets reported as
        # a complete figure by whoever quotes the dashboard next.
        if data_reductions:
            summary_msg += (
                ". Input data was too large for one artifact and was reduced first: "
                + "; ".join(
                    f"'{r.get('title') or r.get('visualization_id')}' {r.get('notice')}"
                    for r in data_reductions
                )
            )
        if data.mode == "slides":
            # Name the design system in the sentence the agent reads. A deck
            # whose theme is never stated is a deck whose theme nobody can hold
            # to account — and the model needs to know whether the id it sent
            # was honoured, because a silent fallback teaches it nothing.
            if theme_choice is not None:
                _theme_label = theme_choice.get("name") or theme_choice.get("id")
                _method = theme_choice.get("method")
                if _method == "model":
                    summary_msg += (
                        f". Deck theme: {_theme_label} ({theme_choice.get('id')}) — "
                        "the theme you named."
                    )
                elif _method == "resolved_unknown_id":
                    summary_msg += (
                        f". Deck theme: {_theme_label} ({theme_choice.get('id')}) — "
                        f"'{theme_choice.get('requested_id')}' is not a theme in the index, "
                        "so the theme was resolved from the request instead. "
                        "Copy the id exactly as the index prints it to choose one."
                    )
                else:
                    summary_msg += (
                        f". Deck theme: {_theme_label} ({theme_choice.get('id')}) — "
                        "resolved from the request (no theme_id was given)."
                    )
            if pptx_repair_attempts:
                summary_msg += f". PPTX execution passed after {pptx_repair_attempts} in-tool repair attempt(s)."
            if preview_images:
                summary_msg += f" Generated {len(preview_images)} slide preview images."
            else:
                summary_msg += (
                    " The deck was built, but slide preview images could not be generated on this "
                    "server — the PPTX file is still downloadable."
                )
        # ★Ours, and it must stay AHEAD of upstream's page branch. Upstream only
        # words the passing case, so on a dirty render its branch appends
        # "Render validation passed." to the very message the agent reads to
        # decide whether to fix anything. `render_clean` is the authoritative
        # signal, not `render_errors` being non-empty — since 0.0.526's repair
        # loop that list also carries NON-FATAL console errors on a clean render.
        elif data.mode == "page" and not render_clean:
            _first = render_errors[0] if render_errors else "no error text captured"
            summary_msg += f". RENDER FAILED with {len(render_errors)} error(s): {_first}"
            if len(render_errors) > 1:
                summary_msg += f" (and {len(render_errors) - 1} more)"
            summary_msg += ". The dashboard code has a bug — use edit_artifact to fix the specific error."
        elif data.mode == "page":
            if repair_attempts:
                summary_msg += f". Render validation passed after {repair_attempts} in-tool repair attempt(s)."
            else:
                summary_msg += ". Render validation passed."
            console_warnings = [e for e in render_errors if e.startswith("[console.error]")]
            if console_warnings:
                summary_msg += f" {len(console_warnings)} non-fatal console error(s) were logged."
        if layout_notice:
            summary_msg += ". " + layout_notice

        # Screenshot attachment is gated on privacy + vision (the screenshot
        # shows the data); validation itself already ran unconditionally.
        # ★`_attach_screenshot` is referenced further down, outside this hunk —
        # dropping upstream's side here would be a NameError, not a lost message.
        _model = runtime_ctx.get("model")
        _attach_screenshot = bool(
            screenshot_base64 and allow_llm_see_data and _model and getattr(_model, "supports_vision", False)
        )
        if _attach_screenshot:
            summary_msg += " Screenshot of the rendered dashboard is attached — review it for visual correctness."

        observation: Dict[str, Any] = {
            "summary": summary_msg,
            "artifact_id": str(artifact.id),
            "mode": data.mode,
            "visualization_count": len(visualizations),
            "visualization_ids": included_viz_ids,
        }
        if render_errors:
            observation["render_errors"] = render_errors
        # Orthogonal: DEF-009's data_reduction disclosure (ours) and 0.0.526's
        # repair_attempts (theirs) describe different things. Both are kept.
        if data_reductions:
            observation["data_reduction"] = data_reductions
            output["data_reduction"] = data_reductions
        if repair_attempts:
            observation["repair_attempts"] = repair_attempts

        # Add preview screenshot for planner reflection (page mode)
        if _attach_screenshot:
            observation["images"] = [{
                "data": screenshot_base64,
                "media_type": "image/png",
                "source_type": "base64",
            }]

        # Add slides-specific info
        if data.mode == "slides":
            if preview_images:
                observation["preview_images"] = preview_images
                observation["slide_count"] = len(preview_images)
            if pptx_path:
                observation["pptx_path"] = pptx_path
            # WHICH theme, and HOW it was chosen — visible in the console next
            # to the deck it produced. Mirrors how layout_check rides both
            # `observation` (for the agent) and `output` (for anything reading
            # the raw tool result).
            if theme_choice is not None:
                observation["deck_theme"] = content.get("deck_theme")
                output["deck_theme"] = content.get("deck_theme")
            # DEF-011: the layout check's verdict, mirroring how data_reduction
            # rides both `observation` (for the agent) and `output` (for
            # anything reading the raw tool result) above. `layout_notice` is
            # only set when there is something to say — a clean "checked, no
            # issues" run stays silent here just like DEF-010 does.
            if layout_check_result is not None:
                if layout_notice:
                    observation["layout_notice"] = layout_notice
                observation["layout_check"] = content.get("layout_check")
                output["layout_check"] = content.get("layout_check")
            if pptx_repair_attempts:
                observation["repair_attempts"] = pptx_repair_attempts
            # Attach the first slide preview for planner reflection — same
            # privacy + vision gate as the page-mode screenshot (the preview
            # shows the data).
            _slides_preview_b64 = self._first_preview_base64(preview_images)
            if _slides_preview_b64 and allow_llm_see_data and _model and getattr(_model, "supports_vision", False):
                observation["summary"] += " First slide preview is attached — review it for visual correctness."
                observation["images"] = [{
                    "data": _slides_preview_b64,
                    "media_type": "image/png",
                    "source_type": "base64",
                }]

        if warnings:
            observation["warnings"] = warnings

        yield ToolEndEvent(
            type="tool.end",
            payload={
                "output": output,
                "observation": observation,
            }
        )

    def _trim_none(self, obj: Any) -> Any:
        """Remove None values and empty collections from nested structures."""
        try:
            if isinstance(obj, dict):
                out = {}
                for k, v in obj.items():
                    tv = self._trim_none(v)
                    if tv is None:
                        continue
                    if isinstance(tv, (dict, list)) and len(tv) == 0:
                        continue
                    out[k] = tv
                return out
            if isinstance(obj, list):
                items = [self._trim_none(v) for v in obj]
                return [v for v in items if not (v is None or (isinstance(v, (dict, list)) and len(v) == 0))]
            return obj
        except Exception:
            return obj

    def _build_slides_prompt(
        self,
        user_prompt: str,
        title: str | None,
        viz_profiles: List[Dict[str, Any]],
        instructions_context: str,
        report_title: str | None,
        allow_llm_see_data: bool,
        messages_context: str = "",
        image_count: int = 0,
        organization_settings: Any = None,
        files: Optional[List[Dict[str, Any]]] = None,
        report: Any = None,
        resolved_theme: Any = None,
    ) -> str:
        """Build the prompt for generating slides using python-pptx code.

        `report` is here so the deck's stored theme (`Report.theme_name`) can
        reach the design system; `resolved_theme` lets the caller resolve once
        and hand the same object to the generated code, so the prompt and the
        runtime cannot name different themes.
        """
        viz_json = json.dumps(viz_profiles, indent=2, default=str)

        # ── Design system ────────────────────────────────────────────────────
        # The INDEX of every theme (so the model can pick when nothing
        # deterministic matched the request) plus the FULL spec of the one that
        # was resolved. Deliberately not all specs: 81 of them is ~32k tokens for
        # 80 themes this deck will not use. This prompt is built once per deck,
        # not once per call, so the index's ~1.2k is affordable.
        theme_section = ""
        _themes = _load_pptx_themes()
        if _themes is not None:
            try:
                theme = resolved_theme
                if theme is None:
                    theme = _resolve_deck_theme(
                        user_text=user_prompt,
                        report=report,
                        organization_settings=organization_settings,
                    )
                index = _themes.index_lines()
                spec = _themes.spec_block(theme) if theme is not None else ""
                if spec:
                    theme_section = f"""
═══════════════════════════════════════════════════════════════════════════════
DESIGN SYSTEM — BUILD THIS DECK IN THE THEME BELOW
═══════════════════════════════════════════════════════════════════════════════

{spec}

Every theme available, if the user asked for a look this one does not serve
(`id  category  when to use`). Switch only on an explicit request; otherwise
build the deck in the theme specified above and do not blend two of them:

{index}
"""
            except Exception:
                # A theme is a design improvement, never a precondition for
                # producing a deck. Fall back to the prompt as it was.
                theme_section = ""

        # Embeddable art. Only images: python-pptx cannot place a PDF, and the
        # executor only loads image/* bytes, so advertising anything else would
        # promise an id that image() will reject.
        embeddable = [
            f for f in (files or [])
            if str(f.get("content_type") or "").startswith("image/")
        ]
        if embeddable:
            listed = "\n".join(
                f"  - {f['id']} — {f.get('filename') or 'image'}" for f in embeddable
            )
            embeddable_images = f"\n\n  Available image ids:\n{listed}"
        else:
            embeddable_images = (
                "\n\n  No images are attached to this deck — image_ids is empty. "
                "Build the design from shapes, color and type."
            )

        language_directive = build_language_directive(organization_settings)

        # Build attached images context
        images_context = ""
        if image_count > 0:
            images_context = f"\n**Attached Images:** {image_count} image(s) provided for visual reference. Use these to understand the design intent, branding, color schemes, or layout preferences the user wants to incorporate."

        return f"""Role: presentation author using python-pptx.{language_directive}
Generate python-pptx code to create a polished slide deck.

═══════════════════════════════════════════════════════════════════════════════
DECK CRAFT — decide this BEFORE writing any code
═══════════════════════════════════════════════════════════════════════════════

**1. Settle the storyline first.** Write the argument in one sentence: what
should the audience believe or do after seeing this? Every slide either
supports that sentence or gets cut. Then order slides so each earns the next.
The arc that works for analytical decks:

  1. Title — subject, audience, date.
  2. The headline — the single most important finding, stated outright. Do not
     save the conclusion for the end; executives read the first two slides.
  3. Evidence — one slide per supporting point, each with a chart.
  4. What changed / why — drivers, segments, root cause.
  5. So what — implications, risks, recommended actions.
  6. Appendix — detail, methodology, caveats.

  For a status or review deck, replace 2-5 with: where we are → what moved →
  what's blocked → what's next.

**2. Titles carry the message.** The title is the one line everyone reads.
Make it the finding, not the topic:
  - Weak:   "Revenue by Region"
  - Strong: "EMEA drove all of Q3 growth; every other region was flat"
A reader should page through titles alone and get the whole argument. If a
title could sit on any deck in any quarter, it is a label, not a takeaway.
Keep titles under ~12 words so they fit one line.

**3. One idea per slide.** If a slide needs "and" twice to explain, split it.
  - One chart per slide, unless two are being directly compared.
  - At most 5 bullets, at most 2 lines each, no sub-bullets.
  - No paragraphs. If prose is needed, the deliverable is a document, not a deck.
  - Numbers carry units and periods ("$4.2M, Q3" — not "4200000").
  - Put supporting detail in speaker notes via `slide.notes_slide`, not in
    shrunken body text.

**4. Never invent a number.** Every figure in a title or takeaway must match
what the chart shows. If the data does not support the claim, change the claim.

**5. Hold ONE visual system across the whole deck.** Pick a palette and stick
to it for every slide — same background family, same accent, same type scale,
same margins. A deck where slide 3 is light and slide 5 is dark, or where card
colors change without meaning, reads as broken no matter how good any single
slide is.

**6. A deck does NOT require data.** `visualizations` is often EMPTY — a topic,
narrative or announcement deck ("a deck about the 2026 World Cup") has no
charts at all, and that is a valid deck, not an error.

  - **Never index `visualizations[0]` without checking the list first.** On an
    empty list that raises IndexError and loses the whole deck. Guard every
    data-driven slide with `if visualizations:` and skip it otherwise.
  - With no data, carry the design with type, color, shapes and images: a
    full-bleed title, a section divider, a numbered-point layout, a quote, a
    stat stated as large type (only if the user supplied the number).
  - Do NOT invent charts, metrics or figures to fill the space. A confident
    typographic slide beats a fabricated bar chart.
  - Requested slide count is a hard constraint: "2 slides" means exactly 2.

**7. The DATA SHAPE picks the slide, not taste.** Look at what the slide has to
show and take the archetype from that — never rotate through layouts for
variety:

  - ONE metric, or one metric against its target → a stat slide: the number as
    large type, one line saying what it is and over what period. No chart.
  - 2-4 categories → labelled horizontal bars, value printed at the end of each
    bar. A pie of four slices is harder to read than four bars.
  - A time series → a line, with the axis in real time order and the period
    named. Never a bar chart of dates.
  - 3 comparable entities (units, options, periods) → a three-up: one column
    card each, same fields in the same order in every column.
  - Two measures on DIFFERENT scales (a count and a percentage, revenue and
    margin) → two exhibits, side by side or on two slides.
    NEVER two series on one axis: the smaller one flattens to the baseline and
    reads as zero.
  - More than ~10 rows → cut to the top 8-10 and say so in the subtitle, or move
    the full table to the appendix.
  - No data at all → a narrative slide, carried by type, shapes and colour. Do
    not invent figures to justify a chart.
{theme_section}
═══════════════════════════════════════════════════════════════════════════════
AVAILABLE IN NAMESPACE (already provided — do not import)
═══════════════════════════════════════════════════════════════════════════════

Python-pptx classes and functions:
- Presentation, Inches, Pt, Emu, RGBColor
- PP_ALIGN, MSO_ANCHOR, MSO_SHAPE
- XL_CHART_TYPE, XL_LEGEND_POSITION
- CategoryChartData, ChartData

Note: Inches, Pt, Emu are functions, not methods.
   Use: Inches(1), Pt(24), Emu(914400)
   Not: 1.inches, 24.pt, value.inches

Note: a line is not a fill. shape.fill.solid() exists; shape.line.solid() does not.
   Use: shape.line.fill.solid()  or  shape.line.color.rgb = RGBColor(...)
   Not: shape.line.solid()
   LineFormat has only: color, dash_style, fill, width.

Data variables:
- visualizations: List[Dict] — each has 'title', 'columns', 'rows'
- report: Dict with 'id', 'title', 'theme'. `report['theme']` is None or a dict
  carrying 'name' and, when a design system resolved, 'id', 'fonts' and
  'palette'. It is a reference, not a substitute for writing the colours out:
  set every colour and font name explicitly on the shapes you draw.

Images (only present when the user attached or generated some):
- image_ids: List[str] — the embeddable image ids available to this deck
- image(file_id) -> stream — pass straight to add_picture. Returns a fresh
  stream per call, so the same image may be placed on several slides:
    pic = slide.shapes.add_picture(image(image_ids[0]), Inches(0), Inches(0),
                                   width=Inches(13.333))
  Cover the slide for a hero/background, or inset it in a content column.
  There is no filesystem access — `image()` is the only way to place art.
  When an image sits behind text, draw a translucent scrim rectangle between
  them or the text becomes unreadable; send the picture to the back by
  inserting it first.{embeddable_images}

Output:
- _pptx_output_path: str — path to save the presentation to

═══════════════════════════════════════════════════════════════════════════════
YOUR VISUALIZATIONS
═══════════════════════════════════════════════════════════════════════════════

{viz_json}

{"(Full sample data included above)" if allow_llm_see_data else "(Data samples hidden for privacy - use column names and row_count)"}

═══════════════════════════════════════════════════════════════════════════════
TASK
═══════════════════════════════════════════════════════════════════════════════

**Report Title:** {report_title or title or 'Presentation'}
**User Request:** {user_prompt}
{images_context}
{f"**Organization Instructions:** {instructions_context}" if instructions_context else ""}

═══════════════════════════════════════════════════════════════════════════════
PYTHON-PPTX QUICK REFERENCE
═══════════════════════════════════════════════════════════════════════════════

**Setup (16:9 widescreen):**
```python
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
```

**Add blank slide with dark background:**
```python
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 23, 42)  # slate-900
```

**Add text box:**
```python
txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Title Text"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
```

**Add bar chart (use this pattern for charts):**
```python
chart_data = CategoryChartData()
chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
chart_data.add_series('Revenue', (1.2, 1.5, 1.8, 2.1))

x, y, cx, cy = Inches(1), Inches(2), Inches(11), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
).chart

# Style the chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
plot = chart.plots[0]
plot.has_data_labels = True
```

**★ CHART TEXT IS BLACK BY DEFAULT — RETHEME IT OR IT VANISHES ON A DARK SLIDE**

python-pptx does NOT inherit your slide colours. Every text element of a chart
starts black, so on a dark background the chart title reads as an empty gap and
the axis labels disappear. Set `chart.font.color.rgb` FIRST — it cascades to the
title, both axes and the legend — then override anything you want brighter:

```python
# One line that covers title + axes + legend. Do this on EVERY chart.
chart.font.color.rgb = TEXT_MUTED
chart.font.size = Pt(12)

# The chart title specifically (it is the one most often left black):
if chart.has_title:
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

# Data labels sit on top of the bars, so they need the bright colour:
plot.has_data_labels = True
plot.data_labels.font.color.rgb = TEXT_LIGHT
plot.data_labels.font.size = Pt(11)
```

**★ SERIES COLOURS ALSO DEFAULT — to Office blue/red/green/purple**

Untouched, a pie or multi-series chart renders in Office's stock palette, which
will not match the palette you chose for the deck. Assign your own:

```python
# Single-series bar/column/line — one colour from your palette:
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = SECONDARY

# Pie/doughnut — colour each SLICE, not the series:
slice_colors = [PRIMARY, SECONDARY, ACCENT, TEXT_MUTED]
for i, point in enumerate(plot.series[0].points):
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = slice_colors[i % len(slice_colors)]
```

**Other chart types:**
- XL_CHART_TYPE.COLUMN_CLUSTERED - vertical bars
- XL_CHART_TYPE.LINE - line chart
- XL_CHART_TYPE.PIE - pie chart
- XL_CHART_TYPE.AREA - area chart

**Dark background (slate-900 = RGB(15, 23, 42)):**
```python
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 23, 42)
```

**Access visualization data** (only inside an `if visualizations:` guard — the
list is empty for a narrative deck):
```python
viz = visualizations[0]
# Each entry of viz['columns'] is a DICT, e.g. {{'field': 'Revenue', 'headerName': 'Revenue'}}.
# Extract the row keys from the 'field' values first:
fields = [c['field'] if isinstance(c, dict) else str(c) for c in viz['columns']]  # e.g. ['AlbumTitle', 'Revenue']
rows = viz['rows']        # list of dicts like {{'AlbumTitle': 'Greatest Hits', 'Revenue': 1500.0}}

# Get categories and values for a chart:
categories = [str(row.get(fields[0], '')) for row in rows]  # First column as labels
values = [float(row.get(fields[1]) or 0) for row in rows]   # Second column as values

# IMPORTANT: never pass a column dict itself to a text property or a row index —
# use fields[i] (a string like 'Revenue'): row[fields[1]] is row['Revenue']
```

═══════════════════════════════════════════════════════════════════════════════
DESIGN PHILOSOPHY - CREATE BEAUTIFUL, PROFESSIONAL SLIDES
═══════════════════════════════════════════════════════════════════════════════

**COLOR STRATEGY - Be Topic-Specific:**
Choose colors that feel designed for THIS topic. If your colors would work for any presentation, you haven't made specific enough choices.

Structure: One DOMINANT color (60-70% visual weight), 1-2 supporting tones, one accent.

Example palettes (pick one that fits the topic):
- **Midnight Executive**: Navy (0,31,63), Steel (119,136,153), Gold accent (212,175,55)
- **Forest & Moss**: Deep green (34,87,76), Sage (138,154,91), Cream (245,245,220)
- **Coral Energy**: Coral (255,127,80), Teal (0,128,128), Sand (244,232,214)
- **Ocean Depths**: Deep blue (0,51,102), Aqua (0,180,180), Pearl (240,248,255)
- **Sunset Warm**: Burgundy (128,0,32), Orange (255,140,0), Cream (255,253,240)
- **Modern Minimal**: Charcoal (54,69,79), Light gray (220,220,220), Teal accent (0,150,136)

**Choose the palette from THIS deck's subject, and commit to it.**
Two decks on different subjects must not arrive looking identical — a retail
review and a risk review should be recognisably different documents. Pick before
you write any slide, name your choice in a comment at the top of the code, and
then use those exact three colours everywhere: dominant for titles, panels and
chart series; the pale tint for cards; the accent for every numeral and every
small label. Do not introduce a fourth colour, and do not drift to a different
palette halfway through the deck.
Semantic colours (a red for a decline, a green for growth) are the one
exception and are not part of the palette count.

**Layout variety — vary between slides:**
Every slide should have visual elements — charts, shapes, or decorative elements. Avoid text-only slides.

**SLIDE ARCHETYPES — pick the one that matches what the slide is DOING.**
Do not cycle through them for variety's sake; a ranking wants bars, a sequence
wants badges. Never use the same archetype twice in a row.

1. COVER — exactly one, first, and no footer on it.
   ★ FULL BLEED. Start by drawing a rectangle over the ENTIRE canvas
   (0, 0, 13.333in, 7.5in) filled with the dominant colour, then place
   everything on top of it in light text. A pale cover with a coloured accent
   bar is the default a model reaches for and it is exactly what makes a deck
   look ordinary — the opening slide is the one place to spend the colour.
   Over that ground: two large circles in a slightly lighter shade of the
   dominant, bled off the right edge so they are half outside the canvas, plus
   one or two small accent-coloured circles. Eyebrow in the accent colour, then
   a 44-60pt serif title in white, then a serif sub-line, then one muted
   sentence of scope. A small muted line low-left ("Leadership working deck").
   No data, no KPI cards, no footer on this slide — it sets the tone, nothing else.
   The "avoid text-only slides" rule does NOT apply here: the shapes are the
   visual element.

2. SECTION DIVIDER — full-bleed dominant colour, a 44-60pt serif title centred
   left, a short accent rule under it, one line of what the section covers.
   Use between parts of a long deck. No footer.

3. METRICS — 3 or 4 cards across, each: a big accent serif number, a bold
   label, a muted unit line. Never more than 4 in a row; a fifth means two rows
   or a second slide. Inverted conclusion panel below them.

4. CHART + INSIGHT — chart on the left ~60%, a card on the right ~35% holding
   2-3 numbered takeaways. The takeaways say what the chart MEANS, they do not
   restate the bars.
   ★ For a ranking, `XL_CHART_TYPE.BAR_CLUSTERED` draws categories BOTTOM-UP, so
   the largest lands at the bottom and the ranking reads upside down. Reverse
   your category and value lists, or use `COLUMN_CLUSTERED`.

5. ACTIONS — a numbered row per recommendation: a filled square on the left in
   the dominant colour with an accent serif numeral, then a serif headline and
   one grey sentence of reasoning. Three or four rows, ordered by expected
   return, not by convenience.

6. PROCESS — 4-5 circular badges in a row, each with an accent numeral, a serif
   step name and a one-line caption. A full-width band underneath carrying the
   outcome of the whole sequence.

7. COMPARISON — 3-5 column cards, each with a filled header in the dominant
   colour and light text, then a short bulleted list. For business units,
   options, or periods being weighed against each other.

8. TWO-PANEL — two cards side by side, each with its own small uppercase label
   and a list of term + one-line definition. For "who does what" or
   "what it costs / what it returns". Statement band underneath.

Every slide carries visual structure. Avoid text-only slides.

**TYPEFACE — set it explicitly on every run, never leave it to the default:**
Pair a SERIF for display with a SANS for body. That pairing is the single
strongest signal a person designed the deck; one default sans throughout is the
strongest signal a machine generated it.
- Headings, titles, KPI numbers, card headers: `Cambria` (serif)
- Body text, labels, captions, axis labels, footers: `Calibri` (sans)
- Set `run.font.name` on EVERY run. python-pptx does not inherit a deck font —
  a run with no name falls back to the viewer's default and the pairing is lost.

**TYPE SCALE — wide jumps, because the gaps are what read as hierarchy:**
| Use | Size |
|---|---|
| Source lines, footers, page numbers | 9-11pt |
| Body, card descriptions, bullets | 14-16pt |
| Subtitle under a title, card headers | 18-22pt |
| Slide title | 28-36pt |
| Cover / section-divider title | 44-60pt |
| A single hero number | 60-96pt |
Body at 18-24pt against a 36pt title is nearly flat and reads as generated.

**SLIDE FURNITURE — the same three lines on EVERY content slide.**
This repetition is what makes a set of slides read as one document rather than
a folder of pictures. It is not decoration; leave it out and the deck falls apart.
1. EYEBROW: 9-11pt, bold, UPPERCASE, letter-spaced, muted colour, above the
   title. Names the section — "PERFORMANCE", "CATEGORY MIX", "ACTIONS".
   (python-pptx has no letter-spacing property: emulate it by joining the
   characters with spaces, e.g. "P E R F O R M A N C E".)
2. TITLE: serif, 28-36pt, in the dominant colour. A statement, not a label —
   "Where the money comes from" beats "Category analysis".
3. SUBTITLE: one grey sentence, 13-15pt, stating the measure and its scope —
   "Net sales = gross less discount, all outlets, full history."
Then, pinned to the bottom of every content slide (NOT the cover):
- FOOTER left: what the deck is ABOUT, 9pt, muted, letter-spaced — the subject
  and the period, e.g. "CITY MART RETAIL · SALES REVIEW · FY2025".
  ★ NEVER put the report's internal title in the footer. Report titles are
  working labels a person typed to find the thing again ("FT dz2-crm",
  "test 3", "copy of Q3") and one landed in a footer on every slide of a real
  deck. Write the footer from the DATA and the QUESTION, not from `report.title`.
  If you cannot name the subject confidently, use the data source's display
  name alone. Never invent a company name that is not in the data.
- PAGE NUMBER right: same size and colour, zero-padded — 02, 03.

**ONE INVERTED PANEL PER SLIDE — where the conclusion goes.**
Every content slide carries exactly one filled rectangle in the DOMINANT colour
with light text on it, holding the single thing the reader should take away. A
small accent-coloured label above it ("WHAT IT MEANS", "OUTCOME", "SO WHAT").
A conclusion buried as the fourth bullet is a conclusion nobody reads. One per
slide — two competing dark panels and neither carries weight.

**CARDS — the workhorse container:**
- Fill: a very pale tint of the dominant colour (near-white, not grey).
- Border: hairline, ~0.75pt, a shade darker than the fill. No drop shadows.
- Optional 4-6pt accent-coloured rule along the top edge.
- Generous inner padding — roughly 0.25 inch on every side.

**NUMERALS — always the accent colour, always the serif.**
Every KPI value, ranking number, step number and card index. Numbers are the
reason the deck exists; make them the thing the eye lands on.

**VISUAL ELEMENTS:**
- Shapes: rectangles, rounded rectangles, and OVALS. Large circles bled off the
  slide edge (partly outside the canvas) are the cheapest way to make a cover
  look designed — two big ones in a darker shade of the background, one or two
  small ones in the accent colour.
- Full-width bands for statements; column cards with filled headers for
  comparisons; circular numbered badges for a process.
- Never an accent line directly under a title — the hallmark of generic slides.

**Common mistakes to avoid:**
- Using `value.inches` instead of `Inches(value)` — Inches/Pt/Emu are functions.
- Repeating the same layout across slides — vary it.
- Center-aligning body text — use left alignment.
- Using only blue without topic-specific reasoning.
- Text-only slides without visual elements.
- Accent lines directly under titles (hallmark of generic slides).
- Cramming too much data — limit charts to top 8-10 items.
- **Leaving chart text black on a dark slide** — see "Rendering defects to
  prevent" below for the properties to set.
- **Leaving series colours at the Office default** (blue/red/green/purple). They
  will not match the palette you picked, and a pie chart makes it obvious.
- **Showing raw column or table names to the reader** — `net_amount`,
  `fact_sales.net_amount`. Titles, axis labels, series names, KPI captions and
  source footers all get a human label.
- **Truncating text that had room to fit.** Side panels and callout boxes are
  wider than an axis label; do not reuse a tight `[:20]` cap everywhere.
- Adding a chart without checking its rows first — `CategoryChartData` raises
  "chart data contains no categories" on an empty list and that failure loses
  the whole deck, not just the slide.

**Technical requirements:**
1. Define `generate_slides(visualizations, report)` returning a Presentation.
2. Use 16:9 widescreen: Inches(13.333) x Inches(7.5).
3. Create real charts with slide.shapes.add_chart() + CategoryChartData.
4. Use visualization data from the visualizations list. Read rows with
   `viz.get('rows', [])` and ALWAYS guard before charting:
   `if rows:` — build the chart; otherwise render the slide without it.
5. Margins: start shapes at Inches(0.75) to Inches(1) from edges.

**Rendering defects to prevent (these are what make a deck look broken):**
- **Charts on dark backgrounds render unreadable.** python-pptx defaults every
  chart label to near-black, which disappears on a dark slide. On a dark
  background set them explicitly — category and value tick labels, data
  labels, and the chart title:
  `chart.font.color.rgb = LIGHT` plus
  `chart.category_axis.tick_labels.font.color.rgb = LIGHT` and
  `chart.value_axis.tick_labels.font.color.rgb = LIGHT`.
  The chart TITLE is the one that gets forgotten, and it disappears completely.
  Pick series colors that contrast with the background too.
- **Content running off the slide.** The canvas is 7.5in tall. Keep every TEXT
  shape between Inches(0.4) and Inches(7.1) vertically, and its text inside
  Inches(0.75) from the left and right edges. Before emitting a card grid,
  check `top + height` for the LAST row fits.
  ★ This is about text that overflows by ACCIDENT. It does not override the
  archetypes above: the COVER and SECTION DIVIDER grounds are full bleed by
  design (0, 0, 13.333in, 7.5in), the cover's circles are meant to run off the
  right edge, and statement bands are meant to be full width. A background
  shape may leave the canvas; the words on top of it may not.
- **Titles colliding with what follows.** A long title wraps to 2-3 lines. Give
  the title box enough height for the wrap and start the next element BELOW it;
  never overlap a subtitle, accent line or chart with the title block.

═══════════════════════════════════════════════════════════════════════════════
OUTPUT FORMAT - Example with Design Principles Applied
═══════════════════════════════════════════════════════════════════════════════

```python
def generate_slides(visualizations, report):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette - choose colors that fit the topic
    PRIMARY = RGBColor(0, 51, 102)      # Deep blue
    SECONDARY = RGBColor(0, 128, 128)   # Teal
    ACCENT = RGBColor(255, 140, 0)      # Orange accent
    BG_DARK = RGBColor(15, 23, 42)      # Dark background
    TEXT_LIGHT = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184)

    def set_background(slide, color=BG_DARK):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_accent_shape(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def humanize(name):
        # ★Column names are for the database, not the reader. `net_amount`
        # becomes "Net Amount"; `fact_sales.net_amount` loses the table too.
        # A deck that says "Revenue = net_amount" has leaked its plumbing.
        return str(name).split('.')[-1].replace('_', ' ').strip().title()

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 1: Title with accent shape
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    # Accent shape behind title
    add_accent_shape(slide, Inches(0), Inches(2.5), Inches(5), Inches(2.5), PRIMARY)

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(3), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = report.get('title', 'Presentation')
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 2: KPI Cards Row (if we have numeric data)
    # ═══════════════════════════════════════════════════════════════
    if visualizations and visualizations[0].get('rows'):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)

        viz = visualizations[0]
        rows = viz.get('rows', [])
        columns = viz.get('columns', [])

        # Create 3 KPI cards across the slide
        card_width = Inches(3.5)
        card_height = Inches(2.5)
        start_x = Inches(1)
        card_y = Inches(2.5)
        gap = Inches(0.5)

        for i, col in enumerate(columns[:3]):
            if i >= 3:
                break
            x = start_x + i * (card_width + gap)

            # Card background
            card = add_accent_shape(slide, x, card_y, card_width, card_height, PRIMARY)

            # Value (large number)
            val = rows[0].get(col, 0) if rows else 0
            val_box = slide.shapes.add_textbox(x + Inches(0.3), card_y + Inches(0.5), card_width - Inches(0.6), Inches(1.2))
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = "{{:,.0f}}".format(float(val)) if isinstance(val, (int, float)) else str(val)
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = TEXT_LIGHT

            # Label
            label_box = slide.shapes.add_textbox(x + Inches(0.3), card_y + Inches(1.7), card_width - Inches(0.6), Inches(0.6))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = humanize(col)
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_MUTED

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 3: Chart with title (different layout)
    # ═══════════════════════════════════════════════════════════════
    if visualizations:
        viz = visualizations[0]
        columns = viz.get('columns', [])
        rows = viz.get('rows', [])

        if len(columns) >= 2 and rows:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_background(slide)

            # Title on left side
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(5), Inches(1))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = viz.get('title', 'Data Analysis')
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = TEXT_LIGHT

            # Extract data
            col_label = columns[0]
            col_value = columns[1]
            # Category labels: keep them whole. Truncate only what genuinely
            # cannot fit, and never at a tight fixed width — "City Value
            # Vitamins & Supplemen..." helps nobody.
            categories = [str(row.get(col_label, '')) for row in rows[:8]]
            values = [float(row.get(col_value, 0) or 0) for row in rows[:8]]

            # Series name is shown to the reader, so give it a human label
            # rather than the raw column name.
            series_label = humanize(col_value)

            # Chart (full width below title)
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series(series_label, tuple(values))

            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5),
                chart_data
            ).chart
            chart.has_legend = False

            # ★Retheme the chart text — without this it renders black on the
            # dark background and the title is effectively invisible.
            chart.font.color.rgb = TEXT_MUTED
            chart.font.size = Pt(12)
            if chart.has_title:
                tp = chart.chart_title.text_frame.paragraphs[0]
                tp.font.color.rgb = TEXT_LIGHT
                tp.font.size = Pt(16)
                tp.font.bold = True

            # ★Series colour from the palette, not Office's default blue.
            plot = chart.plots[0]
            plot.series[0].format.fill.solid()
            plot.series[0].format.fill.fore_color.rgb = SECONDARY
            plot.has_data_labels = True
            plot.data_labels.font.color.rgb = TEXT_LIGHT
            plot.data_labels.font.size = Pt(11)

    return prs

# Execute and save
prs = generate_slides(visualizations, report)
prs.save(_pptx_output_path)
```

Create a beautiful, varied presentation following these design principles. Each slide should look DIFFERENT from the others. Use visual elements, accent shapes, and thoughtful color choices:"""

    def _build_page_system_prompt(self) -> str:
        """Static system prompt for page/dashboard generation.

        Contains only stable reference material (sandbox runtime docs,
        component contract, filtering rules, design guidance, output format).
        Kept free of per-call state so provider prompt caching can reuse it
        across every create/edit call.
        """
        return f"""Role: frontend developer and data visualization engineer. You build React dashboard artifacts from the user's design request and visualization data (both provided in the user message), following this reference.

═══════════════════════════════════════════════════════════════════════════════
REFERENCE — TOOLS, COMPONENTS & DATA
═══════════════════════════════════════════════════════════════════════════════

{SANDBOX_RUNTIME_PROMPT}

CHARTING:

**`<EChart height={{N}} option={{{{...}}}} />`** — chart wrapper. Supports ALL ECharts chart types. 'dash' theme pre-configures colors, tooltip, grid, axes. For standard charts, only write data mapping:
```jsx
<EChart height={{300}} option={{{{ xAxis: {{ type: 'category', data: rows.map(r => r.name) }}, yAxis: {{ type: 'value' }}, series: [{{ type: 'bar', data: rows.map(r => r.val) }}] }}}} />
<EChart height={{300}} option={{{{ tooltip: {{ trigger: 'item' }}, series: [{{ type: 'pie', radius: ['45%','75%'], data: rows.map(r => ({{ value: r.amt, name: r.lbl }})) }}] }}}} />
<EChart height={{300}} option={{{{ xAxis: {{ type: 'category', data: rows.map(r => r.date) }}, yAxis: {{ type: 'value' }}, series: [{{ type: 'line', data: rows.map(r => r.val), areaStyle: {{ opacity: 0.15 }} }}] }}}} />
```
For advanced charts (radar, gauge, treemap, sunburst, funnel, sankey, calendar heatmap, parallel coordinates, graph), pass the full ECharts option — the theme still provides colors and tooltip:
```jsx
<EChart height={{300}} option={{{{ radar: {{ indicator: indicators }}, series: [{{ type: 'radar', data: radarData }}] }}}} />
<EChart height={{250}} option={{{{ series: [{{ type: 'gauge', data: [{{ value: 72 }}], detail: {{ formatter: '{{value}}%' }} }}] }}}} />
<EChart height={{400}} option={{{{ series: [{{ type: 'treemap', data: treeData }}] }}}} />
```

AVAILABLE COMPONENTS (convenience shortcuts — not requirements):
- `<BowKpi title="" value={{fmt(n, {{currency:true}})}} subtitle="" color="#3B82F6" className="" titleClassName="" subtitleClassName="" style={{{{}}}} />` — **REQUIRED for every KPI / stat / big-number tile.** It measures the card's real width and scales the value's type to fit (wrapping rather than shrinking past legibility), is `min-w-0` so it shrinks inside a grid, and uses tabular numerals. A hand-rolled `<div className="text-3xl">{{value}}</div>` cuts the last digits off a long value with no error and no visual sign — the user reads a wrong number. If a design truly needs custom tile markup, still render the VALUE through `<BowFitText className="font-semibold">{{value}}</BowFitText>` inside a `min-w-0` parent. `<KPICard>` is an alias of `<BowKpi>`. Theme these to match your color story:
  - Dark: `className="bg-slate-900 border-slate-700 text-white" titleClassName="text-slate-400"`
  - Colored: `className="bg-indigo-50 border-indigo-200 text-indigo-900" titleClassName="text-indigo-600"`
- `<SectionCard title="" subtitle="" className="" titleClassName="" subtitleClassName="" style={{{{}}}}>...children...</SectionCard>` — same theming: `className` replaces defaults, `titleClassName`/`subtitleClassName` for text. Theme to match.
- `<FilterSelect label="" options={{arr}} selected={{arr}} onChange={{fn}} searchable={{bool}} className="" style={{{{}}}} />` — multi-select dropdown (portaled). Built-in search at 8+ options. `className` replaces default theme (bg-white border-slate-200 text-slate-900) — pass e.g. `className="bg-slate-900 border-slate-700 text-slate-100"` for dark.
- `<FilterSearch label="" value={{str}} onChange={{e => setFilter(field, e.target.value)}} placeholder="Search..." className="" style={{{{}}}} />` — text search. `className` replaces default theme.
- `<FilterDateRange label="" value={{filters[field] || {{}}}} onChange={{val => setFilter(field, val)}} type="date" className="" style={{{{}}}} />` — date range picker. `className` replaces default theme.
- `fmt(n, opts)` — `{{currency:true}}`, `{{pct:true}}`, auto K/M/B
- `<LoadingSpinner size={{32}} />`

All components are fully themeable via `className`/`titleClassName`/`subtitleClassName`/`style`. Don't leave default white/slate styling when your design calls for something different. If the design needs something these can't express — build custom React + Tailwind.

**INFO POPOVER (required):** Pass `viz={{viz[N]}}` to every `<KPICard>` and `<SectionCard>` you build from a visualization. This renders a small built-in "ⓘ" button that lets users inspect the data behind each component (Data tab with rows, Code tab with the query). Use the index of the visualization the card is derived from (the primary one if it combines several). When a card renders FILTERED rows (you called `filterRows(viz[N].rows)`), ALSO pass `rows={{<those filtered rows>}}` so the popover shows the filtered view that matches the component, not the full dataset. When a card AGGREGATES or derives its value client-side, ALSO pass `calc="<formula>"` describing the math with real column names, e.g. `calc="SUM(UnitPrice × Quantity) grouped by GenreName"` or `calc="COUNT(DISTINCT CustomerId)"` — the popover shows it as a "Calculation" line. If you render a chart with a bare `<EChart>` that is NOT inside a `<SectionCard>`, pass `viz={{viz[N]}}` (and `rows`/`calc` if relevant) to the `<EChart>` itself so it still gets the popover.

**CUSTOM MARKUP — add `data-dash-*` attributes (required):** Whenever you build your OWN containers instead of `<BowKpi>`/`<SectionCard>`/`<EChart>` (custom `<div>` KPI tiles — whose VALUE must still go through `<BowFitText>`, chart wrappers, tables), annotate each item's outer element with `data-dash-viz="N"` (source visualization index) and `data-dash-calc="<formula>"` when the value is derived. A global overlay then renders the same Data/Code/Calc popover on each item. Example: `<div data-dash-viz={{0}} data-dash-calc="SUM(UnitPrice × Quantity)">...custom tile...</div>`. EVERY metric, chart, and table must be reachable via either a prebuilt component's `viz` prop OR a `data-dash-viz` attribute — never leave an item with no way to inspect its data.

DATA ACCESS:

```javascript
const data = useArtifactData(); // Returns null while loading
// data = {{ report: {{id, title}}, visualizations: [...] }}
```

Each visualization:
```js
{{
  id: "uuid",
  title: "Visualization Title",
  columns: [{{ "headerName": "Album Title", "field": "AlbumTitle", "dtype": "object", "unique_count": 150 }}, ...],
  rows: [{{ "AlbumTitle": "Battlestar Galactica", "total_revenue": 35.82 }}, ...],
  row_count: 4321,   // TRUE dataset size
  view: {{ /* chart config hints */ }},
  dataModel: {{ /* series/axis config */ }}
}}
```

- Use `column.field` to access row values: `row[column.field]`
- Use `column.headerName` for display labels
- Column metadata includes `dtype` (pandas type) and `unique_count` — use these for filter/format decisions
- **Do not hardcode data** — all values should come from `data.visualizations[N].rows`
- **Sample vs full data:** the `rows`/`sample_rows` shown in the user message are a SAMPLE (capped at 100 rows per visualization) for generation and preview. At runtime the dashboard receives the FULL dataset — `row_count` rows. `row_count` is the true dataset size; `sample_row_count` is the sample size. Write code that works on the full dataset (aggregate with reduce/Map, paginate long tables) and NEVER hardcode workarounds for the sample size.
- **Defensive coding**: Row values and properties can be `null`/`undefined`. Use optional chaining or fallbacks before calling `.includes()`, `.toLowerCase()`, `.startsWith()`, `.split()`, etc. Example: `(row.name || '').includes('x')` or `String(val ?? '').toLowerCase()`. Do not call string methods on a value that could be nullish.

View hints — honor the viz config:
The `view_config` on each visualization describes how the author wants the data rendered. Follow it when generating code.

- `view_config.aggregation` (`"sum" | "avg" | "count" | "min" | "max"`): the raw rows are granular, so aggregate the relevant value column before rendering (especially for `count`, `metric_card`, `pie_chart`, `heatmap`). Use `rows.reduce(...)`. Example for a metric card with aggregation=sum:
  ```js
  const total = useMemo(
    () => viz[0].rows.reduce((s, r) => s + (Number(r.revenue) || 0), 0),
    [viz]
  );
  ```
  For pie/heatmap/bar charts that group by a category, group first and aggregate the value per group rather than using the first matching row.

- `view_config.series_aggregations` (array of `{{key, aggregation}}`): apply the given aggregation per series when building multi-series bar/line/area charts.

- `view_config.default_filters` (array of `{{column, operator, value}}`): the author wants the dashboard to open with these filters already applied. Seed them on first mount so the initial view matches the intent, for example:
  ```js
  const {{ filters, setFilter, filterRows }} = useFilters();
  useEffect(() => {{
    // Seed defaults once — operators follow the useFilters contract.
    {{/* for each entry in view_config.default_filters */}}
    setFilter('column_name', value);
  }}, []);
  ```
  If the underlying runtime uses richer operators (`equals`, `greater_than`, etc.), either call `setFilter` with the operator-aware object it expects, or compute the filtered rows directly via `filterRows(viz[N].rows)` once the filter is seeded. Render the filtered view when defaults are present so the initial numbers match the author's intent.

FILTERING:
- Use `useFilters()` hook for cross-visualization filtering — returns `{{ filters, setFilter, resetFilters, filterRows }}`
- YOU choose which columns to filter — use `dtype` and `unique_count` from the column metadata:
  - `<FilterSelect>` for low-cardinality columns (`unique_count` < ~50, dtype "object"/"int64" with few values)
  - `<FilterSearch>` for high-cardinality text columns (`unique_count` > 50, dtype "object")
  - `<FilterDateRange>` for date/time columns (dtype contains "datetime" or values are date strings)
- Get unique values directly: `[...new Set(viz[N].rows.map(r => r[field]))]`

FILTER FEASIBILITY AUDIT — DO THIS FIRST, BEFORE WRITING CODE:
Before wiring any cross-viz filter, verify it will actually work. A filter that looks wired but silently leaves some vizs untouched is a broken dashboard, not a partial one.

For each dimension you intend to filter by:
1. **Enumerate participating vizs** — which vizs should this filter affect? (Usually: any viz whose topic logically shares the dimension, e.g. a "customer" filter should affect every viz about customers, payments, orders, etc.)
2. **Check column presence** — does each participating viz have the filter column (directly, or via a rename you can handle with `fieldMap`)? Check the `columns` array in YOUR VISUALIZATIONS below.
3. **Decide per dimension**:
   - ALL participants have the column → wire the global filter, use `fieldMap` for renames.
   - SOME participants lack the column but the gap is genuine (no join key in the source data) → make the filter LOCAL to the vizs that support it; do not pretend it affects others.
   - SOME participants lack the column but they should have it (the underlying data supports it, the query just didn't project the column) → **do not wire the filter; do not build the dashboard with a dead filter.** End your response by reporting the gap so the planner can recreate the offending queries before you try again. Example: "Cannot wire `customer_id` filter — `payments` viz lacks `customer_id` but `payments.customer_id` exists in schema. Recreate the payments query with `customer_id` projected, then retry create_artifact."

FILTER PLACEMENT — global vs local:
- **Global filter** (column present in 2+ vizs AFTER the audit above): place in a top-level filter bar above all content. Use one shared filter + `fieldMap` for renames, not duplicates.
- **Local filter** (column present in only 1 viz): place INSIDE that viz's `<SectionCard>`, visually next to the chart/table it affects.
- When a filter affects multiple vizs, add visible UI indication that they're linked.

FILTER DATA FLOW:
- Every viz that passes the feasibility audit for a filter should use `filterRows()` as its data source — for charts, tables, and any KPI/summary derived from that viz.
- KPI cards that summarize filtered data (sum, count, avg) should be computed from filtered rows, not from raw `viz[N].rows`.
- Do not call `filterRows` on a viz that doesn't have the filter column just to "be safe" — silently passing rows through makes the filter look active when it isn't. Audit first, wire second.

EXAMPLE 1 — Global "region" filter affecting KPIs + bar chart + table:
  const {{ filters, setFilter, resetFilters, filterRows }} = useFilters();
  const regions = useMemo(() => [...new Set(vizSales.rows.map(r => r.region))], [vizSales]);
  // ALL downstream from vizSales uses filtered:
  const filteredSales = filterRows(vizSales.rows);
  const totalRevenue = useMemo(() => filteredSales.reduce((s, r) => s + r.revenue, 0), [filteredSales]);
  const chartData = useMemo(() => ({{ labels: filteredSales.map(r => r.month), values: filteredSales.map(r => r.revenue) }}), [filteredSales]);
  // Cross-viz filtering with field mapping:
  const filteredDetails = filterRows(vizDetails.rows, {{ region: 'RegionName' }});
  // Layout: <FilterSelect> in top bar, KPIs below, charts below that

EXAMPLE 2 — Local filter inside a SectionCard:
  const {{ filters, setFilter, filterRows }} = useFilters();
  const filtered = filterRows(vizProducts.rows);
  // Layout: <SectionCard title="Products"><FilterSelect .../><EChart ... /></SectionCard>

- Include a Reset button when any filters are active (`Object.keys(filters).length > 0`)
- After filtering, if a visualization has zero matching rows, display "No data matches current filters"

═══════════════════════════════════════════════════════════════════════════════
DESIGN GUIDANCE (use when the user hasn't specified styling)
═══════════════════════════════════════════════════════════════════════════════

If the user specified a theme/style/colors above, follow that — skip this section.
Otherwise, design a visually striking, publication-quality dashboard — not a generic template.

COLOR & IDENTITY:
- Pick a cohesive color story that fits the data topic. A finance dashboard should feel different from a music dashboard, which should feel different from a healthcare dashboard.
- Choose one dominant color (60-70%), 1-2 supporting tones, and one accent for highlights/CTAs.
- Do NOT default to generic blue. Blue is fine if it fits the topic — but earn it, don't default to it.
- Theme ALL components (KPICard, SectionCard, filters) to match — use `className`, `titleClassName`, `subtitleClassName` props. Default white/slate is only appropriate for a clean/minimal design intent.

LAYOUT & HIERARCHY:
- Lead with the most important insight — KPIs or headline metric at the top.
- Create clear visual hierarchy: primary chart large, secondary charts smaller, supporting data compact.
- Use intentional whitespace — not "fill every pixel" but not "float in empty space" either.
- Vary card sizes and chart heights to create rhythm. A grid of same-sized boxes is boring.

TYPOGRAPHY & POLISH:
- Clean, modern typography. Titles concise and descriptive, not generic ("Revenue by Region" not "Chart 1").
- Subtle shadows, rounded corners, light borders — enough depth to feel crafted, not flat.
- Light mode default. Dark mode only if the topic or user suggests it.

CHART SELECTION:
- Choose the best visualization for the data shape — don't default to bar charts for everything.
- Standard charts (bar, line, pie, area) for simple relationships. Advanced charts (radar, gauge, treemap, funnel, sankey, heatmap) when the data structure rewards it.
- Show data from different angles without redundancy. Each chart should reveal something the others don't.

The goal: it should look like a designer built it for this specific dataset, not like a template was filled in.

═══════════════════════════════════════════════════════════════════════════════
RESPONSIVE LAYOUT (REQUIRED — always applies, even when the user specified a theme/style)
═══════════════════════════════════════════════════════════════════════════════

The dashboard is embedded in an iframe whose width is NOT fixed — the SAME code renders in a narrow chat side-panel (~360–480px), a normal report view (~900px), and a full-screen / published view (up to ~1920px). It MUST reflow gracefully at every width with NO horizontal page scroll and NO clipped or squished content. Build it fluid and mobile-first; only deviate if the user EXPLICITLY asked for a fixed width.

Concrete rules — follow all of them:
- **Outer container:** fluid width, never a fixed pixel width. Use `w-full min-h-full` (add `max-w-screen-2xl mx-auto` only if you want to cap width on huge screens). Responsive padding: `p-4 md:p-6 lg:p-8`. NEVER `w-[1200px]`, `min-w-[...]`, or any fixed-pixel width on layout containers.
- **KPI / stat rows:** use a responsive grid that collapses on narrow screens, e.g. `grid grid-cols-2 md:grid-cols-4 gap-4` (2-up on mobile → 4-up on desktop). Do NOT use a flex row of fixed-width cards that overflows.
- **Chart grids:** start single-column and add columns at breakpoints, e.g. `grid grid-cols-1 lg:grid-cols-2 gap-6`. A primary/feature chart can stay full-width (`col-span-full` or its own row). Never lock a multi-column grid with no single-column fallback.
- **Charts:** give each `<EChart>` a `w-full` container and a fixed `height` (px) — it auto-resizes to its container via ResizeObserver, so width takes care of itself. Do not set a pixel width on charts.
- **Tables & wide content:** wrap in `<div className="overflow-x-auto">` so a wide table scrolls inside its card instead of blowing out the page width. Use `min-w-full` on the `<table>`, not a fixed width.
- **Filter bars:** `flex flex-wrap gap-3` so filters wrap to the next line on narrow widths instead of overflowing.
- **Text & numbers:** allow large KPI numbers to scale (e.g. `text-2xl md:text-3xl`) and use `truncate`/`break-words` where labels can be long, so nothing overflows its card.
- **Sanity check:** before finishing, mentally render at ~380px wide — every row must wrap to 1–2 columns, no element wider than the viewport, no horizontal scrollbar on the body.

═══════════════════════════════════════════════════════════════════════════════
OUTPUT FORMAT
═══════════════════════════════════════════════════════════════════════════════

```
<script type="text/babel">
function App() {{
  const data = useArtifactData();
  if (!data) return <div className="flex items-center justify-center h-screen text-gray-400"><LoadingSpinner size={{32}} /></div>;
  const viz = data.visualizations;
  // ... concise dashboard code
}}
ReactDOM.createRoot(document.getElementById('root')).render(<App />);
</script>
```

Structure: all code should be inside `function App() {{ ... }}` with `ReactDOM.createRoot(document.getElementById('root')).render(<App />);` at the end. Do not put return statements outside a function.

Rules: `<script type="text/babel">` wrapper. `useArtifactData()` for data. `<EChart option={{...}} />` for charts. Pass `viz={{viz[N]}}` to every KPICard/SectionCard so the built-in info popover shows the data behind it. RESPONSIVE — fluid width, responsive grids (`grid-cols-1 md:grid-cols-2 lg:grid-cols-N`), no fixed-pixel widths, no horizontal page scroll at any width (see RESPONSIVE LAYOUT section above); required unless the user asked for a fixed width. Handle zero rows. No hardcoded data. No UUIDs/branding/emoji. Guard nullish values before string methods (use `(val || '')` or `String(val ?? '')`).

**Code size:** Write compact code — no unnecessary variables, comments, or verbose JSX. Omit default props. Don't repeat theme styling the 'dash' theme already provides. Prefer inline expressions over separate variables when used once. For simple dashboards target under 8K characters. For detailed/specific user requests, use as much space as needed to faithfully implement their design — fidelity to the user's request is more important than brevity."""

    def _build_page_prompt(
        self,
        user_prompt: str,
        title: str | None,
        viz_profiles: List[Dict[str, Any]],
        instructions_context: str,
        report_title: str | None,
        allow_llm_see_data: bool,
        messages_context: str = "",
        image_count: int = 0,
        organization_settings: Any = None,
        files: Optional[List[Dict[str, Any]]] = None,
    ) -> str:
        """Build the dynamic user prompt for page/dashboard generation.

        Static reference material lives in _build_page_system_prompt (cached
        as the system prompt); this carries only per-call state.
        """
        viz_json = json.dumps(viz_profiles, indent=2, default=str)

        language_directive = build_language_directive(organization_settings)

        # Build attached images context
        images_context = ""
        if image_count > 0:
            images_context = f"\n**Attached Images:** {image_count} image(s) provided for visual reference. Use these to understand the design intent, branding, color schemes, or layout preferences the user wants to incorporate."

        # Build embedded-files context (generated images / uploaded images or PDFs).
        # These are rendered via the <BowFile> sandbox component by file id.
        files_context = ""
        if files:
            lines = [
                f'- id="{f["id"]}"  type={f.get("content_type", "")}  name={f.get("filename", "")}'
                for f in files
            ]
            files_context = (
                "\n**Embedded Files (render with `<BowFile>`):** You MUST place each of these "
                "files in the layout using `<BowFile id=\"<id>\" />` (see the BowFile entry in the "
                "sandbox runtime docs). Images render inline; PDFs render in a viewer. Do NOT inline "
                "base64 and do NOT use a raw <img src>. Available files:\n" + "\n".join(lines)
            )

        # Note: Previous artifact code is now available via observation context (from create_artifact/read_artifact)
        # The planner can call read_artifact if needed to load previous code into context

        return f"""═══════════════════════════════════════════════════════════════════════════════
Design request (primary specification — takes precedence when it conflicts with reference defaults)
═══════════════════════════════════════════════════════════════════════════════

**Report Title:** {report_title or title or 'Dashboard'}
**User Request:** {user_prompt}
{images_context}
{files_context}
{f"**Organization Instructions:**{chr(10)}{instructions_context}" if instructions_context else ""}

{f"**Conversation History:**{chr(10)}{messages_context}" if messages_context else ""}
{language_directive}

If the user specified a theme, layout, colors, or style above — follow that exactly.
If the user did not specify styling, use the design guidance from the reference instructions.

YOUR VISUALIZATIONS:

{viz_json}

{"(Sample data included above — rows are a sample capped at 100; row_count is the true dataset size)" if allow_llm_see_data else "(Data samples hidden for privacy - use column names and row_count to understand the data structure)"}

Now create the dashboard:"""

    async def _build_file_datauris(self, db, included_files: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Read embedded files from disk and return them as data: URIs.

        Shared with the PDF export path — see _artifact_images.build_file_datauris.
        """
        from ._artifact_images import build_file_datauris

        return await build_file_datauris(db, included_files)

    def _build_prompt(
        self,
        user_prompt: str,
        title: str | None,
        mode: str,
        viz_profiles: List[Dict[str, Any]],
        instructions_context: str,
        report_title: str | None,
        allow_llm_see_data: bool,
        messages_context: str = "",
        image_count: int = 0,
        organization_settings: Any = None,
        files: Optional[List[Dict[str, Any]]] = None,
        report: Any = None,
        resolved_theme: Any = None,
    ) -> str:
        """Build the prompt for generating artifact code. Dispatches to mode-specific builders."""
        if mode == "slides":
            return self._build_slides_prompt(
                user_prompt=user_prompt,
                title=title,
                viz_profiles=viz_profiles,
                instructions_context=instructions_context,
                report_title=report_title,
                allow_llm_see_data=allow_llm_see_data,
                messages_context=messages_context,
                image_count=image_count,
                organization_settings=organization_settings,
                files=files,
                report=report,
                resolved_theme=resolved_theme,
            )
        return self._build_page_prompt(
            user_prompt=user_prompt,
            title=title,
            viz_profiles=viz_profiles,
            instructions_context=instructions_context,
            report_title=report_title,
            allow_llm_see_data=allow_llm_see_data,
            messages_context=messages_context,
            image_count=image_count,
            organization_settings=organization_settings,
            files=files,
        )

    def _extract_code(self, response: str, mode: str = "page") -> str:
        """Extract the code from the LLM response.

        For 'page' mode: Extract React code from <script type="text/babel"> tags
        For 'slides' mode: Extract python-pptx code from python code blocks
        """
        if mode == "slides":
            return self._extract_slides_python(response)

        # Dashboard mode - extract React code from script tags
        start_marker = "<script type=\"text/babel\">"
        end_marker = "</script>"

        start_idx = response.find(start_marker)
        if start_idx == -1:
            # Try alternative markers
            start_marker = "<script type='text/babel'>"
            start_idx = response.find(start_marker)

        if start_idx != -1:
            end_idx = response.find(end_marker, start_idx)
            if end_idx != -1:
                code = response[start_idx:end_idx + len(end_marker)]
                return self._sanitize_code(self._ensure_app_wrapper(code))

        # If no script tags found, wrap the response
        code = response.strip()
        if not code.startswith("<script"):
            code = f'<script type="text/babel">\n{code}\n</script>'

        return self._sanitize_code(self._ensure_app_wrapper(code))

    @staticmethod
    def _sanitize_code(code: str) -> str:
        """Fix common LLM code generation artifacts deterministically."""
        import re

        # Fix double-brace pattern: function App() {\n{ ... }\n}
        # The LLM sometimes wraps the function body in an extra block scope.
        # Match: function App() {\n{ at the start, and }\n} at the end (before render call)
        code = re.sub(
            r'(function\s+\w+\s*\([^)]*\)\s*\{)\s*\n\s*\{',
            r'\1',
            code,
        )
        # Remove the matching trailing extra }
        # Look for }\n}\n before ReactDOM.createRoot
        code = re.sub(
            r'\}\s*\n\s*\}\s*\n(\s*ReactDOM\.createRoot)',
            r'}\n\1',
            code,
        )

        return code

    @staticmethod
    def _ensure_app_wrapper(code: str) -> str:
        """Ensure code has a proper App component wrapper.

        LLM sometimes outputs bare return statements outside a function.
        Detect and fix by wrapping the inner code in function App() + ReactDOM.createRoot.
        """
        import re

        # Check if code already has an App function/component
        if re.search(r'function\s+App\s*\(', code) or re.search(r'(?:const|let|var)\s+App\s*=', code):
            return code

        # Extract inner code between script tags
        inner_match = re.search(
            r'<script\s+type=["\']text/babel["\']>\s*([\s\S]*?)\s*</script>',
            code
        )
        if not inner_match:
            return code

        inner = inner_match.group(1).strip()

        # Strip any existing broken ReactDOM.createRoot/render calls
        inner = re.sub(r'ReactDOM\.createRoot\(.*?\)\.render\(.*?\);?\s*$', '', inner, flags=re.DOTALL).strip()

        # DEF-008: this wrapper exists for a real case — the model emits a bare
        # `return (<div>…)` with no function around it. But it wrapped WHATEVER
        # it found, without asking whether that was code. When the model replied
        # with a status sentence instead of a component, the wrapper produced:
        #
        #     function App() {
        #     Creating the full-totals CFC Sales Dashboard from the compact data.
        #     }
        #
        # — a syntactically perfect shell around a syntactically impossible body.
        # The tool then returned success and the artifact stored as `completed`,
        # so the failure was discovered by the BROWSER, in front of the user
        # ("Missing semicolon (3:8)"). Wrapping prose cannot produce a component;
        # returning the code untouched lets the caller's guard see prose for what
        # it is and ask for a real answer.
        if not _looks_like_component_code(inner):
            logger.warning(
                "_ensure_app_wrapper: refusing to wrap non-code (%d chars) — "
                "reply looks like prose, not a component", len(inner)
            )
            return code

        logger.warning("_ensure_app_wrapper: LLM output missing function App() wrapper — auto-wrapping")

        wrapped = (
            '<script type="text/babel">\n'
            'function App() {\n'
            f'{inner}\n'
            '}\n'
            "ReactDOM.createRoot(document.getElementById('root')).render(<App />);\n"
            '</script>'
        )
        return wrapped

    def _extract_slides_python(self, response: str) -> str:
        """Extract python-pptx code for slides mode.

        ★Two defects lived in the version this replaces, and both surfaced as
        the same unhelpful message — ``invalid syntax (<string>, line 1)``:

        1. The generic fence pattern was ``r'```\\s*([\\s\\S]*?)```'``. ``\\s*``
           does not match a language tag, so a block opened with anything other
           than the exact lowercase ``python`` — ``py``, ``Python`` — kept that
           word as the first line of the extracted "code". Line 1 was then a
           bare identifier, and the deck died before python-pptx was reached.

        2. Both patterns required a CLOSING fence. When the model ran out of
           output tokens mid-deck there was no closer, both searches missed,
           and the `def generate_slides` fallback returned everything from the
           def to wherever the text stopped — code cut off mid-statement. The
           deck that prompted this fix ended at ``for row in`` on line 881 of
           an otherwise valid 37,720-character function.

        The second one is why this returns rather than repairs: a truncated
        deck cannot be salvaged by cleverness here, only regenerated. What this
        can do is hand the caller something that names the real problem, which
        is `slides_code_problem`.
        """
        import re

        # An opening fence with an OPTIONAL language tag, and an optional
        # closer. `(?:```|\Z)` is the whole point — a run that was cut off has
        # no closing fence and must still yield its code.
        fenced = re.search(
            r'```[ \t]*[A-Za-z0-9_+-]*[ \t]*\r?\n([\s\S]*?)(?:```|\Z)',
            response,
        )
        if fenced:
            body = fenced.group(1).strip()
            if body:
                return body

        # No fence at all. Anchor on the function and, when the deck finished
        # properly, stop at its final save.
        func_start = response.find('def generate_slides')
        if func_start != -1:
            save_end = response.rfind('prs.save(')
            if save_end != -1:
                end_idx = response.find(')', save_end)
                if end_idx != -1:
                    return response[func_start:end_idx + 1].strip()
            return response[func_start:].strip()

        # Fallback: return the response as-is
        return response.strip()

    @staticmethod
    def _slides_code_problem(code: str) -> Optional[str]:
        """Name what is wrong with generated deck code, before running it.

        ★Returns a sentence a model can act on, or None. The point is the
        difference between the two things that were being reported identically:
        a deck the model got WRONG, which a retry can fix from the interpreter's
        message, and a deck the model did not FINISH, where the interpreter's
        message describes the cut rather than the mistake and telling the model
        to "fix that specific error" sends it to repair a line it wrote
        correctly.
        """
        if not (code or "").strip():
            return "The generated deck code was empty."

        if "def generate_slides" not in code:
            return (
                "The generated code has no `def generate_slides(visualizations, report):` "
                "function, so there was nothing to run."
            )

        try:
            compile(code, "<string>", "exec")
            return None
        except SyntaxError as err:
            # ★Bound outside the handler on purpose: Python deletes the `as`
            # name when the except block ends, so reading it below would raise
            # UnboundLocalError — a crash inside the diagnostic, replacing the
            # message it exists to produce.
            syntax_error = err

        # It does not parse. Distinguish "cut off" from "written wrong": a deck
        # that never reaches its own save call stopped early.
        if "prs.save(" not in code:
            return (
                "The generated deck code was CUT OFF before it finished — it never "
                "reaches `prs.save(_pptx_output_path)`. This is a length problem, not "
                "a mistake in the code: produce a SHORTER deck (fewer slides, less "
                "repeated layout code, factor repeated shape-drawing into helpers) so "
                "that it completes."
            )
        return (
            f"The generated deck code does not parse: {syntax_error.msg} "
            f"(line {syntax_error.lineno}). Return the complete corrected code."
        )

