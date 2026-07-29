"""Four presentation properties that must hold for data from any connector.

1. No unit is rendered that the data did not supply (no hardcoded currency).
2. A value is never cut off by the box it is drawn in.
3. Value axes read the same in the browser and in every export.
4. A category axis never presents two different things under one label.

None of these tests names a table, column, connector, currency or figure.
"""
from __future__ import annotations

import re
from pathlib import Path

import pytest

from app.services.number_format import (
    ABBREVIATION_STEPS,
    PPTX_PLAIN_NUMBER_FORMAT,
    abbreviate_number,
    axis_label_formatter_js,
    normalize_currency_code,
    pptx_axis_number_format,
    qualify_duplicate_labels,
)

REPO = Path(__file__).resolve().parents[4]
FRONTEND = REPO / "frontend"


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------

def _strip_js_comments(source: str) -> str:
    """Remove // and /* */ comments so an assertion can never be satisfied by
    a comment that merely talks about the thing being asserted."""
    source = re.sub(r"/\*.*?\*/", "", source, flags=re.S)
    return re.sub(r"(?m)//.*$", "", source)


def _code_of(rel_path: str) -> str:
    text = (FRONTEND / rel_path).read_text(encoding="utf-8")
    assert text.strip(), f"{rel_path} is empty"
    return _strip_js_comments(text)


def test_comment_stripper_actually_strips():
    """Guard the guard: if this stopped working, every source assertion below
    would start passing on comments alone."""
    stripped = _strip_js_comments("var a = 1; // currency: 'USD'\n/* currency: 'USD' */\nvar b;")
    assert "USD" not in stripped
    assert "var a = 1;" in stripped and "var b;" in stripped


# ---------------------------------------------------------------------------
# Property 1 — never render a unit the data did not supply
# ---------------------------------------------------------------------------

@pytest.mark.parametrize("value", [True, 1, "", "  ", "US", "USDX", "12A", None, {}, ["USD"]])
def test_non_code_never_names_a_currency(value):
    assert normalize_currency_code(value) is None


@pytest.mark.parametrize("value,expected", [("usd", "USD"), (" ils ", "ILS"), ("MMK", "MMK")])
def test_three_letter_code_is_honored(value, expected):
    assert normalize_currency_code(value) == expected


def test_sandbox_formatter_has_no_default_currency():
    code = _code_of("public/libs/artifact-globals.js")
    # The defect: `opts.currency === true ? 'USD' : opts.currency`.
    assert "'USD'" not in code and '"USD"' not in code, (
        "artifact-globals must not carry a fallback currency code"
    )
    assert "function currencyCode(" in code, "the ISO-4217 guard must exist"
    assert "style: 'currency', currency: code" in code, (
        "currency style must be driven by the validated code, not by opts directly"
    )


def test_metric_card_has_no_default_currency():
    code = _code_of("components/dashboard/kpi/MetricCard.vue")
    assert "'USD'" not in code and '"USD"' not in code


def test_count_widget_has_no_default_currency():
    code = _code_of("components/RenderCount.vue")
    assert "currency: 'USD'" not in code
    assert "currencyCode" in code, "the widget must read a code from the view config"


def test_sandbox_prompt_does_not_advertise_a_unitless_currency_symbol():
    text = (
        REPO / "backend/app/ai/tools/implementations/_sandbox_context.py"
    ).read_text(encoding="utf-8")
    fmt_line = next(line for line in text.splitlines() if "`fmt(n" in line)
    assert "ISO-4217" in fmt_line, "the model must be told a code is required"
    assert "{currency: true}" in fmt_line and "plain number" in fmt_line


# ---------------------------------------------------------------------------
# Property 2 — a value is never cut off by its own card
# ---------------------------------------------------------------------------

def test_metric_card_value_is_sized_to_content_not_fixed():
    code = _code_of("components/dashboard/kpi/MetricCard.vue")
    assert "text-4xl" not in code, "a fixed type size cannot fit an arbitrary value"
    assert "valueStyle" in code
    for prop in ("wordBreak", "overflowWrap", "fontSize"):
        assert prop in code, f"the value must be able to {prop}"


def test_sandbox_kpi_card_value_is_sized_to_content_not_fixed():
    """The sandbox tile sizes its value to the box it was actually given.

    ★This originally asserted `fitValueStyle(props.value…)`, a helper that
    scaled the type off `String(value).length`. That is still better than a
    fixed size, but it was measured clipping a real dashboard: 105,150,299,753
    needed 285px in a 222px card, because character counts do not know how wide
    a character is or how wide the card is. The tile now MEASURES
    (`BowFitText`), so the rule is asserted at that level — see
    test_artifact_kpi_value_not_clipped.py for the full guard.
    """
    code = _code_of("public/libs/artifact-globals.js")
    assert "text-2xl font-semibold" not in code, "fixed type size on a clipping card"
    assert "window.BowFitText =" in code, "the measuring value-fitter is gone"
    assert "scrollWidth" in code, (
        "the tile no longer measures the rendered text against its box — it is "
        "guessing, and a guess loses digits at some card width"
    )
    assert re.search(r"h\(window\.BowFitText,", code), (
        "the KPI tile must actually render its value through BowFitText"
    )


def test_no_render_path_hides_overflow_without_wrapping():
    """Both cards clip their own overflow; both must therefore wrap."""
    for rel in ("public/libs/artifact-globals.js", "components/dashboard/kpi/MetricCard.vue"):
        code = _code_of(rel)
        if "overflow-hidden" in code:
            assert "overflowWrap" in code or "break-words" in code, (
                f"{rel} clips overflow but never allows the value to wrap"
            )


# ---------------------------------------------------------------------------
# Property 3 — one abbreviation rule, every render path
# ---------------------------------------------------------------------------

@pytest.mark.parametrize(
    "value,expected",
    [
        (0, "0"),
        (999, "999"),
        (-999, "-999"),
        (1_000, "1.0K"),
        (8_000_000_000, "8.0B"),
        (70_000_000_000, "70.0B"),
        (4_300_000_000, "4.3B"),
        (-4_300_000_000, "-4.3B"),
        (2_500_000_000_000, "2.5T"),
    ],
)
def test_abbreviation_rule(value, expected):
    assert abbreviate_number(value) == expected


def test_abbreviation_never_emits_a_unit():
    for value in (1, 1e3, 1e6, 1e9, 1e12):
        text = abbreviate_number(value)
        assert not any(sym in text for sym in "$€£¥₪")


@pytest.mark.parametrize(
    "max_abs,expected_suffix",
    [(999, None), (5_000, "K"), (5_000_000, "M"), (8_000_000_000, "B"), (3e12, "T")],
)
def test_pptx_number_format_matches_the_same_buckets(max_abs, expected_suffix):
    code = pptx_axis_number_format(max_abs)
    if expected_suffix is None:
        assert code == PPTX_PLAIN_NUMBER_FORMAT
    else:
        assert code.endswith(f'"{expected_suffix}"')
        # One trailing comma per 1000 of scaling — B must divide by 1e9.
        commas = code.count(",", code.index("0"))
        assert commas == {"K": 1, "M": 2, "B": 3, "T": 4}[expected_suffix]


def test_pptx_number_format_survives_junk():
    for bad in (None, "abc", float("nan")):
        assert pptx_axis_number_format(bad) == PPTX_PLAIN_NUMBER_FORMAT


def test_js_formatter_is_self_contained():
    """The headless renderer behind the Word export loads ECharts and nothing
    else — a formatter that called a sandbox global would silently do nothing
    in exactly the path that was broken."""
    js = axis_label_formatter_js()
    assert "fmt(" not in js and "window." not in js
    for threshold, suffix in ABBREVIATION_STEPS:
        assert f"'{suffix}'" in js
        assert f"{threshold:.0f}" in js


def test_generated_chart_code_formats_every_value_axis():
    from app.services.artifact_codegen import generate_echart_option_code

    for dm in (
        {"type": "bar_chart", "series": [{"key": "k", "value": "v", "name": "V"}]},
        {"type": "line_chart", "series": [{"key": "k", "value": "v"}], "group_by": "g"},
        {"type": "bar_chart", "series": [{"key": "k", "value": "v"}], "horizontal": True},
        {"type": "scatter_plot", "series": [{"x": "a", "y": "b"}]},
    ):
        code = generate_echart_option_code(dm, 0)
        value_axes = re.findall(r"[xy]Axis: \{ type: 'value'[^\n]*", code)
        assert value_axes, f"no value axis emitted for {dm['type']}"
        for axis in value_axes:
            assert "formatter:" in axis, f"unformatted value axis in {dm['type']}: {axis}"


def test_pptx_export_normalizes_a_real_chart(tmp_path):
    """python-pptx never errors on a badly formatted axis, so assert on the
    file it wrote."""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    from app.ai.code_execution.pptx_executor import normalize_chart_axes

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    # Two categories deliberately share a label.
    data.categories = ["dup", "dup", "other"]
    data.add_series("s", (70_000_000_000, 8_000_000_000, 4_300_000_000))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6), Inches(4), data
    )
    path = tmp_path / "deck.pptx"
    prs.save(str(path))

    # Before: python-pptx leaves the axis on the general format.
    before = Presentation(str(path)).slides[0].shapes[0].chart
    assert before.value_axis.tick_labels.number_format in ("General", None)

    assert normalize_chart_axes(path) == 1

    after = Presentation(str(path)).slides[0].shapes[0].chart
    assert after.value_axis.tick_labels.number_format == pptx_axis_number_format(7e10)
    assert after.value_axis.tick_labels.number_format_is_linked is False

    # Read the labels back out of the SAVED file. Asserting only on the return
    # value of normalize_chart_axes passes while the label rewrite silently
    # does nothing — that is exactly what happened, and only a rendered slide
    # showed it.
    labels = [str(c) for plot in after.plots for c in plot.categories]
    assert len(set(labels)) == len(labels), f"duplicate labels survived: {labels}"
    assert labels[0].startswith("dup") and labels[1].startswith("dup")
    assert labels[0] != labels[1]


def test_pptx_normalize_never_raises_on_a_bad_file(tmp_path):
    from app.ai.code_execution.pptx_executor import normalize_chart_axes

    junk = tmp_path / "not-a-deck.pptx"
    junk.write_bytes(b"nonsense")
    assert normalize_chart_axes(junk) == 0


# ---------------------------------------------------------------------------
# Property 4 — a category axis never labels two things the same
# ---------------------------------------------------------------------------

def test_unique_labels_are_untouched():
    assert qualify_duplicate_labels(["a", "b", "c"]) == ["a", "b", "c"]


def test_collision_is_qualified_by_the_column_that_differs():
    out = qualify_duplicate_labels(
        ["Common", "Common", "Rare"], ["ParentA", "ParentB", "ParentA"]
    )
    assert out == ["Common (ParentA)", "Common (ParentB)", "Rare"]
    assert len(set(out)) == len(out)


def test_collision_with_no_qualifier_states_the_ambiguity():
    out = qualify_duplicate_labels(["Common", "Common"])
    assert out == ["Common (1 of 2)", "Common (2 of 2)"]


def test_qualifier_that_also_repeats_still_yields_unique_labels():
    out = qualify_duplicate_labels(["x", "x", "x"], ["p", "p", "q"])
    assert len(set(out)) == 3


def test_generated_chart_code_disambiguates_categories():
    from app.services.artifact_codegen import generate_echart_option_code

    code = generate_echart_option_code(
        {"type": "bar_chart", "series": [{"key": "k", "value": "v"}]}, 0
    )
    assert "_qualCol" in code and "_catIds" in code
    # The old shape looked categories up by the label alone, which merges two
    # different things whose labels collide.
    assert "rows.find(r => r['k'] === c)" not in code


def test_frontend_chart_builder_disambiguates_categories():
    code = _code_of("components/dashboard/charts/EChartsVisual.vue")
    assert "pickQualifierColumn(" in code
    assert "qualifyDuplicateLabels(" in code
    assert "catIdOf" in code
    assert 'const cat = String(r[categoryKey] ?? \'\')' not in code, (
        "buckets must key on category identity, not on the label"
    )
