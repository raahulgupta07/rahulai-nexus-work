"""FIX-A (2026-08-21) — the theme's ground sheet buried the model's dark slides.

Measured on 0.0.543.18, four decks generated through the chat API: the
boardroom furniture pass painted a full-bleed WHITE ``ground:flat`` rectangle
as the first spTree shape on EVERY slide — including slides where the
generated deck code had set its own dark background via ``<p:bg>``
(0B0F19 cover, 0E3A53 closer). The rectangle sits over the background and
under the text, and the pass never recolors text, so the white headlines those
slides were designed with became white-on-white at exactly 1.00:1 — invisible.
Victims: the Fabric deck's cover and closing titles, and the main heading of
every content slide in the Power BI deck (the model made that whole deck dark).

The fix: ``_slide_sets_own_ground`` checks the slide XML for ``<p:bg>``. When
the model gave a slide its own background, the pass skips the GROUND layers
only (ground sheet, paper grain, ruling, margin rule) and records
``ground:kept_slide_own``. Overlays (masthead, tracker, footer, chip) still
paint — they sit above the text with their own colors.

★These tests build real pptx files and run the real pass — no schema, safe in
`tests/unit/fork/` where the migration fixture is a no-op.
"""
import pathlib

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches


BOARDROOM_LAYOUT = {
    "ground": "flat",
    "ground_color": "#FFFFFF",
    "footer_rule": "hairline",
    "footer": {"left": "City Mart", "right": "{page} / {pages}"},
}
THEME = {"palette": {"muted_text": "#64748B", "heading_text": "#1F3A5F"}}


def _deck(tmp_path, dark_first: bool) -> pathlib.Path:
    """Two-slide deck. Slide 1 optionally sets its OWN dark background —
    exactly what generated deck code does with ``slide.background.fill``."""
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    for i in range(2):
        slide = prs.slides.add_slide(blank)
        if dark_first and i == 0:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x0B, 0x0F, 0x19)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.text = "State of City Mart: Together We Win"
    path = tmp_path / ("dark.pptx" if dark_first else "plain.pptx")
    prs.save(str(path))
    return path


def _paint(path):
    from app.ai.decks.motifs import paint_theme_furniture
    return paint_theme_furniture(path, THEME, dict(BOARDROOM_LAYOUT))


def _first_shape_is_full_bleed_white(slide) -> bool:
    """The defect's signature: a white rectangle at 0,0 covering the canvas,
    first in the spTree draw order."""
    shapes = list(slide.shapes)
    if not shapes:
        return False
    s = shapes[0]
    try:
        covers = (s.left or 0) <= 0 and (s.top or 0) <= 0 and \
            (s.width or 0) >= 12000000 and (s.height or 0) >= 6700000
        white = s.fill.fore_color.rgb == RGBColor(0xFF, 0xFF, 0xFF)
        return covers and white
    except Exception:
        return False


class TestASlideWithItsOwnBackgroundKeepsIt:
    def test_no_white_sheet_over_a_dark_slide(self, tmp_path):
        """★The defect exactly as measured: dark ``<p:bg>``, white sheet on top."""
        path = _deck(tmp_path, dark_first=True)
        _paint(path)
        prs = Presentation(str(path))
        assert not _first_shape_is_full_bleed_white(prs.slides[0])

    def test_the_slides_own_bg_element_survives(self, tmp_path):
        path = _deck(tmp_path, dark_first=True)
        _paint(path)
        prs = Presentation(str(path))
        assert prs.slides[0]._element.cSld.find(qn("p:bg")) is not None

    def test_the_skip_is_recorded_not_silent(self, tmp_path):
        """A clean run must be distinguishable from a run where the guard never
        fired — the painted list carries the decision."""
        path = _deck(tmp_path, dark_first=True)
        result = _paint(path)
        assert "ground:kept_slide_own" in result["painted"]

    def test_overlays_still_paint_on_a_dark_slide(self, tmp_path):
        """★The guard skips GROUND layers only. A fix that skipped the whole
        pass would delete the footer and tracker from every dark slide."""
        path = _deck(tmp_path, dark_first=True)
        prs_before = Presentation(str(path))
        count_before = len(prs_before.slides[0].shapes)
        _paint(path)
        prs = Presentation(str(path))
        assert len(prs.slides[0].shapes) > count_before


class TestAPlainSlideStillGetsTheThemeGround:
    def test_the_ground_sheet_still_paints(self, tmp_path):
        """★Positive control. A change that stopped painting grounds everywhere
        would satisfy every test above and delete the theme."""
        path = _deck(tmp_path, dark_first=False)
        result = _paint(path)
        assert "ground:flat" in result["painted"]

    def test_a_mixed_deck_treats_each_slide_on_its_own(self, tmp_path):
        """Slide 1 dark (kept), slide 2 plain (themed) — one deck, both rules."""
        path = _deck(tmp_path, dark_first=True)
        result = _paint(path)
        assert "ground:kept_slide_own" in result["painted"]
        assert "ground:flat" in result["painted"]
        prs = Presentation(str(path))
        assert _first_shape_is_full_bleed_white(prs.slides[1])


class TestTheDefaultThemeIsReadable:
    def test_boardroom_muted_text_clears_the_readability_bar(self):
        """FIX-B. boardroom is DEFAULT_THEME_ID, and its muted grey (#8E9BAA,
        2.83:1 on white) was used for every subtitle, caption and axis label of
        every un-styled deck. 4.5:1 is the bar."""
        from app.ai.decks import pptx_themes as pt

        def lum(h):
            h = h.lstrip("#")
            c = [int(h[i:i + 2], 16) / 255 for i in (0, 2, 4)]
            c = [x / 12.92 if x <= 0.04045 else ((x + 0.055) / 1.055) ** 2.4 for x in c]
            return 0.2126 * c[0] + 0.7152 * c[1] + 0.0722 * c[2]

        palette = pt._THEMES["boardroom"].palette
        bg, muted, body = palette["background"], palette["muted_text"], palette["body_text"]
        hi, lo = sorted((lum(bg), lum(muted)), reverse=True)
        assert (hi + 0.05) / (lo + 0.05) >= 4.5
        # muted must stay LIGHTER than body, or the hierarchy inverts
        assert lum(muted) > lum(body)


class TestPreviewsAreRenderedForRetina:
    def test_the_default_dpi_covers_a_2x_display(self):
        """FIX-C. 150 dpi → 2000px-wide PNGs, upsampled on retina — the user's
        "blur". 220 dpi ≈ 2933px, ≥2× the ~1400px display width of the viewer."""
        import inspect
        from app.ai.code_execution.pptx_executor import PptxPreviewService
        sig = inspect.signature(PptxPreviewService.generate_previews)
        assert sig.parameters["dpi"].default >= 220
