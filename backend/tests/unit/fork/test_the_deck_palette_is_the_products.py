"""A deck's palette belongs to the product, not to whatever the model typed.

Generated deck code writes ``RGBColor(0x1F, 0x4E, 0x79)`` literals because that
is what python-pptx examples look like. Two things follow, and they are
deliberately kept apart:

* the palette is handed IN — ``theme`` is bound into the exec namespace, so
  generated code can name a role instead of inventing a hex; and
* the finished file is READ BACK — a near-miss of a role is snapped onto it,
  and a colour that is nowhere near any role is REPORTED, never guessed at and
  never grounds for refusing the deck.

That last sentence is the important one. This tree already shipped a guard that
refused valid generated code (``_SANCTIONED_FILE_COLLECTIONS`` bound the upload
list by NAME where the runtime binds it by POSITION), and users saw the product
reject correct work. ``off_palette`` is a number to read, not a gate to fail.
"""

import pytest

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from app.ai.code_execution.pptx_executor import (
    PptxCodeExecutor,
    apply_theme_palette,
    palette_from_theme,
)

# A small brand palette, in the shape a caller would actually pass: colour
# roles, a series list, and non-colour entries that must be ignored rather than
# rejected.
THEME = {
    "PRIMARY": "#1F4E79",
    "ACCENT": "E8710A",
    "SERIES": ["#1F4E79", "#7F7F7F"],
    "FONT": "Inter",
    "TITLE_SIZE": 32,
}

FAR_AWAY = RGBColor(0xC0, 0x00, 0x00)          # a deliberate red, nowhere near a role
NEAR_MISS = RGBColor(0x20, 0x4E, 0x7A)         # PRIMARY, rounded differently
ON_PALETTE = RGBColor(0x1F, 0x4E, 0x79)        # PRIMARY exactly


def _deck(path, colors):
    """One slide, one textbox per colour, coloured fill and coloured text."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, color in enumerate(colors):
        box = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(4), Inches(0.8))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = f"row {i}"
        run.font.size = Pt(18)
        run.font.color.rgb = color
    prs.save(str(path))
    return path


def _colors_in(path):
    """Every run colour in a saved deck, read back independently of the pass."""
    found = []
    prs = Presentation(str(path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    try:
                        found.append(str(run.font.color.rgb).upper())
                    except Exception:
                        pass
    return found


# ---------------------------------------------------------------------------
# The palette reaches generated code
# ---------------------------------------------------------------------------

CODE = """
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
run = box.text_frame.paragraphs[0].add_run()
run.text = "hello"
run.font.color.rgb = RGBColor.from_string(theme['PRIMARY'].lstrip('#'))
prs.save(_pptx_output_path)
"""

NO_THEME_CODE = """
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
prs.save(_pptx_output_path)
"""


def test_generated_code_can_read_the_theme(tmp_path):
    out = tmp_path / "themed.pptx"
    PptxCodeExecutor().execute_pptx_code(
        code=CODE,
        visualizations=[],
        report={"id": "r1", "title": "t"},
        output_path=out,
        theme=THEME,
    )
    assert _colors_in(out) == ["1F4E79"]


def test_a_theme_can_be_read_with_get(tmp_path):
    """`theme` is a plain dict — dict access and .get() both have to work."""
    out = tmp_path / "get.pptx"
    code = CODE.replace("theme['PRIMARY']", "theme.get('PRIMARY', '#000000')")
    PptxCodeExecutor().execute_pptx_code(
        code=code,
        visualizations=[],
        report={"id": "r1", "title": "t"},
        output_path=out,
        theme=THEME,
    )
    assert _colors_in(out) == ["1F4E79"]


def test_without_a_theme_the_namespace_is_what_it_always_was(tmp_path):
    """Absent theme must change nothing: the name is not bound at all.

    A deck that does not ask for a theme still builds — that is the half that
    matters, and it is asserted first so this cannot degrade into a test that
    only proves an exception.
    """
    out = tmp_path / "plain.pptx"
    path, _ = PptxCodeExecutor().execute_pptx_code(
        code=NO_THEME_CODE,
        visualizations=[],
        report={"id": "r1", "title": "t"},
        output_path=out,
    )
    assert path.exists()

    with pytest.raises(NameError):
        PptxCodeExecutor().execute_pptx_code(
            code=CODE,
            visualizations=[],
            report={"id": "r1", "title": "t"},
            output_path=tmp_path / "unthemed.pptx",
        )


def test_execute_takes_theme_as_an_optional_keyword():
    import inspect

    sig = inspect.signature(PptxCodeExecutor.execute_pptx_code)
    param = sig.parameters["theme"]
    assert param.default is None
    assert param.kind is inspect.Parameter.KEYWORD_ONLY


# ---------------------------------------------------------------------------
# Reading a palette out of a theme dict
# ---------------------------------------------------------------------------

def test_a_palette_takes_colours_and_ignores_everything_else():
    palette = palette_from_theme(THEME)
    assert set(palette) == {"1F4E79", "E8710A", "7F7F7F"}
    assert palette["1F4E79"] == "PRIMARY"
    assert palette["7F7F7F"] == "SERIES[1]"


def test_junk_in_a_theme_is_ignored_not_fatal():
    assert palette_from_theme(None) == {}
    assert palette_from_theme("#1F4E79") == {}
    assert palette_from_theme({"A": None, "B": ["nope", 7], "C": "#abc"}) == {"AABBCC": "C"}


# ---------------------------------------------------------------------------
# What the post-pass reports
# ---------------------------------------------------------------------------

def test_off_palette_colours_are_counted(tmp_path):
    deck = _deck(tmp_path / "off.pptx", [ON_PALETTE, FAR_AWAY])
    summary = apply_theme_palette(deck, THEME)

    assert summary["off_palette"] == 1
    assert summary["off_palette_colors"] == ["C00000"]
    assert "1F4E79" in summary["colors_seen"]
    assert summary["palette"] == ["1F4E79", "7F7F7F", "E8710A"]


def test_a_colour_that_is_off_palette_is_left_alone(tmp_path):
    """Report it; do not guess at it. A deliberate red must survive the pass."""
    deck = _deck(tmp_path / "keep.pptx", [FAR_AWAY])
    summary = apply_theme_palette(deck, THEME)

    assert summary["remapped"] == 0
    assert _colors_in(deck) == ["C00000"]


def test_a_near_miss_is_snapped_onto_its_role(tmp_path):
    deck = _deck(tmp_path / "near.pptx", [NEAR_MISS])
    summary = apply_theme_palette(deck, THEME)

    assert summary["remapped"] == 1
    assert summary["off_palette"] == 0
    assert _colors_in(deck) == ["1F4E79"]


def test_an_on_palette_deck_is_neither_remapped_nor_flagged(tmp_path):
    deck = _deck(tmp_path / "clean.pptx", [ON_PALETTE, ON_PALETTE])
    summary = apply_theme_palette(deck, THEME)

    assert summary["remapped"] == 0
    assert summary["off_palette"] == 0
    assert summary["colors_seen"] == ["1F4E79"]


def test_without_a_palette_it_only_reports(tmp_path):
    """No theme means nothing to compare against — observe, never invent."""
    deck = _deck(tmp_path / "report.pptx", [FAR_AWAY, NEAR_MISS])
    summary = apply_theme_palette(deck)

    assert summary["remapped"] == 0
    assert summary["off_palette"] == 0
    assert summary["colors_seen"] == ["204E7A", "C00000"]
    assert summary["palette"] == []


def test_shape_fills_are_seen_too(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = FAR_AWAY
    deck = tmp_path / "fill.pptx"
    prs.save(str(deck))

    summary = apply_theme_palette(deck, THEME)
    assert "C00000" in summary["off_palette_colors"]


# ---------------------------------------------------------------------------
# It never raises — same contract as normalize_chart_axes
# ---------------------------------------------------------------------------

def test_a_missing_file_is_not_an_error(tmp_path):
    summary = apply_theme_palette(tmp_path / "nothing-here.pptx", THEME)
    assert summary["remapped"] == 0
    assert summary["colors_seen"] == []


def test_a_file_that_is_not_a_deck_is_not_an_error(tmp_path):
    junk = tmp_path / "broken.pptx"
    junk.write_bytes(b"this is not a zip, let alone a presentation")
    summary = apply_theme_palette(junk, THEME)
    assert summary == {
        "remapped": 0,
        "off_palette": 0,
        "off_palette_colors": [],
        "colors_seen": [],
        "palette": ["1F4E79", "7F7F7F", "E8710A"],
    }


def test_a_deck_that_cannot_be_saved_still_reports(tmp_path, monkeypatch):
    """A failed save costs the remap, not the deck and not the report."""
    deck = _deck(tmp_path / "readonly.pptx", [NEAR_MISS])

    import app.ai.code_execution.pptx_executor as mod

    original = mod.Presentation

    class _NoSave:
        def __init__(self, inner):
            self._inner = inner

        @property
        def slides(self):
            return self._inner.slides

        def save(self, _path):
            raise OSError("read-only volume")

    monkeypatch.setattr(mod, "Presentation", lambda p: _NoSave(original(p)))

    summary = apply_theme_palette(deck, THEME)
    assert summary["remapped"] == 0
    assert "204E7A" in summary["colors_seen"]
    assert _colors_in(deck) == ["204E7A"]
