"""A theme's prohibitions are enforced on the saved file, not asked for.

Across live deck generations the model honoured every theme's palette and
typography and ignored every theme's "strictly avoid" list — shadows and
rounded corners on every deck. These tests pin the mechanical half of that list
to the file: build a deck that violates it, run the pass, look at the XML.

The negative control is the important one. An enforcer that squares the corners
of a theme which permits rounded corners is a new bug, not a fix.
"""

import pytest

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from app.ai.code_execution.pptx_executor import enforce_theme_rules


ALL_TOKENS = ["shadows", "rounded_corners", "gradients", "multiple_accents", "boxes"]


def _rounded_shadowed_gradient_deck(path: Path) -> Path:
    """One slide, one rounded shape with an inherited shadow and a gradient fill."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(4), Inches(2)
    )
    shape.text_frame.text = "Revenue"
    shape.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    fill = shape.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = RGBColor.from_string("1F4E79")
    stops[0].position = 0.0
    stops[1].color.rgb = RGBColor.from_string("E8F1F8")
    stops[1].position = 1.0

    prs.save(str(path))
    return path


def _shape(path: Path):
    prs = Presentation(str(path))
    return list(prs.slides)[0].shapes[0]


def _sp_pr(shape):
    return shape._element.spPr


@pytest.fixture()
def deck(tmp_path: Path) -> Path:
    return _rounded_shadowed_gradient_deck(tmp_path / "deck.pptx")


# ---------------------------------------------------------------------------
# The three mechanical prohibitions
# ---------------------------------------------------------------------------


def test_shadow_is_turned_off_in_the_file(deck: Path):
    """A shape with no effectLst inherits the theme's shadow; an empty one does not."""
    assert _sp_pr(_shape(deck)).find(qn("a:effectLst")) is None  # inherits, i.e. shadowed

    result = enforce_theme_rules(deck, {"avoid": ["shadows"]})

    effect_lst = _sp_pr(_shape(deck)).find(qn("a:effectLst"))
    assert effect_lst is not None, "shadow left inherited"
    assert len(effect_lst) == 0, "an effect survived the pass"
    assert result["shadows_cleared"] == 1
    assert result["violations"]["shadows"] == 1


def test_explicit_drop_shadow_is_removed(tmp_path: Path):
    """The shadow the generating code wrote by hand goes too."""
    path = _rounded_shadowed_gradient_deck(tmp_path / "explicit.pptx")
    prs = Presentation(str(path))
    shape = list(prs.slides)[0].shapes[0]
    spPr = shape._element.spPr
    spPr.append(
        spPr.makeelement(qn("a:effectLst"), {})
    )
    effect_lst = spPr.find(qn("a:effectLst"))
    effect_lst.append(effect_lst.makeelement(qn("a:outerShdw"), {"blurRad": "40000"}))
    prs.save(str(path))

    result = enforce_theme_rules(path, {"avoid": ["shadows"]})

    effect_lst = _sp_pr(_shape(path)).find(qn("a:effectLst"))
    assert len(effect_lst) == 0
    assert result["shadows_cleared"] == 1


def test_rounded_corners_are_squared(deck: Path):
    assert _sp_pr(_shape(deck)).find(qn("a:prstGeom")).get("prst") == "roundRect"

    result = enforce_theme_rules(deck, {"avoid": ["rounded_corners"]})

    prst_geom = _sp_pr(_shape(deck)).find(qn("a:prstGeom"))
    assert prst_geom.get("prst") == "rect"
    av_lst = prst_geom.find(qn("a:avLst"))
    assert av_lst is None or len(av_lst) == 0, "adjustment outlived its geometry"
    assert result["corners_squared"] == 1
    assert result["violations"]["rounded_corners"] == 1


def test_gradient_is_flattened_to_its_dominant_stop(deck: Path):
    from pptx.enum.dml import MSO_FILL

    result = enforce_theme_rules(deck, {"avoid": ["gradients"]})

    fill = _shape(deck).fill
    assert fill.type == MSO_FILL.SOLID
    assert str(fill.fore_color.rgb) in {"1F4E79", "E8F1F8"}
    assert result["gradients_flattened"] == 1
    assert result["violations"]["gradients"] == 1


def test_dominant_stop_is_the_one_covering_most_of_the_shape(tmp_path: Path):
    """A hairline light sliver at the top does not make the shape light."""
    path = tmp_path / "lopsided.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(4), Inches(2)
    )
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor.from_string("FFFFFF")
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = RGBColor.from_string("1F4E79")
    fill.gradient_stops[1].position = 0.05
    prs.save(str(path))

    enforce_theme_rules(path, {"avoid": ["gradients"]})

    assert str(_shape(path).fill.fore_color.rgb) == "1F4E79"


def test_all_three_are_enforced_in_one_pass(deck: Path):
    result = enforce_theme_rules(deck, {"avoid": ALL_TOKENS})

    assert result["shadows_cleared"] == 1
    assert result["corners_squared"] == 1
    assert result["gradients_flattened"] == 1

    shape = _shape(deck)
    assert _sp_pr(shape).find(qn("a:prstGeom")).get("prst") == "rect"
    assert len(_sp_pr(shape).find(qn("a:effectLst"))) == 0


# ---------------------------------------------------------------------------
# The negative control — only what the theme actually forbids
# ---------------------------------------------------------------------------


def test_empty_avoid_list_changes_nothing(deck: Path):
    before = deck.read_bytes()

    result = enforce_theme_rules(deck, {"avoid": []})

    assert deck.read_bytes() == before, "file rewritten for a theme that forbids nothing"
    assert result == {
        "shadows_cleared": 0,
        "corners_squared": 0,
        "gradients_flattened": 0,
        "violations": {},
        "reported_only": [],
    }


def test_theme_without_an_avoid_key_changes_nothing(deck: Path):
    before = deck.read_bytes()

    result = enforce_theme_rules(deck, {"primary": "#1F4E79", "font": "Inter"})

    assert deck.read_bytes() == before
    assert result["corners_squared"] == 0
    assert result["violations"] == {}


def test_no_theme_at_all_changes_nothing(deck: Path):
    before = deck.read_bytes()

    assert enforce_theme_rules(deck)["shadows_cleared"] == 0
    assert deck.read_bytes() == before


def test_forbidding_shadows_does_not_square_corners(deck: Path):
    """The prohibitions are independent; enforcing one must not enforce another."""
    result = enforce_theme_rules(deck, {"avoid": ["shadows"]})

    assert _sp_pr(_shape(deck)).find(qn("a:prstGeom")).get("prst") == "roundRect"
    assert result["corners_squared"] == 0
    assert "rounded_corners" not in result["violations"]
    assert result["gradients_flattened"] == 0


def test_a_deck_already_compliant_is_not_rewritten(tmp_path: Path):
    """Nothing to do means no save, and honest zeroes."""
    path = tmp_path / "clean.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(1)
    )
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string("1F4E79")
    prs.save(str(path))
    before = path.read_bytes()

    result = enforce_theme_rules(path, {"avoid": ALL_TOKENS})

    assert path.read_bytes() == before
    assert result["shadows_cleared"] == 0
    assert result["corners_squared"] == 0
    assert result["gradients_flattened"] == 0


# ---------------------------------------------------------------------------
# Reported, not acted on
# ---------------------------------------------------------------------------


def test_multiple_accents_is_counted_and_never_repainted(tmp_path: Path):
    """Colour belongs to apply_theme_palette; this pass only says what it saw."""
    path = tmp_path / "accents.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, hexval in enumerate(("1F4E79", "C0504D", "9BBB59")):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1 + i), Inches(1), Inches(1), Inches(1)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(hexval)
        shape.shadow.inherit = False
    prs.save(str(path))

    result = enforce_theme_rules(path, {"avoid": ["multiple_accents"]})

    assert result["violations"]["multiple_accents"] == 2  # three accents, one allowed
    assert "multiple_accents" in result["reported_only"]
    colors = {
        str(s.fill.fore_color.rgb) for s in list(Presentation(str(path)).slides)[0].shapes
    }
    assert colors == {"1F4E79", "C0504D", "9BBB59"}, "a colour was repainted"


def test_boxes_are_counted_not_converted(deck: Path):
    result = enforce_theme_rules(deck, {"avoid": ["boxes"]})

    assert result["violations"]["boxes"] == 1
    assert result["reported_only"] == ["boxes"]
    assert _shape(deck).has_text_frame, "the card was taken apart"


def test_reported_only_lists_every_unenforceable_token(deck: Path):
    result = enforce_theme_rules(deck, {"avoid": ALL_TOKENS + ["legends"]})

    assert result["reported_only"] == ["boxes", "legends", "multiple_accents"]


def test_spelling_variants_are_the_same_prohibition(deck: Path):
    result = enforce_theme_rules(deck, {"avoid": ["Drop Shadows", "border-radius"]})

    assert result["shadows_cleared"] == 1
    assert result["corners_squared"] == 1
    assert set(result["violations"]) == {"shadows", "rounded_corners"}


# ---------------------------------------------------------------------------
# It never raises, and it never gates
# ---------------------------------------------------------------------------


def test_missing_file_returns_cleanly(tmp_path: Path):
    result = enforce_theme_rules(tmp_path / "nope.pptx", {"avoid": ALL_TOKENS})

    assert result["shadows_cleared"] == 0
    assert result["violations"] == {}
    assert result["reported_only"] == ["boxes", "multiple_accents"]


def test_corrupt_file_returns_cleanly(tmp_path: Path):
    path = tmp_path / "corrupt.pptx"
    path.write_bytes(b"this is not a pptx")

    result = enforce_theme_rules(path, {"avoid": ["shadows"]})

    assert result["shadows_cleared"] == 0
    assert result["violations"] == {}


def test_a_violating_deck_is_never_rejected(deck: Path):
    """Report, do not gate — a violation is a number to read, not a failure."""
    result = enforce_theme_rules(deck, {"avoid": ALL_TOKENS})

    assert isinstance(result, dict)
    assert deck.exists()
    assert Presentation(str(deck)).slides  # still a deck anyone can open


def test_junk_theme_shapes_do_not_raise(deck: Path):
    for junk in (None, {}, {"avoid": None}, {"avoid": "shadows"}, {"avoid": [None, 7]}, "nope", 42):
        result = enforce_theme_rules(deck, junk)
        assert isinstance(result["violations"], dict)


def test_a_group_hides_nothing(tmp_path: Path):
    """Shapes inside a group carry the corners a reader actually sees."""
    path = tmp_path / "grouped.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
    )
    b = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(1), Inches(2), Inches(1)
    )
    try:
        slide.shapes.add_group_shape([a, b])
    except AttributeError:  # pragma: no cover - older python-pptx has no grouping
        pytest.skip("python-pptx too old to build a group shape")
    prs.save(str(path))

    result = enforce_theme_rules(path, {"avoid": ["rounded_corners"]})

    assert result["corners_squared"] == 2
    prs = Presentation(str(path))
    presets = [
        el.get("prst")
        for el in list(prs.slides)[0].shapes._spTree.iter(qn("a:prstGeom"))
    ]
    assert "roundRect" not in presets
