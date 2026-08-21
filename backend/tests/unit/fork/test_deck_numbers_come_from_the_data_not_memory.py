"""FIX-D / FIX-E (2026-08-21) — decks stated figures their own data never held.

Measured on 0.0.543.18: the slide-building step ran real queries, got correct
answers, and then wrote python-pptx code containing numbers TYPED FROM MEMORY
of the conversation. The excel deck put Riverside at 1,995 against a real
1,459 (flipping a branch ranking); the Fabric deck invented 7 of 12 months of
a revenue chart, the generated code even carrying the comment "Sample seasonal
monthly values". The invented numbers were self-consistent, so nothing inside
the deck ever looked wrong.

FIX-D: `deck_grounding.check_deck_code` runs the figure_grounding engine over
the generated code's string literals and hardcoded numeric series, BEFORE the
code executes. FIX-E: `data_coverage_label` + `stamp_data_coverage` write the
data's newest period onto the cover deterministically, because on the same
measured decks "this year" was silently answered with 2023 rows.

★These are the replay tests the phase plan demanded: the excel fabrication
must be caught, and a faithful transcription of the same data must pass.
"""
import pathlib

import pytest


# The excel deck's real branch totals (thousands) — Riverside is 1,459, and the
# deck shipped 1,995. Column sum is 12,473; the deck shipped 12,827.
BRANCH_VIZ = {
    "id": "viz-1",
    "title": "Branch totals",
    "rows": [
        {"branch": "Downtown", "revenue": 3120},
        {"branch": "Airport Rd", "revenue": 2610},
        {"branch": "Junction Sq", "revenue": 1877},
        {"branch": "Hledan", "revenue": 1485},
        {"branch": "Riverside", "revenue": 1459},
        {"branch": "North Point", "revenue": 1922},
    ],
}


def _check(code, vizzes):
    from app.ai.decks.deck_grounding import check_deck_code
    return check_deck_code(code, vizzes)


class TestTheGateCatchesTheMeasuredFabrications:
    def test_riverside_1995_is_caught(self):
        """★The excel deck's exact defect: a figure in slide text that appears
        nowhere in the data."""
        code = (
            "def generate_slides(visualizations, report):\n"
            "    title = 'Riverside led the network at 1,995'\n"
        )
        result = _check(code, [BRANCH_VIZ])
        assert result["status"] == "checked"
        tokens = [v["token"] for v in result["violations"]]
        assert "1,995" in tokens

    def test_an_invented_network_total_is_caught(self):
        code = (
            "def generate_slides(visualizations, report):\n"
            "    tile = 'Network revenue: 12,827'\n"
        )
        result = _check(code, [BRANCH_VIZ])
        assert [v["token"] for v in result["violations"]] == ["12,827"]

    def test_a_hardcoded_chart_series_is_caught(self):
        """★The Fabric deck's defect: invented months fed straight to
        add_series. None of these values is derivable from the rows."""
        code = (
            "def generate_slides(visualizations, report):\n"
            "    values = [4180.5, 3975.2, 5190.1, 6210.7]\n"
        )
        result = _check(code, [BRANCH_VIZ])
        assert result["violations"], "invented series must be flagged"
        assert all(v["where"] == "series" for v in result["violations"])


class TestFaithfulCodePasses:
    def test_the_true_figures_pass(self):
        """★Positive control — the phase plan's own bar: correct transcription
        (cells, the real column total, the real mean) must NOT be flagged."""
        code = (
            "def generate_slides(visualizations, report):\n"
            "    a = 'Riverside closed at 1,459'\n"
            "    b = 'Network revenue: 12,473'\n"
            "    c = 'Average branch: 2,078.8'\n"
        )
        result = _check(code, [BRANCH_VIZ])
        assert result["status"] == "checked"
        assert result["violations"] == []

    def test_code_that_computes_from_rows_has_nothing_to_flag(self):
        """The desired behaviour produces no literals at all."""
        code = (
            "def generate_slides(visualizations, report):\n"
            "    rows = visualizations[0]['rows']\n"
            "    total = sum(r['revenue'] for r in rows)\n"
            "    title = f'Network revenue: {total:,}'\n"
            "    series = [float(r['revenue']) for r in rows]\n"
        )
        result = _check(code, [BRANCH_VIZ])
        assert result["violations"] == []

    def test_a_series_in_display_scale_passes(self):
        """Chart code writes 202.3 meaning 202.3M — the gate must try the
        scales charts are drawn in before crying fabrication."""
        viz = {"rows": [
            {"month": "Mar", "revenue": 202300000},
            {"month": "Apr", "revenue": 187300000},
            {"month": "May", "revenue": 154800000},
        ]}
        code = (
            "def generate_slides(visualizations, report):\n"
            "    values = [202.3, 187.3, 154.8]\n"
        )
        assert _check(code, [viz])["violations"] == []

    def test_an_identifier_stored_as_a_string_cell_passes(self):
        """★Measured false positive on the real Power BI deck: article numbers
        are stored as TEXT cells, and naming one verbatim was flagged as
        invented. A value copied from a cell is grounded whatever the column's
        storage type."""
        viz = {"rows": [
            {"article": "1000000345139", "revenue": 6700000},
            {"article": "1000000381190", "revenue": 5100000},
            {"article": "1000000370261", "revenue": 4300000},
        ]}
        code = (
            "def generate_slides(visualizations, report):\n"
            "    t = 'Article 1000000345139 led at 6.7M'\n"
        )
        assert _check(code, [viz])["violations"] == []

    def test_a_rounded_row_count_claim_passes(self):
        """★Measured false positive: "300,000+ line records" cites the query's
        row_count, rounded to its own trailing zeros. Honest rounding is held
        to the precision the notation shows, not to the unit."""
        viz = dict(BRANCH_VIZ, row_count=301245)
        code = (
            "def generate_slides(visualizations, report):\n"
            "    t = '300,000+ line records analysed'\n"
        )
        assert _check(code, [viz])["violations"] == []

    def test_rounding_slack_does_not_resurrect_the_fabrications(self):
        """The relaxed rule must not wave the measured defects through: none
        of them ends in zeros at its stated precision."""
        code = (
            "def generate_slides(visualizations, report):\n"
            "    a = 'Riverside 1,995'\n"
            "    b = 'Network 12,827'\n"
        )
        tokens = [v["token"] for v in _check(code, [BRANCH_VIZ])["violations"]]
        assert "1,995" in tokens and "12,827" in tokens

    def test_layout_numbers_and_percentages_never_trip_it(self):
        """Inches/Pt/RGB args are scalar call arguments (never collected),
        small integers are structural, and percentages are computed values the
        engine passes by design."""
        code = (
            "def generate_slides(visualizations, report):\n"
            "    x = Inches(1); y = Pt(24); c = RGBColor(15, 23, 42)\n"
            "    t = 'Top 10 branches drove 78% of growth'\n"
            "    shares = [41.3, 58.7, 22.1]\n"
        )
        assert _check(code, [BRANCH_VIZ])["violations"] == []


class TestTheGateFailsOpenNotClosed:
    def test_no_numeric_data_means_skipped(self):
        """A narrative deck must not be blocked by its own emptiness."""
        result = _check("def generate_slides(v, r):\n    t = 'Launch 5,000'\n", [])
        assert result["status"] == "skipped"

    def test_unparseable_code_means_unavailable_not_a_crash(self):
        result = _check("def generate_slides(:", [BRANCH_VIZ])
        assert result["status"] == "unavailable"


class TestDataCoverage:
    def test_the_newest_month_wins_across_visualizations(self):
        from app.ai.decks.deck_grounding import data_coverage_label
        vizzes = [
            {"rows": [{"period": "2023-04", "v": 1}, {"period": "2025-01", "v": 2}]},
            {"rows": [{"period": "2025-12-01", "v": 3}, {"period": "2024-06-15", "v": 4}]},
        ]
        assert data_coverage_label(vizzes) == "Dec 2025"

    def test_a_bare_year_column_gives_a_year(self):
        from app.ai.decks.deck_grounding import data_coverage_label
        assert data_coverage_label(
            [{"rows": [{"year": 2023, "v": 1}, {"year": 2025, "v": 2}]}]
        ) == "2025"

    def test_no_period_column_means_no_label(self):
        from app.ai.decks.deck_grounding import data_coverage_label
        assert data_coverage_label(
            [{"rows": [{"branch": "Downtown", "revenue": 3120}]}]
        ) is None

    def test_the_stamp_lands_on_the_cover_and_reads_against_a_dark_ground(self, tmp_path):
        """★Build a real deck with a dark cover (the FIX-A case), stamp it,
        reopen it, and find the stamp — in the light colour."""
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Emu

        from app.ai.decks.deck_grounding import stamp_data_coverage

        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0B, 0x0F, 0x19)
        path = tmp_path / "deck.pptx"
        prs.save(str(path))

        result = stamp_data_coverage(str(path), "Dec 2025")
        assert result["status"] == "stamped"
        assert result["dark_cover"] is True

        reopened = Presentation(str(path))
        texts = [
            s.text_frame.text
            for s in reopened.slides[0].shapes
            if s.has_text_frame
        ]
        assert any("Data through Dec 2025" in t for t in texts)


class TestTheProfileDisclosesAnEmptyResult:
    def test_zero_rows_carries_a_warning_the_model_reads(self):
        """FIX-E: two real "2026" queries returned zero rows, were recorded as
        successful, and the deck silently answered with 2023 data."""
        from app.ai.tools.implementations.create_artifact import CreateArtifactTool

        profile = CreateArtifactTool._build_viz_profile(
            None,
            {"id": "v", "title": "2026 revenue", "rows": [], "row_count": 0, "columns": []},
            True,
        )
        assert "ZERO rows" in profile.get("empty_result_warning", "")

    def test_a_populated_result_carries_none(self):
        from app.ai.tools.implementations.create_artifact import CreateArtifactTool

        profile = CreateArtifactTool._build_viz_profile(
            None,
            {"id": "v", "rows": BRANCH_VIZ["rows"], "row_count": 6,
             "columns": [{"field": "branch"}, {"field": "revenue"}]},
            True,
        )
        assert "empty_result_warning" not in profile
