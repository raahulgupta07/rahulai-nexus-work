"""The bottom of a slide has ONE author, and the prompt says which one.

WHY THIS FILE EXISTS
--------------------
Two writers put furniture at the bottom of a content slide and neither knew
about the other:

  * `_build_slides_prompt` told the model to draw a source footer and a
    zero-padded page number on every content slide;
  * `decks/motifs.py` reopens the SAVED deck afterwards and paints the THEME's
    tracker, footer rule, footer and page number — default ON.

Measured on a real deck: two progress trackers stacked, and two source footers
overlapping. The fix is in two halves and this file covers both, because either
half alone still ships an overlapping deck:

  FIX 1  the prompt is written from what the painter will ACTUALLY do — the
         resolved theme's own layout entry — so the model is told to leave the
         footer alone only where a footer is genuinely coming.
  FIX 2  the painter strips model-drawn TEXT out of the footer band before it
         paints, because the prompt is a request and a request is not a
         guarantee. The strip is the half that does not depend on the model
         listening.

HOW THIS IS MEASURED, AND WHY
-----------------------------
For fix 1: every assertion runs against the STRING `_build_slides_prompt`
returns, never against the text of `create_artifact.py`. A rule sitting in a
branch that did not run is not a rule, and — the reason this repo insists on it
— ★a source-scanning test here has matched its own explanatory comment at least
four times. There is nothing to strip when you read a return value.

For fix 2: the painter is driven against a REAL `python-pptx` Presentation
written to a real file. `_strip_model_footer_text` is XML surgery on the shape
tree; a mock cannot tell you whether the element left the tree, and the whole
defect is about what is on the slide afterwards.

★WHICH THEME PAINTS A FOOTER IS PROBED AT RUNTIME, never hardcoded. There are
81 registered design systems and ~24 legitimately paint nothing (`ted-style`
forbids footers and page numbers outright). Naming one would make this file
fail the day a layout is re-read from its prompt — a change to product data,
not to product behaviour. If NO theme is found in either class the test skips
with a reason, and that skip is itself news: it would mean the layouts registry
had emptied out.
"""
import pytest

from app.ai.tools.implementations import create_artifact as ca
from app.ai.tools.implementations.create_artifact import CreateArtifactTool

# The never-draw instruction, verbatim. This string is the contract — it is what
# the model reads — so a reword is a product decision and must break this file.
NEVER_DRAW = "Never draw your own"

# The pre-fix instructions, verbatim. They are still correct for a theme that
# paints nothing, so their SURVIVAL is asserted just as hard as their absence.
DRAW_THE_PAGE_NUMBER = "PAGE NUMBER right:"
DRAW_THE_FOOTER = "FOOTER left:"

# `_build_slides_prompt` never touches `self` — pinned by
# `test_a_deck_layout_earns_its_place.py`, whose calling pattern this reuses —
# so it is called as a plain function. No tool construction, no session, no
# database, which is what `tests/unit/fork` requires.
_build = CreateArtifactTool.__dict__["_build_slides_prompt"]


def build_prompt(**overrides) -> str:
    """The prompt a real deck request produces."""
    kwargs = dict(
        user_prompt="quarterly revenue review for the board",
        title="Q3 revenue",
        viz_profiles=[],
        instructions_context="",
        report_title=None,
        allow_llm_see_data=False,
    )
    kwargs.update(overrides)
    return _build(object(), **kwargs)


def _themes_by_what_they_paint():
    """`(paints_footer, paints_nothing)` — two real theme objects, or Nones.

    Probed through the product's OWN pair of readers (`_load_pptx_themes` for
    the registry, `_theme_layout_for` for the layout), so this cannot disagree
    with what `_build_slides_prompt` will see when it asks the same questions.
    """
    themes = ca._load_pptx_themes()
    if themes is None:
        return None, None
    paints_footer = paints_nothing = None
    for theme in themes.all_themes():
        layout = ca._theme_layout_for(getattr(theme, "id", None))
        if not isinstance(layout, dict):
            continue
        footer = layout.get("footer")
        if paints_footer is None and isinstance(footer, dict) and footer:
            paints_footer = theme
        if (
            paints_nothing is None
            and not isinstance(footer, dict)
            and not layout.get("footer_rule")
            and not layout.get("tracker")
        ):
            paints_nothing = theme
        if paints_footer is not None and paints_nothing is not None:
            break
    return paints_footer, paints_nothing


@pytest.fixture(scope="module")
def theme_that_paints_a_footer():
    theme, _ = _themes_by_what_they_paint()
    if theme is None:
        pytest.skip(
            "no registered theme has a layout entry defining a footer. Nothing "
            "in this file is broken — but that is news: the whole point of "
            "FIX 1 is that some themes paint one. Read decks/theme_layouts.py "
            "before deleting this test."
        )
    return theme


@pytest.fixture(scope="module")
def theme_that_paints_nothing():
    _, theme = _themes_by_what_they_paint()
    if theme is None:
        pytest.skip(
            "every registered theme paints footer, footer rule or tracker, so "
            "the pre-fix instructions can never be reached. ~24 layouts are "
            "deliberately empty; if that is no longer true, the else-branch in "
            "_build_slides_prompt is dead code and should be removed on purpose."
        )
    return theme


@pytest.fixture(autouse=True)
def painter_on(monkeypatch):
    """The furniture pass is ON, which is its shipped default.

    Pinned rather than assumed: `_deck_theme_furniture_enabled` reads a runtime
    setting, and a test whose outcome depends on an org's flag is a test that
    reports on the flag. The switch itself gets its own case below.
    """
    monkeypatch.setattr(ca, "_deck_theme_furniture_enabled", lambda: True)


# ═══════════════════════════════════════════════════════════════════════════
# FIX 1 — the prompt says who draws the bottom of the slide
# ═══════════════════════════════════════════════════════════════════════════


def test_a_theme_that_paints_a_footer_tells_the_model_not_to(
    theme_that_paints_a_footer,
):
    """The model is told the furniture arrives after its code runs."""
    prompt = build_prompt(resolved_theme=theme_that_paints_a_footer)

    assert NEVER_DRAW in prompt


def test_a_theme_that_paints_a_footer_withdraws_the_draw_it_instruction(
    theme_that_paints_a_footer,
):
    """And the OLD instruction is gone, which is the half that matters.

    Telling the model both things is worse than telling it neither: it drew a
    page number because it was asked to, and the theme drew one too. An
    absence assertion, so it is paired with the presence assertion above and
    with the survival case below — a prompt that had lost BOTH instructions
    would satisfy this alone.
    """
    prompt = build_prompt(resolved_theme=theme_that_paints_a_footer)

    assert DRAW_THE_PAGE_NUMBER not in prompt
    assert NEVER_DRAW in prompt, (
        "the draw-it instruction is gone and nothing replaced it — the deck "
        "now gets no footer from either writer"
    )


def test_a_theme_that_paints_nothing_still_gets_the_old_instructions(
    theme_that_paints_nothing,
):
    """★THE POSITIVE CONTROL for the two absence assertions above.

    ~24 layouts paint no footer, no footer rule and no tracker. Nothing else
    will draw them, so the model still must — and the pre-fix wording is
    exactly right for that case. A fix that simply deleted the instructions
    passes every absence assertion in this file and fails here.
    """
    prompt = build_prompt(resolved_theme=theme_that_paints_nothing)

    assert DRAW_THE_PAGE_NUMBER in prompt
    assert DRAW_THE_FOOTER in prompt
    assert NEVER_DRAW not in prompt


def test_the_two_instructions_are_never_in_the_same_prompt(
    theme_that_paints_a_footer, theme_that_paints_nothing
):
    """Whatever the theme, the model is told once and told one thing.

    This is the property the defect violated, stated directly rather than
    inferred from the two cases above.
    """
    for theme in (theme_that_paints_a_footer, theme_that_paints_nothing):
        prompt = build_prompt(resolved_theme=theme)
        assert (NEVER_DRAW in prompt) != (DRAW_THE_PAGE_NUMBER in prompt), (
            f"theme {getattr(theme, 'id', theme)!r} produced a prompt carrying "
            "both instructions or neither"
        )


def test_switching_the_furniture_pass_off_hands_the_footer_back_to_the_model(
    monkeypatch, theme_that_paints_a_footer
):
    """`hybrid_deck_theme_furniture=false` is a real deployment.

    The painter never runs there, so a prompt that still said "never draw your
    own footer" would produce decks with no footer at all — the same class of
    defect as the overlap, arrived at from the other side. The instruction is
    written from what the painter WILL do, and with the pass off it will do
    nothing.
    """
    monkeypatch.setattr(ca, "_deck_theme_furniture_enabled", lambda: False)

    prompt = build_prompt(resolved_theme=theme_that_paints_a_footer)

    assert NEVER_DRAW not in prompt
    assert DRAW_THE_PAGE_NUMBER in prompt


def test_a_theme_that_cannot_be_read_still_gets_a_footer_from_somebody(
    theme_that_paints_a_footer,
):
    """A bogus theme costs the deck its design system, never its footer.

    `_theme_layout_for` answers `{}` for anything it cannot resolve, so the
    painter will paint nothing — and the prompt must therefore ask the model
    for the footer. Handing it an object that is not a Theme drives that path
    for real rather than asserting it from the source.
    """
    prompt = build_prompt(resolved_theme=object())

    assert DRAW_THE_PAGE_NUMBER in prompt
    assert NEVER_DRAW not in prompt


# ═══════════════════════════════════════════════════════════════════════════
# FIX 2 — the painter clears the band before it paints
# ═══════════════════════════════════════════════════════════════════════════

_BAND_TEXT = "MODEL DREW THIS FOOTER"
_BODY_TEXT = "MODEL DREW THIS BODY LINE"
_SILENT_SHAPE = "an empty rectangle in the band"

# A layout that paints a footer, and one that paints none. Written here rather
# than taken from the registry so the two differ in exactly one thing.
_LAYOUT_WITH_FOOTER = {
    "ground": "flat",
    "ground_color": "#FFFFFF",
    "ornament": (),
    "footer": {"left": "SUBJECT", "right": "{page} / {pages}", "color": "#777777", "size_px": 10},
}
_LAYOUT_WITHOUT_FOOTER = {"ground": "flat", "ground_color": "#FFFFFF", "ornament": ()}

_THEME = {"palette": {"muted_text": "#777777"}, "fonts": ["Cambria", "Calibri"]}


def _deck_with_a_model_drawn_footer(path):
    """A one-slide deck carrying the three shapes the strip must tell apart.

    Kept to one slide and four shapes: this runs in the inner-loop suite, and
    the rule under test is per-shape, so a bigger deck measures nothing more.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    height = prs.slide_height
    width = prs.slide_width

    # In the band (the painter's band starts at 648 of a 720px canvas).
    box = slide.shapes.add_textbox(
        Emu(int(width * 0.05)), Emu(int(height * 0.94)), Emu(int(width * 0.5)), Emu(int(height * 0.04))
    )
    box.text_frame.text = _BAND_TEXT

    # Body height — a title, a bullet, anything the deck is actually about.
    body = slide.shapes.add_textbox(
        Emu(int(width * 0.05)), Emu(int(height * 0.30)), Emu(int(width * 0.5)), Emu(int(height * 0.08))
    )
    body.text_frame.text = _BODY_TEXT

    # In the band and NOT text: a rule, a swatch, a chart's baseline. Deleting
    # one of these to fix a cosmetic overlap costs the slide its content.
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(width * 0.60)),
        Emu(int(height * 0.95)),
        Emu(int(width * 0.20)),
        Emu(int(height * 0.01)),
    )
    shape.name = _SILENT_SHAPE

    prs.save(str(path))
    return path


def _texts_and_names(path):
    """`(set of non-empty texts, set of shape names)` on the saved deck's slide 1."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slide = prs.slides[0]
    texts, names = set(), set()
    for shape in slide.shapes:
        names.add(shape.name)
        try:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.add(shape.text_frame.text.strip())
        except Exception:
            continue
    return texts, names


@pytest.fixture
def painted_with_a_footer(tmp_path):
    from app.ai.decks.motifs import paint_theme_furniture

    deck = _deck_with_a_model_drawn_footer(tmp_path / "deck.pptx")
    result = paint_theme_furniture(deck, _THEME, _LAYOUT_WITH_FOOTER)
    return deck, result


def test_the_models_own_footer_is_removed_before_the_theme_paints_one(
    painted_with_a_footer,
):
    """The overlap, measured on the file: only one footer survives."""
    deck, _ = painted_with_a_footer
    texts, _names = _texts_and_names(deck)

    assert _BAND_TEXT not in texts


def test_the_strip_says_it_happened(painted_with_a_footer):
    """A row that vanishes silently is the same defect in the other direction.

    `painted` is what reaches the step record, so the audit marker is the only
    way anyone later can tell a stripped deck from one the model laid out
    correctly. The prefix is asserted, not the exact suffix — the marker is a
    diagnostic, and pinning its full wording buys nothing.
    """
    _deck, result = painted_with_a_footer

    assert any(str(name).startswith("stripped:") for name in result["painted"]), (
        f"nothing in {result['painted']!r} records the removal"
    )


def test_text_at_body_height_is_left_exactly_where_it_was(painted_with_a_footer):
    """★THE POSITIVE CONTROL. A strip that ate the slide passes the case above.

    This is the assertion that keeps FIX 2 a narrowing rather than a purge: the
    band is a band, not "the lower half", and the deck's own content has to
    come through untouched.
    """
    deck, _ = painted_with_a_footer
    texts, _names = _texts_and_names(deck)

    assert _BODY_TEXT in texts


def test_a_shape_with_no_text_in_the_band_survives(painted_with_a_footer):
    """A picture, a rule or a chart dipping into the bottom is not a footer.

    Deleting one costs the slide its content to fix a cosmetic overlap, which
    is a worse trade than the overlap.
    """
    deck, _ = painted_with_a_footer
    _texts, names = _texts_and_names(deck)

    assert _SILENT_SHAPE in names


def test_a_theme_that_paints_no_footer_leaves_the_models_footer_alone(tmp_path):
    """★THE OTHER POSITIVE CONTROL, and the one that pairs with FIX 1.

    With no footer and no footer rule coming, the model's own footer is the
    slide's ONLY footer. Removing it would leave the deck with none — and this
    is exactly the theme class whose prompt still ASKS the model to draw one.
    """
    from app.ai.decks.motifs import paint_theme_furniture

    deck = _deck_with_a_model_drawn_footer(tmp_path / "deck.pptx")
    result = paint_theme_furniture(deck, _THEME, _LAYOUT_WITHOUT_FOOTER)

    texts, names = _texts_and_names(deck)
    assert _BAND_TEXT in texts
    assert _BODY_TEXT in texts
    assert _SILENT_SHAPE in names
    assert not any(str(name).startswith("stripped:") for name in result["painted"])


def test_the_deck_is_still_a_deck_afterwards(painted_with_a_footer):
    """The painter's whole contract: never raise, always leave a valid file.

    Cheap, and it is the assertion that fails first if the XML surgery ever
    removes an element the file format requires.
    """
    from pptx import Presentation

    deck, result = painted_with_a_footer

    assert result["slides"] == 1
    assert len(Presentation(str(deck)).slides) == 1
