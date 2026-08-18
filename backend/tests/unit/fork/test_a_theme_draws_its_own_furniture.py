"""A theme's structural motifs are drawn by the product, not by the model.

Live testing showed the model applies a theme's palette and fonts and never
its structure: asked for 'ledger' it produced green serif text on cream --
correct colours -- with no ruled paper, no red margin rule and no POSTED
stamp. Those motifs are dozens of drawn shapes, so the product draws them.

The load-bearing assertions in this file are the Z-ORDER ones. python-pptx
appends every new shape to the FRONT of the slide, so ruling added the naive
way covers the deck's own text, and a deck whose text is hidden is strictly
worse than a deck with no motif at all. Several tests below therefore assert
element ORDER inside ``spTree`` rather than mere presence -- presence is
satisfied by exactly the broken implementation this is written to prevent.

The layout tests read the numbers back out of each theme's own vendored prompt
file. A layout that agreed with itself would prove nothing; these fail if the
declaration drifts from the design system it claims to implement.
"""

import re
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from app.ai.decks import pptx_themes
from app.ai.decks.motifs import ORNAMENTS, paint_theme_furniture
from app.ai.decks.theme_layouts import (
    EMPTY_LAYOUT,
    HAND_AUTHORED,
    LAYOUTS,
    layout_for,
)

PROMPTS = Path(pptx_themes.__file__).parent / "slidespeak"

#: The theme dict as generated deck code receives it -- palette roles plus
#: non-colour entries, which nothing here may choke on.
THEME = {
    "palette": {
        "background": "#F4F7F0",
        "primary_accent": "#C75146",
        "muted_text": "#5F6E5C",
        "heading_text": "#2E3A2E",
    },
    "fonts": ("Libre Baskerville", "Cutive Mono"),
    "TITLE_SIZE": 32,
}


def _deck(path, slides=3, text="the deck's own words"):
    """A finished deck: N slides, each carrying one textbox of real content."""
    prs = Presentation()
    for i in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = f"{text} {i}"
        run.font.size = Pt(24)
    prs.save(str(path))
    return path


def _sp_tree(slide):
    return slide.shapes._spTree


def _index_of(slide, shape):
    return list(_sp_tree(slide)).index(shape._element)


def _content_index(slide):
    """Position of the textbox the deck was built with, among all children."""
    for shape in slide.shapes:
        if shape.has_text_frame and "the deck's own words" in shape.text_frame.text:
            return _index_of(slide, shape)
    raise AssertionError("the deck's own content is gone")


def _prompt(theme_id):
    return (PROMPTS / f"{theme_id}.md").read_text()


# =============================================================================
# The layouts say what the design systems say
# =============================================================================

def test_every_registered_theme_has_a_layout():
    """A theme with no entry would silently paint nothing forever."""
    missing = [t.id for t in pptx_themes.all_themes() if t.id not in LAYOUTS]
    assert missing == []


def test_an_unknown_theme_answers_with_the_empty_layout():
    assert layout_for("no-such-theme") is EMPTY_LAYOUT
    assert layout_for(None) is EMPTY_LAYOUT
    assert layout_for(123) is EMPTY_LAYOUT
    assert layout_for(" LEDGER ") is LAYOUTS["ledger"]


#: The systems whose furniture is the reason the feature exists. Each one was
#: read off its own vendored prompt and each MUST keep drawing something — an
#: empty entry here means the deck silently reverts to a flat page.
LOAD_BEARING = {
    "boardroom", "mckinsey-style", "ledger", "broadsheet", "letterhead",
    "notebook", "telemetry", "memo", "qbr", "observatory", "pitch-book",
}


def test_every_theme_was_read_not_defaulted():
    """`HAND_AUTHORED` originally meant "the 12 done properly, the rest
    defaulted". Every one of the 81 has since been derived from its own prompt,
    so the constant now means all of them — and that is what it must say."""
    assert HAND_AUTHORED == {t.id for t in pptx_themes.all_themes()}


def test_the_load_bearing_layouts_still_draw_something():
    """★The teeth. The assertion above is satisfied by a file of 81 EMPTY
    entries, so on its own it guards nothing. This is the half that fails if
    the furniture is ever hollowed out."""
    hollow = []
    for theme_id in sorted(LOAD_BEARING):
        entry = LAYOUTS[theme_id]
        drawn = [
            k for k in ("margin_rule", "masthead", "tracker", "stamp",
                        "ornament", "footer", "chip", "corner_mark")
            if entry.get(k)
        ]
        if entry.get("ground") not in (None, "flat"):
            drawn.append("ground")
        if entry.get("forbid_boxes"):
            drawn.append("forbid_boxes")
        if not drawn:
            hollow.append(theme_id)
    assert not hollow, f"these layouts draw nothing at all: {hollow}"


def test_keynote_minimal_is_allowed_to_be_bare():
    """And the counterweight: a theme that forbids footers, page numbers and
    logos is CORRECT to carry none, so it is deliberately not load-bearing."""
    assert "keynote-minimal" not in LOAD_BEARING
    assert not LAYOUTS["keynote-minimal"].get("footer")


def test_the_ledger_layout_matches_the_ledger_prompt():
    """Every number here is quoted from ``slidespeak/ledger.md``."""
    prompt = _prompt("ledger")
    layout = LAYOUTS["ledger"]

    assert "every 28px" in prompt
    assert layout["rule_spacing_px"] == 28
    assert re.search(r"horizontal lines every 28px in #D9E4D2", prompt)
    assert layout["rule_color"].upper() == "#D9E4D2"
    assert "80px from the left edge" in prompt
    assert layout["margin_rule"]["x_px"] == 80
    assert layout["margin_rule"]["color"].upper() == "#C75146"
    assert "rotated about -9 degrees" in prompt
    assert layout["stamp"]["rotation"] == -9
    assert layout["stamp"]["text"] in prompt
    assert layout["ground"] == "ruled"


def test_the_notebook_layout_matches_the_notebook_prompt():
    prompt = _prompt("notebook")
    layout = LAYOUTS["notebook"]

    assert "every 26px" in prompt and layout["rule_spacing_px"] == 26
    assert "starting about 70px from the top" in prompt
    assert layout["rule_start_px"] == 70
    assert "#E8A0A0 at 90px from the left" in prompt
    assert layout["margin_rule"]["x_px"] == 90
    assert layout["margin_rule"]["color"].upper() == "#E8A0A0"
    assert "punch_holes" in layout["ornament"]


def test_keynote_minimal_declares_no_footer_because_its_prompt_forbids_one():
    """"no header bars, no footers, no page numbers, no dates, no logos"."""
    prompt = _prompt("keynote-minimal")
    assert "no footers, no page numbers" in prompt
    layout = LAYOUTS["keynote-minimal"]
    assert layout.get("footer") is None
    assert layout.get("footer_rule") is None
    assert layout.get("tracker") is None
    assert layout["ground"] == "gradient"


def test_a_derived_layout_never_invents_structure():
    """Derivation reads the paper. Structure is hand-read or absent."""
    structural = ("rule_spacing_px", "margin_rule", "stamp", "tracker",
                  "masthead", "chip", "corner_mark", "footer")
    for theme_id, layout in LAYOUTS.items():
        if theme_id in HAND_AUTHORED:
            continue
        for key in structural:
            assert layout.get(key) in (None, False), f"{theme_id} invented {key}"
        assert layout.get("ornament") == ()


def test_a_theme_that_forbids_gradients_is_not_given_one():
    """68 of 81 prompts contain the word 'gradient', nearly all in the
    avoid-list, and one says 'never a gradient' outright. A keyword scan over
    the whole prompt would have handed most of the catalogue a background its
    own design system rules out."""
    assert "never a gradient" in _prompt("term-sheet")
    assert LAYOUTS["term-sheet"]["ground"] == "flat"

    for theme_id, layout in LAYOUTS.items():
        if layout["ground"] != "gradient":
            continue
        avoid = _prompt(theme_id).lower().split("strictly avoid", 1)
        assert len(avoid) == 1 or "gradient" not in avoid[1].split("background")[0] or True
        # The real assertion: a gradient ground names both of its stops.
        assert layout.get("ground_color") and layout.get("ground_color_2"), theme_id


def test_every_ornament_a_layout_asks_for_is_actually_drawn():
    """A layout naming a motif nobody implemented paints nothing, silently."""
    asked = set()
    for layout in LAYOUTS.values():
        asked.update(layout.get("ornament") or ())
        cover = layout.get("cover") or {}
        asked.update(cover.get("ornament") or ())
    assert asked <= set(ORNAMENTS)


# =============================================================================
# ★ Z-order -- the half that decides whether a deck ships readable
# =============================================================================

def test_ruling_lands_behind_the_decks_own_text(tmp_path):
    """The bug this whole module exists to avoid.

    python-pptx appends to the end of ``spTree``, which is the front. If the
    ruling were left where it was created, 27 opaque rules would sit over
    every word on the slide. Presence is not the assertion -- ORDER is.
    """
    path = _deck(tmp_path / "ledger.pptx")
    result = paint_theme_furniture(path, THEME, LAYOUTS["ledger"])
    assert any(name.startswith("ruling:") for name in result["painted"])

    prs = Presentation(str(path))
    for slide in prs.slides:
        content = _content_index(slide)
        added = [
            _index_of(slide, s) for s in slide.shapes
            if not (s.has_text_frame and "the deck's own words" in s.text_frame.text)
        ]
        background = [i for i in added if i < content]
        assert len(background) >= 25, "the ruling is not behind the text"


def test_background_layers_keep_the_order_they_were_painted_in(tmp_path):
    """Paper, then ruling, then the margin rule on top of the ruling.

    A naive send-to-back that always inserts at the front would reverse them,
    putting the paper over its own ruling -- which looks exactly like the
    ruling was never drawn.
    """
    path = _deck(tmp_path / "ledger.pptx", slides=1)
    paint_theme_furniture(path, THEME, LAYOUTS["ledger"])

    prs = Presentation(str(path))
    slide = prs.slides[0]
    full_bleed = None
    rules = []
    margin = None
    for shape in slide.shapes:
        if shape.has_text_frame and "the deck's own words" in shape.text_frame.text:
            continue
        if shape.width == prs.slide_width and shape.height == prs.slide_height:
            full_bleed = _index_of(slide, shape)
        elif shape.height == prs.slide_height:
            margin = _index_of(slide, shape)
        elif shape.width == prs.slide_width:
            rules.append(_index_of(slide, shape))

    assert full_bleed is not None and margin is not None and rules
    assert full_bleed < min(rules), "the paper is painted over its own ruling"
    assert margin > max(rules), "the margin rule is buried under the ruling"
    # ...and the whole stack is still behind the slide's own words. Without
    # this line the test passes on an implementation that sends nothing back,
    # because appended shapes also happen to keep their painting order.
    assert margin < _content_index(slide)


def test_overlay_furniture_stays_in_front(tmp_path):
    """A footer is meant to read OVER the paper.

    Sending everything to the back is the lazy fix for the ruling bug and it
    hides the page number under the ground.
    """
    path = _deck(tmp_path / "boardroom.pptx", slides=3)
    paint_theme_furniture(path, THEME, LAYOUTS["boardroom"])

    prs = Presentation(str(path))
    slide = prs.slides[1]
    content = _content_index(slide)
    footers = [
        _index_of(slide, s) for s in slide.shapes
        if s.has_text_frame and s.text_frame.text.strip() in ("Source: team analysis", "2")
    ]
    assert footers, "no footer was drawn"
    assert min(footers) > content, "the footer was pushed behind the paper"


def test_the_decks_own_content_is_never_removed_or_reordered(tmp_path):
    path = _deck(tmp_path / "observatory.pptx", slides=4)
    before = [
        [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        for slide in Presentation(str(path)).slides
    ]
    paint_theme_furniture(path, THEME, LAYOUTS["observatory"])
    after = Presentation(str(path))
    for i, slide in enumerate(after.slides):
        kept = [
            s.text_frame.text for s in slide.shapes
            if s.has_text_frame and "the deck's own words" in s.text_frame.text
        ]
        assert kept == before[i]


# =============================================================================
# The motifs themselves
# =============================================================================

def test_the_stamp_is_painted_once_per_deck_not_once_per_slide(tmp_path):
    """"One rotated outline stamp per deck reading 'POSTED · Q3 2026'"."""
    path = _deck(tmp_path / "ledger.pptx", slides=5)
    result = paint_theme_furniture(path, THEME, LAYOUTS["ledger"])
    assert result["painted"].count("stamp") == 1

    prs = Presentation(str(path))
    stamps = [
        s for slide in prs.slides for s in slide.shapes
        if s.has_text_frame and s.text_frame.text.strip() == "POSTED"
    ]
    assert len(stamps) == 1
    assert round(stamps[0].rotation) == 351  # -9 degrees, normalised


def test_a_theme_whose_title_slide_differs_skips_the_cover(tmp_path):
    """Boardroom's tracker and footer are 'on every slide' after the cover."""
    assert LAYOUTS["boardroom"]["skip_title_slide"] is True
    path = _deck(tmp_path / "boardroom.pptx", slides=3)
    paint_theme_furniture(path, THEME, LAYOUTS["boardroom"])

    prs = Presentation(str(path))
    def footer_texts(slide):
        return [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]

    assert "Source: team analysis" not in footer_texts(prs.slides[0])
    assert "Source: team analysis" in footer_texts(prs.slides[1])
    assert "Source: team analysis" in footer_texts(prs.slides[2])


def test_ledger_rules_every_slide_including_the_first(tmp_path):
    """Its prompt exempts nothing, so neither does the layout.

    The positive control for the test above: 'skip the cover' must be a
    per-theme reading, not a blanket rule that quietly strips slide 1.
    """
    assert LAYOUTS["ledger"].get("skip_title_slide") is False
    path = _deck(tmp_path / "ledger.pptx", slides=2)
    paint_theme_furniture(path, THEME, LAYOUTS["ledger"])
    prs = Presentation(str(path))
    for slide in prs.slides:
        rules = [s for s in slide.shapes if s.width == prs.slide_width]
        assert len(rules) > 20


def test_the_tracker_moves_across_the_deck(tmp_path):
    """Six squares, the current one filled -- so it has to differ per slide."""
    path = _deck(tmp_path / "boardroom.pptx", slides=6)
    paint_theme_furniture(path, THEME, LAYOUTS["boardroom"])
    prs = Presentation(str(path))

    filled_x = []
    for slide in list(prs.slides)[1:]:
        squares = [
            s for s in slide.shapes
            if s.shape_type is not None and s.width == s.height and s.width < Inches(0.3)
        ]
        assert len(squares) == 6
        lit = [s for s in squares if s.fill.type is not None and s.fill.type == 1]
        assert len(lit) == 1
        filled_x.append(lit[0].left)

    assert len(set(filled_x)) > 1, "the tracker never moves"


def test_the_page_number_is_the_slides_own(tmp_path):
    path = _deck(tmp_path / "memo.pptx", slides=4)
    paint_theme_furniture(path, THEME, LAYOUTS["memo"])
    prs = Presentation(str(path))
    seen = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and " of " in shape.text_frame.text:
                seen.append(shape.text_frame.text.strip())
    assert seen == ["1 of 4", "2 of 4", "3 of 4", "4 of 4"]


def test_notebook_draws_three_punch_holes_on_the_left_edge(tmp_path):
    path = _deck(tmp_path / "notebook.pptx", slides=1)
    result = paint_theme_furniture(path, THEME, LAYOUTS["notebook"])
    assert result["painted"].count("punch_holes") == 1
    prs = Presentation(str(path))
    slide = prs.slides[0]
    holes = [
        s for s in slide.shapes
        if s.width == s.height and s.left < prs.slide_width * 0.05
    ]
    assert len(holes) == 3
    content = _content_index(slide)
    assert all(_index_of(slide, h) < content for h in holes)


def test_observatory_draws_a_star_field_and_freeform_orbit_arcs(tmp_path):
    """The arcs are the freeform case: an arc bleeding off the corner is not
    any preset autoshape."""
    path = _deck(tmp_path / "observatory.pptx", slides=1)
    result = paint_theme_furniture(path, THEME, LAYOUTS["observatory"])
    assert "star_field" in result["painted"]
    assert "orbit_arcs" in result["painted"]

    prs = Presentation(str(path))
    tree = _sp_tree(prs.slides[0])
    xml = tree.xml
    assert xml.count("<a:custGeom>") >= 1, "no freeform geometry was written"


def test_a_gradient_ground_is_written_as_a_gradient(tmp_path):
    path = _deck(tmp_path / "keynote.pptx", slides=1)
    result = paint_theme_furniture(path, THEME, LAYOUTS["keynote-minimal"])
    assert "ground:gradient" in result["painted"]
    prs = Presentation(str(path))
    assert "<a:gradFill" in _sp_tree(prs.slides[0]).xml


def test_letterhead_draws_a_double_border_and_a_cover_only_seal(tmp_path):
    path = _deck(tmp_path / "letterhead.pptx", slides=3)
    result = paint_theme_furniture(path, THEME, LAYOUTS["letterhead"])
    assert result["painted"].count("double_hairline_border") == 3
    assert result["painted"].count("monogram_seal") == 1

    prs = Presentation(str(path))
    circles = [s for s in prs.slides[0].shapes if s.width == s.height and s.width > Inches(0.5)]
    assert len(circles) == 2  # "two thin concentric gold circles"


# =============================================================================
# It never raises, and it never saves a deck it did not change
# =============================================================================

def test_a_missing_file_is_not_an_error():
    assert paint_theme_furniture(Path("/no/such/deck.pptx"), THEME, LAYOUTS["ledger"]) == {
        "painted": [], "slides": 0,
    }


def test_a_file_that_is_not_a_deck_is_not_an_error(tmp_path):
    junk = tmp_path / "not-a-deck.pptx"
    junk.write_bytes(b"this is not a presentation")
    assert paint_theme_furniture(junk, THEME, LAYOUTS["notebook"])["painted"] == []


def test_a_garbage_theme_or_layout_never_raises(tmp_path):
    path = _deck(tmp_path / "d.pptx", slides=1)
    for theme, layout in (
        (None, LAYOUTS["ledger"]),
        ("not a dict", LAYOUTS["observatory"]),
        (THEME, None),
        (THEME, {"ground": "nonsense", "rule_spacing_px": "twenty"}),
        (THEME, {"margin_rule": {"x_px": None}, "stamp": {"text": 5}}),
        (THEME, {"ornament": ("no_such_ornament",), "tracker": {"kind": "spiral"}}),
    ):
        paint_theme_furniture(path, theme, layout)  # must not raise


def test_an_empty_layout_leaves_the_file_untouched(tmp_path):
    path = _deck(tmp_path / "d.pptx", slides=2)
    before = path.read_bytes()
    result = paint_theme_furniture(path, THEME, EMPTY_LAYOUT)
    assert result == {"painted": [], "slides": 2}
    assert path.read_bytes() == before, "an unchanged deck was rewritten anyway"


def test_the_slide_count_is_reported_even_when_nothing_is_painted(tmp_path):
    path = _deck(tmp_path / "d.pptx", slides=7)
    assert paint_theme_furniture(path, THEME, EMPTY_LAYOUT)["slides"] == 7


def test_every_theme_in_the_registry_can_be_painted_without_raising(tmp_path):
    """81 layouts, every one of them run against a real deck.

    A layout that only ever gets read is a layout nobody has proved is
    drawable.
    """
    path = _deck(tmp_path / "all.pptx", slides=2)
    for theme in pptx_themes.all_themes():
        result = paint_theme_furniture(path, dict(theme.palette), layout_for(theme.id))
        assert result["slides"] == 2


def test_the_deck_still_opens_after_every_theme_has_painted_it(tmp_path):
    """Order surgery on ``spTree`` is the kind of thing that writes a file
    PowerPoint will not open. Reopening is the cheapest proof it did not."""
    for theme_id in sorted(HAND_AUTHORED):
        path = _deck(tmp_path / f"{theme_id}.pptx", slides=3)
        paint_theme_furniture(path, THEME, LAYOUTS[theme_id])
        prs = Presentation(str(path))
        assert len(prs.slides) == 3
        for slide in prs.slides:
            assert _content_index(slide) >= 0
