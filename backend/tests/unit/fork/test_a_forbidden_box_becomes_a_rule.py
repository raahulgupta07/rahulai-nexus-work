"""A theme that forbids boxes gets the reference's left-rule rows.

Measured against the vendored reference renders, our decks and the design
systems they claim differ structurally, not chromatically: ours are filled,
bordered panels; the reference for the SAME theme is four rows separated by
nothing but a thin left rule. `enforce_theme_rules` clears the mechanical
prohibitions and deliberately only COUNTS `boxes`; `convert_boxes_to_rules`
is the layout half.

★The load-bearing tests in this file are the REFUSALS. Converting a genuine
card row is the easy half and one assertion covers it. The expensive failures
all look like over-eagerness: a takeaway panel several themes REQUIRE, a chart
plot ground, a row of KPI tiles and a deck whose theme never forbade boxes at
all are every one of them a filled rectangle, and turning any of them into a
rule breaks the deck in a way no exception reports. A wrong conversion is
worse than an unconverted card, so most of what follows asserts that nothing
happened -- and each of those is paired with the positive control that fails
if the converter is simply disabled.
"""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from app.ai.decks import motifs
from app.ai.decks.motifs import convert_boxes_to_rules
from app.ai.decks.theme_layouts import CANVAS_H, CANVAS_W, EMPTY_LAYOUT

#: The theme dict as generated deck code receives it.
THEME = {
    "palette": {
        "background": "#FFFFFF",
        "primary_accent": "#C75146",
        "muted_text": "#5F6E5C",
        "heading_text": "#1F3A5F",
    },
    "fonts": ("Inter", "Inter"),
}

STARK = {"forbid_boxes": True}
LENIENT = {"forbid_boxes": False}

CARD_FILL = "F2F4F7"      # a light tint of the paper -- what a card is
PANEL_FILL = "1F3A5F"     # the loud shape -- what a takeaway panel is


# =============================================================================
# Building decks the way generated code builds them
# =============================================================================

def _x(prs, px):
    return int(px * prs.slide_width / CANVAS_W)


def _y(prs, px):
    return int(px * prs.slide_height / CANVAS_H)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(prs, slide, x, y, w, h, fill=CARD_FILL, outline=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(_x(prs, x)), Emu(_y(prs, y)), Emu(_x(prs, w)), Emu(_y(prs, h))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    if outline:
        shape.line.color.rgb = RGBColor.from_string("D0D5DD")
        shape.line.width = Pt(1)
    return shape


def _text(prs, slide, x, y, w, h, words):
    box = slide.shapes.add_textbox(
        Emu(_x(prs, x)), Emu(_y(prs, y)), Emu(_x(prs, w)), Emu(_y(prs, h))
    )
    run = box.text_frame.paragraphs[0].add_run()
    run.text = words
    run.font.size = Pt(12)
    return box


def _card_rows(prs, slide, count=4, fill=CARD_FILL, top=140, x=80, w=560, h=90, step=110):
    """The shape the model actually produces: filled panels with text on top."""
    for i in range(count):
        _rect(prs, slide, x, top + i * step, w, h, fill=fill)
        _text(prs, slide, x + 20, top + i * step + 22, w - 60, 40, f"row {i} says something")


def _deck(tmp_path, build, name="d.pptx"):
    prs = Presentation()
    slide = _blank(prs)
    build(prs, slide)
    path = tmp_path / name
    prs.save(str(path))
    return path


def _shapes(path):
    return list(Presentation(str(path)).slides[0].shapes)


def _text_boxes(path):
    """Every run of real text, with the position it was drawn at."""
    found = []
    for shape in _shapes(path):
        try:
            if shape.has_text_frame and shape.text_frame.text.strip():
                found.append(
                    (shape.text_frame.text, shape.left, shape.top, shape.width, shape.height)
                )
        except Exception:
            continue
    return sorted(found)


def _rules(path):
    """The thin vertical accent rules a conversion leaves behind."""
    prs = Presentation(str(path))
    limit = _x(prs, 8)
    return [
        s for s in prs.slides[0].shapes
        if s.width is not None and s.width <= limit and s.height is not None
        and s.height > s.width and s.fill.type == MSO_FILL.SOLID
    ]


# =============================================================================
# The conversion itself
# =============================================================================

def test_a_row_of_cards_becomes_a_row_of_rules(tmp_path):
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))

    result = convert_boxes_to_rules(path, THEME, STARK)

    assert result["cards_converted"] == 4
    assert result["cards_left"] == 0
    assert result["why_left"] == []


def test_the_converted_card_keeps_no_fill_and_no_border(tmp_path):
    """The prohibition is the point: a card that keeps its fill is unconverted."""
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))
    convert_boxes_to_rules(path, THEME, STARK)

    prs = Presentation(str(path))
    panels = [
        s for s in prs.slides[0].shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.width > _x(prs, 100)
    ]
    assert len(panels) == 4
    for shape in panels:
        assert shape.fill.type == MSO_FILL.BACKGROUND, "a card kept its fill"
        assert shape.line.fill.type == MSO_FILL.BACKGROUND, "a card kept its border"


def test_each_converted_card_gains_one_thin_rule_at_its_left_edge(tmp_path):
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))
    prs_before = Presentation(str(path))
    card_boxes = sorted(
        (s.left, s.top, s.height) for s in prs_before.slides[0].shapes
        if s.has_text_frame and not s.text_frame.text.strip()
    )

    convert_boxes_to_rules(path, THEME, STARK)

    rules = sorted((r.left, r.top, r.height) for r in _rules(path))
    assert len(rules) == 4
    assert rules == card_boxes, "a rule is not at its card's left edge and full height"


def test_the_rule_is_painted_in_the_theme_accent(tmp_path):
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))
    convert_boxes_to_rules(path, THEME, STARK)

    for rule in _rules(path):
        assert str(rule.fill.fore_color.rgb) == "C75146"


def test_no_word_on_the_slide_moves(tmp_path):
    """Text sits ON the card, so a conversion that reflows anything is wrong."""
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))
    before = _text_boxes(path)

    convert_boxes_to_rules(path, THEME, STARK)

    assert _text_boxes(path) == before


def test_the_rule_is_drawn_behind_the_text_it_belongs_beside(tmp_path):
    """★Z-order is the difficulty in this module and a rule over its own row's
    words is worse than no rule at all. Presence is satisfied by exactly the
    broken implementation this asserts against."""
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))
    convert_boxes_to_rules(path, THEME, STARK)

    prs = Presentation(str(path))
    slide = prs.slides[0]
    # `slide.shapes` iterates in `spTree` order, which is back to front.
    rule_positions, text_positions = [], []
    for index, shape in enumerate(slide.shapes):
        if shape.width <= _x(prs, 8) and shape.height > shape.width:
            rule_positions.append(index)
        elif shape.has_text_frame and shape.text_frame.text.strip():
            text_positions.append(index)

    assert len(rule_positions) == 4 and len(text_positions) == 4
    assert max(rule_positions) < max(text_positions), "a rule was painted over the words"
    for rule, words in zip(sorted(rule_positions), sorted(text_positions)):
        assert rule < words, "a row's rule sits in front of that row's own text"


def test_the_deck_still_opens_and_keeps_every_shape(tmp_path):
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))
    before = len(_shapes(path))

    convert_boxes_to_rules(path, THEME, STARK)

    after = _shapes(path)
    assert len(after) == before + 4, "a conversion added or lost a shape"


def test_several_slides_are_each_converted(tmp_path):
    prs = Presentation()
    for _ in range(3):
        _card_rows(prs, _blank(prs))
    path = tmp_path / "many.pptx"
    prs.save(str(path))

    assert convert_boxes_to_rules(path, THEME, STARK)["cards_converted"] == 12


# =============================================================================
# ★What must never be converted
# =============================================================================

def test_a_takeaway_panel_is_left_alone(tmp_path):
    """Several themes REQUIRE a bottom-line panel. It is dark and it is alone."""
    def build(prs, slide):
        _card_rows(prs, slide, count=0)
        _rect(prs, slide, 80, 560, 700, 90, fill=PANEL_FILL)
        _text(prs, slide, 100, 585, 600, 40, "So what: the number moved.")

    path = _deck(tmp_path, build)
    before = path.read_bytes()

    result = convert_boxes_to_rules(path, THEME, STARK)

    assert result["cards_converted"] == 0
    assert result["cards_left"] == 1
    assert "dark or strongly coloured" in result["why_left"][0]
    assert path.read_bytes() == before, "an unchanged deck was rewritten anyway"


def test_a_light_takeaway_panel_beside_a_card_row_survives_it(tmp_path):
    """The positive control for the refusal above: the row still converts.

    A refusal-only test passes on a converter that does nothing at all.
    """
    def build(prs, slide):
        _card_rows(prs, slide)
        _rect(prs, slide, 80, 600, 700, 60, fill="F7F7F5")
        _text(prs, slide, 100, 618, 500, 40, "So what: the number moved.")

    path = _deck(tmp_path, build)
    result = convert_boxes_to_rules(path, THEME, STARK)

    assert result["cards_converted"] == 4
    assert result["cards_left"] == 1
    assert "near-identical rows" in result["why_left"][0]


def test_a_chart_plot_background_is_left_alone(tmp_path):
    def build(prs, slide):
        _rect(prs, slide, 100, 120, 700, 400, fill="FAFAFA")
        data = CategoryChartData()
        data.categories = ["a", "b", "c"]
        data.add_series("s", (1.0, 2.0, 3.0))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Emu(_x(prs, 100)), Emu(_y(prs, 120)), Emu(_x(prs, 700)), Emu(_y(prs, 400)), data,
        )

    path = _deck(tmp_path, build)
    before = path.read_bytes()

    result = convert_boxes_to_rules(path, THEME, STARK)

    assert result["cards_converted"] == 0
    assert path.read_bytes() == before
    assert any(getattr(s, "has_chart", False) for s in _shapes(path)), "the chart is gone"


def test_a_row_of_kpi_tiles_is_left_alone(tmp_path):
    """Tiles are laid out ACROSS. A rule down the side of one rules off nothing."""
    def build(prs, slide):
        for i in range(4):
            x = 80 + i * 280
            _rect(prs, slide, x, 200, 240, 160)
            _text(prs, slide, x + 20, 240, 200, 60, f"{i}42%")

    path = _deck(tmp_path, build)
    result = convert_boxes_to_rules(path, THEME, STARK)

    assert result["cards_converted"] == 0
    assert result["cards_left"] == 4
    assert all("near-identical rows" in why for why in result["why_left"])


def test_a_full_bleed_decorative_shape_is_left_alone(tmp_path):
    def build(prs, slide):
        _rect(prs, slide, 0, 0, 1280, 720, fill="F5F5F5")
        _text(prs, slide, 100, 300, 600, 60, "a title over the bleed")

    path = _deck(tmp_path, build)
    assert convert_boxes_to_rules(path, THEME, STARK)["cards_converted"] == 0


def test_a_full_width_band_is_too_big_to_be_a_row(tmp_path):
    """Even light, even textual, even in pairs: a band that spans the slide is
    a header, a footer strip or a takeaway, never one of several rows."""
    def build(prs, slide):
        for i in range(2):
            _rect(prs, slide, 40, 120 + i * 300, 1200, 200)
            _text(prs, slide, 80, 160 + i * 300, 900, 60, f"band {i}")

    path = _deck(tmp_path, build)
    assert convert_boxes_to_rules(path, THEME, STARK)["cards_converted"] == 0


def test_a_small_swatch_or_divider_is_left_alone(tmp_path):
    def build(prs, slide):
        for i in range(4):
            _rect(prs, slide, 80, 200 + i * 40, 60, 12)
            _text(prs, slide, 90, 200 + i * 40, 40, 12, f"{i}")

    path = _deck(tmp_path, build)
    assert convert_boxes_to_rules(path, THEME, STARK)["cards_converted"] == 0


def test_a_panel_with_nothing_on_it_is_left_alone(tmp_path):
    """No text means a plot ground or an ornament, not a row of content."""
    def build(prs, slide):
        for i in range(4):
            _rect(prs, slide, 80, 140 + i * 110, 560, 90)

    path = _deck(tmp_path, build)
    result = convert_boxes_to_rules(path, THEME, STARK)

    assert result["cards_converted"] == 0
    assert result["cards_left"] == 4
    assert all("no text sits on it" in why for why in result["why_left"])


def test_an_oval_is_never_a_card(tmp_path):
    def build(prs, slide):
        for i in range(4):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Emu(_x(prs, 80)), Emu(_y(prs, 140 + i * 110)),
                Emu(_x(prs, 560)), Emu(_y(prs, 90)),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(CARD_FILL)
            _text(prs, slide, 100, 160 + i * 110, 400, 40, f"row {i}")

    path = _deck(tmp_path, build)
    assert convert_boxes_to_rules(path, THEME, STARK)["cards_converted"] == 0


def test_cards_inside_a_group_are_not_reached(tmp_path):
    """Deliberate: a group is composed geometry and this pass does not open it."""
    def build(prs, slide):
        _card_rows(prs, slide)

    path = _deck(tmp_path, build)
    prs = Presentation(str(path))
    slide = prs.slides[0]
    tree = slide.shapes._spTree
    group = tree.add_grpSp()
    for element in [el for el in list(tree) if el.tag == qn("p:sp")]:
        tree.remove(element)
        group.append(element)
    prs.save(str(path))

    assert convert_boxes_to_rules(path, THEME, STARK)["cards_converted"] == 0


def test_a_table_is_never_opened(tmp_path):
    def build(prs, slide):
        slide.shapes.add_table(
            3, 3, Emu(_x(prs, 100)), Emu(_y(prs, 150)), Emu(_x(prs, 800)), Emu(_y(prs, 300))
        )

    path = _deck(tmp_path, build)
    before = path.read_bytes()

    assert convert_boxes_to_rules(path, THEME, STARK)["cards_converted"] == 0
    assert path.read_bytes() == before


# =============================================================================
# ★The negative control: a theme that never forbade boxes
# =============================================================================

def test_a_theme_without_forbid_boxes_is_untouched(tmp_path):
    """An over-eager converter is a new bug, not a fix.

    Byte equality, not a count: a pass that opened and re-saved the deck
    without converting anything still rewrites the file, and that is the
    shape of a change nobody asked for.
    """
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))
    before = path.read_bytes()

    result = convert_boxes_to_rules(path, THEME, LENIENT)

    assert result == {"cards_converted": 0, "cards_left": 0, "why_left": []}
    assert path.read_bytes() == before


def test_a_layout_that_says_nothing_about_boxes_is_untouched(tmp_path):
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))
    before = path.read_bytes()

    assert convert_boxes_to_rules(path, THEME, EMPTY_LAYOUT)["cards_converted"] == 0
    assert path.read_bytes() == before


def test_the_same_deck_converts_when_the_theme_does_forbid_boxes(tmp_path):
    """The control for both tests above -- they are equally satisfied by a
    converter that never runs."""
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))
    assert convert_boxes_to_rules(path, THEME, STARK)["cards_converted"] == 4


# =============================================================================
# It never raises
# =============================================================================

def test_a_missing_file_is_not_an_error():
    assert convert_boxes_to_rules(Path("/no/such/deck.pptx"), THEME, STARK) == {
        "cards_converted": 0, "cards_left": 0, "why_left": [],
    }


def test_a_file_that_is_not_a_deck_is_not_an_error(tmp_path):
    junk = tmp_path / "not-a-deck.pptx"
    junk.write_bytes(b"this is not a presentation")

    assert convert_boxes_to_rules(junk, THEME, STARK)["cards_converted"] == 0


def test_a_truncated_deck_is_not_an_error(tmp_path):
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))
    path.write_bytes(path.read_bytes()[: len(path.read_bytes()) // 2])

    assert convert_boxes_to_rules(path, THEME, STARK)["cards_converted"] == 0


@pytest.mark.parametrize("theme,layout", [
    (None, STARK),
    ("not a dict", STARK),
    (THEME, None),
    (THEME, "not a dict"),
    (THEME, {"forbid_boxes": "yes please"}),
    (THEME, {"forbid_boxes": True, "accent": 17}),
    ({"palette": "not a dict"}, STARK),
])
def test_a_garbage_theme_or_layout_never_raises(tmp_path, theme, layout):
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))

    result = convert_boxes_to_rules(path, theme, layout)  # must not raise

    assert set(result) == {"cards_converted", "cards_left", "why_left"}


def test_an_empty_deck_is_not_an_error(tmp_path):
    prs = Presentation()
    path = tmp_path / "empty.pptx"
    prs.save(str(path))

    assert convert_boxes_to_rules(path, THEME, STARK)["cards_converted"] == 0


def test_a_deck_that_cannot_be_saved_reports_nothing_converted(tmp_path, monkeypatch):
    """Nothing persisted means nothing was done, and those cards are still cards."""
    path = _deck(tmp_path, lambda prs, slide: _card_rows(prs, slide))
    real = motifs.Presentation

    class _Unsaveable:
        def __init__(self, inner):
            self._inner = inner

        def __getattr__(self, name):
            return getattr(self._inner, name)

        def save(self, *args, **kwargs):
            raise OSError("disk full")

    monkeypatch.setattr(motifs, "Presentation", lambda p: _Unsaveable(real(p)))

    result = convert_boxes_to_rules(path, THEME, STARK)

    assert result["cards_converted"] == 0
    assert result["cards_left"] == 4
    assert any("could not be saved" in why for why in result["why_left"])


def test_the_why_list_is_bounded(tmp_path):
    """A summary, not a log: 200 refusals must not become 200 sentences."""
    def build(prs, slide):
        for i in range(60):
            _rect(prs, slide, 80, 40 + i * 11, 300, 40, fill=PANEL_FILL)
            _text(prs, slide, 90, 45 + i * 11, 200, 30, f"{i}")

    path = _deck(tmp_path, build)
    result = convert_boxes_to_rules(path, THEME, STARK)

    assert result["cards_left"] == 60
    assert len(result["why_left"]) <= 41
    assert "more left alone" in result["why_left"][-1]
