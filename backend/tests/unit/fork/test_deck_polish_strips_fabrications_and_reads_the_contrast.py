"""FIX-F / FIX-G / FIX-H (2026-08-21) — the polish passes on the shipped deck.

Measured on the phase-2 verification deck (numerically correct, gate-clean):
every slide still credited "Source: team analysis" — an attribution the model
invents as furniture — and the charts labelled their bars 86.438351 because
computing from rows ships raw float precision. And the original RCA's cause 5
was never fixed: nothing MEASURES readability on the shipped file, so a white
title on a white ground passes every existing check at exactly 1.00:1.

FIX-G deletes the fabricated credit (deletion, not replacement — no connector
name is available here, and an absent credit is honest). FIX-H gives shown
data labels a display number format. FIX-F measures text-vs-backdrop contrast
and reports — it never recolours.

★These build real pptx files and run the real passes — no schema, safe in
`tests/unit/fork/` where the migration fixture is a no-op.
"""
import pathlib

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt


def _blank_deck(tmp_path, name="deck.pptx"):
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    return prs, tmp_path / name


def _add_textbox(slide, text, color=None, left=1, top=1, w=8, h=1):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    box.text_frame.text = text
    if color is not None:
        box.text_frame.paragraphs[0].runs[0].font.color.rgb = color
    return box


class TestFabricatedAttributionsAreStripped:
    def test_the_measured_credit_line_is_deleted(self, tmp_path):
        """★"Source: team analysis" — the exact string shipped on every slide
        of the phase-2 verification deck."""
        from app.ai.decks.deck_polish import strip_fabricated_attribution

        prs, path = _blank_deck(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, "Revenue by branch")
        _add_textbox(slide, "Source: team analysis", top=6)
        prs.save(str(path))

        result = strip_fabricated_attribution(str(path))
        assert result["status"] == "checked"
        assert result["removed"] == 1

        texts = [
            s.text_frame.text
            for s in Presentation(str(path)).slides[0].shapes
            if s.has_text_frame
        ]
        assert "Source: team analysis" not in texts
        assert "Revenue by branch" in texts

    def test_the_variants_die_too(self, tmp_path):
        from app.ai.decks.deck_polish import strip_fabricated_attribution

        prs, path = _blank_deck(tmp_path)
        for line in (
            "Source: internal analysis",
            "Source: internal data analysis",
            "Source: analyst estimates",
            "Source - our analysis",
        ):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_textbox(slide, line)
        prs.save(str(path))

        assert strip_fabricated_attribution(str(path))["removed"] == 4

    def test_a_real_attribution_survives(self, tmp_path):
        """★The whole point of requiring a generic-analysis terminal word: a
        credit that names an actual source must never be touched."""
        from app.ai.decks.deck_polish import strip_fabricated_attribution

        prs, path = _blank_deck(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, "Source: Microsoft Fabric sales lakehouse")
        _add_textbox(slide, "Source: POS transaction data, FY2025", top=2)
        prs.save(str(path))

        assert strip_fabricated_attribution(str(path))["removed"] == 0
        texts = [
            s.text_frame.text
            for s in Presentation(str(path)).slides[0].shapes
            if s.has_text_frame
        ]
        assert len(texts) == 2

    def test_a_matched_paragraph_inside_a_bigger_frame_goes_alone(self, tmp_path):
        from app.ai.decks.deck_polish import strip_fabricated_attribution

        prs, path = _blank_deck(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = _add_textbox(slide, "Riverside closed at 1,459")
        box.text_frame.add_paragraph().text = "Source: team analysis"
        prs.save(str(path))

        assert strip_fabricated_attribution(str(path))["removed"] == 1
        reopened = Presentation(str(path))
        frame_texts = [
            s.text_frame.text
            for s in reopened.slides[0].shapes
            if s.has_text_frame
        ]
        assert frame_texts == ["Riverside closed at 1,459"]


class TestDataLabelsGetADisplayFormat:
    def _chart_deck(self, tmp_path, values):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        prs, path = _blank_deck(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        data = CategoryChartData()
        data.categories = [f"c{i}" for i in range(len(values))]
        data.add_series("s", values)
        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1), Inches(8), Inches(5), data,
        )
        plot = frame.chart.plots[0]
        plot.has_data_labels = True
        prs.save(str(path))
        return path

    def test_raw_float_labels_are_rounded(self, tmp_path):
        """★86.438351 — the exact label measured on the phase-2 deck."""
        from app.ai.decks.deck_polish import round_chart_data_labels

        path = self._chart_deck(tmp_path, (86.438351, 92.7431, 78.09992))
        result = round_chart_data_labels(str(path))
        assert result["status"] == "checked"
        assert result["charts_fixed"] == 1

        reopened = Presentation(str(path))
        chart = next(
            s.chart for s in reopened.slides[0].shapes if s.has_chart
        )
        assert chart.plots[0].data_labels.number_format == "#,##0.0"

    def test_thousands_read_as_integers(self, tmp_path):
        from app.ai.decks.deck_polish import round_chart_data_labels

        path = self._chart_deck(tmp_path, (1711.2989, 2540.77, 1102.4))
        round_chart_data_labels(str(path))
        chart = next(
            s.chart for s in Presentation(str(path)).slides[0].shapes if s.has_chart
        )
        assert chart.plots[0].data_labels.number_format == "#,##0"

    def test_a_chart_with_clean_values_is_left_alone(self, tmp_path):
        """No sub-tenth precision → nothing to fix, file untouched."""
        from app.ai.decks.deck_polish import round_chart_data_labels

        path = self._chart_deck(tmp_path, (3120.0, 2610.0, 1459.0))
        assert round_chart_data_labels(str(path))["charts_fixed"] == 0


class TestTheContrastCheckSeesWhatTheLayoutCheckCannot:
    def test_white_on_white_is_caught(self, tmp_path):
        """★RCA cause 5, the exact defect class: white title over a white
        ground, geometrically perfect, invisible."""
        from app.ai.decks.deck_polish import check_deck_contrast

        prs, path = _blank_deck(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(
            slide, "State of City Mart: Together We Win",
            color=RGBColor(0xFF, 0xFF, 0xFF),
        )
        prs.save(str(path))

        result = check_deck_contrast(str(path))
        assert result["status"] == "checked"
        assert len(result["issues"]) == 1
        assert result["issues"][0]["ratio"] < 1.1

    def test_light_text_on_a_dark_slide_bg_passes(self, tmp_path):
        """★The phase-1 case: the model set its own dark <p:bg> and used white
        text. A check that read the default white ground would flag every
        correct dark slide."""
        from app.ai.decks.deck_polish import check_deck_contrast

        prs, path = _blank_deck(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0B, 0x0F, 0x19)
        _add_textbox(slide, "Together We Win", color=RGBColor(0xFF, 0xFF, 0xFF))
        prs.save(str(path))

        assert check_deck_contrast(str(path))["issues"] == []

    def test_text_reads_against_its_own_card_not_the_slide(self, tmp_path):
        """White text INSIDE a dark card on a white slide is readable; judged
        against the slide it would be a false positive."""
        from pptx.enum.shapes import MSO_SHAPE

        from app.ai.decks.deck_polish import check_deck_contrast

        prs, path = _blank_deck(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(5), Inches(2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
        box = _add_textbox(
            slide, "3,120 — Downtown", color=RGBColor(0xFF, 0xFF, 0xFF),
            left=1.2, top=1.3, w=4, h=1,
        )
        assert box is not None
        prs.save(str(path))

        assert check_deck_contrast(str(path))["issues"] == []

    def test_an_inherited_colour_is_skipped_not_guessed(self, tmp_path):
        """★A run with no explicit colour resolves per master. Guessing black
        (or white) would flag readable decks and teach the agent to distrust
        the check — unresolved runs are skipped."""
        from app.ai.decks.deck_polish import check_deck_contrast

        prs, path = _blank_deck(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide, "No explicit colour anywhere here")
        prs.save(str(path))

        assert check_deck_contrast(str(path))["issues"] == []

    def test_a_missing_file_fails_open(self, tmp_path):
        from app.ai.decks.deck_polish import check_deck_contrast

        result = check_deck_contrast(str(tmp_path / "nope.pptx"))
        assert result["status"] == "unavailable"
