"""Phase 6.1 — build the file-support fixture corpus.

Every fixture is a REAL file of its format, produced by the same libraries or
by LibreOffice, not a renamed text file. A renamed fixture measures the
extension registry and nothing else; the point of this corpus is to exercise
the actual readers, so `.doc` has to be a real binary Word document and `.xls`
a real BIFF workbook.

Four of them are adversarial, chosen because each defeats a *different* guard:
an image-only PDF (nothing to extract), a one-line docx (trips the
MIN_USABLE_DOC_CHARS floor), a glyph-soup PDF (extraction "succeeds" and
returns nonsense), and a truncated docx (the reader raises).
"""
import json
import os
import shutil
import subprocess
import sys
import zipfile

OUT = sys.argv[1] if len(sys.argv) > 1 else "/out"
os.makedirs(OUT, exist_ok=True)

# The one sentence every text-bearing fixture carries. Probes assert on it, so
# a reader that returns *something* but not this is caught as garbage rather
# than counted as a pass.
MARK = "CityAgent fixture marker ALPHA-7731 quarterly revenue report."
ROWS = [("region", "revenue"), ("North", "1200"), ("South", "980"), ("East", "1440")]

made = []


def record(name, how):
    p = os.path.join(OUT, name)
    made.append({"file": name, "bytes": os.path.getsize(p) if os.path.exists(p) else 0,
                 "made_by": how, "exists": os.path.exists(p)})


def w(name, text, mode="w"):
    with open(os.path.join(OUT, name), mode) as fh:
        fh.write(text)


# ---------------------------------------------------------------- plain text
w("sample.txt", MARK + "\n")
record("sample.txt", "literal")

w("sample.md", f"# Report\n\n{MARK}\n\n| region | revenue |\n|---|---|\n| North | 1200 |\n")
record("sample.md", "literal")

w("sample.log", f"2026-08-03 12:00:00 INFO {MARK}\n2026-08-03 12:00:01 WARN disk 91%\n")
record("sample.log", "literal")

w("sample.csv", "\n".join(",".join(r) for r in ROWS) + "\n")
record("sample.csv", "literal")

w("sample.tsv", "\n".join("\t".join(r) for r in ROWS) + "\n")
record("sample.tsv", "literal")

w("sample.json", json.dumps([{"region": r[0], "revenue": int(r[1])} for r in ROWS[1:]], indent=2))
record("sample.json", "literal")

w("sample.ndjson", "".join(json.dumps({"region": r[0], "revenue": int(r[1])}) + "\n" for r in ROWS[1:]))
record("sample.ndjson", "literal")

w("sample.html", f"<html><body><h1>Report</h1><p>{MARK}</p>"
                 "<table><tr><th>region</th><th>revenue</th></tr>"
                 "<tr><td>North</td><td>1200</td></tr></table></body></html>")
record("sample.html", "literal")

w("sample.xml", f"<?xml version='1.0'?><report><note>{MARK}</note>"
                "<row region='North' revenue='1200'/></report>")
record("sample.xml", "literal")

w("sample.yaml", f"title: Report\nnote: {MARK}\nrows:\n  - region: North\n    revenue: 1200\n")
record("sample.yaml", "literal")

# ---------------------------------------------------------------- tabular bin
import pandas as pd
df = pd.DataFrame([{"region": r[0], "revenue": int(r[1])} for r in ROWS[1:]])

df.to_excel(os.path.join(OUT, "sample.xlsx"), index=False)
record("sample.xlsx", "pandas/openpyxl")

df.to_parquet(os.path.join(OUT, "sample.parquet"), index=False)
record("sample.parquet", "pandas/pyarrow")

# ---------------------------------------------------------------- ooxml docs
import docx
d = docx.Document()
d.add_heading("Quarterly Report", 0)
d.add_paragraph(MARK)
t = d.add_table(rows=1, cols=2)
t.rows[0].cells[0].text = "region"
t.rows[0].cells[1].text = "revenue"
d.save(os.path.join(OUT, "sample.docx"))
record("sample.docx", "python-docx")

# Adversarial: one short line. Non-whitespace length is deliberately under
# MIN_USABLE_DOC_CHARS (16) so the floor fires if it is applied to OOXML.
d2 = docx.Document()
d2.add_paragraph("Rev up 4%.")
d2.save(os.path.join(OUT, "adv_oneline.docx"))
record("adv_oneline.docx", "python-docx (adversarial: 10 chars)")

from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
s = prs.slides.add_slide(prs.slide_layouts[5])
s.shapes.title.text = "Quarterly Report"
s.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1)).text_frame.text = MARK
prs.save(os.path.join(OUT, "sample.pptx"))
record("sample.pptx", "python-pptx")

# ---------------------------------------------------------------- images
from PIL import Image, ImageDraw
img = Image.new("RGB", (96, 32), "white")
ImageDraw.Draw(img).text((2, 10), MARK[:20], fill="black")
for ext, kw in (("png", {}), ("jpg", {"quality": 92}), ("bmp", {}),
                ("tiff", {}), ("webp", {})):
    img.save(os.path.join(OUT, f"sample.{ext}"), **kw)
    record(f"sample.{ext}", "PIL")

# Adversarial: a PDF whose only content is that raster. Text extraction has
# literally nothing to find, which is the scanned-document case.
img.save(os.path.join(OUT, "adv_imageonly.pdf"), "PDF", resolution=150)
record("adv_imageonly.pdf", "PIL (adversarial: image-only PDF)")

# ---------------------------------------------------------------- email / zip
from email.message import EmailMessage
m = EmailMessage()
m["From"] = "finance@cityagent.io"
m["To"] = "ops@cityagent.io"
m["Subject"] = "Quarterly numbers"
m.set_content(MARK)
w("sample.eml", m.as_string())
record("sample.eml", "email.message")

with zipfile.ZipFile(os.path.join(OUT, "sample.zip"), "w") as z:
    z.write(os.path.join(OUT, "sample.csv"), "sample.csv")
    z.write(os.path.join(OUT, "sample.txt"), "sample.txt")
record("sample.zip", "zipfile")

# ---------------------------------------------------------------- LibreOffice
# The legacy binary formats have no pure-python writer worth trusting, so they
# are converted from the OOXML fixtures above. This also produces the PDF, so
# the text-bearing PDF is a real typeset document rather than a raster.
def convert(src, target, outname=None):
    subprocess.run(["soffice", "--headless", "--convert-to", target,
                    "--outdir", OUT, os.path.join(OUT, src)],
                   check=False, capture_output=True, timeout=180)
    if outname:
        record(outname, f"libreoffice {src} -> {target}")


convert("sample.docx", "pdf", "sample.pdf")
convert("sample.docx", "doc", "sample.doc")
convert("sample.docx", "rtf", "sample.rtf")
convert("sample.docx", "odt", "sample.odt")
convert("sample.pptx", "ppt", "sample.ppt")
convert("sample.pptx", "odp", "sample.odp")
convert("sample.xlsx", "xls", "sample.xls")

# ---------------------------------------------------------------- corrupt
# Adversarial: a real docx with its tail cut off. zipfile raises rather than
# returning nothing, which is a different code path from an empty extraction.
src = os.path.join(OUT, "sample.docx")
data = open(src, "rb").read()
with open(os.path.join(OUT, "adv_corrupt.docx"), "wb") as fh:
    fh.write(data[: len(data) // 2])
record("adv_corrupt.docx", "truncated docx (adversarial)")

# Adversarial: glyph soup. A real PDF with every ToUnicode CMap removed still
# renders, but pypdf then has no way to map glyph codes back to characters —
# the signature failure doc_text_looks_garbled() exists for.
try:
    from pypdf import PdfReader, PdfWriter
    r = PdfReader(os.path.join(OUT, "sample.pdf"))
    wtr = PdfWriter()
    for page in r.pages:
        res = page.get("/Resources")
        fonts = res.get("/Font") if res else None
        if fonts:
            for k in list(fonts.keys()):
                fobj = fonts[k].get_object()
                fobj.pop("/ToUnicode", None)
                for df in (fobj.get("/DescendantFonts") or []):
                    df.get_object().pop("/ToUnicode", None)
        wtr.add_page(page)
    with open(os.path.join(OUT, "adv_garbled.pdf"), "wb") as fh:
        wtr.write(fh)
    record("adv_garbled.pdf", "pypdf: ToUnicode stripped (adversarial)")
except Exception as e:
    print("garbled pdf FAILED:", e, file=sys.stderr)

# Housekeeping: soffice leaves lock files behind.
for f in os.listdir(OUT):
    if f.startswith("."):
        os.remove(os.path.join(OUT, f))

print(json.dumps(made, indent=2))
missing = [m["file"] for m in made if not m["exists"] or m["bytes"] == 0]
print(f"\nTOTAL {len(made)}  MISSING/EMPTY {len(missing)}: {missing}", file=sys.stderr)
